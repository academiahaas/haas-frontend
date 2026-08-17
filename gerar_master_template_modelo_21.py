import os
import math
import urllib.request
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# PALETA OFICIAL HAAS LANGUAGE - BENTO POP EDITION
BRANCO = RGBColor(0xFF, 0xFF, 0xFF)
ROXO = RGBColor(0x4F, 0x35, 0xE2)
AZUL = RGBColor(0x1E, 0x3A, 0x8A)
SLATE = RGBColor(0x33, 0x41, 0x55)
TURQUESA = RGBColor(0x14, 0xC8, 0xB0)
ROXO_CLARO = RGBColor(0xED, 0xE9, 0xFE)
TURQUESA_CLARO = RGBColor(0xE0, 0xF2, 0xF1)
FUNDO_BG = RGBColor(0xF1, 0xF5, 0xF9) # Off-white para destacar o Bento Grid

LOGO_URL = "https://jdppxfokfhqjudwfwckd.supabase.co/storage/v1/object/public/haas-academy/assignments/Design%20sem%20nome%20(4).png"
LOGO_PATH = "/tmp/haas-logo.png"

if not os.path.exists(LOGO_PATH):
    try:
        urllib.request.urlretrieve(LOGO_URL, LOGO_PATH)
    except Exception:
        pass

W, H = Inches(13.333), Inches(7.5)
TOTAL_SLIDES = 22

prs = Presentation()
prs.slide_width = W
prs.slide_height = H

def base_slide(num, hide_logo=False):
    sl = prs.slides.add_slide(prs.slide_layouts[6])
    sl.background.fill.solid()
    sl.background.fill.fore_color.rgb = FUNDO_BG

    if os.path.exists(LOGO_PATH) and not hide_logo:
        sl.shapes.add_picture(LOGO_PATH, Inches(10.8), Inches(0.2), height=Inches(0.65))
        
    tb = sl.shapes.add_textbox(Inches(0.5), Inches(7.12), Inches(8), Inches(0.35))
    p = tb.text_frame.paragraphs[0]
    p.text = "HAAS LANGUAGE  •  BENTO POP COLOR GRID EDITION 📱"
    p.font.size = Pt(10); p.font.color.rgb = ROXO; p.font.bold = True
    
    tb2 = sl.shapes.add_textbox(Inches(10.8), Inches(7.12), Inches(2.0), Inches(0.35))
    p2 = tb2.text_frame.paragraphs[0]
    p2.text = f"{num:02d} / {TOTAL_SLIDES:02d}"
    p2.font.size = Pt(10); p2.font.color.rgb = AZUL; p2.font.bold = True; p2.alignment = PP_ALIGN.RIGHT
    return sl

def add_header(sl, tit, sub=""):
    tb = sl.shapes.add_textbox(Inches(0.5), Inches(0.25), Inches(10), Inches(0.65))
    tf = tb.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]; p.text = tit
    p.font.size = Pt(30); p.font.color.rgb = ROXO; p.font.bold = True
    
    if sub:
        tb_s = sl.shapes.add_textbox(Inches(0.5), Inches(0.95), Inches(10), Inches(0.4))
        p_s = tb_s.text_frame.paragraphs[0]
        p_s.text = sub
        p_s.font.size = Pt(16); p_s.font.color.rgb = AZUL; p_s.font.bold = True

def add_bento_tile(sl, tit, desc, left, top, width, height, bg_c=BRANCO, tit_c=ROXO, desc_c=SLATE, tit_sz=16, desc_sz=14, border_c=None, badge=None, align_center_v=False):
    """Cria um Card Bento Tile para o painel pop"""
    s = sl.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    s.fill.solid(); s.fill.fore_color.rgb = bg_c
    if border_c:
        s.line.color.rgb = border_c; s.line.width = Pt(1.5)
    else:
        s.line.fill.background()
        
    tf = s.text_frame; tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE if align_center_v else MSO_ANCHOR.TOP
    
    margin_v = Inches(0.12) if height < Inches(1.2) else Inches(0.2)
    tf.margin_left = margin_v; tf.margin_top = margin_v; tf.margin_right = margin_v; tf.margin_bottom = margin_v
    
    if badge:
        p_b = tf.paragraphs[0]
        p_b.text = badge.upper()
        p_b.font.size = Pt(10); p_b.font.color.rgb = tit_c; p_b.font.bold = True
        
        if tit:
            p = tf.add_paragraph()
            p.text = tit; p.font.size = Pt(tit_sz); p.font.color.rgb = tit_c; p.font.bold = True; p.space_before = Pt(4)
    else:
        if tit:
            p = tf.paragraphs[0]
            p.text = tit; p.font.size = Pt(tit_sz); p.font.color.rgb = tit_c; p.font.bold = True
            
    if desc:
        p2 = tf.add_paragraph() if (badge or tit) else tf.paragraphs[0]
        p2.text = desc; p2.font.size = Pt(desc_sz); p2.font.color.rgb = desc_c
        if badge or tit:
            p2.space_before = Pt(4)

def img_bento_box(sl, left, top, width, height, label="[IMAGEM BENTO]"):
    s = sl.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    s.fill.solid(); s.fill.fore_color.rgb = ROXO_CLARO
    s.line.color.rgb = TURQUESA; s.line.width = Pt(2)
    tf = s.text_frame; tf.word_wrap = True; tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER
    p = tf.paragraphs[0]; p.text = label; p.font.size = Pt(16); p.font.color.rgb = ROXO; p.font.bold = True

def add_bento_flashcard(sl, num_str, desc, left, top, width, height, bg_c=BRANCO, b_c=ROXO):
    s = sl.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    s.fill.solid(); s.fill.fore_color.rgb = bg_c
    s.line.fill.background()
    
    chip = sl.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left + Inches(0.12), top + Inches(0.12), Inches(0.65), Inches(0.35))
    chip.fill.solid(); chip.fill.fore_color.rgb = b_c; chip.line.fill.background()
    tf_chip = chip.text_frame; tf_chip.vertical_anchor = MSO_ANCHOR.MIDDLE
    p_chip = tf_chip.paragraphs[0]; p_chip.text = f"#{num_str}"; p_chip.alignment = PP_ALIGN.CENTER
    p_chip.font.size = Pt(11); p_chip.font.color.rgb = BRANCO; p_chip.font.bold = True
    
    tf = s.text_frame; tf.word_wrap = True; tf.vertical_anchor = MSO_ANCHOR.TOP
    tf.margin_left = Inches(0.88); tf.margin_top = Inches(0.15); tf.margin_right = Inches(0.15); tf.margin_bottom = Inches(0.12)
    p = tf.paragraphs[0]; p.text = desc; p.font.size = Pt(14); p.font.color.rgb = SLATE if bg_c == BRANCO else ROXO; p.font.bold = False

# ==========================================================
# CONSTRUÇÃO DOS 22 SLIDES - MODELO 21
# ==========================================================

# SLIDE 1: CAPA (Mosaico Bento de 4 Tiles)
sl1 = base_slide(1, hide_logo=True)
add_bento_tile(sl1, "[Título da Aula / Tema Principal]", "📅 Data da Aula: [DD/MM/AAAA]", Inches(0.5), Inches(1.2), Inches(8.0), Inches(3.2), bg_c=ROXO, tit_c=BRANCO, desc_c=TURQUESA_CLARO, tit_sz=36, desc_sz=18, badge="LESSON TITLE")
add_bento_tile(sl1, "HAAS BENTO", "POP GRID v2.1", Inches(8.62), Inches(1.2), Inches(4.21), Inches(1.5), bg_c=TURQUESA, tit_c=AZUL, desc_c=AZUL, tit_sz=20, desc_sz=14, align_center_v=True)
add_bento_tile(sl1, "", "[Subtítulo / Descrição da Aula]", Inches(0.5), Inches(4.52), Inches(8.0), Inches(2.3), bg_c=AZUL, tit_c=BRANCO, desc_c=BRANCO, desc_sz=18, align_center_v=True)
img_bento_box(sl1, Inches(8.62), Inches(2.82), Inches(4.21), Inches(4.0), "[Imagem Hero Bento]")

# SLIDE 2: PRESENÇA (Painel Pop 4 Widgets)
sl2 = base_slide(2)
add_header(sl2, "Presença & Integrantes", "Painel da Turma")
add_bento_tile(sl2, "Professor(a):", "[Nome do(a) Professor(a)]", Inches(0.5), Inches(1.5), Inches(6.0), Inches(2.5), bg_c=ROXO, tit_c=TURQUESA, desc_c=BRANCO, tit_sz=20, desc_sz=16, badge="DOCENTE")
add_bento_tile(sl2, "Aluno(s):", "[Nome do Aluno / Lista da Turma]", Inches(6.62), Inches(1.5), Inches(6.21), Inches(5.3), bg_c=TURQUESA, tit_c=AZUL, desc_c=AZUL, tit_sz=20, desc_sz=16, badge="ALUNOS")
add_bento_tile(sl2, "Sessão Ativa", "Haas High-Level Education", Inches(0.5), Inches(4.12), Inches(6.0), Inches(2.68), bg_c=AZUL, tit_c=BRANCO, desc_c=TURQUESA_CLARO, tit_sz=18, desc_sz=14, badge="STATUS")

# SLIDE 3: WARM-UP (3 Bento Widgets)
sl3 = base_slide(3)
add_header(sl3, "Aquecimento / Warm-up", "Foco 100% na Fala (3 Minutos)")
add_bento_tile(sl3, "Pergunta Provocativa Rápida:", "[Insira aqui a pergunta provocativa de aquecimento]", Inches(0.5), Inches(1.5), Inches(7.5), Inches(2.6), bg_c=ROXO, tit_c=BRANCO, desc_c=BRANCO, tit_sz=20, desc_sz=16, badge="PROMPT")
add_bento_tile(sl3, "Fundamentação Neurocognitiva:", "Ativa o córtex pré-frontal e reduz a inibição linguística inicial.", Inches(0.5), Inches(4.22), Inches(7.5), Inches(2.58), bg_c=TURQUESA_CLARO, border_c=TURQUESA, tit_c=AZUL, desc_c=SLATE, tit_sz=18, desc_sz=15, align_center_v=True)
img_bento_box(sl3, Inches(8.12), Inches(1.5), Inches(4.71), Inches(5.3), "[Imagem Warm-up]")

# SLIDE 4: OBJETIVOS (Bento Tiles Goals 1, 2, 3)
sl4 = base_slide(4)
add_header(sl4, "Objetivos da Aula", "Metas da Sessão")
add_bento_tile(sl4, "Goal #01", "[Objetivo 1]", Inches(0.5), Inches(1.5), Inches(7.0), Inches(1.6), bg_c=ROXO, tit_c=TURQUESA, desc_c=BRANCO, tit_sz=16, desc_sz=15, align_center_v=True)
add_bento_tile(sl4, "Goal #02", "[Objetivo 2]", Inches(0.5), Inches(3.22), Inches(7.0), Inches(1.6), bg_c=AZUL, tit_c=TURQUESA, desc_c=BRANCO, tit_sz=16, desc_sz=15, align_center_v=True)
add_bento_tile(sl4, "Goal #03", "[Objetivo 3]", Inches(0.5), Inches(4.94), Inches(7.0), Inches(1.86), bg_c=TURQUESA, tit_c=AZUL, desc_c=AZUL, tit_sz=16, desc_sz=15, align_center_v=True)
add_bento_tile(sl4, "Nota de Escrita:", "Nesta aula daremos atenção à escrita para reforçar a memorização de estruturas/vocabulário e reduzir a interferência da língua materna.", Inches(7.62), Inches(1.5), Inches(5.21), Inches(5.3), bg_c=ROXO_CLARO, border_c=ROXO, tit_c=ROXO, desc_c=SLATE, tit_sz=18, desc_sz=15)

# SLIDE 5: METODOLOGIA (Mosaico Bento Pop 3x2)
sl5 = base_slide(5)
add_header(sl5, "A Nossa Metodologia", "Estrutura Haas Language")
add_bento_tile(sl5, "1. Fala Ativa", "Conversação inicial sem julgamentos.", Inches(0.5), Inches(1.5), Inches(3.9), Inches(2.5), bg_c=ROXO, tit_c=TURQUESA, desc_c=BRANCO, badge="PASSO 01")
add_bento_tile(sl5, "2. Exposição", "Apresentação leve da gramática.", Inches(4.52), Inches(1.5), Inches(3.9), Inches(2.5), bg_c=TURQUESA, tit_c=AZUL, desc_c=AZUL, badge="PASSO 02")
add_bento_tile(sl5, "3. Prática", "Exercícios com suporte.", Inches(8.54), Inches(1.5), Inches(4.29), Inches(2.5), bg_c=AZUL, tit_c=TURQUESA, desc_c=BRANCO, badge="PASSO 03")
add_bento_tile(sl5, "4. Conversação", "Aplicação prática das estruturas em contexto real.", Inches(0.5), Inches(4.12), Inches(6.0), Inches(2.68), bg_c=ROXO_CLARO, border_c=ROXO, tit_c=ROXO, desc_c=SLATE, badge="PASSO 04")
add_bento_tile(sl5, "5. Escrita", "Consolidação visual e redução da tradução automática.", Inches(6.62), Inches(4.12), Inches(6.21), Inches(2.68), bg_c=TURQUESA_CLARO, border_c=TURQUESA, tit_c=AZUL, desc_c=SLATE, badge="PASSO 05")

# SLIDES 6 A 8: EXPOSIÇÃO (Mosaico Bento Triplo)
for s_num in [6, 7, 8]:
    sl = base_slide(s_num)
    add_header(sl, f"Explicação do Tema — Parte {s_num - 5}")
    add_bento_tile(sl, "Conteúdo Principal:", f"[Explicação e exemplo central da estrutura - Bloco {s_num - 5}]", Inches(0.5), Inches(1.5), Inches(12.333), Inches(2.5), bg_c=ROXO, tit_c=TURQUESA, desc_c=BRANCO, tit_sz=20, desc_sz=15, badge="FOCUS")
    add_bento_tile(sl, "Exemplos Práticos:", "[Exemplos adicionais e frases do dia a dia]", Inches(0.5), Inches(4.12), Inches(6.0), Inches(2.68), bg_c=TURQUESA, tit_c=AZUL, desc_c=AZUL, badge="EXAMPLES")
    add_bento_tile(sl, "Notas & Detalhes:", "[Exceções, nuances gramaticais e contexto de uso]", Inches(6.62), Inches(4.12), Inches(6.21), Inches(2.68), bg_c=AZUL, tit_c=TURQUESA, desc_c=BRANCO, badge="RULES")

# SLIDE 9: SÍNTESE + CULTURA
sl9 = base_slide(9)
add_header(sl9, "Síntese do Conteúdo & Sugestão Cultural")
add_bento_tile(sl9, "Resumo do Tema:", "[Resumo sintético dos pontos chave ensinados na aula]", Inches(0.5), Inches(1.5), Inches(6.0), Inches(5.3), bg_c=AZUL, tit_c=TURQUESA, desc_c=BRANCO, tit_sz=18, desc_sz=15, badge="SUMMARY")
add_bento_tile(sl9, "Sugestão Cultural:", "🎵 [Nome da Música / Vídeo recomendado sobre o assunto].\n\nO professor pode rodar o áudio/vídeo ou deixar como recomendação extra para o aluno praticar.", Inches(6.62), Inches(1.5), Inches(6.21), Inches(5.3), bg_c=ROXO, tit_c=TURQUESA, desc_c=BRANCO, tit_sz=18, desc_sz=15, badge="MEDIA")

# SLIDES 10 E 11: PRÁTICA GUIADA
sl10 = base_slide(10)
add_header(sl10, "Prática Guiada — Blocos 1 a 5", "Correção de Vícios")
add_bento_flashcard(sl10, "01", "[Lacunas contextualizadas com opções]", Inches(0.5), Inches(1.5), Inches(12.333), Inches(1.2), bg_c=BRANCO, b_c=ROXO)
add_bento_flashcard(sl10, "02", "[Lacunas contextualizadas com opções]", Inches(0.5), Inches(2.82), Inches(6.0), Inches(1.8), bg_c=TURQUESA_CLARO, b_c=TURQUESA)
add_bento_flashcard(sl10, "03", "[Preenchimento direto de verbos]", Inches(6.62), Inches(2.82), Inches(6.21), Inches(1.8), bg_c=TURQUESA_CLARO, b_c=TURQUESA)
add_bento_flashcard(sl10, "04", "[Reestruturação / Transformação de frase]", Inches(0.5), Inches(4.74), Inches(6.0), Inches(2.06), bg_c=ROXO_CLARO, b_c=ROXO)
add_bento_flashcard(sl10, "05", "[Reestruturação / Transformação de frase]", Inches(6.62), Inches(4.74), Inches(6.21), Inches(2.06), bg_c=ROXO_CLARO, b_c=ROXO)

sl11 = base_slide(11)
add_header(sl11, "Prática Guiada — Blocos 6 a 10", "Desafios Práticos")
add_bento_flashcard(sl11, "06", "[Preposição / Conector correto]", Inches(0.5), Inches(1.5), Inches(12.333), Inches(1.2), bg_c=BRANCO, b_c=TURQUESA)
add_bento_flashcard(sl11, "07", "[Escolha da opção correta]", Inches(0.5), Inches(2.82), Inches(6.0), Inches(1.8), bg_c=ROXO_CLARO, b_c=ROXO)
add_bento_flashcard(sl11, "08", "[Substituição por sinônimo formal]", Inches(6.62), Inches(2.82), Inches(6.21), Inches(1.8), bg_c=ROXO_CLARO, b_c=ROXO)
add_bento_flashcard(sl11, "09", "[Identificação de erro e correção]", Inches(0.5), Inches(4.74), Inches(6.0), Inches(2.06), bg_c=TURQUESA_CLARO, b_c=TURQUESA)
add_bento_flashcard(sl11, "10", "[Tradução de sentido sem uso literal]", Inches(6.62), Inches(4.74), Inches(6.21), Inches(2.06), bg_c=TURQUESA_CLARO, b_c=TURQUESA)

# ==========================================================
# SLIDES 12 A 19: CONVERSAÇÃO (INOVAÇÃO BENTO WIDGET WALL INÉDITA!)
# ESQUERDA: PERGUNTA ATIVA BENTO (X = 0.5", W = 6.9") + PROMPTS (Y = 4.3")
# DIREITA: PAREDE DE WIDGETS DE HISTÓRICO BENTO GRID COLORIDO (X = 7.6", W = 5.2")
# ==========================================================
perguntas_config = [
    "[Pergunta principal 1 sobre o tema]", "[Pergunta principal 2 sobre o tema]", "[Pergunta principal 3 sobre o tema]",
    "[Pergunta principal 4 sobre o tema]", "[Pergunta principal 5 sobre o tema]", "[Pergunta principal 6 sobre o tema]",
    "[Pergunta principal 7 sobre o tema]", "[Pergunta principal 8 sobre o tema]"
]

def add_bento_conversation_widgets(sl, history_list):
    """Gera uma parede de Widgets Bento Grid Coloridos na direita com as perguntas passadas"""
    count = len(history_list)
    if count == 0: return
    
    cols = 2 if count <= 4 else 3
    rows = math.ceil(count / cols)
    
    gap = 0.12
    avail_w = 5.21 - (gap * (cols - 1))
    avail_h = 5.3 - (gap * (rows - 1))
    
    w_tile = avail_w / cols
    h_tile = avail_h / rows
    
    bento_colors = [
        (TURQUESA, AZUL),
        (ROXO, BRANCO),
        (AZUL, BRANCO),
        (ROXO_CLARO, ROXO),
        (TURQUESA_CLARO, AZUL),
        (ROXO, TURQUESA),
        (AZUL, TURQUESA)
    ]
    
    for i, h_txt in enumerate(history_list):
        r = i // cols
        c = i % cols
        
        x = 7.62 + c * (w_tile + gap)
        y = 1.5 + r * (h_tile + gap)
        
        bg_c, txt_c = bento_colors[i % len(bento_colors)]
        clean_t = h_txt[:28] + "..." if len(h_txt) > 28 else h_txt
        d_sz = 11 if count <= 3 else (10 if count <= 5 else 9)
        
        add_bento_tile(sl, f"Q{i+1:02d}", clean_t, Inches(x), Inches(y), Inches(w_tile), Inches(h_tile),
                       bg_c=bg_c, tit_c=txt_c, desc_c=txt_c, tit_sz=11, desc_sz=d_sz, align_center_v=True)

for idx in range(8):
    s_num = 12 + idx
    sl = base_slide(s_num)
    add_header(sl, "Conversação", f"Sessão de Fala [{idx + 1}/8]")
    p_desc = perguntas_config[idx]
    
    if idx == 0:
        # Pergunta 1 Ampla Hero Bento
        add_bento_tile(sl, f"PERGUNTA ATIVA Q0{idx+1}:", p_desc, Inches(0.5), Inches(1.5), Inches(12.333), Inches(2.7), bg_c=ROXO, tit_c=TURQUESA, desc_c=BRANCO, tit_sz=18, desc_sz=18, align_center_v=True)
        add_bento_tile(sl, "PROMPTS DE APOIO:", "• Para expandir: 'Dê um exemplo recente...'\n• Para concordar/discordar: 'Qual é o outro lado dessa moeda?'", Inches(0.5), Inches(4.32), Inches(12.333), Inches(2.48), bg_c=TURQUESA_CLARO, border_c=TURQUESA, tit_c=AZUL, desc_c=SLATE, tit_sz=16, desc_sz=14, align_center_v=True)
    else:
        # LADO ESQUERDO: Pergunta Ativa + Prompts
        add_bento_tile(sl, f"PERGUNTA ATIVA Q0{idx+1}:", p_desc, Inches(0.5), Inches(1.5), Inches(6.9), Inches(2.7), bg_c=ROXO, tit_c=TURQUESA, desc_c=BRANCO, tit_sz=18, desc_sz=16, align_center_v=True)
        add_bento_tile(sl, "PROMPTS DE APOIO:", "• Para expandir: 'Dê um exemplo recente...'\n• Para discordar: 'Qual é o outro lado?'", Inches(0.5), Inches(4.32), Inches(6.9), Inches(2.48), bg_c=TURQUESA_CLARO, border_c=TURQUESA, tit_c=AZUL, desc_c=SLATE, tit_sz=15, desc_sz=14, align_center_v=True)
        
        # LADO DIREITO: Parede de Widgets Bento Coloridos de Histórico!
        history_list = [perguntas_config[i] for i in range(idx)]
        add_bento_conversation_widgets(sl, history_list)

# SLIDE 20: ESCRITA
sl20 = base_slide(20)
add_header(sl20, "Atividade de Escrita", "Consolidação e Memorização")
add_bento_tile(sl20, "Tarefa de Escrita Recomendada:", "[Instrução detalhada da tarefa: escrever um texto, frases ou diálogo]", Inches(0.5), Inches(1.5), Inches(12.333), Inches(2.7), bg_c=ROXO, tit_c=TURQUESA, desc_c=BRANCO, tit_sz=18, desc_sz=15, badge="WRITING TASK")
add_bento_tile(sl20, "Fundamentação Pedagógica:", "Identificamos que a escrita nesta etapa é fundamental para você fixar as estruturas aprendidas, expandir vocabulário ativo e diminuir a tradução mental automática do seu idioma nativo.", Inches(0.5), Inches(4.32), Inches(12.333), Inches(2.48), bg_c=TURQUESA_CLARO, border_c=TURQUESA, tit_c=AZUL, desc_c=SLATE, desc_sz=14, align_center_v=True, badge="RATIONALE")

# SLIDE 21: TAREFA DE CASA
sl21 = base_slide(21)
add_header(sl21, "Encerramento — Tarefa de Casa", "Ação e Prática Contínua")
add_bento_tile(sl21, "Tarefa de Casa:", "• [Atividade prática 1 para realizar fora da aula]\n• [Atividade de fixação de vocabulário ou escuta]", Inches(0.5), Inches(1.5), Inches(12.333), Inches(5.3), bg_c=AZUL, tit_c=TURQUESA, desc_c=BRANCO, tit_sz=18, desc_sz=15, badge="HOMEWORK")

# SLIDE 22: ENCERRAMENTO & FEEDBACK
sl22 = base_slide(22)
add_header(sl22, "Encerramento & Feedback Rápido", "Considerações Finais")
add_bento_tile(sl22, "Agradecimento:", "Obrigado pela dedicação na aula de hoje!", Inches(0.5), Inches(1.5), Inches(12.333), Inches(2.2), bg_c=ROXO, tit_c=TURQUESA, desc_c=BRANCO, tit_sz=20, desc_sz=16, badge="THANKS")
add_bento_tile(sl22, "Espaço para Feedback Rápido:", "[Espaço destinado para o professor coletar dúvidas rápidas e considerações do aluno]", Inches(0.5), Inches(3.82), Inches(12.333), Inches(2.98), bg_c=TURQUESA_CLARO, border_c=TURQUESA, tit_c=AZUL, desc_c=SLATE, tit_sz=18, desc_sz=15, badge="FEEDBACK")

out_path_local = "template-haas-modelo-21-bento-pop.pptx"
prs.save(out_path_local)

server_dir = "/var/www/haas-frontend-desk-mobile-oficial/public/temp-miniaturas/"
os.makedirs(server_dir, exist_ok=True)
out_path_server = os.path.join(server_dir, "template-haas-modelo-21-bento-pop.pptx")

import shutil
shutil.copyfile(out_path_local, out_path_server)
print(f"🎉 MODELO 21 BENTO POP COLOR GRID GERADO EM: {out_path_server}")
