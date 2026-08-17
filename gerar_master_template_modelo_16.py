import os
import urllib.request
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# PALETA OFICIAL HAAS LANGUAGE - SWISS TYPOGRAPHIC STYLE
BRANCO = RGBColor(0xFF, 0xFF, 0xFF)
ROXO = RGBColor(0x4F, 0x35, 0xE2)
AZUL = RGBColor(0x1E, 0x3A, 0x8A)
SLATE = RGBColor(0x33, 0x41, 0x55)
TURQUESA = RGBColor(0x14, 0xC8, 0xB0)
BORDA = RGBColor(0x1E, 0x3A, 0x8A)
FUNDO_SWISS = RGBColor(0xF8, 0xFA, 0xFC)
ROXO_CLARO = RGBColor(0xED, 0xE9, 0xFE)
TURQUESA_CLARO = RGBColor(0xE0, 0xF2, 0xF1)

LOGO_URL = "https://jdppxfokfhqjudwfwckd.supabase.co/storage/v1/object/public/haas-academy/assignments/Design%20sem%20nome%20(4).png"
LOGO_PATH = "/tmp/haas-logo.png"

if not os.path.exists(LOGO_PATH):
    try:
        urllib.request.urlretrieve(LOGO_URL, LOGO_PATH)
    except Exception:
        pass

W, H = Inches(13.333), Inches(7.5)
TOTAL_SLIDES = 22

prs = Presentation()
prs.slide_width = W
prs.slide_height = H

def base_slide(num, hide_logo=False):
    sl = prs.slides.add_slide(prs.slide_layouts[6])
    sl.background.fill.solid()
    sl.background.fill.fore_color.rgb = BRANCO

    # EIXO VERTICAL SUÍÇO
    grid_axis = sl.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.8), Inches(0), Pt(2), H)
    grid_axis.fill.solid(); grid_axis.fill.fore_color.rgb = ROXO; grid_axis.line.fill.background()

    num_box = sl.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0.3), Inches(0.8), Inches(0.5))
    num_box.fill.solid(); num_box.fill.fore_color.rgb = ROXO; num_box.line.fill.background()
    tf_nb = num_box.text_frame; tf_nb.vertical_anchor = MSO_ANCHOR.MIDDLE
    p_nb = tf_nb.paragraphs[0]; p_nb.text = f"{num:02d}"
    p_nb.font.size = Pt(12); p_nb.font.color.rgb = BRANCO; p_nb.font.bold = True; p_nb.alignment = PP_ALIGN.CENTER

    if os.path.exists(LOGO_PATH) and not hide_logo:
        sl.shapes.add_picture(LOGO_PATH, Inches(10.8), Inches(0.2), height=Inches(0.65))
        
    tb = sl.shapes.add_textbox(Inches(1.0), Inches(7.1), Inches(8), Inches(0.35))
    p = tb.text_frame.paragraphs[0]
    p.text = "HAAS LANGUAGE  •  INTERNATIONAL TYPOGRAPHIC STYLE"
    p.font.size = Pt(10); p.font.color.rgb = AZUL; p.font.bold = True
    
    tb2 = sl.shapes.add_textbox(Inches(10.8), Inches(7.1), Inches(2.0), Inches(0.35))
    p2 = tb2.text_frame.paragraphs[0]
    p2.text = f"{num} / {TOTAL_SLIDES}"
    p2.font.size = Pt(10); p2.font.color.rgb = ROXO; p2.font.bold = True; p2.alignment = PP_ALIGN.RIGHT
    return sl

def add_header(sl, tit, sub=""):
    tb = sl.shapes.add_textbox(Inches(1.0), Inches(0.3), Inches(9.5), Inches(0.65))
    tf = tb.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]; p.text = tit.upper()
    p.font.size = Pt(28); p.font.color.rgb = ROXO; p.font.bold = True
    
    if sub:
        tb_s = sl.shapes.add_textbox(Inches(1.0), Inches(0.95), Inches(9.5), Inches(0.4))
        p_s = tb_s.text_frame.paragraphs[0]
        p_s.text = sub
        p_s.font.size = Pt(16); p_s.font.color.rgb = AZUL; p_s.font.bold = True

def add_swiss_card(sl, tit, desc, left, top, width, height, bg_color=FUNDO_SWISS, top_border_color=ROXO, tit_color=ROXO, desc_color=SLATE, tit_size=16, desc_size=14, align_center_v=False):
    s = sl.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    s.fill.solid(); s.fill.fore_color.rgb = bg_color
    s.line.color.rgb = BORDA; s.line.width = Pt(1)

    if top_border_color:
        tb_line = sl.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, Inches(0.08))
        tb_line.fill.solid(); tb_line.fill.fore_color.rgb = top_border_color; tb_line.line.fill.background()

    tf = s.text_frame; tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE if align_center_v else MSO_ANCHOR.TOP
    tf.margin_left = Inches(0.25); tf.margin_top = Inches(0.2); tf.margin_right = Inches(0.25); tf.margin_bottom = Inches(0.2)
    
    if tit:
        p = tf.paragraphs[0]
        p.text = tit; p.font.size = Pt(tit_size); p.font.color.rgb = tit_color; p.font.bold = True
        if desc:
            p2 = tf.add_paragraph()
            p2.text = desc; p2.font.size = Pt(desc_size); p2.font.color.rgb = desc_color; p2.space_before = Pt(6)
    else:
        if desc:
            p = tf.paragraphs[0]
            p.text = desc; p.font.size = Pt(desc_size); p.font.color.rgb = desc_color

def img_swiss_box(sl, left, top, width, height, label="[IMAGEM SUÍÇA]"):
    s = sl.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    s.fill.solid(); s.fill.fore_color.rgb = ROXO_CLARO
    s.line.color.rgb = ROXO; s.line.width = Pt(2)
    tf = s.text_frame; tf.word_wrap = True; tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER
    p = tf.paragraphs[0]; p.text = label; p.font.size = Pt(16); p.font.color.rgb = ROXO; p.font.bold = True

def add_swiss_flashcard(sl, num_str, desc, left, top, width, height, top_c=ROXO):
    s = sl.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    s.fill.solid(); s.fill.fore_color.rgb = BRANCO; s.line.color.rgb = BORDA; s.line.width = Pt(1)
    
    tb_line = sl.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, Inches(0.06))
    tb_line.fill.solid(); tb_line.fill.fore_color.rgb = top_c; tb_line.line.fill.background()

    num_tag = sl.shapes.add_shape(MSO_SHAPE.RECTANGLE, left + Inches(0.15), top + Inches(0.15), Inches(0.5), Inches(0.35))
    num_tag.fill.solid(); num_tag.fill.fore_color.rgb = top_c; num_tag.line.fill.background()
    tf_tag = num_tag.text_frame; tf_tag.vertical_anchor = MSO_ANCHOR.MIDDLE
    p_tag = tf_tag.paragraphs[0]; p_tag.text = num_str; p_tag.alignment = PP_ALIGN.CENTER
    p_tag.font.size = Pt(11); p_tag.font.color.rgb = BRANCO; p_tag.font.bold = True

    tf = s.text_frame; tf.word_wrap = True; tf.vertical_anchor = MSO_ANCHOR.TOP
    tf.margin_left = Inches(0.75); tf.margin_top = Inches(0.15); tf.margin_right = Inches(0.15); tf.margin_bottom = Inches(0.1)
    p = tf.paragraphs[0]; p.text = desc; p.font.size = Pt(14); p.font.color.rgb = SLATE; p.font.bold = False

# SLIDE 1
sl1 = base_slide(1, hide_logo=True)
tb_c = sl1.shapes.add_textbox(Inches(1.0), Inches(1.5), Inches(6.5), Inches(2.2))
tf_c = tb_c.text_frame; tf_c.word_wrap = True
p = tf_c.paragraphs[0]; p.text = "[Título da Aula / Tema Principal]"
p.font.size = Pt(40); p.font.color.rgb = ROXO; p.font.bold = True

tb_sub = sl1.shapes.add_textbox(Inches(1.0), Inches(3.8), Inches(6.5), Inches(1.2))
tf_sub = tb_sub.text_frame; tf_sub.word_wrap = True
p_sub = tf_sub.paragraphs[0]; p_sub.text = "[Subtítulo / Descrição da Aula]"
p_sub.font.size = Pt(22); p_sub.font.color.rgb = AZUL; p_sub.font.bold = True

tb_date = sl1.shapes.add_textbox(Inches(1.0), Inches(5.5), Inches(6.5), Inches(0.6))
tf_date = tb_date.text_frame; tf_date.word_wrap = True
p_date = tf_date.paragraphs[0]; p_date.text = "📅 Data da Aula: [DD/MM/AAAA]"
p_date.font.size = Pt(16); p_date.font.color.rgb = TURQUESA; p_date.font.bold = True

img_swiss_box(sl1, Inches(7.8), Inches(0.8), Inches(5.0), Inches(5.9), "[IMAGEM ASMETRICA SUÍÇA]")

# SLIDE 2
sl2 = base_slide(2)
add_header(sl2, "Presença & Integrantes", "Grade de Registro")
add_swiss_card(sl2, "PROFESSOR(A):", "[Nome do(a) Professor(a)]", Inches(1.0), Inches(1.5), Inches(5.8), Inches(5.2), top_border_color=ROXO, tit_size=18, desc_size=16)
add_swiss_card(sl2, "ALUNO(S):", "[Nome do Aluno / Lista da Turma]", Inches(7.1), Inches(1.5), Inches(5.7), Inches(5.2), top_border_color=TURQUESA, tit_size=18, desc_size=16)

# SLIDE 3
sl3 = base_slide(3)
add_header(sl3, "Aquecimento / Warm-up", "Foco 100% na Fala (3 Minutos)")
add_swiss_card(sl3, "PERGUNTA PROVOCATIVA RÁPIDA:", "[Insira aqui a pergunta provocativa de aquecimento]", Inches(1.0), Inches(1.5), Inches(7.0), Inches(2.8), top_border_color=ROXO, tit_size=18, desc_size=16)
add_swiss_card(sl3, "FUNDAMENTAÇÃO NEUROCOGNITIVA:", "Ativa o córtex pré-frontal e reduz a inibição linguística inicial.", Inches(1.0), Inches(4.5), Inches(7.0), Inches(2.2), top_border_color=TURQUESA, tit_size=15, desc_size=14)
img_swiss_box(sl3, Inches(8.3), Inches(1.5), Inches(4.5), Inches(5.2), "[IMAGEM SUÍÇA]")

# SLIDE 4
sl4 = base_slide(4)
add_header(sl4, "Objetivos da Aula", "Metas e Foco de Aprendizado")
add_swiss_card(sl4, "METAS DA AULA:", "• [Objetivo 1]\n• [Objetivo 2]\n• [Objetivo 3]", Inches(1.0), Inches(1.5), Inches(7.5), Inches(5.2), top_border_color=ROXO, tit_size=18, desc_size=16)
add_swiss_card(sl4, "NOTA DE ESCRITA:", "Nesta aula daremos atenção à escrita para reforçar a memorização de estruturas/vocabulário e reduzir a interferência da língua materna.", Inches(8.8), Inches(1.5), Inches(4.0), Inches(5.2), top_border_color=AZUL, tit_size=16, desc_size=14)

# SLIDE 5
sl5 = base_slide(5)
add_header(sl5, "A Nossa Metodologia", "Estrutura Fixa Haas Language")
met_left = [
    ("1. FALA ATIVA", "Conversação inicial sem julgamentos.", ROXO),
    ("2. EXPOSIÇÃO", "Apresentação leve da gramática.", TURQUESA),
    ("3. PRÁTICA GUIADA", "Exercícios com suporte do professor.", AZUL)
]
top_y = 1.5
for t, d, c in met_left:
    add_swiss_card(sl5, t, d, Inches(1.0), Inches(top_y), Inches(6.0), Inches(1.6), top_border_color=c, tit_size=15, desc_size=13)
    top_y += 1.8

add_swiss_card(sl5, "4. CONVERSAÇÃO", "Aplicação prática das estruturas em contexto real.", Inches(7.3), Inches(1.5), Inches(5.5), Inches(2.5), top_border_color=ROXO, tit_size=18, desc_size=15)
add_swiss_card(sl5, "5. ESCRITA", "Consolidação visual e redução da tradução mental automática.", Inches(7.3), Inches(4.2), Inches(5.5), Inches(2.5), top_border_color=TURQUESA, tit_size=18, desc_size=15)

# SLIDES 6 A 8
for s_num in [6, 7, 8]:
    sl = base_slide(s_num)
    add_header(sl, f"Explicação do Tema — Parte {s_num - 5}")
    add_swiss_card(sl, "CONCEITO PRINCIPAL:", f"[Explicação e exemplo central da estrutura - Bloco {s_num - 5}]", Inches(1.0), Inches(1.5), Inches(7.5), Inches(5.2), top_border_color=ROXO, tit_size=20, desc_size=15)
    add_swiss_card(sl, "EXEMPLOS PRÁTICOS:", "[Exemplos adicionais e frases do dia a dia]", Inches(8.8), Inches(1.5), Inches(4.0), Inches(2.5), top_border_color=TURQUESA, tit_size=16, desc_size=13)
    add_swiss_card(sl, "NOTAS & DETALHES:", "[Exceções, nuances gramaticais e contexto]", Inches(8.8), Inches(4.2), Inches(4.0), Inches(2.5), top_border_color=AZUL, tit_size=16, desc_size=13)

# SLIDE 9
sl9 = base_slide(9)
add_header(sl9, "Síntese do Conteúdo & Sugestão Cultural")
add_swiss_card(sl9, "RESUMO DO TEMA:", "[Resumo sintético dos pontos chave ensinados na aula]", Inches(1.0), Inches(1.5), Inches(5.7), Inches(5.2), top_border_color=ROXO, tit_size=18, desc_size=15)
add_swiss_card(sl9, "SUGESTÃO CULTURAL:", "🎵 [Nome da Música / Vídeo recomendado sobre o assunto].\n\nO professor pode rodar o áudio/vídeo ou deixar como recomendação extra para o aluno praticar.", Inches(7.0), Inches(1.5), Inches(5.8), Inches(5.2), top_border_color=TURQUESA, tit_size=18, desc_size=15)

# SLIDES 10 E 11
sl10 = base_slide(10)
add_header(sl10, "Prática Guiada — Blocos 1 a 5", "Correção de Vícios")
exs_1_5 = [
    ("01", "[Lacunas contextualizadas com opções]"),
    ("02", "[Lacunas contextualizadas com opções]"),
    ("03", "[Preenchimento direto de verbos]"),
    ("04", "[Reestruturação / Transformação de frase]"),
    ("05", "[Reestruturação / Transformação de frase]")
]
for idx, (num_s, d_e) in enumerate(exs_1_5):
    col = idx % 2; row = idx // 2
    l = 1.0 if col == 0 else 7.0
    w = 5.7 if (idx < 4 or len(exs_1_5) % 2 == 0) else 11.7
    if idx == 4: l = 1.0
    t = 1.5 + row * 1.7
    add_swiss_flashcard(sl10, num_s, d_e, Inches(l), Inches(t), Inches(w), Inches(1.5), top_c=ROXO if idx%2==0 else TURQUESA)

sl11 = base_slide(11)
add_header(sl11, "Prática Guiada — Blocos 6 a 10", "Desafios Práticos")
exs_6_10 = [
    ("06", "[Preposição / Conector correto]"),
    ("07", "[Escolha da opção correta]"),
    ("08", "[Substituição por sinônimo formal]"),
    ("09", "[Identificação de erro e correção]"),
    ("10", "[Tradução de sentido sem uso literal]")
]
for idx, (num_s, d_e) in enumerate(exs_6_10):
    col = idx % 2; row = idx // 2
    l = 1.0 if col == 0 else 7.0
    w = 5.7 if (idx < 4 or len(exs_6_10) % 2 == 0) else 11.7
    if idx == 4: l = 1.0
    t = 1.5 + row * 1.7
    add_swiss_flashcard(sl11, num_s, d_e, Inches(l), Inches(t), Inches(w), Inches(1.5), top_c=TURQUESA if idx%2==0 else ROXO)

# SLIDES 12 A 19: CONVERSAÇÃO (ESCADARIA SUÍÇA NA DIREITA)
perguntas_config = [
    "[Pergunta principal 1 sobre o tema]", "[Pergunta principal 2 sobre o tema]", "[Pergunta principal 3 sobre o tema]",
    "[Pergunta principal 4 sobre o tema]", "[Pergunta principal 5 sobre o tema]", "[Pergunta principal 6 sobre o tema]",
    "[Pergunta principal 7 sobre o tema]", "[Pergunta principal 8 sobre o tema]"
]

def add_swiss_vertical_staircase_right(sl, history_list):
    count = len(history_list)
    if count == 0: return
    gap = 0.1
    avail_h = 5.2 - (gap * (count - 1))
    card_h = avail_h / count
    
    for i, h_txt in enumerate(history_list):
        y_pos = 1.5 + i * (card_h + gap)
        clean_t = h_txt[:35] + "..." if len(h_txt) > 35 else h_txt
        d_sz = 11 if count <= 3 else (10 if count <= 5 else 9)
        bg_col = ROXO_CLARO if i % 2 == 0 else TURQUESA_CLARO
        top_c = ROXO if i % 2 == 0 else TURQUESA
        add_swiss_card(sl, f"{i+1:02d} // HISTÓRICO", clean_t, Inches(8.3), Inches(y_pos), Inches(4.5), Inches(card_h), bg_color=bg_col, top_border_color=top_c, tit_size=10, desc_size=d_sz)

for idx in range(8):
    s_num = 12 + idx
    sl = base_slide(s_num)
    add_header(sl, "Conversação", f"Thread {idx + 1} / 8")
    p_desc = perguntas_config[idx]
    
    if idx == 0:
        add_swiss_card(sl, f"PERGUNTA ATIVA {idx+1:02d}:", p_desc, Inches(1.0), Inches(1.5), Inches(11.8), Inches(2.7), bg_color=ROXO_CLARO, top_border_color=ROXO, tit_size=18, desc_size=20, align_center_v=True)
        add_swiss_card(sl, "NOTAS DE APOIO AO PROFESSOR:", "• Para expandir: 'Dê um exemplo recente...'\n• Para concordar/discordar: 'Qual é o outro lado dessa moeda?'", Inches(1.0), Inches(4.4), Inches(11.8), Inches(2.3), bg_color=FUNDO_SWISS, top_border_color=TURQUESA, tit_size=16, desc_size=14, align_center_v=True)
    else:
        add_swiss_card(sl, f"PERGUNTA ATIVA {idx+1:02d}:", p_desc, Inches(1.0), Inches(1.5), Inches(7.0), Inches(2.7), bg_color=ROXO_CLARO, top_border_color=ROXO, tit_size=18, desc_size=18, align_center_v=True)
        add_swiss_card(sl, "NOTAS DE APOIO:", "• Para expandir: 'Dê um exemplo recente...'\n• Para discordar: 'Qual é o outro lado?'", Inches(1.0), Inches(4.4), Inches(7.0), Inches(2.3), bg_color=FUNDO_SWISS, top_border_color=TURQUESA, tit_size=15, desc_size=14, align_center_v=True)
        add_swiss_vertical_staircase_right(sl, [perguntas_config[i] for i in range(idx)])

# SLIDE 20
sl20 = base_slide(20)
add_header(sl20, "Atividade de Escrita", "Consolidação e Memorização")
add_swiss_card(sl20, "TAREFA DE ESCRITA RECOMENDADA:", "[Instrução detalhada da tarefa: escrever um texto, frases ou diálogo]", Inches(1.0), Inches(1.5), Inches(11.8), Inches(2.8), top_border_color=ROXO, tit_size=18, desc_size=15)
add_swiss_card(sl20, "FUNDAMENTAÇÃO PEDAGÓGICA:", "Identificamos que a escrita nesta etapa é fundamental para você fixar as estruturas aprendidas, expandir vocabulário ativo e diminuir a tradução mental automática do seu idioma nativo.", Inches(1.0), Inches(4.5), Inches(11.8), Inches(2.2), top_border_color=TURQUESA, desc_size=14, align_center_v=True)

# SLIDE 21
sl21 = base_slide(21)
add_header(sl21, "Encerramento — Tarefa de Casa", "Ação e Prática Contínua")
add_swiss_card(sl21, "TAREFA DE CASA:", "• [Atividade prática 1 para realizar fora da aula]\n• [Atividade de fixação de vocabulário ou escuta]", Inches(1.0), Inches(1.5), Inches(11.8), Inches(5.2), top_border_color=ROXO, tit_size=18, desc_size=15)

# SLIDE 22
sl22 = base_slide(22)
add_header(sl22, "Encerramento & Feedback Rápido", "Considerações Finais")
add_swiss_card(sl22, "AGRADECIMENTO:", "Obrigado pela dedicação na aula de hoje!", Inches(1.0), Inches(1.5), Inches(11.8), Inches(2.0), top_border_color=TURQUESA, tit_size=20, desc_size=16)
add_swiss_card(sl22, "ESPAÇO PARA FEEDBACK RÁPIDO:", "[Espaço destinado para o professor coletar dúvidas rápidas e considerações do aluno]", Inches(1.0), Inches(3.8), Inches(11.8), Inches(2.9), top_border_color=ROXO, tit_size=18, desc_size=15)

out_path_server = "/var/www/haas-frontend-desk-mobile-oficial/public/temp-miniaturas/template-haas-modelo-16-swiss.pptx"
prs.save(out_path_server)
prs.save("template-haas-modelo-16-swiss.pptx")
print(f"🎉 MODELO 16 SWISS TYPOGRAPHIC GERADO EM: {out_path_server}")
