import os
import shutil
import urllib.request
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml import parse_xml

# PALETA OFICIAL HAAS LANGUAGE - OVERPRINT PRINTING EDITION
BRANCO = RGBColor(0xFF, 0xFF, 0xFF)
ROXO = RGBColor(0x4F, 0x35, 0xE2)
AZUL = RGBColor(0x1E, 0x3A, 0x8A)
SLATE = RGBColor(0x33, 0x41, 0x55)
TURQUESA = RGBColor(0x14, 0xC8, 0xB0)
ROXO_CLARO = RGBColor(0xED, 0xE9, 0xFE)
TURQUESA_CLARO = RGBColor(0xE0, 0xF2, 0xF1)
FUNDO_BG = RGBColor(0xF8, 0xFA, 0xFC)

LOGO_URL = "https://jdppxfokfhqjudwfwckd.supabase.co/storage/v1/object/public/haas-academy/assignments/Design%20sem%20nome%20(4).png"
LOGO_PATH = "/tmp/haas-logo.png"

if not os.path.exists(LOGO_PATH):
    try:
        urllib.request.urlretrieve(LOGO_URL, LOGO_PATH)
    except Exception:
        pass

W, H = Inches(13.333), Inches(7.5)
TOTAL_SLIDES = 22

prs = Presentation()
prs.slide_width = W
prs.slide_height = H

def apply_translucent_fill(shape, hex_color="4F35E2", alpha_pct=40):
    spPr = shape.element.spPr
    solidFill = spPr.find('{http://schemas.openxmlformats.org/drawingml/2006/main}solidFill')
    if solidFill is not None:
        spPr.remove(solidFill)
    
    alpha_val = int(alpha_pct * 1000)
    xml = f"""
    <a:solidFill xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main">
      <a:srgbClr val="{hex_color}">
        <a:alpha val="{alpha_val}"/>
      </a:srgbClr>
    </a:solidFill>
    """
    spPr.append(parse_xml(xml))

def base_slide(num, hide_logo=False):
    sl = prs.slides.add_slide(prs.slide_layouts[6])
    sl.background.fill.solid()
    sl.background.fill.fore_color.rgb = FUNDO_BG

    if os.path.exists(LOGO_PATH) and not hide_logo:
        sl.shapes.add_picture(LOGO_PATH, Inches(10.8), Inches(0.2), height=Inches(0.65))
        
    tb = sl.shapes.add_textbox(Inches(0.5), Inches(7.12), Inches(8), Inches(0.35))
    p = tb.text_frame.paragraphs[0]
    p.text = "HAAS LANGUAGE  •  OVERPRINT PRINTING EDITION 🖨️"
    p.font.size = Pt(10); p.font.color.rgb = ROXO; p.font.bold = True
    
    tb2 = sl.shapes.add_textbox(Inches(10.8), Inches(7.12), Inches(2.0), Inches(0.35))
    p2 = tb2.text_frame.paragraphs[0]
    p2.text = f"{num:02d} / {TOTAL_SLIDES:02d}"
    p2.font.size = Pt(10); p2.font.color.rgb = AZUL; p2.font.bold = True; p2.alignment = PP_ALIGN.RIGHT
    return sl

def add_header(sl, tit, sub=""):
    tb = sl.shapes.add_textbox(Inches(0.5), Inches(0.25), Inches(10), Inches(0.65))
    tf = tb.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]; p.text = tit
    p.font.size = Pt(30); p.font.color.rgb = ROXO; p.font.bold = True
    
    if sub:
        tb_s = sl.shapes.add_textbox(Inches(0.5), Inches(0.95), Inches(10), Inches(0.4))
        p_s = tb_s.text_frame.paragraphs[0]
        p_s.text = sub
        p_s.font.size = Pt(16); p_s.font.color.rgb = AZUL; p_s.font.bold = True

def add_overprint_card(sl, tit, desc, left, top, width, height, bg_hex="4F35E2", alpha=30, border_hex="14C8B0", tit_c=ROXO, desc_c=SLATE, tit_sz=16, desc_sz=14, watermark_text="", align_center_v=False):
    base = sl.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(left), Inches(top), Inches(width), Inches(height))
    apply_translucent_fill(base, hex_color=bg_hex, alpha_pct=alpha)
    base.line.color.rgb = TURQUESA if border_hex=="14C8B0" else ROXO
    base.line.width = Pt(1.5)

    if watermark_text:
        wm_box = sl.shapes.add_textbox(Inches(left + 0.1), Inches(top + 0.05), Inches(width - 0.2), Inches(height - 0.1))
        tf_wm = wm_box.text_frame; tf_wm.word_wrap = True; tf_wm.vertical_anchor = MSO_ANCHOR.TOP
        p_wm = tf_wm.paragraphs[0]
        p_wm.text = watermark_text.upper()
        p_wm.font.size = Pt(36)
        p_wm.font.bold = True
        p_wm.font.color.rgb = ROXO_CLARO if bg_hex=="4F35E2" else TURQUESA_CLARO

    tf = base.text_frame; tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE if align_center_v else MSO_ANCHOR.TOP
    tf.margin_left = Inches(0.25); tf.margin_top = Inches(0.2); tf.margin_right = Inches(0.25); tf.margin_bottom = Inches(0.2)
    
    if tit:
        p = tf.paragraphs[0]
        p.text = tit; p.font.size = Pt(tit_sz); p.font.color.rgb = tit_c; p.font.bold = True
        if desc:
            p2 = tf.add_paragraph()
            p2.text = desc; p2.font.size = Pt(desc_sz); p2.font.color.rgb = desc_c; p2.space_before = Pt(6)
    else:
        if desc:
            p = tf.paragraphs[0]
            p.text = desc; p.font.size = Pt(desc_sz); p.font.color.rgb = desc_c

def img_overprint_box(sl, left, top, width, height, label="[POSTER OVERPRINT]"):
    c1 = sl.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(left - 0.15), Inches(top + 0.15), Inches(width), Inches(height))
    apply_translucent_fill(c1, "14C8B0", 40)
    c1.line.fill.background()

    s = sl.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(left), Inches(top), Inches(width), Inches(height))
    apply_translucent_fill(s, "4F35E2", 30)
    s.line.color.rgb = ROXO; s.line.width = Pt(2.0)
    tf = s.text_frame; tf.word_wrap = True; tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER
    p = tf.paragraphs[0]; p.text = label; p.font.size = Pt(16); p.font.color.rgb = ROXO; p.font.bold = True

def add_overprint_flashcard(sl, num_str, desc, left, top, width, height, bg_hex="4F35E2"):
    s = sl.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(left), Inches(top), Inches(width), Inches(height))
    s.fill.solid(); s.fill.fore_color.rgb = BRANCO
    s.line.color.rgb = ROXO if bg_hex=="4F35E2" else TURQUESA; s.line.width = Pt(1.5)

    tag = sl.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(left - 0.1), Inches(top + 0.12), Inches(0.75), Inches(0.35))
    apply_translucent_fill(tag, bg_hex, 80)
    tag.line.fill.background()
    tf_tag = tag.text_frame; tf_tag.vertical_anchor = MSO_ANCHOR.MIDDLE
    p_tag = tf_tag.paragraphs[0]; p_tag.text = f"#{num_str}"; p_tag.alignment = PP_ALIGN.CENTER
    p_tag.font.size = Pt(11); p_tag.font.color.rgb = BRANCO; p_tag.font.bold = True

    tf = s.text_frame; tf.word_wrap = True; tf.vertical_anchor = MSO_ANCHOR.TOP
    tf.margin_left = Inches(0.85); tf.margin_top = Inches(0.12); tf.margin_right = Inches(0.15); tf.margin_bottom = Inches(0.1)
    p = tf.paragraphs[0]; p.text = desc; p.font.size = Pt(14); p.font.color.rgb = SLATE; p.font.bold = False

# SLIDE 1
sl1 = base_slide(1, hide_logo=True)
add_overprint_card(sl1, "[Título da Aula / Tema Principal]", "[Subtítulo / Descrição da Aula]\n\n📅 Data da Aula: [DD/MM/AAAA]", 0.5, 1.4, 7.2, 5.1, bg_hex="4F35E2", alpha=25, border_hex="14C8B0", tit_c=ROXO, desc_c=SLATE, tit_sz=36, desc_sz=18, watermark_text="OVERPRINT")
img_overprint_box(sl1, 8.1, 1.4, 4.6, 5.1, "[Poster Overprint]")

# SLIDE 2
sl2 = base_slide(2)
add_header(sl2, "Presença & Integrantes", "Sobreposição de Registro")
add_overprint_card(sl2, "Professor(a):", "[Nome do(a) Professor(a)]", 0.5, 1.5, 12.0, 2.3, bg_hex="4F35E2", alpha=20, border_hex="4F35E2", tit_sz=20, desc_sz=16, watermark_text="DOCENTE")
add_overprint_card(sl2, "Aluno(s):", "[Nome do Aluno / Lista da Turma]", 0.5, 4.1, 12.0, 2.5, bg_hex="14C8B0", alpha=25, border_hex="14C8B0", tit_sz=20, desc_sz=16, watermark_text="STUDENTS")

# SLIDE 3
sl3 = base_slide(3)
add_header(sl3, "Aquecimento / Warm-up", "Foco 100% na Fala (3 Minutos)")
add_overprint_card(sl3, "Pergunta Provocativa Rápida:", "[Insira aqui a pergunta provocativa de aquecimento]", 0.5, 1.5, 7.0, 2.5, bg_hex="4F35E2", alpha=20, tit_sz=20, desc_sz=16, watermark_text="SPEAK")
add_overprint_card(sl3, "Fundamentação Neurocognitiva:", "Ativa o córtex pré-frontal e reduz a inibição linguística inicial.", 0.5, 4.3, 7.0, 2.3, bg_hex="14C8B0", alpha=30, tit_sz=18, desc_sz=15, align_center_v=True, watermark_text="NEURO")
img_overprint_box(sl3, 7.9, 1.5, 4.8, 5.1, "[Imagem Overprint]")

# SLIDE 4
sl4 = base_slide(4)
add_header(sl4, "Objetivos da Aula", "Sobreposição de Metas")
add_overprint_card(sl4, "Goal #01", "[Objetivo 1]", 0.5, 1.5, 7.0, 1.5, bg_hex="4F35E2", alpha=20, align_center_v=True)
add_overprint_card(sl4, "Goal #02", "[Objetivo 2]", 0.5, 3.2, 7.0, 1.5, bg_hex="14C8B0", alpha=25, align_center_v=True)
add_overprint_card(sl4, "Goal #03", "[Objetivo 3]", 0.5, 4.9, 7.0, 1.7, bg_hex="1E3A8A", alpha=20, align_center_v=True)
add_overprint_card(sl4, "Nota de Escrita:", "Nesta aula daremos atenção à escrita para reforçar a memorização de estruturas/vocabulário e reduzir a interferência da língua materna.", 7.9, 1.5, 4.8, 5.1, bg_hex="4F35E2", alpha=25, tit_sz=18, desc_sz=15, watermark_text="WRITE")

# SLIDE 5
sl5 = base_slide(5)
add_header(sl5, "A Nossa Metodologia", "Estrutura Haas Language")
add_overprint_card(sl5, "1. Fala Ativa", "Conversação inicial sem julgamentos.", 0.5, 1.5, 5.8, 2.4, bg_hex="4F35E2", alpha=25, watermark_text="PASSO 01")
add_overprint_card(sl5, "2. Exposição", "Apresentação leve da gramática.", 6.7, 1.5, 6.0, 2.4, bg_hex="14C8B0", alpha=30, watermark_text="PASSO 02")
add_overprint_card(sl5, "3. Prática", "Exercícios com suporte.", 0.5, 4.2, 3.8, 2.5, bg_hex="1E3A8A", alpha=20, watermark_text="PASSO 03")
add_overprint_card(sl5, "4. Conversa", "Aplicação prática real.", 4.6, 4.2, 4.0, 2.5, bg_hex="4F35E2", alpha=25, watermark_text="PASSO 04")
add_overprint_card(sl5, "5. Escrita", "Consolidação visual.", 8.9, 4.2, 3.8, 2.5, bg_hex="14C8B0", alpha=30, watermark_text="PASSO 05")

# SLIDES 6 A 8
for s_num in [6, 7, 8]:
    sl = base_slide(s_num)
    add_header(sl, f"Explicação do Tema — Parte {s_num - 5}")
    add_overprint_card(sl, "Conteúdo Principal:", f"[Explicação e exemplo central da estrutura - Bloco {s_num - 5}]", 0.5, 1.5, 7.0, 5.1, bg_hex="4F35E2", alpha=25, tit_sz=20, desc_sz=15, watermark_text="FOCUS")
    add_overprint_card(sl, "Exemplos Práticos:", "[Exemplos adicionais e frases do dia a dia]", 7.9, 1.5, 4.8, 2.4, bg_hex="14C8B0", alpha=30, watermark_text="EXAMP")
    add_overprint_card(sl, "Notas & Detalhes:", "[Exceções, nuances gramaticais e contexto de uso]", 7.9, 4.2, 4.8, 2.4, bg_hex="1E3A8A", alpha=20, watermark_text="RULES")

# SLIDE 9
sl9 = base_slide(9)
add_header(sl9, "Síntese do Conteúdo & Sugestão Cultural")
add_overprint_card(sl9, "Resumo do Tema:", "[Resumo sintético dos pontos chave ensinados na aula]", 0.5, 1.5, 5.8, 5.1, bg_hex="4F35E2", alpha=25, tit_sz=18, desc_sz=15, watermark_text="SUMMARY")
add_overprint_card(sl9, "Sugestão Cultural:", "🎵 [Nome da Música / Vídeo recomendado sobre o assunto].\n\nO professor pode rodar o áudio/vídeo ou deixar como recomendação extra para o aluno praticar.", 6.7, 1.5, 6.0, 5.1, bg_hex="14C8B0", alpha=30, tit_sz=18, desc_sz=15, watermark_text="MEDIA")

# SLIDES 10 E 11
sl10 = base_slide(10)
add_header(sl10, "Prática Guiada — Blocos 1 a 5", "Correção de Vícios")
add_overprint_flashcard(sl10, "01", "[Lacunas contextualizadas com opções]", 0.5, 1.5, 12.0, 1.2, bg_hex="4F35E2")
add_overprint_flashcard(sl10, "02", "[Lacunas contextualizadas com opções]", 0.5, 2.9, 5.8, 1.8, bg_hex="14C8B0")
add_overprint_flashcard(sl10, "03", "[Preenchimento direto de verbos]", 6.7, 2.9, 5.8, 1.8, bg_hex="14C8B0")
add_overprint_flashcard(sl10, "04", "[Reestruturação / Transformação de frase]", 0.5, 4.9, 5.8, 1.8, bg_hex="4F35E2")
add_overprint_flashcard(sl10, "05", "[Reestruturação / Transformação de frase]", 6.7, 4.9, 5.8, 1.8, bg_hex="4F35E2")

sl11 = base_slide(11)
add_header(sl11, "Prática Guiada — Blocos 6 a 10", "Desafios Práticos")
add_overprint_flashcard(sl11, "06", "[Preposição / Conector correto]", 0.5, 1.5, 12.0, 1.2, bg_hex="14C8B0")
add_overprint_flashcard(sl11, "07", "[Escolha da opção correta]", 0.5, 2.9, 5.8, 1.8, bg_hex="4F35E2")
add_overprint_flashcard(sl11, "08", "[Substituição por sinônimo formal]", 6.7, 2.9, 5.8, 1.8, bg_hex="4F35E2")
add_overprint_flashcard(sl11, "09", "[Identificação de erro e correção]", 0.5, 4.9, 5.8, 1.8, bg_hex="14C8B0")
add_overprint_flashcard(sl11, "10", "[Tradução de sentido sem uso literal]", 6.7, 4.9, 5.8, 1.8, bg_hex="14C8B0")

# SLIDES 12 A 19
perguntas_config = [
    "[Pergunta principal 1 sobre o tema]", "[Pergunta principal 2 sobre o tema]", "[Pergunta principal 3 sobre o tema]",
    "[Pergunta principal 4 sobre o tema]", "[Pergunta principal 5 sobre o tema]", "[Pergunta principal 6 sobre o tema]",
    "[Pergunta principal 7 sobre o tema]", "[Pergunta principal 8 sobre o tema]"
]

def add_overprint_history_accordion(sl, history_list):
    count = len(history_list)
    if count == 0: return
    
    gap = 0.12
    avail_h = 5.1 - (gap * (count - 1))
    card_h = avail_h / count
    bg_hexes = ["4F35E2", "14C8B0", "1E3A8A", "4F35E2", "14C8B0", "1E3A8A", "4F35E2"]
    
    for i, h_txt in enumerate(history_list):
        y_pos = 1.5 + i * (card_h + gap)
        clean_t = h_txt[:30] + "..." if len(h_txt) > 30 else h_txt
        d_sz = 11 if count <= 3 else (10 if count <= 5 else 9)
        hex_c = bg_hexes[i % len(bg_hexes)]
        add_overprint_card(sl, f"Q{i+1:02d}", clean_t, 7.9, y_pos, 4.8, card_h, bg_hex=hex_c, alpha=25, tit_sz=11, desc_sz=d_sz, align_center_v=True, watermark_text=f"Q0{i+1}")

for idx in range(8):
    s_num = 12 + idx
    sl = base_slide(s_num)
    add_header(sl, "Conversação", f"Sessão Ativa [{idx + 1}/8]")
    p_desc = perguntas_config[idx]
    
    if idx == 0:
        add_overprint_card(sl, f"PERGUNTA ATIVA Q0{idx+1}:", p_desc, 0.5, 1.5, 12.0, 2.6, bg_hex="4F35E2", alpha=25, tit_sz=18, desc_sz=18, watermark_text="SPEAKING HERO", align_center_v=True)
        add_overprint_card(sl, "PROMPTS DE APOIO:", "• Para expandir: 'Dê um exemplo recente...'\n• Para concordar/discordar: 'Qual é o outro lado dessa moeda?'", 0.5, 4.3, 12.0, 2.3, bg_hex="14C8B0", alpha=30, tit_sz=16, desc_sz=14, watermark_text="TEACHER SUPPORT", align_center_v=True)
    else:
        add_overprint_card(sl, f"PERGUNTA ATIVA Q0{idx+1}:", p_desc, 0.5, 1.5, 7.0, 2.6, bg_hex="4F35E2", alpha=25, tit_sz=18, desc_sz=16, watermark_text=f"ACTIVE Q0{idx+1}", align_center_v=True)
        add_overprint_card(sl, "PROMPTS DE APOIO:", "• Para expandir: 'Dê um exemplo recente...'\n• Para discordar: 'Qual é o outro lado?'", 0.5, 4.3, 7.0, 2.3, bg_hex="14C8B0", alpha=30, tit_sz=15, desc_sz=14, watermark_text="TEACHER SUPPORT", align_center_v=True)
        
        history_list = [perguntas_config[i] for i in range(idx)]
        add_overprint_history_accordion(sl, history_list)

# SLIDE 20
sl20 = base_slide(20)
add_header(sl20, "Atividade de Escrita", "Consolidação e Memorização")
add_overprint_card(sl20, "Tarefa de Escrita Recomendada:", "[Instrução detalhada da tarefa: escrever um texto, frases ou diálogo]", 0.5, 1.5, 12.0, 2.6, bg_hex="4F35E2", alpha=25, tit_sz=18, desc_sz=15, watermark_text="WRITING TASK")
add_overprint_card(sl20, "Fundamentação Pedagógica:", "Identificamos que a escrita nesta etapa é fundamental para você fixar as estruturas aprendidas, expandir vocabulário ativo e diminuir a tradução mental automática do seu idioma nativo.", 0.5, 4.3, 12.0, 2.3, bg_hex="14C8B0", alpha=30, desc_sz=14, watermark_text="NEURO RATIONALE", align_center_v=True)

# SLIDE 21
sl21 = base_slide(21)
add_header(sl21, "Encerramento — Tarefa de Casa", "Ação e Prática Contínua")
add_overprint_card(sl21, "Tarefa de Casa:", "• [Atividade prática 1 para realizar fora da aula]\n• [Atividade de fixação de vocabulário ou escuta]", 0.5, 1.5, 12.0, 5.1, bg_hex="4F35E2", alpha=25, tit_sz=18, desc_sz=15, watermark_text="HOMEWORK")

# SLIDE 22
sl22 = base_slide(22)
add_header(sl22, "Encerramento & Feedback Rápido", "Considerações Finais")
add_overprint_card(sl22, "Agradecimento:", "Obrigado pela dedicação na aula de hoje!", 0.5, 1.5, 12.0, 2.2, bg_hex="4F35E2", alpha=25, tit_sz=20, desc_sz=16, watermark_text="THANKS")
add_overprint_card(sl22, "Espaço para Feedback Rápido:", "[Espaço destinado para o professor coletar dúvidas rápidas e considerações do aluno]", 0.5, 3.9, 12.0, 2.7, bg_hex="14C8B0", alpha=30, tit_sz=18, desc_sz=15, watermark_text="FEEDBACK")

out_path_local = "template-haas-modelo-23-overprint.pptx"
prs.save(out_path_local)

server_dir = "/var/www/haas-frontend-desk-mobile-oficial/public/temp-miniaturas/"
os.makedirs(server_dir, exist_ok=True)
out_path_server = os.path.join(server_dir, "template-haas-modelo-23-overprint.pptx")

shutil.copyfile(out_path_local, out_path_server)
print(f"SUCCESSFULLY GENERATED OVERPRINT MODEL 23: {out_path_server}")
