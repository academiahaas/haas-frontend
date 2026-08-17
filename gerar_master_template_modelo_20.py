import os
import urllib.request
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml import parse_xml

# PALETA OFICIAL HAAS LANGUAGE - LIQUID GRADIENT EDITION
BRANCO = RGBColor(0xFF, 0xFF, 0xFF)
ROXO = RGBColor(0x4F, 0x35, 0xE2)
AZUL = RGBColor(0x1E, 0x3A, 0x8A)
SLATE = RGBColor(0x33, 0x41, 0x55)
TURQUESA = RGBColor(0x14, 0xC8, 0xB0)
BORDA = RGBColor(0xCB, 0xD5, 0xE1)
ROXO_CLARO = RGBColor(0xED, 0xE9, 0xFE)
TURQUESA_CLARO = RGBColor(0xE0, 0xF2, 0xF1)

LOGO_URL = "https://jdppxfokfhqjudwfwckd.supabase.co/storage/v1/object/public/haas-academy/assignments/Design%20sem%20nome%20(4).png"
LOGO_PATH = "/tmp/haas-logo.png"

if not os.path.exists(LOGO_PATH):
    try:
        urllib.request.urlretrieve(LOGO_URL, LOGO_PATH)
    except Exception:
        pass

W, H = Inches(13.333), Inches(7.5)
TOTAL_SLIDES = 22

prs = Presentation()
prs.slide_width = W
prs.slide_height = H

def apply_liquid_gradient(sl):
    """Aplica o Fundo Gradiente Fluido Liquid/Aurora (Navy -> Purple -> Turquoise)"""
    bg = sl.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), W, H)
    spPr = bg.element.spPr
    grad_xml = """
    <a:gradFill xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main">
      <a:gsLst>
        <a:gs pos="0"><a:srgbClr val="1E3A8A"/></a:gs>
        <a:gs pos="50000"><a:srgbClr val="4F35E2"/></a:gs>
        <a:gs pos="100000"><a:srgbClr val="14C8B0"/></a:gs>
      </a:gsLst>
      <a:lin ang="3150000" scaled="1"/>
    </a:gradFill>
    """
    gradFill = parse_xml(grad_xml)
    solidFill = spPr.find('{http://schemas.openxmlformats.org/drawingml/2006/main}solidFill')
    if solidFill is not None:
        spPr.remove(solidFill)
    spPr.append(gradFill)
    bg.line.fill.background()

def base_slide(num, hide_logo=False):
    sl = prs.slides.add_slide(prs.slide_layouts[6])
    apply_liquid_gradient(sl)

    if os.path.exists(LOGO_PATH) and not hide_logo:
        sl.shapes.add_picture(LOGO_PATH, Inches(10.8), Inches(0.2), height=Inches(0.65))
        
    tb = sl.shapes.add_textbox(Inches(0.8), Inches(7.1), Inches(8), Inches(0.35))
    p = tb.text_frame.paragraphs[0]
    p.text = "HAAS LANGUAGE  •  LIQUID GRADIENT EDITION 🌊"
    p.font.size = Pt(10); p.font.color.rgb = BRANCO; p.font.bold = True
    
    tb2 = sl.shapes.add_textbox(Inches(10.8), Inches(7.1), Inches(2.0), Inches(0.35))
    p2 = tb2.text_frame.paragraphs[0]
    p2.text = f"{num:02d} / {TOTAL_SLIDES:02d}"
    p2.font.size = Pt(10); p2.font.color.rgb = TURQUESA; p2.font.bold = True; p2.alignment = PP_ALIGN.RIGHT
    return sl

def add_header(sl, tit, sub=""):
    tb = sl.shapes.add_textbox(Inches(0.8), Inches(0.25), Inches(9.8), Inches(0.75))
    tf = tb.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]; p.text = tit
    p.font.size = Pt(32); p.font.color.rgb = BRANCO; p.font.bold = True
    
    if sub:
        tb_s = sl.shapes.add_textbox(Inches(0.8), Inches(0.95), Inches(9.8), Inches(0.4))
        p_s = tb_s.text_frame.paragraphs[0]
        p_s.text = sub
        p_s.font.size = Pt(16); p_s.font.color.rgb = TURQUESA; p_s.font.bold = True

def add_crisp_card(sl, tit, desc, left, top, width, height, bg_color=BRANCO, border_color=None, tit_color=ROXO, desc_color=SLATE, tit_size=16, desc_size=14, align_center_v=False):
    s = sl.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    s.fill.solid(); s.fill.fore_color.rgb = bg_color
    if border_color:
        s.line.color.rgb = border_color; s.line.width = Pt(1.5)
    else:
        s.line.fill.background()
        
    tf = s.text_frame; tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE if align_center_v else MSO_ANCHOR.TOP
    tf.margin_left = Inches(0.25); tf.margin_top = Inches(0.2); tf.margin_right = Inches(0.25); tf.margin_bottom = Inches(0.2)
    
    if tit:
        p = tf.paragraphs[0]
        p.text = tit; p.font.size = Pt(tit_size); p.font.color.rgb = tit_color; p.font.bold = True
        if desc:
            p2 = tf.add_paragraph()
            p2.text = desc; p2.font.size = Pt(desc_size); p2.font.color.rgb = desc_color; p2.space_before = Pt(6)
    else:
        if desc:
            p = tf.paragraphs[0]
            p.text = desc; p.font.size = Pt(desc_size); p.font.color.rgb = desc_color

def img_liquid_box(sl, left, top, width, height, label="[IMAGEM LIQUID]"):
    s = sl.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    s.fill.solid(); s.fill.fore_color.rgb = ROXO_CLARO
    s.line.color.rgb = TURQUESA; s.line.width = Pt(2)
    tf = s.text_frame; tf.word_wrap = True; tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER
    p = tf.paragraphs[0]; p.text = label; p.font.size = Pt(16); p.font.color.rgb = ROXO; p.font.bold = True

def add_liquid_flashcard(sl, num_str, desc, left, top, width, height, b_c=ROXO):
    s = sl.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    s.fill.solid(); s.fill.fore_color.rgb = BRANCO; s.line.color.rgb = b_c; s.line.width = Pt(1.5)
    
    chip = sl.shapes.add_shape(MSO_SHAPE.OVAL, left + Inches(0.15), top + (height/2) - Inches(0.25), Inches(0.5), Inches(0.5))
    chip.fill.solid(); chip.fill.fore_color.rgb = b_c; chip.line.fill.background()
    tf_chip = chip.text_frame; tf_chip.vertical_anchor = MSO_ANCHOR.MIDDLE
    p_chip = tf_chip.paragraphs[0]; p_chip.text = num_str; p_chip.alignment = PP_ALIGN.CENTER
    p_chip.font.size = Pt(11); p_chip.font.color.rgb = BRANCO; p_chip.font.bold = True
    
    tf = s.text_frame; tf.word_wrap = True; tf.vertical_anchor = MSO_ANCHOR.TOP
    tf.margin_left = Inches(0.8); tf.margin_top = Inches(0.15); tf.margin_right = Inches(0.15); tf.margin_bottom = Inches(0.12)
    p = tf.paragraphs[0]; p.text = desc; p.font.size = Pt(14); p.font.color.rgb = SLATE; p.font.bold = False

# SLIDE 1
sl1 = base_slide(1, hide_logo=True)
add_crisp_card(sl1, "[Título da Aula / Tema Principal]", "[Subtítulo / Descrição da Aula]\n\n📅 Data da Aula: [DD/MM/AAAA]", Inches(0.8), Inches(1.2), Inches(11.733), Inches(3.2), bg_color=BRANCO, border_color=TURQUESA, tit_size=38, desc_size=20, align_center_v=True)
img_liquid_box(sl1, Inches(0.8), Inches(4.6), Inches(11.733), Inches(2.2), "[Imagem da Capa - Panorama Fluid]")

# SLIDE 2
sl2 = base_slide(2)
add_header(sl2, "Presença & Integrantes", "Sessão de Aprendizado")
add_crisp_card(sl2, "Professor(a):", "[Nome do(a) Professor(a)]", Inches(0.8), Inches(1.5), Inches(7.5), Inches(2.4), bg_color=BRANCO, border_color=ROXO, tit_size=20, desc_size=16)
add_crisp_card(sl2, "Aluno(s):", "[Nome do Aluno / Lista da Turma]", Inches(5.0), Inches(4.1), Inches(7.533), Inches(2.7), bg_color=BRANCO, border_color=TURQUESA, tit_size=20, desc_size=16)

# SLIDE 3
sl3 = base_slide(3)
add_header(sl3, "Aquecimento / Warm-up", "Foco 100% na Fala (3 Minutos)")
add_crisp_card(sl3, "Pergunta Provocativa Rápida:", "[Insira aqui a pergunta provocativa de aquecimento]", Inches(0.8), Inches(1.5), Inches(11.733), Inches(2.2), bg_color=BRANCO, border_color=ROXO, tit_size=20, desc_size=16)
img_liquid_box(sl3, Inches(0.8), Inches(3.9), Inches(5.5), Inches(2.9), "[Imagem Warm-up]")
add_crisp_card(sl3, "Fundamentação Neurocognitiva:", "Ativa o córtex pré-frontal e reduz a inibição linguística inicial.", Inches(6.5), Inches(3.9), Inches(6.033), Inches(2.9), bg_color=TURQUESA_CLARO, border_color=TURQUESA, tit_size=18, desc_size=15, align_center_v=True)

# SLIDE 4
sl4 = base_slide(4)
add_header(sl4, "Objetivos da Aula", "Metas da Sessão")
add_crisp_card(sl4, "Goal 01:", "[Objetivo 1]", Inches(0.8), Inches(1.5), Inches(7.2), Inches(1.5), bg_color=BRANCO, border_color=ROXO, tit_size=14, desc_size=15, align_center_v=True)
add_crisp_card(sl4, "Goal 02:", "[Objetivo 2]", Inches(0.8), Inches(3.2), Inches(7.2), Inches(1.5), bg_color=BRANCO, border_color=ROXO, tit_size=14, desc_size=15, align_center_v=True)
add_crisp_card(sl4, "Goal 03:", "[Objetivo 3]", Inches(0.8), Inches(4.9), Inches(7.2), Inches(1.5), bg_color=BRANCO, border_color=ROXO, tit_size=14, desc_size=15, align_center_v=True)
add_crisp_card(sl4, "Nota de Escrita:", "Nesta aula daremos atenção à escrita para reforçar a memorização de estruturas/vocabulário e reduzir a interferência da língua materna.", Inches(8.3), Inches(1.5), Inches(4.233), Inches(4.9), bg_color=ROXO_CLARO, border_color=ROXO, tit_size=18, desc_size=15)

# SLIDE 5
sl5 = base_slide(5)
add_header(sl5, "A Nossa Metodologia", "Estrutura Haas Language")
add_crisp_card(sl5, "1. Fala Ativa", "Conversação inicial sem julgamentos.", Inches(0.8), Inches(1.5), Inches(3.7), Inches(2.5), bg_color=BRANCO, border_color=ROXO)
add_crisp_card(sl5, "2. Exposição", "Apresentação leve da gramática.", Inches(4.8), Inches(1.5), Inches(3.7), Inches(2.5), bg_color=BRANCO, border_color=TURQUESA)
add_crisp_card(sl5, "3. Prática", "Exercícios com suporte.", Inches(8.8), Inches(1.5), Inches(3.733), Inches(2.5), bg_color=BRANCO, border_color=AZUL)
add_crisp_card(sl5, "4. Conversação", "Aplicação prática das estruturas em contexto real.", Inches(0.8), Inches(4.2), Inches(5.7), Inches(2.6), bg_color=BRANCO, border_color=TURQUESA)
add_crisp_card(sl5, "5. Escrita", "Consolidação visual e redução da tradução automática.", Inches(6.8), Inches(4.2), Inches(5.733), Inches(2.6), bg_color=BRANCO, border_color=ROXO)

# SLIDES 6 A 8
for s_num in [6, 7, 8]:
    sl = base_slide(s_num)
    add_header(sl, f"Explicação do Tema — Parte {s_num - 5}")
    add_crisp_card(sl, "Conteúdo Principal:", f"[Explicação e exemplo central da estrutura - Bloco {s_num - 5}]", Inches(0.8), Inches(1.5), Inches(11.733), Inches(2.5), bg_color=BRANCO, border_color=ROXO, tit_size=20, desc_size=15)
    add_crisp_card(sl, "Exemplos Práticos:", "[Exemplos adicionais e frases do dia a dia]", Inches(0.8), Inches(4.2), Inches(5.7), Inches(2.6), bg_color=BRANCO, border_color=TURQUESA)
    add_crisp_card(sl, "Notas & Detalhes:", "[Exceções, nuances gramaticais e contexto de uso]", Inches(6.8), Inches(4.2), Inches(5.733), Inches(2.6), bg_color=BRANCO, border_color=AZUL)

# SLIDE 9
sl9 = base_slide(9)
add_header(sl9, "Síntese do Conteúdo & Sugestão Cultural")
add_crisp_card(sl9, "Resumo do Tema:", "[Resumo sintético dos pontos chave ensinados na aula]", Inches(0.8), Inches(1.5), Inches(11.733), Inches(2.5), bg_color=BRANCO, border_color=ROXO, tit_size=18, desc_size=15)
add_crisp_card(sl9, "Sugestão Cultural:", "🎵 [Nome da Música / Vídeo recomendado sobre o assunto].\n\nO professor pode rodar o áudio/vídeo ou deixar como recomendação extra para o aluno praticar.", Inches(0.8), Inches(4.2), Inches(11.733), Inches(2.6), bg_color=BRANCO, border_color=TURQUESA, tit_size=18, desc_size=15)

# SLIDES 10 E 11
sl10 = base_slide(10)
add_header(sl10, "Prática Guiada — Blocos 1 a 5", "Correção de Vícios")
add_liquid_flashcard(sl10, "01", "[Lacunas contextualizadas com opções]", Inches(0.8), Inches(1.5), Inches(11.733), Inches(1.2), b_c=ROXO)
add_liquid_flashcard(sl10, "02", "[Lacunas contextualizadas com opções]", Inches(0.8), Inches(2.9), Inches(5.7), Inches(1.8), b_c=TURQUESA)
add_liquid_flashcard(sl10, "03", "[Preenchimento direto de verbos]", Inches(6.8), Inches(2.9), Inches(5.733), Inches(1.8), b_c=TURQUESA)
add_liquid_flashcard(sl10, "04", "[Reestruturação / Transformação de frase]", Inches(0.8), Inches(4.9), Inches(5.7), Inches(1.9), b_c=ROXO)
add_liquid_flashcard(sl10, "05", "[Reestruturação / Transformação de frase]", Inches(6.8), Inches(4.9), Inches(5.733), Inches(1.9), b_c=ROXO)

sl11 = base_slide(11)
add_header(sl11, "Prática Guiada — Blocos 6 a 10", "Desafios Práticos")
add_liquid_flashcard(sl11, "06", "[Preposição / Conector correto]", Inches(0.8), Inches(1.5), Inches(11.733), Inches(1.2), b_c=TURQUESA)
add_liquid_flashcard(sl11, "07", "[Escolha da opção correta]", Inches(0.8), Inches(2.9), Inches(5.7), Inches(1.8), b_c=ROXO)
add_liquid_flashcard(sl11, "08", "[Substituição por sinônimo formal]", Inches(6.8), Inches(2.9), Inches(5.733), Inches(1.8), b_c=ROXO)
add_liquid_flashcard(sl11, "09", "[Identificação de erro e correção]", Inches(0.8), Inches(4.9), Inches(5.7), Inches(1.9), b_c=TURQUESA)
add_liquid_flashcard(sl11, "10", "[Tradução de sentido sem uso literal]", Inches(6.8), Inches(4.9), Inches(5.733), Inches(1.9), b_c=TURQUESA)

# SLIDES 12 A 19: CONVERSAÇÃO (AURORA TIMELINE WAVE)
perguntas_config = [
    "[Pergunta principal 1 sobre o tema]", "[Pergunta principal 2 sobre o tema]", "[Pergunta principal 3 sobre o tema]",
    "[Pergunta principal 4 sobre o tema]", "[Pergunta principal 5 sobre o tema]", "[Pergunta principal 6 sobre o tema]",
    "[Pergunta principal 7 sobre o tema]", "[Pergunta principal 8 sobre o tema]"
]

def add_aurora_top_timeline(sl, history_list):
    count = len(history_list)
    if count == 0: return
    gap = 0.12
    avail_w = 11.733 - (gap * (count - 1))
    pill_w = avail_w / count
    
    for i, h_txt in enumerate(history_list):
        left_x = 0.8 + i * (pill_w + gap)
        clean_t = h_txt[:26] + "..." if len(h_txt) > 26 else h_txt
        d_sz = 11 if count <= 3 else (10 if count <= 5 else 9)
        acc_c = TURQUESA if i % 2 == 0 else ROXO
        add_crisp_card(sl, f"Q{i+1:02d}", clean_t, Inches(left_x), Inches(1.4), Inches(pill_w), Inches(0.95), bg_color=BRANCO, border_color=acc_c, tit_color=acc_c, desc_size=d_sz, align_center_v=True)

for idx in range(8):
    s_num = 12 + idx
    sl = base_slide(s_num)
    add_header(sl, "Conversação", f"Sessão de Fala [{idx + 1}/8]")
    p_desc = perguntas_config[idx]
    
    if idx == 0:
        add_crisp_card(sl, "Pergunta Ativa:", p_desc, Inches(0.8), Inches(1.5), Inches(11.733), Inches(3.2), bg_color=BRANCO, border_color=ROXO, tit_size=18, desc_size=20, align_center_v=True)
        add_crisp_card(sl, "Apoio ao Professor:", "• Para expandir: 'Dê um exemplo recente...'\n• Para concordar/discordar: 'Qual é o outro lado dessa moeda?'", Inches(0.8), Inches(4.9), Inches(11.733), Inches(1.9), bg_color=TURQUESA_CLARO, border_color=TURQUESA, tit_size=16, desc_size=14, align_center_v=True)
    else:
        history_list = [perguntas_config[i] for i in range(idx)]
        add_aurora_top_timeline(sl, history_list)
        
        add_crisp_card(sl, f"Pergunta Ativa {idx+1:02d}:", p_desc, Inches(0.8), Inches(2.5), Inches(11.733), Inches(2.5), bg_color=BRANCO, border_color=ROXO, tit_size=18, desc_size=18, align_center_v=True)
        add_crisp_card(sl, "Apoio ao Professor:", "• Para expandir: 'Dê um exemplo recente...'\n• Para discordar: 'Qual é o outro lado?'", Inches(0.8), Inches(5.2), Inches(11.733), Inches(1.6), bg_color=TURQUESA_CLARO, border_color=TURQUESA, tit_size=15, desc_size=14, align_center_v=True)

# SLIDE 20
sl20 = base_slide(20)
add_header(sl20, "Atividade de Escrita", "Consolidação e Memorização")
add_crisp_card(sl20, "Tarefa de Escrita Recomendada:", "[Instrução detalhada da tarefa: escrever um texto, frases ou diálogo]", Inches(0.8), Inches(1.5), Inches(11.733), Inches(2.8), bg_color=BRANCO, border_color=ROXO, tit_size=18, desc_size=15)
add_crisp_card(sl20, "Fundamentação Pedagógica:", "Identificamos que a escrita nesta etapa é fundamental para você fixar as estruturas aprendidas, expandir vocabulário ativo e diminuir a tradução mental automática do seu idioma nativo.", Inches(0.8), Inches(4.5), Inches(11.733), Inches(2.3), bg_color=TURQUESA_CLARO, border_color=TURQUESA, desc_size=14, align_center_v=True)

# SLIDE 21
sl21 = base_slide(21)
add_header(sl21, "Encerramento — Tarefa de Casa", "Ação e Prática Contínua")
add_crisp_card(sl21, "Tarefa de Casa:", "• [Atividade prática 1 para realizar fora da aula]\n• [Atividade de fixação de vocabulário ou escuta]", Inches(0.8), Inches(1.5), Inches(11.733), Inches(5.2), bg_color=BRANCO, border_color=ROXO, tit_size=18, desc_size=15)

# SLIDE 22
sl22 = base_slide(22)
add_header(sl22, "Encerramento & Feedback Rápido", "Considerações Finais")
add_crisp_card(sl22, "Agradecimento:", "Obrigado pela dedicação na aula de hoje!", Inches(0.8), Inches(1.5), Inches(11.733), Inches(2.0), bg_color=TURQUESA_CLARO, border_color=TURQUESA, tit_size=20, desc_size=16)
add_crisp_card(sl22, "Espaço para Feedback Rápido:", "[Espaço destinado para o professor coletar dúvidas rápidas e considerações do aluno]", Inches(0.8), Inches(3.7), Inches(11.733), Inches(3.0), bg_color=BRANCO, border_color=ROXO, tit_size=18, desc_size=15)

out_path_local = "template-haas-modelo-20-liquid.pptx"
prs.save(out_path_local)

server_dir = "/var/www/haas-frontend-desk-mobile-oficial/public/temp-miniaturas/"
os.makedirs(server_dir, exist_ok=True)
out_path_server = os.path.join(server_dir, "template-haas-modelo-20-liquid.pptx")

import shutil
shutil.copyfile(out_path_local, out_path_server)
print(f"🎉 MODELO 20 LIQUID GRADIENT GERADO EM: {out_path_server}")
