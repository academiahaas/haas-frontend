import os
import urllib.request
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# PALETA OFICIAL HAAS LANGUAGE - SCALE CONTRAST EDITION
BRANCO = RGBColor(0xFF, 0xFF, 0xFF)
ROXO = RGBColor(0x4F, 0x35, 0xE2)
AZUL = RGBColor(0x1E, 0x3A, 0x8A)
SLATE = RGBColor(0x33, 0x41, 0x55)
TURQUESA = RGBColor(0x14, 0xC8, 0xB0)
BORDA = RGBColor(0xCB, 0xD5, 0xE1)
FUNDO_SCALE = RGBColor(0xF8, 0xFA, 0xFC)
ROXO_CLARO = RGBColor(0xED, 0xE9, 0xFE)
TURQUESA_CLARO = RGBColor(0xE0, 0xF2, 0xF1)

LOGO_URL = "https://jdppxfokfhqjudwfwckd.supabase.co/storage/v1/object/public/haas-academy/assignments/Design%20sem%20nome%20(4).png"
LOGO_PATH = "/tmp/haas-logo.png"

if not os.path.exists(LOGO_PATH):
    try:
        urllib.request.urlretrieve(LOGO_URL, LOGO_PATH)
    except Exception:
        pass

W, H = Inches(13.333), Inches(7.5)
TOTAL_SLIDES = 22

prs = Presentation()
prs.slide_width = W
prs.slide_height = H

def base_slide(num, hide_logo=False):
    sl = prs.slides.add_slide(prs.slide_layouts[6])
    sl.background.fill.solid()
    sl.background.fill.fore_color.rgb = FUNDO_SCALE

    axis = sl.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.5), Inches(0.2), Pt(3), Inches(6.8))
    axis.fill.solid(); axis.fill.fore_color.rgb = ROXO; axis.line.fill.background()

    if os.path.exists(LOGO_PATH) and not hide_logo:
        sl.shapes.add_picture(LOGO_PATH, Inches(10.8), Inches(0.2), height=Inches(0.65))
        
    tb = sl.shapes.add_textbox(Inches(0.7), Inches(7.1), Inches(8), Inches(0.35))
    p = tb.text_frame.paragraphs[0]
    p.text = "HAAS LANGUAGE  •  SCALE CONTRAST EDITION"
    p.font.size = Pt(10); p.font.color.rgb = AZUL; p.font.bold = True
    
    tb2 = sl.shapes.add_textbox(Inches(10.8), Inches(7.1), Inches(2.0), Inches(0.35))
    p2 = tb2.text_frame.paragraphs[0]
    p2.text = f"{num:02d} / {TOTAL_SLIDES:02d}"
    p2.font.size = Pt(10); p2.font.color.rgb = ROXO; p2.font.bold = True; p2.alignment = PP_ALIGN.RIGHT
    return sl

def add_header(sl, tit, sub=""):
    tb = sl.shapes.add_textbox(Inches(0.7), Inches(0.25), Inches(9.8), Inches(0.75))
    tf = tb.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]; p.text = tit
    p.font.size = Pt(32); p.font.color.rgb = ROXO; p.font.bold = True
    
    if sub:
        tb_s = sl.shapes.add_textbox(Inches(0.7), Inches(0.98), Inches(9.8), Inches(0.4))
        p_s = tb_s.text_frame.paragraphs[0]
        p_s.text = sub
        p_s.font.size = Pt(15); p_s.font.color.rgb = AZUL; p_s.font.bold = False

def add_scale_card(sl, giant_label, tit, desc, left, top, width, height, bg_color=BRANCO, border_color=BORDA, accent_color=ROXO, tit_size=16, desc_size=14, align_center_v=False):
    s = sl.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    s.fill.solid(); s.fill.fore_color.rgb = bg_color
    if border_color:
        s.line.color.rgb = border_color; s.line.width = Pt(1.5)
    else:
        s.line.fill.background()
        
    tf = s.text_frame; tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE if align_center_v else MSO_ANCHOR.TOP
    tf.margin_left = Inches(0.25); tf.margin_top = Inches(0.2); tf.margin_right = Inches(0.25); tf.margin_bottom = Inches(0.2)
    
    if giant_label:
        p_g = tf.paragraphs[0]
        p_g.text = giant_label
        p_g.font.size = Pt(38); p_g.font.color.rgb = accent_color; p_g.font.bold = True
        
        if tit:
            p = tf.add_paragraph()
            p.text = tit; p.font.size = Pt(tit_size); p.font.color.rgb = ROXO; p.font.bold = True; p.space_before = Pt(2)
    else:
        if tit:
            p = tf.paragraphs[0]
            p.text = tit; p.font.size = Pt(tit_size); p.font.color.rgb = ROXO; p.font.bold = True
            
    if desc:
        p2 = tf.add_paragraph() if (giant_label or tit) else tf.paragraphs[0]
        p2.text = desc; p2.font.size = Pt(desc_size); p2.font.color.rgb = SLATE; p2.font.bold = False
        if giant_label or tit:
            p2.space_before = Pt(4)

def img_scale_box(sl, left, top, width, height, label="[IMAGEM EMPÁTICA]"):
    s = sl.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    s.fill.solid(); s.fill.fore_color.rgb = ROXO_CLARO
    s.line.color.rgb = TURQUESA; s.line.width = Pt(2)
    tf = s.text_frame; tf.word_wrap = True; tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER
    p = tf.paragraphs[0]; p.text = label; p.font.size = Pt(16); p.font.color.rgb = ROXO; p.font.bold = True

def add_scale_flashcard(sl, num_str, desc, left, top, width, height, accent_c=ROXO):
    s = sl.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    s.fill.solid(); s.fill.fore_color.rgb = BRANCO; s.line.color.rgb = BORDA; s.line.width = Pt(1)
    
    num_block = sl.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, Inches(0.9), height)
    num_block.fill.solid(); num_block.fill.fore_color.rgb = accent_c; num_block.line.fill.background()
    tf_nb = num_block.text_frame; tf_nb.vertical_anchor = MSO_ANCHOR.MIDDLE
    p_nb = tf_nb.paragraphs[0]; p_nb.text = num_str; p_nb.alignment = PP_ALIGN.CENTER
    p_nb.font.size = Pt(26); p_nb.font.color.rgb = BRANCO; p_nb.font.bold = True
    
    tf = s.text_frame; tf.word_wrap = True; tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    tf.margin_left = Inches(1.1); tf.margin_top = Inches(0.12); tf.margin_right = Inches(0.2); tf.margin_bottom = Inches(0.12)
    p = tf.paragraphs[0]; p.text = desc; p.font.size = Pt(14); p.font.color.rgb = SLATE; p.font.bold = False

# SLIDE 1: CAPA
sl1 = base_slide(1, hide_logo=True)
tb_c = sl1.shapes.add_textbox(Inches(0.7), Inches(1.2), Inches(7.0), Inches(2.5))
tf_c = tb_c.text_frame; tf_c.word_wrap = True
p = tf_c.paragraphs[0]; p.text = "[TÍTULO]"
p.font.size = Pt(64); p.font.color.rgb = ROXO; p.font.bold = True

tb_full = sl1.shapes.add_textbox(Inches(0.7), Inches(3.6), Inches(7.0), Inches(1.8))
tf_full = tb_full.text_frame; tf_full.word_wrap = True
p_f = tf_full.paragraphs[0]; p_f.text = "[Título Completo / Tema Principal da Aula]"
p_f.font.size = Pt(22); p_f.font.color.rgb = AZUL; p_f.font.bold = True

p_s = tf_full.add_paragraph(); p_s.text = "[Subtítulo / Descrição da Aula]"
p_s.font.size = Pt(16); p_s.font.color.rgb = SLATE; p_s.space_before = Pt(8)

p_d = tf_full.add_paragraph(); p_d.text = "📅 Data da Aula: [DD/MM/AAAA]"
p_d.font.size = Pt(15); p_d.font.color.rgb = TURQUESA; p_d.font.bold = True; p_d.space_before = Pt(12)

img_scale_box(sl1, Inches(8.0), Inches(1.2), Inches(4.8), Inches(5.5), "[IMAGEM EM ESCALA]")

# SLIDE 2: PRESENÇA
sl2 = base_slide(2)
add_header(sl2, "Presença & Integrantes", "Registro da Sessão")
add_scale_card(sl2, "DOCENTE", "Professor(a):", "[Nome do(a) Professor(a)]", Inches(0.7), Inches(1.5), Inches(5.8), Inches(5.3), border_color=ROXO, accent_color=ROXO, desc_size=16)
add_scale_card(sl2, "TURMA", "Aluno(s):", "[Nome do Aluno / Lista da Turma]", Inches(6.8), Inches(1.5), Inches(6.0), Inches(5.3), border_color=TURQUESA, accent_color=TURQUESA, desc_size=16)

# SLIDE 3: WARM-UP
sl3 = base_slide(3)
add_header(sl3, "Aquecimento / Warm-up", "Foco 100% na Fala (3 Minutos)")
add_scale_card(sl3, "SPEAK", "Pergunta Provocativa Rápida:", "[Insira aqui a pergunta provocativa de aquecimento]", Inches(0.7), Inches(1.5), Inches(6.8), Inches(2.6), border_color=ROXO, accent_color=ROXO, desc_size=16)
add_scale_card(sl3, "NEURO", "Fundamentação Neurocognitiva:", "Ativa o córtex pré-frontal e reduz a inibição linguística inicial.", Inches(0.7), Inches(4.3), Inches(6.8), Inches(2.5), border_color=TURQUESA, accent_color=TURQUESA, desc_size=14)
img_scale_box(sl3, Inches(7.8), Inches(1.5), Inches(5.0), Inches(5.3), "[IMAGEM WARM-UP]")

# SLIDE 4: OBJETIVOS
sl4 = base_slide(4)
add_header(sl4, "Objetivos da Aula", "Metas e Foco")
add_scale_card(sl4, "METAS", "O que vamos alcançar hoje:", "• [Objetivo 1]\n• [Objetivo 2]\n• [Objetivo 3]", Inches(0.7), Inches(1.5), Inches(7.2), Inches(5.3), border_color=ROXO, accent_color=ROXO, desc_size=16)
add_scale_card(sl4, "ESCRITA", "Nota de Escrita:", "Nesta aula daremos atenção à escrita para reforçar a memorização de estruturas/vocabulário e reduzir a interferência da língua materna.", Inches(8.2), Inches(1.5), Inches(4.6), Inches(5.3), border_color=TURQUESA, accent_color=TURQUESA, desc_size=14)

# SLIDE 5: METODOLOGIA
sl5 = base_slide(5)
add_header(sl5, "A Nossa Metodologia", "Estrutura Fixa Haas Language")
metodologias = [
    ("01", "Fala Ativa", "Conversação inicial sem julgamentos.", ROXO),
    ("02", "Exposição", "Apresentação leve da gramática.", TURQUESA),
    ("03", "Prática", "Exercícios com suporte.", AZUL),
    ("04", "Conversa", "Aplicação prática real.", ROXO),
    ("05", "Escrita", "Consolidação e memorização.", TURQUESA)
]
left_m = 0.7
for num_g, tit_m, desc_m, acc_c in metodologias:
    add_scale_card(sl5, num_g, tit_m, desc_m, Inches(left_m), Inches(1.5), Inches(2.28), Inches(5.3), border_color=BORDA, accent_color=acc_c, tit_size=15, desc_size=12)
    left_m += 2.45

# SLIDES 6 A 8: EXPOSIÇÃO
for s_num in [6, 7, 8]:
    sl = base_slide(s_num)
    add_header(sl, f"Explicação do Tema — Parte {s_num - 5}")
    add_scale_card(sl, "FOCUS", "Conceito Principal:", f"[Explicação e exemplo central da estrutura - Bloco {s_num - 5}]", Inches(0.7), Inches(1.5), Inches(12.1), Inches(2.6), border_color=ROXO, accent_color=ROXO, tit_size=18, desc_size=16)
    add_scale_card(sl, "EXEMPLOS", "Aplicações Práticas:", "[Exemplos adicionais e frases do dia a dia]", Inches(0.7), Inches(4.3), Inches(5.9), Inches(2.5), border_color=TURQUESA, accent_color=TURQUESA, tit_size=16, desc_size=14)
    add_scale_card(sl, "REGRAS", "Notas & Detalhes:", "[Exceções, nuances gramaticais e contexto]", Inches(6.9), Inches(4.3), Inches(5.9), Inches(2.5), border_color=AZUL, accent_color=AZUL, tit_size=16, desc_size=14)

# SLIDE 9: SÍNTESE + CULTURA
sl9 = base_slide(9)
add_header(sl9, "Síntese do Conteúdo & Sugestão Cultural")
add_scale_card(sl9, "RESUMO", "Pontos Chave:", "[Resumo sintético dos pontos chave ensinados na aula]", Inches(0.7), Inches(1.5), Inches(5.9), Inches(5.3), border_color=ROXO, accent_color=ROXO, tit_size=18, desc_size=15)
add_scale_card(sl9, "CULTURA", "Mídia Recomendada:", "🎵 [Nome da Música / Vídeo recomendado sobre o assunto].\n\nO professor pode rodar o áudio/vídeo ou deixar como recomendação extra para o aluno praticar.", Inches(6.9), Inches(1.5), Inches(5.9), Inches(5.3), border_color=TURQUESA, accent_color=TURQUESA, tit_size=18, desc_size=15)

# SLIDES 10 E 11: PRÁTICA GUIADA
sl10 = base_slide(10)
add_header(sl10, "Prática Guiada — Blocos 1 a 5", "Correção de Vícios")
exs_1_5 = [
    ("01", "[Lacunas contextualizadas com opções]"),
    ("02", "[Lacunas contextualizadas com opções]"),
    ("03", "[Preenchimento direto de verbos]"),
    ("04", "[Reestruturação / Transformação de frase]"),
    ("05", "[Reestruturação / Transformação de frase]")
]
for idx, (num_s, d_e) in enumerate(exs_1_5):
    col = idx % 2; row = idx // 2
    l = 0.7 if col == 0 else 6.9
    w = 5.9 if (idx < 4 or len(exs_1_5) % 2 == 0) else 12.1
    if idx == 4: l = 0.7
    t = 1.5 + row * 1.7
    add_scale_flashcard(sl10, num_s, d_e, Inches(l), Inches(t), Inches(w), Inches(1.5), accent_c=ROXO if idx%2==0 else TURQUESA)

sl11 = base_slide(11)
add_header(sl11, "Prática Guiada — Blocos 6 a 10", "Desafios Práticos")
exs_6_10 = [
    ("06", "[Preposição / Conector correto]"),
    ("07", "[Escolha da opção correta]"),
    ("08", "[Substituição por sinônimo formal]"),
    ("09", "[Identificação de erro e correção]"),
    ("10", "[Tradução de sentido sem uso literal]")
]
for idx, (num_s, d_e) in enumerate(exs_6_10):
    col = idx % 2; row = idx // 2
    l = 0.7 if col == 0 else 6.9
    w = 5.9 if (idx < 4 or len(exs_6_10) % 2 == 0) else 12.1
    if idx == 4: l = 0.7
    t = 1.5 + row * 1.7
    add_scale_flashcard(sl11, num_s, d_e, Inches(l), Inches(t), Inches(w), Inches(1.5), accent_c=TURQUESA if idx%2==0 else ROXO)

# SLIDES 12 A 19: CONVERSAÇÃO SCALE CONTRAST (Pt 72)
perguntas_config = [
    "[Pergunta principal 1 sobre o tema]", "[Pergunta principal 2 sobre o tema]", "[Pergunta principal 3 sobre o tema]",
    "[Pergunta principal 4 sobre o tema]", "[Pergunta principal 5 sobre o tema]", "[Pergunta principal 6 sobre o tema]",
    "[Pergunta principal 7 sobre o tema]", "[Pergunta principal 8 sobre o tema]"
]

def add_scale_history_ticker(sl, history_list):
    count = len(history_list)
    if count == 0: return
    gap = 0.12
    avail_w = 8.1 - (gap * (count - 1))
    card_w = avail_w / count
    
    for i, h_txt in enumerate(history_list):
        left_x = 4.7 + i * (card_w + gap)
        clean_t = h_txt[:26] + "..." if len(h_txt) > 26 else h_txt
        d_sz = 11 if count <= 3 else (10 if count <= 5 else 9)
        bg_col = ROXO_CLARO if i % 2 == 0 else TURQUESA_CLARO
        acc_c = ROXO if i % 2 == 0 else TURQUESA
        add_scale_card(sl, f"{i+1:02d}", "", clean_t, Inches(left_x), Inches(5.3), Inches(card_w), Inches(1.5), bg_color=bg_col, border_color=acc_c, accent_color=acc_c, desc_size=d_sz)

for idx in range(8):
    s_num = 12 + idx
    sl = base_slide(s_num)
    add_header(sl, "Conversação", f"Thread {idx + 1} de 8")
    p_desc = perguntas_config[idx]
    
    giant_num_box = sl.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.7), Inches(1.5), Inches(3.7), Inches(5.3))
    giant_num_box.fill.solid(); giant_num_box.fill.fore_color.rgb = ROXO_CLARO
    giant_num_box.line.color.rgb = ROXO; giant_num_box.line.width = Pt(2)
    tf_gn = giant_num_box.text_frame; tf_gn.vertical_anchor = MSO_ANCHOR.MIDDLE
    
    p_g = tf_gn.paragraphs[0]; p_g.text = f"Q{idx+1:02d}"
    p_g.font.size = Pt(72); p_g.font.color.rgb = ROXO; p_g.font.bold = True; p_g.alignment = PP_ALIGN.CENTER
    p_g_sub = tf_gn.add_paragraph(); p_g_sub.text = "PERGUNTA ATIVA"
    p_g_sub.font.size = Pt(12); p_g_sub.font.color.rgb = AZUL; p_g_sub.font.bold = True; p_g_sub.alignment = PP_ALIGN.CENTER; p_g_sub.space_before = Pt(8)

    if idx == 0:
        add_scale_card(sl, "", "Pergunta Principal:", p_desc, Inches(4.7), Inches(1.5), Inches(8.1), Inches(2.2), border_color=ROXO, tit_size=18, desc_size=18, align_center_v=True)
        add_scale_card(sl, "", "Notas de Apoio ao Professor:", "• Para expandir: 'Dê um exemplo recente...'\n• Para concordar/discordar: 'Qual é o outro lado dessa moeda?'", Inches(4.7), Inches(3.9), Inches(8.1), Inches(2.9), bg_color=TURQUESA_CLARO, border_color=TURQUESA, tit_size=16, desc_size=14, align_center_v=True)
    else:
        add_scale_card(sl, "", "Pergunta Principal:", p_desc, Inches(4.7), Inches(1.5), Inches(8.1), Inches(2.0), border_color=ROXO, tit_size=18, desc_size=16, align_center_v=True)
        add_scale_card(sl, "", "Apoio ao Professor:", "• Para expandir: 'Dê um exemplo recente...'\n• Para discordar: 'Qual é o outro lado?'", Inches(4.7), Inches(3.6), Inches(8.1), Inches(1.5), bg_color=TURQUESA_CLARO, border_color=TURQUESA, tit_size=14, desc_size=13, align_center_v=True)
        history_list = [perguntas_config[i] for i in range(idx)]
        add_scale_history_ticker(sl, history_list)

# SLIDE 20: ESCRITA
sl20 = base_slide(20)
add_header(sl20, "Atividade de Escrita", "Consolidação e Memorização")
add_scale_card(sl20, "TASK", "Tarefa de Escrita Recomendada:", "[Instrução detalhada da tarefa: escrever um texto, frases ou diálogo]", Inches(0.7), Inches(1.5), Inches(12.1), Inches(2.8), border_color=ROXO, accent_color=ROXO, tit_size=18, desc_size=15)
add_scale_card(sl20, "WHY", "Fundamentação Pedagógica:", "Identificamos que a escrita nesta etapa é fundamental para você fixar as estruturas aprendidas, expandir vocabulário ativo e diminuir a tradução mental automática do seu idioma nativo.", Inches(0.7), Inches(4.5), Inches(12.1), Inches(2.3), border_color=TURQUESA, accent_color=TURQUESA, desc_size=14, align_center_v=True)

# SLIDE 21: TAREFA DE CASA
sl21 = base_slide(21)
add_header(sl21, "Encerramento — Tarefa de Casa", "Ação e Prática Contínua")
add_scale_card(sl21, "HOMEWORK", "Tarefa de Casa:", "• [Atividade prática 1 para realizar fora da aula]\n• [Atividade de fixação de vocabulário ou escuta]", Inches(0.7), Inches(1.5), Inches(12.1), Inches(5.3), border_color=ROXO, accent_color=ROXO, tit_size=18, desc_size=15)

# SLIDE 22: ENCERRAMENTO & FEEDBACK
sl22 = base_slide(22)
add_header(sl22, "Encerramento & Feedback Rápido", "Considerações Finais")
add_scale_card(sl22, "THANKS", "Agradecimento:", "Obrigado pela dedicação na aula de hoje!", Inches(0.7), Inches(1.5), Inches(12.1), Inches(2.0), border_color=TURQUESA, accent_color=TURQUESA, tit_size=20, desc_size=16)
add_scale_card(sl22, "FEEDBACK", "Espaço para Feedback Rápido:", "[Espaço destinado para o professor coletar dúvidas rápidas e considerações do aluno]", Inches(0.7), Inches(3.7), Inches(12.1), Inches(3.1), border_color=ROXO, accent_color=ROXO, tit_size=18, desc_size=15)

out_path_server = "/var/www/haas-frontend-desk-mobile-oficial/public/temp-miniaturas/template-haas-modelo-18-scale.pptx"
prs.save(out_path_server)
prs.save("template-haas-modelo-18-scale.pptx")
print(f"🎉 MODELO 18 SCALE CONTRAST GERADO EM: {out_path_server}")
