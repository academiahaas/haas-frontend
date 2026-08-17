import os
import math
import urllib.request
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# PALETA OFICIAL HAAS LANGUAGE - MODELO 13 (ORBITAL SPHERES)
BRANCO = RGBColor(0xFF, 0xFF, 0xFF)
ROXO = RGBColor(0x4F, 0x35, 0xE2)
AZUL = RGBColor(0x1E, 0x3A, 0x8A) # Fundo Principal
SLATE = RGBColor(0x33, 0x41, 0x55)
TURQUESA = RGBColor(0x14, 0xC8, 0xB0)
TURQUESA_CLARO = RGBColor(0xE0, 0xF2, 0xF1)
BORDA = RGBColor(0xCB, 0xD5, 0xE1)
ROXO_CLARO = RGBColor(0xED, 0xE9, 0xFE)

LOGO_URL = "https://jdppxfokfhqjudwfwckd.supabase.co/storage/v1/object/public/haas-academy/assignments/Design%20sem%20nome%20(4).png"
LOGO_PATH = "/tmp/haas-logo.png"

if not os.path.exists(LOGO_PATH):
    try:
        urllib.request.urlretrieve(LOGO_URL, LOGO_PATH)
    except Exception:
        pass

W, H = Inches(13.333), Inches(7.5)
TOTAL_SLIDES = 22
prs = Presentation()
prs.slide_width = W
prs.slide_height = H

def base_slide(num, hide_logo=False):
    sl = prs.slides.add_slide(prs.slide_layouts[6])
    sl.background.fill.solid()
    sl.background.fill.fore_color.rgb = AZUL

    # Esferas de Fundo Decorativas
    c1 = sl.shapes.add_shape(MSO_SHAPE.OVAL, Inches(9.5), Inches(-2.5), Inches(6), Inches(6))
    c1.fill.solid(); c1.fill.fore_color.rgb = ROXO; c1.line.fill.background()
    
    c2 = sl.shapes.add_shape(MSO_SHAPE.OVAL, Inches(11.0), Inches(-0.5), Inches(4), Inches(4))
    c2.fill.solid(); c2.fill.fore_color.rgb = TURQUESA; c2.line.fill.background()

    c3 = sl.shapes.add_shape(MSO_SHAPE.OVAL, Inches(-2.0), Inches(5.0), Inches(5), Inches(5))
    c3.fill.solid(); c3.fill.fore_color.rgb = ROXO; c3.line.fill.background()

    c4 = sl.shapes.add_shape(MSO_SHAPE.OVAL, Inches(-0.5), Inches(6.0), Inches(3), Inches(3))
    c4.fill.solid(); c4.fill.fore_color.rgb = TURQUESA; c4.line.fill.background()

    if os.path.exists(LOGO_PATH) and not hide_logo:
        sl.shapes.add_picture(LOGO_PATH, Inches(10.8), Inches(0.2), height=Inches(0.65))
        
    tb = sl.shapes.add_textbox(Inches(0.5), Inches(7.12), Inches(8), Inches(0.35))
    p = tb.text_frame.paragraphs[0]
    p.text = "HAAS LANGUAGE  •  ORBITAL SPHERES EDITION"
    p.font.size = Pt(10); p.font.color.rgb = BRANCO; p.font.bold = True
    
    tb2 = sl.shapes.add_textbox(Inches(10.8), Inches(7.12), Inches(2.0), Inches(0.35))
    p2 = tb2.text_frame.paragraphs[0]
    p2.text = f"SLIDE {num} / {TOTAL_SLIDES}"
    p2.font.size = Pt(10); p2.font.color.rgb = TURQUESA; p2.font.bold = True; p2.alignment = PP_ALIGN.RIGHT
    return sl

def add_header(sl, tit, sub=""):
    tb = sl.shapes.add_textbox(Inches(0.5), Inches(0.35), Inches(10), Inches(0.65))
    tf = tb.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]; p.text = tit
    p.font.size = Pt(32); p.font.color.rgb = BRANCO; p.font.bold = True
    if sub:
        tb_s = sl.shapes.add_textbox(Inches(0.5), Inches(0.98), Inches(10), Inches(0.4))
        p_s = tb_s.text_frame.paragraphs[0]
        p_s.text = sub
        p_s.font.size = Pt(18); p_s.font.color.rgb = TURQUESA; p_s.font.bold = True

def add_rounded_card(sl, tit, desc, left, top, width, height, bg_color=BRANCO, border_color=None, tit_color=ROXO, desc_color=SLATE, tit_size=16, desc_size=14, align_center_v=False):
    s = sl.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    s.fill.solid(); s.fill.fore_color.rgb = bg_color
    if border_color:
        s.line.color.rgb = border_color; s.line.width = Pt(1.5)
    else:
        s.line.fill.background()
        
    tf = s.text_frame; tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE if align_center_v else MSO_ANCHOR.TOP
    margin_val = Inches(0.12) if height < Inches(1.0) else Inches(0.2)
    tf.margin_left = margin_val; tf.margin_top = margin_val; tf.margin_right = margin_val; tf.margin_bottom = margin_val
    
    if tit:
        p = tf.paragraphs[0]
        p.text = tit; p.font.size = Pt(tit_size); p.font.color.rgb = tit_color; p.font.bold = True
        if desc:
            p2 = tf.add_paragraph()
            p2.text = desc; p2.font.size = Pt(desc_size); p2.font.color.rgb = desc_color; p2.space_before = Pt(6)
    else:
        if desc:
            p = tf.paragraphs[0]
            p.text = desc; p.font.size = Pt(desc_size); p.font.color.rgb = desc_color

def img_circle(sl, left, top, diameter, label="[Imagem]"):
    s = sl.shapes.add_shape(MSO_SHAPE.OVAL, left, top, diameter, diameter)
    s.fill.solid(); s.fill.fore_color.rgb = ROXO_CLARO
    s.line.color.rgb = TURQUESA; s.line.width = Pt(3)
    tf = s.text_frame; tf.word_wrap = True; tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]; p.text = label; p.font.size = Pt(16); p.font.color.rgb = ROXO; p.font.bold = True; p.alignment = PP_ALIGN.CENTER

def add_circle_flashcard(sl, num_str, desc, left, top, width, height, b_c=TURQUESA):
    s = sl.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    s.fill.solid(); s.fill.fore_color.rgb = BRANCO; s.line.color.rgb = b_c; s.line.width = Pt(2)
    
    chip = sl.shapes.add_shape(MSO_SHAPE.OVAL, left - Inches(0.25), top + (height/2) - Inches(0.3), Inches(0.6), Inches(0.6))
    chip.fill.solid(); chip.fill.fore_color.rgb = b_c; chip.line.color.rgb = BRANCO; chip.line.width = Pt(1.5)
    tf_chip = chip.text_frame; tf_chip.vertical_anchor = MSO_ANCHOR.MIDDLE
    p_chip = tf_chip.paragraphs[0]; p_chip.text = num_str; p_chip.alignment = PP_ALIGN.CENTER
    p_chip.font.size = Pt(12); p_chip.font.color.rgb = BRANCO; p_chip.font.bold = True
    
    tf = s.text_frame; tf.word_wrap = True; tf.vertical_anchor = MSO_ANCHOR.TOP
    tf.margin_left = Inches(0.45); tf.margin_top = Inches(0.15); tf.margin_right = Inches(0.15); tf.margin_bottom = Inches(0.12)
    p = tf.paragraphs[0]; p.text = desc; p.font.size = Pt(14); p.font.color.rgb = SLATE; p.font.bold = False

# ==========================================
# MAGIA ORBITAL (SISTEMA SOLAR DE PERGUNTAS)
# ==========================================
def add_main_question_circle(sl, question_text):
    """Círculo Gigante Central contendo a Pergunta Ativa e os Prompts"""
    cx = 6.666
    cy = 4.3
    diameter = 4.8
    left = cx - (diameter / 2)
    top = cy - (diameter / 2)
    
    s = sl.shapes.add_shape(MSO_SHAPE.OVAL, Inches(left), Inches(top), Inches(diameter), Inches(diameter))
    s.fill.solid(); s.fill.fore_color.rgb = BRANCO
    s.line.color.rgb = ROXO; s.line.width = Pt(5)
    
    tf = s.text_frame; tf.word_wrap = True; tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    tf.margin_left = Inches(0.4); tf.margin_right = Inches(0.4)
    
    p = tf.paragraphs[0]; p.text = question_text
    p.font.size = Pt(20); p.font.color.rgb = ROXO; p.font.bold = True; p.alignment = PP_ALIGN.CENTER
    
    p2 = tf.add_paragraph()
    p2.text = "• Para expandir: 'Dê um exemplo recente...'\n• Para discordar: 'Qual é o outro lado?'"
    p2.font.size = Pt(13); p2.font.color.rgb = SLATE; p2.alignment = PP_ALIGN.CENTER; p2.space_before = Pt(14)

def add_orbital_history(sl, history_list):
    """Distribui as 'bolinhas' do histórico em órbita geométrica 360º ao redor do círculo central"""
    cx = 6.666
    cy = 4.3
    r_x = 4.6
    r_y = 2.1
    
    for i, h_txt in enumerate(history_list):
        # Distribui os slots em um arco ao redor do centro usando seno e cosseno!
        angle_deg = 180 + i * (360 / 7)
        angle_rad = math.radians(angle_deg)
        x = cx + r_x * math.cos(angle_rad)
        y = cy + r_y * math.sin(angle_rad)
        
        b_size = 2.0
        left = x - (b_size / 2)
        top = y - (b_size / 2)
        
        clean_t = h_txt[:35] + "..." if len(h_txt) > 35 else h_txt
        bg_col = TURQUESA_CLARO if i % 2 == 0 else ROXO_CLARO
        b_col = TURQUESA if i % 2 == 0 else ROXO
        
        s = sl.shapes.add_shape(MSO_SHAPE.OVAL, Inches(left), Inches(top), Inches(b_size), Inches(b_size))
        s.fill.solid(); s.fill.fore_color.rgb = bg_col
        s.line.color.rgb = b_col; s.line.width = Pt(2)
        
        tf = s.text_frame; tf.word_wrap = True; tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        tf.margin_left = Inches(0.15); tf.margin_right = Inches(0.15)
        p = tf.paragraphs[0]; p.text = f"Q{i+1}"
        p.font.size = Pt(12); p.font.color.rgb = ROXO; p.font.bold = True; p.alignment = PP_ALIGN.CENTER
        p2 = tf.add_paragraph(); p2.text = clean_t
        p2.font.size = Pt(10); p2.font.color.rgb = SLATE; p2.alignment = PP_ALIGN.CENTER; p2.space_before = Pt(2)


# ==========================================================
# CONSTRUÇÃO DOS 22 SLIDES
# ==========================================================
sl1 = base_slide(1, hide_logo=True)
img_circle(sl1, Inches(7.0), Inches(1.0), Inches(5.5), "[Imagem Hero Circular]")
add_rounded_card(sl1, "[Título da Aula / Tema Principal]", "", Inches(0.5), Inches(2.2), Inches(6.0), Inches(1.8), bg_color=BRANCO, border_color=None, tit_size=36, align_center_v=True)
add_rounded_card(sl1, "[Subtítulo / Descrição da Aula]", "", Inches(0.5), Inches(4.2), Inches(6.0), Inches(1.2), bg_color=BRANCO, border_color=None, tit_color=AZUL, tit_size=20, align_center_v=True)
add_rounded_card(sl1, "📅 Data da Aula: [DD/MM/AAAA]", "", Inches(0.5), Inches(5.6), Inches(6.0), Inches(0.8), bg_color=BRANCO, border_color=None, tit_color=TURQUESA, tit_size=16, align_center_v=True)

sl2 = base_slide(2)
add_header(sl2, "Presença & Integrantes")
add_rounded_card(sl2, "Professor(a):", "[Nome do(a) Professor(a)]", Inches(0.5), Inches(1.5), Inches(12.3), Inches(1.5), bg_color=ROXO_CLARO, border_color=ROXO)
add_rounded_card(sl2, "Aluno(s):", "[Nome do Aluno / Lista da Turma]", Inches(0.5), Inches(3.4), Inches(12.3), Inches(3.4), bg_color=BRANCO, border_color=TURQUESA)

sl3 = base_slide(3)
add_header(sl3, "Aquecimento / Warm-up", "Foco 100% na Fala (3 Minutos)")
add_rounded_card(sl3, "Pergunta Provocativa Rápida:", "[Insira aqui a pergunta provocativa de aquecimento]", Inches(0.5), Inches(1.5), Inches(7.0), Inches(2.6), bg_color=BRANCO, tit_size=20, desc_size=16)
add_rounded_card(sl3, "", "Ativa o córtex pré-frontal e reduz a inibição linguística inicial.", Inches(0.5), Inches(4.5), Inches(7.0), Inches(2.0), bg_color=TURQUESA_CLARO, border_color=TURQUESA, align_center_v=True)
img_circle(sl3, Inches(8.0), Inches(1.8), Inches(4.5), "[Imagem Circular]")

sl4 = base_slide(4)
add_header(sl4, "Objetivos da Aula")
add_rounded_card(sl4, "Metas da Aula:", "• [Objetivo 1]\n• [Objetivo 2]\n• [Objetivo 3]", Inches(0.5), Inches(1.5), Inches(7.5), Inches(5.3), bg_color=BRANCO)
add_rounded_card(sl4, "Nota de Escrita:", "Nesta aula daremos atenção à escrita para reforçar a memorização de estruturas/vocabulário e reduzir a interferência da língua materna.", Inches(8.4), Inches(1.5), Inches(4.4), Inches(5.3), bg_color=ROXO_CLARO, border_color=ROXO)

sl5 = base_slide(5)
add_header(sl5, "A Nossa Metodologia", "Estrutura Fixa de Aprendizado Haas Language")
metodologias = [
    ("1. Fala Ativa", "Conversação inicial sem julgamentos."), ("2. Exposição", "Apresentação leve da gramática."),
    ("3. Prática", "Exercícios com suporte."), ("4. Conversação", "Aplicação prática das estruturas."), ("5. Escrita", "Consolidação visual.")
]
left_m = 0.5
for i, (tit_m, desc_m) in enumerate(metodologias):
    top_m = 2.0 if i % 2 == 0 else 3.5
    bg_col = BRANCO if i % 2 == 0 else ROXO_CLARO
    border_col = TURQUESA if i % 2 == 0 else ROXO
    s = sl5.shapes.add_shape(MSO_SHAPE.OVAL, Inches(left_m), Inches(top_m), Inches(2.8), Inches(2.8))
    s.fill.solid(); s.fill.fore_color.rgb = bg_col
    s.line.color.rgb = border_col; s.line.width = Pt(2)
    tf = s.text_frame; tf.word_wrap = True; tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    tf.margin_left = Inches(0.2); tf.margin_right = Inches(0.2)
    p = tf.paragraphs[0]; p.text = tit_m; p.font.size = Pt(15); p.font.color.rgb = ROXO; p.font.bold = True; p.alignment = PP_ALIGN.CENTER
    p2 = tf.add_paragraph(); p2.text = desc_m; p2.font.size = Pt(12); p2.font.color.rgb = SLATE; p2.alignment = PP_ALIGN.CENTER; p2.space_before = Pt(4)
    left_m += 2.4

for s_num in [6, 7, 8]:
    sl = base_slide(s_num)
    add_header(sl, f"Explicação do Tema — Parte {s_num - 5}")
    add_rounded_card(sl, "Conteúdo Principal:", f"[Explicação e exemplo central da estrutura - Bloco {s_num - 5}]", Inches(0.5), Inches(1.5), Inches(12.3), Inches(2.5), bg_color=BRANCO, tit_size=20, desc_size=15)
    add_rounded_card(sl, "Exemplos Práticos:", "[Exemplos adicionais e frases do dia a dia]", Inches(0.5), Inches(4.4), Inches(6.0), Inches(2.4), bg_color=ROXO_CLARO, border_color=ROXO)
    add_rounded_card(sl, "Notas & Detalhes:", "[Exceções, nuances gramaticais e contexto de uso]", Inches(6.8), Inches(4.4), Inches(6.0), Inches(2.4), bg_color=TURQUESA_CLARO, border_color=TURQUESA)

sl9 = base_slide(9)
add_header(sl9, "Síntese do Conteúdo & Sugestão Cultural")
add_rounded_card(sl9, "Resumo do Tema:", "[Resumo sintético dos pontos chave ensinados na aula]", Inches(0.5), Inches(1.5), Inches(6.0), Inches(5.3), bg_color=BRANCO)
add_rounded_card(sl9, "Sugestão Cultural:", "🎵 [Nome da Música / Vídeo recomendado sobre o assunto].\n\nO professor pode rodar o áudio/vídeo ou deixar como recomendação extra para o aluno praticar.", Inches(6.8), Inches(1.5), Inches(6.0), Inches(5.3), bg_color=ROXO_CLARO, border_color=ROXO)

sl10 = base_slide(10)
add_header(sl10, "Prática Guiada — Blocos 1 a 5", "Correção de Vícios")
add_circle_flashcard(sl10, "01", "[Lacunas contextualizadas com opções]", Inches(0.8), Inches(1.5), Inches(12.0), Inches(1.2), b_c=ROXO)
add_circle_flashcard(sl10, "02", "[Lacunas contextualizadas com opções]", Inches(0.8), Inches(2.9), Inches(5.5), Inches(1.8), b_c=TURQUESA)
add_circle_flashcard(sl10, "03", "[Preenchimento direto de verbos]", Inches(7.3), Inches(2.9), Inches(5.5), Inches(1.8), b_c=TURQUESA)
add_circle_flashcard(sl10, "04", "[Reestruturação / Transformação de frase]", Inches(0.8), Inches(4.9), Inches(5.5), Inches(1.9), b_c=ROXO)
add_circle_flashcard(sl10, "05", "[Reestruturação / Transformação de frase]", Inches(7.3), Inches(4.9), Inches(5.5), Inches(1.9), b_c=ROXO)

sl11 = base_slide(11)
add_header(sl11, "Prática Guiada — Blocos 6 a 10", "Desafios Práticos")
add_circle_flashcard(sl11, "06", "[Preposição / Conector correto]", Inches(0.8), Inches(1.5), Inches(12.0), Inches(1.2), b_c=TURQUESA)
add_circle_flashcard(sl11, "07", "[Escolha da opção correta]", Inches(0.8), Inches(2.9), Inches(5.5), Inches(1.8), b_c=ROXO)
add_circle_flashcard(sl11, "08", "[Substituição por sinônimo formal]", Inches(7.3), Inches(2.9), Inches(5.5), Inches(1.8), b_c=ROXO)
add_circle_flashcard(sl11, "09", "[Identificação de erro e correção]", Inches(0.8), Inches(4.9), Inches(5.5), Inches(1.9), b_c=TURQUESA)
add_circle_flashcard(sl11, "10", "[Tradução de sentido sem uso literal]", Inches(7.3), Inches(4.9), Inches(5.5), Inches(1.9), b_c=TURQUESA)

# SLIDES 12 A 19: CONVERSAÇÃO (ÓRBITAS ESTILO SISTEMA SOLAR)
perguntas_config = [
    "[Pergunta principal 1 sobre o tema]", "[Pergunta principal 2 sobre o tema]", "[Pergunta principal 3 sobre o tema]",
    "[Pergunta principal 4 sobre o tema]", "[Pergunta principal 5 sobre o tema]", "[Pergunta principal 6 sobre o tema]",
    "[Pergunta principal 7 sobre o tema]", "[Pergunta principal 8 sobre o tema]"
]
for idx in range(8):
    sl = base_slide(12 + idx)
    add_header(sl, "Conversação")
    
    add_main_question_circle(sl, perguntas_config[idx])
    history_list = [perguntas_config[i] for i in range(idx)]
    add_orbital_history(sl, history_list)

sl20 = base_slide(20)
add_header(sl20, "Atividade de Escrita", "Consolidação e Memorização")
add_rounded_card(sl20, "Tarefa de Escrita Recomendada:", "[Instrução detalhada da tarefa: escrever um texto, frases ou diálogo]", Inches(0.5), Inches(1.5), Inches(12.3), Inches(2.8), bg_color=BRANCO, border_color=ROXO)
add_rounded_card(sl20, "", "Identificamos que a escrita nesta etapa é fundamental para você fixar as estruturas aprendidas, expandir vocabulário ativo e diminuir a tradução mental automática do seu idioma nativo.", Inches(0.5), Inches(4.6), Inches(12.3), Inches(2.2), bg_color=TURQUESA_CLARO, border_color=TURQUESA, align_center_v=True)

sl21 = base_slide(21)
add_header(sl21, "Encerramento — Tarefa de Casa", "Ação e Prática Contínua")
add_rounded_card(sl21, "Tarefa de Casa:", "• [Atividade prática 1 para realizar fora da aula]\n• [Atividade de fixação de vocabulário ou escuta]", Inches(0.5), Inches(1.5), Inches(12.3), Inches(5.2), bg_color=BRANCO, border_color=TURQUESA)

sl22 = base_slide(22)
add_header(sl22, "Encerramento & Feedback Rápido", "Considerações Finais")
add_rounded_card(sl22, "Agradecimento:", "Obrigado pela dedicação na aula de hoje!", Inches(0.5), Inches(1.5), Inches(12.3), Inches(2.0), bg_color=ROXO_CLARO, border_color=ROXO, tit_size=20, desc_size=16)
add_rounded_card(sl22, "Espaço para Feedback Rápido:", "[Espaço destinado para o professor coletar dúvidas rápidas e considerações do aluno]", Inches(0.5), Inches(3.8), Inches(12.3), Inches(2.9), bg_color=BRANCO, border_color=TURQUESA, tit_size=18, desc_size=15)

out_path = "/var/www/haas-frontend-desk-mobile-oficial/public/temp-miniaturas/template-haas-modelo-13-orbital.pptx"
prs.save(out_path)
print(f"🎉 SISTEMA SOLAR GERADO COM SUCESSO: {out_path}")
