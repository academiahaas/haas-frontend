import os
import shutil
import urllib.request
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# PALETA OFICIAL HAAS LANGUAGE - ANATOMIA DO TEXTO (DATA-DRIVEN TYPO / SCHEMATIC)
BRANCO = RGBColor(0xFF, 0xFF, 0xFF)
ROXO = RGBColor(0x4F, 0x35, 0xE2)
AZUL = RGBColor(0x1E, 0x3A, 0x8A)
SLATE = RGBColor(0x33, 0x41, 0x55)
TURQUESA = RGBColor(0x14, 0xC8, 0xB0)
ROXO_CLARO = RGBColor(0xED, 0xE9, 0xFE)
TURQUESA_CLARO = RGBColor(0xE0, 0xF2, 0xF1)
FUNDO_BG = RGBColor(0xF8, 0xFA, 0xFC)

LOGO_URL = "https://jdppxfokfhqjudwfwckd.supabase.co/storage/v1/object/public/haas-academy/assignments/Design%20sem%20nome%20(4).png"
LOGO_PATH = "/tmp/haas-logo.png"

if not os.path.exists(LOGO_PATH):
    try:
        urllib.request.urlretrieve(LOGO_URL, LOGO_PATH)
    except Exception:
        pass

W, H = Inches(13.333), Inches(7.5)
TOTAL_SLIDES = 22

prs = Presentation()
prs.slide_width = W
prs.slide_height = H

def base_slide(num, hide_logo=False):
    sl = prs.slides.add_slide(prs.slide_layouts[6])
    sl.background.fill.solid()
    sl.background.fill.fore_color.rgb = FUNDO_BG

    # Guia Técnico Blueprint no Topo (Linha de Eixo com Nó Círculo)
    axis = sl.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.5), Inches(0.2), W - Inches(1.0), Pt(1.5))
    axis.fill.solid(); axis.fill.fore_color.rgb = ROXO; axis.line.fill.background()

    node = sl.shapes.add_shape(MSO_SHAPE.OVAL, Inches(0.5), Inches(0.125), Inches(0.15), Inches(0.15))
    node.fill.solid(); node.fill.fore_color.rgb = TURQUESA; node.line.fill.background()

    if os.path.exists(LOGO_PATH) and not hide_logo:
        sl.shapes.add_picture(LOGO_PATH, Inches(10.8), Inches(0.2), height=Inches(0.65))
        
    tb = sl.shapes.add_textbox(Inches(0.5), Inches(7.12), Inches(8), Inches(0.35))
    p = tb.text_frame.paragraphs[0]
    p.text = "HAAS LANGUAGE  •  TEXT ANATOMY / DATA-DRIVEN EDITION 📐"
    p.font.size = Pt(10); p.font.color.rgb = ROXO; p.font.bold = True
    
    tb2 = sl.shapes.add_textbox(Inches(10.8), Inches(7.12), Inches(2.0), Inches(0.35))
    p2 = tb2.text_frame.paragraphs[0]
    p2.text = f"SCHEMATIC [{num:02d}/{TOTAL_SLIDES:02d}]"
    p2.font.size = Pt(10); p2.font.color.rgb = AZUL; p2.font.bold = True; p2.alignment = PP_ALIGN.RIGHT
    return sl

def add_header(sl, tit, sub=""):
    tb = sl.shapes.add_textbox(Inches(0.5), Inches(0.25), Inches(10), Inches(0.65))
    tf = tb.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]; p.text = f"// {tit}"
    p.font.size = Pt(30); p.font.color.rgb = ROXO; p.font.bold = True
    
    if sub:
        tb_s = sl.shapes.add_textbox(Inches(0.5), Inches(0.95), Inches(10), Inches(0.4))
        p_s = tb_s.text_frame.paragraphs[0]
        p_s.text = f"[ANATOMY SPEC]: {sub}"
        p_s.font.size = Pt(15); p_s.font.color.rgb = AZUL; p_s.font.bold = True

def add_schematic_card(sl, tit, desc, left, top, width, height, bg_c=BRANCO, border_c=ROXO, tag="[COMPONENT]", tit_c=ROXO, desc_c=SLATE, tit_sz=16, desc_sz=14, align_center_v=False):
    s = sl.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(left), Inches(top), Inches(width), Inches(height))
    s.fill.solid(); s.fill.fore_color.rgb = bg_c
    s.line.color.rgb = border_c; s.line.width = Pt(1.5)

    if tag:
        tag_box = sl.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(left), Inches(top - 0.28), Inches(len(tag) * 0.12 + 0.5), Inches(0.28))
        tag_box.fill.solid(); tag_box.fill.fore_color.rgb = border_c; tag_box.line.fill.background()
        tf_t = tag_box.text_frame; tf_t.vertical_anchor = MSO_ANCHOR.MIDDLE
        p_t = tf_t.paragraphs[0]; p_t.text = tag.upper(); p_t.alignment = PP_ALIGN.CENTER
        p_t.font.size = Pt(9); p_t.font.color.rgb = BRANCO; p_t.font.bold = True

    tf = s.text_frame; tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE if align_center_v else MSO_ANCHOR.TOP
    tf.margin_left = Inches(0.25); tf.margin_top = Inches(0.2); tf.margin_right = Inches(0.25); tf.margin_bottom = Inches(0.2)
    
    if tit:
        p = tf.paragraphs[0]
        p.text = tit; p.font.size = Pt(tit_sz); p.font.color.rgb = tit_c; p.font.bold = True
        if desc:
            p2 = tf.add_paragraph()
            p2.text = desc; p2.font.size = Pt(desc_sz); p2.font.color.rgb = desc_c; p2.space_before = Pt(6)
    else:
        if desc:
            p = tf.paragraphs[0]
            p.text = desc; p.font.size = Pt(desc_sz); p.font.color.rgb = desc_c

def img_schematic_box(sl, left, top, width, height, label="[VISUAL_DIAGRAM_SPEC]"):
    s = sl.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(left), Inches(top), Inches(width), Inches(height))
    s.fill.solid(); s.fill.fore_color.rgb = ROXO_CLARO
    s.line.color.rgb = TURQUESA; s.line.width = Pt(2.0)

    l1 = sl.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(left + 0.1), Inches(top + 0.1), Inches(width - 0.2), Inches(height - 0.2))
    l1.fill.background(); l1.line.color.rgb = ROXO; l1.line.width = Pt(1.0)

    tf = s.text_frame; tf.word_wrap = True; tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER
    p = tf.paragraphs[0]; p.text = label; p.font.size = Pt(15); p.font.color.rgb = ROXO; p.font.bold = True

def add_schematic_flashcard(sl, num_str, desc, left, top, width, height, border_c=ROXO):
    s = sl.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(left), Inches(top), Inches(width), Inches(height))
    s.fill.solid(); s.fill.fore_color.rgb = BRANCO
    s.line.color.rgb = border_c; s.line.width = Pt(1.5)

    node = sl.shapes.add_shape(MSO_SHAPE.OVAL, Inches(left + 0.15), Inches(top + (height/2) - 0.18), Inches(0.36), Inches(0.36))
    node.fill.solid(); node.fill.fore_color.rgb = border_c; node.line.fill.background()
    tf_n = node.text_frame; tf_n.vertical_anchor = MSO_ANCHOR.MIDDLE
    p_n = tf_n.paragraphs[0]; p_n.text = num_str; p_n.alignment = PP_ALIGN.CENTER
    p_n.font.size = Pt(10); p_n.font.color.rgb = BRANCO; p_n.font.bold = True

    tf = s.text_frame; tf.word_wrap = True; tf.vertical_anchor = MSO_ANCHOR.TOP
    tf.margin_left = Inches(0.65); tf.margin_top = Inches(0.15); tf.margin_right = Inches(0.15); tf.margin_bottom = Inches(0.12)
    p = tf.paragraphs[0]; p.text = desc; p.font.size = Pt(14); p.font.color.rgb = SLATE; p.font.bold = False

# SLIDE 1
sl1 = base_slide(1, hide_logo=True)
add_schematic_card(sl1, "[Título da Aula / Tema Principal]", "[Subtítulo / Descrição da Aula]\n\n📅 Data da Aula: [DD/MM/AAAA]", 0.5, 1.8, 7.0, 5.0, bg_c=BRANCO, border_c=ROXO, tag="[STRUCTURE_HEAD]", tit_sz=36, desc_sz=18)
img_schematic_box(sl1, 7.8, 1.8, 5.0, 5.0, "[DIAGRAMA_CAPA_SCHEMATIC]")

# SLIDE 2
sl2 = base_slide(2)
add_header(sl2, "Presença & Integrantes", "Registro de Nó de Acesso")
add_schematic_card(sl2, "Professor(a):", "[Nome do(a) Professor(a)]", 0.5, 1.8, 12.333, 2.2, bg_c=ROXO_CLARO, border_c=ROXO, tag="[INSTRUCTOR_NODE]", tit_sz=20, desc_sz=16)
add_schematic_card(sl2, "Aluno(s):", "[Nome do Aluno / Lista da Turma]", 0.5, 4.5, 12.333, 2.4, bg_c=TURQUESA_CLARO, border_c=TURQUESA, tag="[STUDENTS_NODE]", tit_sz=20, desc_sz=16)

# SLIDE 3
sl3 = base_slide(3)
add_header(sl3, "Aquecimento / Warm-up", "Foco 100% na Fala (3 Minutos)")
add_schematic_card(sl3, "Pergunta Provocativa Rápida:", "[Insira aqui a pergunta provocativa de aquecimento]", 0.5, 1.8, 7.0, 2.4, bg_c=BRANCO, border_c=ROXO, tag="[SPEECH_TRIGGER]", tit_sz=20, desc_sz=16)
add_schematic_card(sl3, "Fundamentação Neurocognitiva:", "Ativa o córtex pré-frontal e reduz a inibição linguística inicial.", 0.5, 4.6, 7.0, 2.2, bg_c=TURQUESA_CLARO, border_c=TURQUESA, tag="[NEURO_ANALYSIS]", tit_sz=18, desc_sz=15, align_center_v=True)
img_schematic_box(sl3, 7.8, 1.8, 5.0, 5.0, "[VISUAL_WARMUP_SPEC]")

# SLIDE 4
sl4 = base_slide(4)
add_header(sl4, "Objetivos da Aula", "Desmontagem Pedagógica")
add_schematic_card(sl4, "Metas da Aula:", "• [Objetivo 1]\n• [Objetivo 2]\n• [Objetivo 3]", 0.5, 1.8, 7.0, 5.0, bg_c=BRANCO, border_c=ROXO, tag="[TARGET_CHECKLIST]", tit_sz=18, desc_sz=16)
add_schematic_card(sl4, "Nota de Escrita:", "Nesta aula daremos atenção à escrita para reforçar a memorização de estruturas/vocabulário e reduzir a interferência da língua materna.", 7.8, 1.8, 5.0, 5.0, bg_c=ROXO_CLARO, border_c=ROXO, tag="[WRITING_SPEC]", tit_sz=18, desc_sz=15)

# SLIDE 5
sl5 = base_slide(5)
add_header(sl5, "A Nossa Metodologia", "Arquitetura do Sistema Haas Language")
nodes = [
    ("01", "Fala Ativa", "Conversação inicial sem julgamentos.", ROXO),
    ("02", "Exposição", "Apresentação da gramática.", TURQUESA),
    ("03", "Prática", "Exercícios com suporte.", AZUL),
    ("04", "Conversa", "Aplicação prática em contexto.", ROXO),
    ("05", "Escrita", "Consolidação e memorização.", TURQUESA)
]
top_y = 1.8
for n_num, tit, desc, col in nodes:
    add_schematic_card(sl5, tit, desc, 0.5, top_y, 12.333, 0.9, bg_c=BRANCO, border_c=col, tag=f"[NODE_{n_num}]", tit_sz=15, desc_sz=13, align_center_v=True)
    top_y += 1.05

# SLIDES 6 A 8
for s_num in [6, 7, 8]:
    sl = base_slide(s_num)
    add_header(sl, f"Explicação do Tema — Parte {s_num - 5}")
    add_schematic_card(sl, "Conceito Principal:", f"[Explicação e exemplo central da estrutura - Bloco {s_num - 5}]", 0.5, 1.8, 7.0, 5.0, bg_c=BRANCO, border_c=ROXO, tag="[GRAMMAR_ANATOMY]", tit_sz=20, desc_sz=15)
    add_schematic_card(sl, "Exemplos Práticos:", "[Exemplos adicionais e frases do dia a dia]", 7.8, 1.8, 5.0, 2.3, bg_c=TURQUESA_CLARO, border_c=TURQUESA, tag="[USAGE_EXAMPLES]")
    add_schematic_card(sl, "Notas & Detalhes:", "[Exceções, nuances gramaticais e contexto de uso]", 7.8, 4.5, 5.0, 2.3, bg_c=ROXO_CLARO, border_c=AZUL, tag="[EXCEPTIONS_SPEC]")

# SLIDE 9
sl9 = base_slide(9)
add_header(sl9, "Síntese do Conteúdo & Sugestão Cultural")
add_schematic_card(sl9, "Resumo do Tema:", "[Resumo sintético dos pontos chave ensinados na aula]", 0.5, 1.8, 5.8, 5.0, bg_c=BRANCO, border_c=ROXO, tag="[CORE_SYNTHESIS]", tit_sz=18, desc_sz=15)
add_schematic_card(sl9, "Sugestão Cultural:", "🎵 [Nome da Música / Vídeo recomendado sobre o assunto].\n\nO professor pode rodar o áudio/vídeo ou deixar como recomendação extra para o aluno praticar.", 6.7, 1.8, 6.1, 5.0, bg_c=TURQUESA_CLARO, border_c=TURQUESA, tag="[CULTURAL_MEDIA]", tit_sz=18, desc_sz=15)

# SLIDES 10 E 11
sl10 = base_slide(10)
add_header(sl10, "Prática Guiada — Blocos 1 a 5", "Correção de Vícios")
add_schematic_flashcard(sl10, "01", "[Lacunas contextualizadas com opções]", 0.5, 1.6, 12.333, 1.1, border_c=ROXO)
add_schematic_flashcard(sl10, "02", "[Lacunas contextualizadas com opções]", 0.5, 2.9, 6.0, 1.8, border_c=TURQUESA)
add_schematic_flashcard(sl10, "03", "[Preenchimento direto de verbos]", 6.8, 2.9, 6.0, 1.8, border_c=TURQUESA)
add_schematic_flashcard(sl10, "04", "[Reestruturação / Transformação de frase]", 0.5, 4.9, 6.0, 1.8, border_c=ROXO)
add_schematic_flashcard(sl10, "05", "[Reestruturação / Transformação de frase]", 6.8, 4.9, 6.0, 1.8, border_c=ROXO)

sl11 = base_slide(11)
add_header(sl11, "Prática Guiada — Blocos 6 a 10", "Desafios Práticos")
add_schematic_flashcard(sl11, "06", "[Preposição / Conector correto]", 0.5, 1.6, 12.333, 1.1, border_c=TURQUESA)
add_schematic_flashcard(sl11, "07", "[Escolha da opção correta]", 0.5, 2.9, 6.0, 1.8, border_c=ROXO)
add_schematic_flashcard(sl11, "08", "[Substituição por sinônimo formal]", 6.8, 2.9, 6.0, 1.8, border_c=ROXO)
add_schematic_flashcard(sl11, "09", "[Identificação de erro e correção]", 0.5, 4.9, 6.0, 1.8, border_c=TURQUESA)
add_schematic_flashcard(sl11, "10", "[Tradução de sentido sem uso literal]", 6.8, 4.9, 6.0, 1.8, border_c=TURQUESA)

# SLIDES 12 A 19: CONVERSAÇÃO (SCHEMATIC PIPELINE NODE)
perguntas_config = [
    "[Pergunta principal 1 sobre o tema]", "[Pergunta principal 2 sobre o tema]", "[Pergunta principal 3 sobre o tema]",
    "[Pergunta principal 4 sobre o tema]", "[Pergunta principal 5 sobre o tema]", "[Pergunta principal 6 sobre o tema]",
    "[Pergunta principal 7 sobre o tema]", "[Pergunta principal 8 sobre o tema]"
]

def add_schematic_node_history(sl, history_list):
    count = len(history_list)
    if count == 0: return
    
    gap = 0.12
    avail_h = 5.0 - (gap * (count - 1))
    card_h = avail_h / count
    
    line_x = 7.85
    pipeline_line = sl.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(line_x), Inches(1.8), Pt(2), Inches(4.8))
    pipeline_line.fill.solid(); pipeline_line.fill.fore_color.rgb = TURQUESA; pipeline_line.line.fill.background()
    
    for i, h_txt in enumerate(history_list):
        y_pos = 1.8 + i * (card_h + gap)
        clean_t = h_txt[:30] + "..." if len(h_txt) > 30 else h_txt
        d_sz = 11 if count <= 3 else (10 if count <= 5 else 9)
        border_c = ROXO if i % 2 == 0 else TURQUESA
        bg_col = ROXO_CLARO if i % 2 == 0 else TURQUESA_CLARO
        
        node_dot = sl.shapes.add_shape(MSO_SHAPE.OVAL, Inches(line_x - 0.08), Inches(y_pos + (card_h/2) - 0.08), Inches(0.18), Inches(0.18))
        node_dot.fill.solid(); node_dot.fill.fore_color.rgb = border_c; node_dot.line.fill.background()
        
        add_schematic_card(sl, "", clean_t, 8.1, y_pos, 4.7, card_h, bg_c=bg_col, border_c=border_c, tag=f"[NODE_0{i+1}]", desc_sz=d_sz, align_center_v=True)

for idx in range(8):
    s_num = 12 + idx
    sl = base_slide(s_num)
    add_header(sl, "Conversação", f"Thread de Pergunta [{idx + 1}/8]")
    p_desc = perguntas_config[idx]
    
    if idx == 0:
        add_schematic_card(sl, f"PERGUNTA ATIVA [Q0{idx+1}]:", p_desc, 0.5, 1.8, 12.333, 2.5, bg_c=BRANCO, border_c=ROXO, tag="[ACTIVE_PROMPT_SCHEMATIC]", tit_sz=18, desc_sz=18, align_center_v=True)
        add_schematic_card(sl, "PROMPTS DE APOIO AO PROFESSOR:", "• Para expandir: 'Dê um exemplo recente...'\n• Para concordar/discordar: 'Qual é o outro lado dessa moeda?'", 0.5, 4.5, 12.333, 2.3, bg_c=TURQUESA_CLARO, border_c=TURQUESA, tag="[TEACHER_SUPPORT_LOG]", tit_sz=16, desc_sz=14, align_center_v=True)
    else:
        add_schematic_card(sl, f"PERGUNTA ATIVA [Q0{idx+1}]:", p_desc, 0.5, 1.8, 7.0, 2.5, bg_c=BRANCO, border_c=ROXO, tag=f"[ACTIVE_PROMPT_Q0{idx+1}]", tit_sz=18, desc_sz=16, align_center_v=True)
        add_schematic_card(sl, "PROMPTS DE APOIO:", "• Para expandir: 'Dê um exemplo recente...'\n• Para discordar: 'Qual é o outro lado?'", 0.5, 4.5, 7.0, 2.3, bg_c=TURQUESA_CLARO, border_c=TURQUESA, tag="[TEACHER_SUPPORT_LOG]", tit_sz=15, desc_sz=14, align_center_v=True)
        
        history_list = [perguntas_config[i] for i in range(idx)]
        add_schematic_node_history(sl, history_list)

# SLIDE 20
sl20 = base_slide(20)
add_header(sl20, "Atividade de Escrita", "Consolidação e Memorização")
add_schematic_card(sl20, "Tarefa de Escrita Recomendada:", "[Instrução detalhada da tarefa: escrever um texto, frases ou diálogo]", 0.5, 1.8, 12.333, 2.5, bg_c=BRANCO, border_c=ROXO, tag="[WRITING_TASK_SPEC]", tit_sz=18, desc_sz=15)
add_schematic_card(sl20, "Fundamentação Pedagógica:", "Identificamos que a escrita nesta etapa é fundamental para você fixar as estruturas aprendidas, expandir vocabulário ativo e diminuir a tradução mental automática do seu idioma nativo.", 0.5, 4.5, 12.333, 2.3, bg_c=TURQUESA_CLARO, border_c=TURQUESA, tag="[PEDAGOGICAL_RATIONALE]", desc_sz=14, align_center_v=True)

# SLIDE 21
sl21 = base_slide(21)
add_header(sl21, "Encerramento — Tarefa de Casa", "Ação e Prática Contínua")
add_schematic_card(sl21, "Tarefa de Casa:", "• [Atividade prática 1 para realizar fora da aula]\n• [Atividade de fixação de vocabulário ou escuta]", 0.5, 1.8, 12.333, 5.0, bg_c=BRANCO, border_c=ROXO, tag="[HOMEWORK_ASSIGNMENT]", tit_sz=18, desc_sz=15)

# SLIDE 22
sl22 = base_slide(22)
add_header(sl22, "Encerramento & Feedback Rápido", "Considerações Finais")
add_schematic_card(sl22, "Agradecimento:", "Obrigado pela dedicação na aula de hoje!", 0.5, 1.8, 12.333, 2.2, bg_c=ROXO_CLARO, border_c=ROXO, tag="[CLOSING_SPEC]", tit_sz=20, desc_sz=16)
add_schematic_card(sl22, "Espaço para Feedback Rápido:", "[Espaço destinado para o professor coletar dúvidas rápidas e considerações do aluno]", 0.5, 4.3, 12.333, 2.5, bg_c=TURQUESA_CLARO, border_c=TURQUESA, tag="[FEEDBACK_LOG]", tit_sz=18, desc_sz=15)

out_path_local = "template-haas-modelo-25-data-driven.pptx"
prs.save(out_path_local)

server_dir = "/var/www/haas-frontend-desk-mobile-oficial/public/temp-miniaturas/"
os.makedirs(server_dir, exist_ok=True)
out_path_server = os.path.join(server_dir, "template-haas-modelo-25-data-driven.pptx")

shutil.copyfile(out_path_local, out_path_server)
print(f"SUCCESSFULLY GENERATED DATA-DRIVEN MODEL 25: {out_path_server}")
