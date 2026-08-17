import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

# PALETA OFICIAL HAAS ACADEMY
BRANCO = RGBColor(0xFF, 0xFF, 0xFF)
ROXO = RGBColor(0x4F, 0x35, 0xE2)
AZUL = RGBColor(0x1E, 0x3A, 0x8A)
SLATE = RGBColor(0x33, 0x41, 0x55)
TURQUESA = RGBColor(0x14, 0xC8, 0xB0)
BORDA = RGBColor(0xCB, 0xD5, 0xE1)
FUNDO_CARD = RGBColor(0xF8, 0xFA, 0xFC)
ROXO_CLARO = RGBColor(0xED, 0xE9, 0xFE)
MUTED = RGBColor(0x94, 0xA3, 0xB8)

W, H = Inches(13.333), Inches(7.5)

prs = Presentation()
prs.slide_width = W
prs.slide_height = H

def base_slide(num):
    sl = prs.slides.add_slide(prs.slide_layouts[6])
    sl.background.fill.solid()
    sl.background.fill.fore_color.rgb = BRANCO
    
    # Faixa lateral
    s = sl.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(0.12), H)
    s.fill.solid(); s.fill.fore_color.rgb = ROXO; s.line.fill.background()
    
    # Rodapé
    b = sl.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(7.1), W, Inches(0.4))
    b.fill.solid(); b.fill.fore_color.rgb = ROXO; b.line.fill.background()
    
    # Detalhe Turquesa
    det = sl.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.5), Inches(1.0), Inches(1.2), Pt(4))
    det.fill.solid(); det.fill.fore_color.rgb = TURQUESA; det.line.fill.background()
    
    # Círculo
    c = sl.shapes.add_shape(MSO_SHAPE.OVAL, Inches(12.5), Inches(0.4), Inches(0.3), Inches(0.3))
    c.fill.solid(); c.fill.fore_color.rgb = TURQUESA; c.line.fill.background()
    
    # Texto Rodapé
    tb = sl.shapes.add_textbox(Inches(0.5), Inches(7.15), Inches(8), Inches(0.3))
    p = tb.text_frame.paragraphs[0]
    p.text = "Haas Academia de Idiomas"
    p.font.size = Pt(10); p.font.color.rgb = BRANCO
    
    tb2 = sl.shapes.add_textbox(Inches(11.5), Inches(7.15), Inches(1.3), Inches(0.3))
    p2 = tb2.text_frame.paragraphs[0]
    p2.text = f"Slide {num} / 18"
    p2.font.size = Pt(10); p2.font.color.rgb = BRANCO; p2.alignment = PP_ALIGN.RIGHT
    return sl

def add_header(sl, tit, sub=""):
    tb = sl.shapes.add_textbox(Inches(0.5), Inches(1.2), Inches(12), Inches(0.6))
    tf = tb.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]; p.text = tit
    p.font.size = Pt(28); p.font.color.rgb = ROXO; p.font.bold = True
    
    if sub:
        tb_s = sl.shapes.add_textbox(Inches(0.5), Inches(1.8), Inches(12), Inches(0.4))
        p_s = tb_s.text_frame.paragraphs[0]
        p_s.text = sub
        p_s.font.size = Pt(16); p_s.font.color.rgb = AZUL; p_s.font.bold = True

def add_card(sl, tit, desc, left, top, width, height, bg_color=FUNDO_CARD, border_color=BORDA):
    s = sl.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    s.fill.solid(); s.fill.fore_color.rgb = bg_color
    s.line.color.rgb = border_color; s.line.width = Pt(1)
    tf = s.text_frame; tf.word_wrap = True; tf.margin_left = Inches(0.15); tf.margin_top = Inches(0.15)
    if tit:
        p = tf.paragraphs[0]; p.text = tit; p.font.size = Pt(13); p.font.color.rgb = ROXO; p.font.bold = True
        if desc:
            p2 = tf.add_paragraph(); p2.text = desc; p2.font.size = Pt(11); p2.font.color.rgb = SLATE; p2.space_before = Pt(4)
    else:
        if desc:
            p = tf.paragraphs[0]; p.text = desc; p.font.size = Pt(11); p.font.color.rgb = SLATE

# ==========================================================
# CONSTRUÇÃO EXATA DOS 18 SLIDES
# ==========================================================

# SLIDE 1: CAPA
sl1 = base_slide(1)
add_header(sl1, "[Título da Aula / Tema Principal]", "Bem-vindo(a) à nossa aula!")
add_card(sl1, "Mensagem de Boas-vindas:", "[Mensagem de boas-vindas no idioma que está sendo aprendido]", Inches(0.5), Inches(3.0), Inches(12.3), Inches(2.0))

# SLIDE 2: INTEGRANTES
sl2 = base_slide(2)
add_header(sl2, "Presença & Integrantes", "Lista de Alunos Previstos para esta Turma")
positions = [(0.5, 2.5), (4.7, 2.5), (8.9, 2.5), (0.5, 4.5), (4.7, 4.5), (8.9, 4.5)]
for i, (l, t) in enumerate(positions, 1):
    add_card(sl2, f"Aluno {i}:", "[Nome do Aluno]", Inches(l), Inches(t), Inches(3.9), Inches(1.5))

# SLIDE 3: WARM-UP
sl3 = base_slide(3)
add_header(sl3, "Aquecimento / Warm-up", "Foco 100% na Fala (3 Minutos)")
add_card(sl3, "Pergunta Provocativa Rápida:", "[Insira aqui a pergunta provocativa de aquecimento]", Inches(0.5), Inches(2.6), Inches(12.3), Inches(2.2))
add_card(sl3, "Fundamentação Neurocognitiva:", "Ativa o córtex pré-frontal e reduz a inibição linguística inicial. Todos devem falar nos primeiros 3 minutos.", Inches(0.5), Inches(5.2), Inches(12.3), Inches(1.2), ROXO_CLARO, ROXO)

# SLIDE 4: OBJETIVOS
sl4 = base_slide(4)
add_header(sl4, "Objetivos da Aula", "O que vamos alcançar hoje")
add_card(sl4, "Metas da Aula:", "• [Objetivo 1]\n• [Objetivo 2]\n• [Objetivo 3]", Inches(0.5), Inches(2.5), Inches(12.3), Inches(2.5))
add_card(sl4, "Nota de Escrita:", "Nesta aula daremos atenção à escrita para reforçar a memorização de estruturas/vocabulário e reduzir a interferência da língua materna.", Inches(0.5), Inches(5.3), Inches(12.3), Inches(1.2), FUNDO_CARD, TURQUESA)

# SLIDE 5: METODOLOGIA
sl5 = base_slide(5)
add_header(sl5, "A Nossa Metodologia", "Estrutura Fixa de Aprendizado")
metodologias = [
    ("1. Fala Ativa Primeiro", "Conversação inicial sem julgamentos."),
    ("2. Exposição do Conteúdo", "Apresentação da gramática/tema."),
    ("3. Prática Guiada", "Exercícios com suporte do professor."),
    ("4. Conversação", "Aplicação prática das estruturas."),
    ("5. Escrita (quando necessária)", "Consolidação visual e redução de tradução.")
]
top_m = 2.4
for tit_m, desc_m in metodologias:
    add_card(sl5, tit_m, desc_m, Inches(0.5), Inches(top_m), Inches(12.3), Inches(0.8))
    top_m += 0.9

# SLIDES 6 A 8: EXPOSIÇÃO MULTINÍVEL
for s_num in [6, 7, 8]:
    sl = base_slide(s_num)
    add_header(sl, f"Explicação do Tema — Parte {s_num - 5}", "Estratégia Multinível")
    add_card(sl, "Conteúdo Principal (Nível Inicial / Básico):", f"[Explicação e exemplos simples para os níveis Inicial e Básico - Bloco {s_num - 5}]", Inches(0.5), Inches(2.5), Inches(8.0), Inches(4.0))
    add_card(sl, "Nuances & Exceções\n(Nível Independente / Avançado):", "[Nuances gramaticais detalhadas para alunos avançados]", Inches(8.8), Inches(2.5), Inches(4.0), Inches(4.0), ROXO_CLARO, ROXO)

# SLIDE 9: SÍNTESE + CULTURA
sl9 = base_slide(9)
add_header(sl9, "Síntese do Conteúdo & Sugestão Cultural", "Resumo e Recomendação")
add_card(sl9, "Resumo do Tema:", "[Resumo sintético dos pontos chave ensinados na aula]", Inches(0.5), Inches(2.5), Inches(12.3), Inches(2.2))
add_card(sl9, "Sugestão Cultural (Música ou Vídeo):", "🎵 [Nome da Música / Vídeo recomendado sobre o assunto].\nO professor pode rodar o áudio/vídeo ou deixar como recomendação extra.", Inches(0.5), Inches(5.0), Inches(12.3), Inches(1.5), FUNDO_CARD, AZUL)

# SLIDES 10 E 11: PRÁTICA GUIADA
sl10 = base_slide(10)
add_header(sl10, "Prática Guiada — Exercícios 1 a 5", "Máximo de 10 Exercícios Guiados (Correção de Vícios)")
exs_1_5 = [
    ("Exercício 1:", "[Lacunas contextualizadas com opções] — Inicial / Descoberta"),
    ("Exercício 2:", "[Lacunas contextualizadas com opções] — Básico / Conexão"),
    ("Exercício 3:", "[Preenchimento direto de verbos] — Básico / Conexão"),
    ("Exercício 4:", "[Reestruturação / Transformação] — Intermediário / Autônomo"),
    ("Exercício 5:", "[Reestruturação / Transformação] — Intermediário / Autônomo")
]
top_e = 2.4
for t_e, d_e in exs_1_5:
    add_card(sl10, t_e, d_e, Inches(0.5), Inches(top_e), Inches(12.3), Inches(0.8))
    top_e += 0.9

sl11 = base_slide(11)
add_header(sl11, "Prática Guiada — Exercícios 6 a 10", "Desafios para Níveis Superiores")
exs_6_10 = [
    ("Exercício 6:", "[Preposição / Conector correto] — Intermediário / Autônomo"),
    ("Exercício 7:", "[Escolha de opção correta] — Independente"),
    ("Exercício 8:", "[Substituição por sinônimo formal] — Independente"),
    ("Exercício 9:", "[Identificação de erro e correção] — Avançado / Especialista"),
    ("Exercício 10:", "[Tradução de sentido sem uso literal] — Avançado / Especialista")
]
top_e = 2.4
for t_e, d_e in exs_6_10:
    add_card(sl11, t_e, d_e, Inches(0.5), Inches(top_e), Inches(12.3), Inches(0.8))
    top_e += 0.9

# SLIDES 12 A 15: CONVERSAÇÃO CIRÚRGICA CUMULATIVA
sl12 = base_slide(12)
add_header(sl12, "Conversação Cirúrgica", "Efeito Cumulativo Progressivo")
add_card(sl12, "Pergunta 1 (Revelada):", "[Pergunta principal 1] — Inicial / Descoberta", Inches(0.5), Inches(2.5), Inches(12.3), Inches(2.0), ROXO_CLARO, ROXO)
add_card(sl12, "Sub-prompts Visuais (Apoio):", "• Para expandir: 'Dê um exemplo recente...'\n• Para concordar/discordar: 'Qual é o outro lado dessa moeda?'", Inches(0.5), Inches(4.8), Inches(12.3), Inches(1.5))

sl13 = base_slide(13)
add_header(sl13, "Conversação Cirúrgica", "Efeito Cumulativo Progressivo")
add_card(sl13, "Histórico:", "1. [Pergunta principal 1]", Inches(0.5), Inches(2.4), Inches(12.3), Inches(0.8), FUNDO_CARD, MUTED)
add_card(sl13, "Pergunta 2 (Nova):", "[Pergunta principal 2] — Básico / Conexão", Inches(0.5), Inches(3.4), Inches(12.3), Inches(1.8), ROXO_CLARO, ROXO)
add_card(sl13, "Sub-prompts Visuais (Apoio):", "• Para expandir: 'Dê um exemplo recente...'\n• Para concordar/discordar: 'Qual é o outro lado dessa moeda?'", Inches(0.5), Inches(5.4), Inches(12.3), Inches(1.2))

sl14 = base_slide(14)
add_header(sl14, "Conversação Cirúrgica", "Efeito Cumulativo Progressivo")
add_card(sl14, "Histórico:", "1. [Pergunta 1]\n2. [Pergunta 2]", Inches(0.5), Inches(2.4), Inches(12.3), Inches(1.2), FUNDO_CARD, MUTED)
add_card(sl14, "Pergunta 3 (Nova):", "[Pergunta principal 3] — Intermediário / Autônomo", Inches(0.5), Inches(3.8), Inches(12.3), Inches(1.8), ROXO_CLARO, ROXO)
add_card(sl14, "Sub-prompts Visuais (Apoio):", "• Prompts de expansão para a Pergunta 3", Inches(0.5), Inches(5.8), Inches(12.3), Inches(0.9))

sl15 = base_slide(15)
add_header(sl15, "Conversação Cirúrgica", "Efeito Cumulativo Progressivo")
add_card(sl15, "Histórico Completo:", "1. [Pergunta 1]\n2. [Pergunta 2]\n3. [Pergunta 3]", Inches(0.5), Inches(2.4), Inches(12.3), Inches(1.5), FUNDO_CARD, MUTED)
add_card(sl15, "Pergunta 4 (Nova):", "[Pergunta principal 4] — Avançado / Especialista", Inches(0.5), Inches(4.1), Inches(12.3), Inches(1.6), ROXO_CLARO, ROXO)
add_card(sl15, "Sub-prompts Visuais (Apoio):", "• Prompts de análise crítica para a Pergunta 4", Inches(0.5), Inches(5.9), Inches(12.3), Inches(0.8))

# SLIDE 16: PARTE ESCRITA
sl16 = base_slide(16)
add_header(sl16, "Atividade de Escrita", "Consolidação e Memorização")
add_card(sl16, "Tarefa de Escrita Recomendada:", "[Instrução detalhada da tarefa: escrever um texto, frases ou diálogo dependendo do nível]", Inches(0.5), Inches(2.5), Inches(12.3), Inches(2.5))
add_card(sl16, "Razão Pedagógica (Rodapé Fixo):", "Identificamos que a escrita nesta etapa é fundamental para você fixar as estruturas aprendidas, expandir vocabulário ativo e diminuir a tradução mental automática do seu idioma nativo.", Inches(0.5), Inches(5.2), Inches(12.3), Inches(1.3), ROXO_CLARO, ROXO)

# SLIDE 17: TAREFA DE CASA
sl17 = base_slide(17)
add_header(sl17, "Encerramento — Tarefa de Casa", "Ação e Prática Contínua")
add_card(sl17, "Tarefa de Casa:", "• [Atividade prática 1 para realizar fora da aula]\n• [Atividade de fixação de vocabulário ou escuta]", Inches(0.5), Inches(2.5), Inches(12.3), Inches(3.8))

# SLIDE 18: ENCERRAMENTO & FEEDBACK
sl18 = base_slide(18)
add_header(sl18, "Encerramento & Feedback Rado", "Considerações Finais")
add_card(sl18, "Agradecimento:", "Obrigado pela dedicação na aula de hoje!", Inches(0.5), Inches(2.5), Inches(12.3), Inches(1.5))
add_card(sl18, "Espaço para Feedback Rápido:", "[Espaço destinado para o professor coletar dúvidas rápidas e considerações do aluno]", Inches(0.5), Inches(4.3), Inches(12.3), Inches(2.2), FUNDO_CARD, TURQUESA)

# SALVAR O TEMPLATE
out_path = "/var/www/haas-frontend-desk-mobile-oficial/public/temp-miniaturas/template-haas-modelo.pptx"
prs.save(out_path)
print(f"✅ TEMPLATE MASTER HAAS CRIADO COM SUCESSO EM: {out_path}")
