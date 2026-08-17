import os
import urllib.request
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# PALETA OFICIAL HAAS LANGUAGE - TERMINAL / COMMAND PROMPT EDITION
BG_DARK = RGBColor(0x0F, 0x17, 0x2A)      # Dark Navy Terminal Background (#0F172A)
CARD_DARK = RGBColor(0x1E, 0x29, 0x3B)    # Lighter Terminal Container (#1E293B)
BRANCO = RGBColor(0xFF, 0xFF, 0xFF)
ROXO = RGBColor(0x4F, 0x35, 0xE2)
AZUL = RGBColor(0x1E, 0x3A, 0x8A)
TURQUESA = RGBColor(0x14, 0xC8, 0xB0)     # Prompt Accent Color
SLATE = RGBColor(0x94, 0xA3, 0xB8)      # Terminal Subtitle Color (#94A3B8)
BORDA_DARK = RGBColor(0x33, 0x41, 0x55) # Terminal Line Color (#334155)

LOGO_URL = "https://jdppxfokfhqjudwfwckd.supabase.co/storage/v1/object/public/haas-academy/assignments/Design%20sem%20nome%20(4).png"
LOGO_PATH = "/tmp/haas-logo.png"

if not os.path.exists(LOGO_PATH):
    try:
        urllib.request.urlretrieve(LOGO_URL, LOGO_PATH)
    except Exception:
        pass

W, H = Inches(13.333), Inches(7.5)
TOTAL_SLIDES = 22

prs = Presentation()
prs.slide_width = W
prs.slide_height = H

def base_slide(num, hide_logo=False):
    sl = prs.slides.add_slide(prs.slide_layouts[6])
    sl.background.fill.solid()
    sl.background.fill.fore_color.rgb = BG_DARK

    # Terminal Status Bar (Linha Superior do Terminal)
    top_line = sl.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), W, Inches(0.08))
    top_line.fill.solid(); top_line.fill.fore_color.rgb = TURQUESA; top_line.line.fill.background()

    if os.path.exists(LOGO_PATH) and not hide_logo:
        sl.shapes.add_picture(LOGO_PATH, Inches(10.8), Inches(0.2), height=Inches(0.65))
        
    # Rodapé Command Line
    tb = sl.shapes.add_textbox(Inches(0.5), Inches(7.1), Inches(8), Inches(0.35))
    p = tb.text_frame.paragraphs[0]
    p.text = "haas-cli  •  COMMAND PROMPT EDITION  v2.0"
    p.font.size = Pt(10); p.font.color.rgb = TURQUESA; p.font.bold = True
    
    tb2 = sl.shapes.add_textbox(Inches(10.8), Inches(7.1), Inches(2.0), Inches(0.35))
    p2 = tb2.text_frame.paragraphs[0]
    p2.text = f"SLOT [{num:02d}/{TOTAL_SLIDES:02d}]"
    p2.font.size = Pt(10); p2.font.color.rgb = SLATE; p2.font.bold = True; p2.alignment = PP_ALIGN.RIGHT
    return sl

def add_header(sl, tit, sub=""):
    tb = sl.shapes.add_textbox(Inches(0.5), Inches(0.35), Inches(10), Inches(0.65))
    tf = tb.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]; p.text = f">_ {tit}"
    p.font.size = Pt(28); p.font.color.rgb = BRANCO; p.font.bold = True
    
    if sub:
        tb_s = sl.shapes.add_textbox(Inches(0.5), Inches(0.95), Inches(10), Inches(0.4))
        p_s = tb_s.text_frame.paragraphs[0]
        p_s.text = f"// {sub}"
        p_s.font.size = Pt(16); p_s.font.color.rgb = TURQUESA; p_s.font.bold = False

def add_terminal_card(sl, tit, desc, left, top, width, height, bg_c=CARD_DARK, border_c=BORDA_DARK, tit_c=TURQUESA, desc_c=BRANCO, tit_size=16, desc_size=14, show_dots=True, align_center_v=False):
    """Janela de Terminal de Código com controles (🔴 🟡 🟢)"""
    s = sl.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    s.fill.solid(); s.fill.fore_color.rgb = bg_c
    if border_c:
        s.line.color.rgb = border_c; s.line.width = Pt(1.5)
    else:
        s.line.fill.background()

    # Controles de Janela de Terminal (🔴 🟡 🟢)
    if show_dots and height > Inches(1.2):
        dot_y = top + Inches(0.12)
        colors = [RGBColor(0xEF, 0x44, 0x44), RGBColor(0xF5, 0x9E, 0x0B), RGBColor(0x10, 0xB9, 0x81)]
        for i, dot_c in enumerate(colors):
            dot = sl.shapes.add_shape(MSO_SHAPE.OVAL, left + Inches(0.15 + i * 0.2), dot_y, Inches(0.12), Inches(0.12))
            dot.fill.solid(); dot.fill.fore_color.rgb = dot_c; dot.line.fill.background()

    tf = s.text_frame; tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE if align_center_v else MSO_ANCHOR.TOP
    margin_top_val = Inches(0.32) if (show_dots and height > Inches(1.2)) else Inches(0.18)
    tf.margin_left = Inches(0.25); tf.margin_top = margin_top_val; tf.margin_right = Inches(0.25); tf.margin_bottom = Inches(0.18)
    
    if tit:
        p = tf.paragraphs[0]
        p.text = tit; p.font.size = Pt(tit_size); p.font.color.rgb = tit_c; p.font.bold = True
        if desc:
            p2 = tf.add_paragraph()
            p2.text = desc; p2.font.size = Pt(desc_size); p2.font.color.rgb = desc_c; p2.space_before = Pt(6)
    else:
        if desc:
            p = tf.paragraphs[0]
            p.text = desc; p.font.size = Pt(desc_size); p.font.color.rgb = desc_c

def img_terminal_box(sl, left, top, width, height, label="[assets/media_file.png]"):
    s = sl.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    s.fill.solid(); s.fill.fore_color.rgb = CARD_DARK
    s.line.color.rgb = ROXO; s.line.width = Pt(1.5)

    # Controles de Janela
    colors = [RGBColor(0xEF, 0x44, 0x44), RGBColor(0xF5, 0x9E, 0x0B), RGBColor(0x10, 0xB9, 0x81)]
    for i, dot_c in enumerate(colors):
        dot = sl.shapes.add_shape(MSO_SHAPE.OVAL, left + Inches(0.15 + i * 0.2), top + Inches(0.12), Inches(0.12), Inches(0.12))
        dot.fill.solid(); dot.fill.fore_color.rgb = dot_c; dot.line.fill.background()

    tf = s.text_frame; tf.word_wrap = True; tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER
    p = tf.paragraphs[0]; p.text = f"🖼️ {label}"; p.font.size = Pt(16); p.font.color.rgb = TURQUESA; p.font.bold = True

def add_terminal_flashcard(sl, num_str, desc, left, top, width, height):
    s = sl.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    s.fill.solid(); s.fill.fore_color.rgb = CARD_DARK; s.line.color.rgb = ROXO; s.line.width = Pt(1.5)
    
    chip = sl.shapes.add_shape(MSO_SHAPE.RECTANGLE, left + Inches(0.12), top + Inches(0.12), Inches(0.8), Inches(0.35))
    chip.fill.solid(); chip.fill.fore_color.rgb = ROXO; chip.line.fill.background()
    tf_chip = chip.text_frame; tf_chip.vertical_anchor = MSO_ANCHOR.MIDDLE
    p_chip = tf_chip.paragraphs[0]; p_chip.text = f"${num_str}"; p_chip.alignment = PP_ALIGN.CENTER
    p_chip.font.size = Pt(11); p_chip.font.color.rgb = TURQUESA; p_chip.font.bold = True
    
    tf = s.text_frame; tf.word_wrap = True; tf.vertical_anchor = MSO_ANCHOR.TOP
    tf.margin_left = Inches(1.0); tf.margin_top = Inches(0.12); tf.margin_right = Inches(0.15); tf.margin_bottom = Inches(0.1)
    p = tf.paragraphs[0]; p.text = desc; p.font.size = Pt(14); p.font.color.rgb = BRANCO; p.font.bold = False

# ==========================================================
# CONSTRUÇÃO DOS 22 SLIDES - MODELO 17 (TERMINAL CLI)
# ==========================================================

# SLIDE 1: CAPA TERMINAL
sl1 = base_slide(1, hide_logo=True)
add_terminal_card(sl1, "$ init_haas_lesson.sh", "[Título da Aula / Tema Principal]\n\n// Subtítulo / Descrição da Aula\n\n>_ SYSTEM_DATE: [DD/MM/AAAA]", Inches(0.5), Inches(1.5), Inches(7.0), Inches(5.2), border_c=ROXO, tit_c=TURQUESA, desc_c=BRANCO, tit_size=20, desc_size=18)
img_terminal_box(sl1, Inches(7.8), Inches(1.5), Inches(5.0), Inches(5.2), "assets/cover_hero.png")

# SLIDE 2: PRESENÇA
sl2 = base_slide(2)
add_header(sl2, "SESSION_USERS", "Lista de Participantes")
add_terminal_card(sl2, "$ cat /etc/teachers.cfg", "Professor(a): [Nome do(a) Professor(a)]", Inches(0.5), Inches(1.5), Inches(12.3), Inches(1.8), border_c=ROXO, tit_c=TURQUESA, desc_c=BRANCO)
add_terminal_card(sl2, "$ cat /etc/students.cfg", "Aluno(s): [Nome do Aluno / Lista da Turma]", Inches(0.5), Inches(3.5), Inches(12.3), Inches(3.2), border_c=TURQUESA, tit_c=TURQUESA, desc_c=BRANCO)

# SLIDE 3: WARM-UP
sl3 = base_slide(3)
add_header(sl3, "WARMUP_THREAD", "Foco 100% na Fala (3 Minutos)")
add_terminal_card(sl3, ">_ PROMPT_QUESTION:", "[Insira aqui a pergunta provocativa de aquecimento]", Inches(0.5), Inches(1.5), Inches(7.0), Inches(2.6), border_c=ROXO, tit_c=TURQUESA, desc_c=BRANCO, tit_size=18, desc_size=16)
add_terminal_card(sl3, "// NEUROSCIENCE_LOG:", "Ativa o córtex pré-frontal e reduz a inibição linguística inicial.", Inches(0.5), Inches(4.3), Inches(7.0), Inches(2.4), border_c=TURQUESA, tit_c=SLATE, desc_c=BRANCO, align_center_v=True)
img_terminal_box(sl3, Inches(7.8), Inches(1.5), Inches(5.0), Inches(5.2), "assets/warmup_media.png")

# SLIDE 4: OBJETIVOS
sl4 = base_slide(4)
add_header(sl4, "LESSON_OBJECTIVES", "Metas da Sessão")
add_terminal_card(sl4, ">_ TASKS_TO_COMPLETE:", "[x] [Objetivo 1]\n[x] [Objetivo 2]\n[x] [Objetivo 3]", Inches(0.5), Inches(1.5), Inches(7.5), Inches(5.2), border_c=ROXO, tit_c=TURQUESA, desc_c=BRANCO)
add_terminal_card(sl4, "// WRITING_FOCUS_NOTE:", "Nesta aula daremos atenção à escrita para reforçar a memorização de estruturas e reduzir a interferência da língua materna.", Inches(8.3), Inches(1.5), Inches(4.5), Inches(5.2), border_c=TURQUESA, tit_c=SLATE, desc_c=BRANCO)

# SLIDE 5: METODOLOGIA (CLI PIPELINE 5 BLOCOS)
sl5 = base_slide(5)
add_header(sl5, "HAAS_METHODOLOGY_PIPELINE", "Fluxo de Aprendizado")
met_logs = [
    ("[01/05] FALA_ATIVA", "Conversação inicial sem julgamentos.", ROXO),
    ("[02/05] EXPOSICAO", "Apresentação leve da gramática.", TURQUESA),
    ("[03/05] PRATICA_GUIADA", "Exercícios com suporte do professor.", ROXO),
    ("[04/05] CONVERSACAO", "Aplicação prática das estruturas.", TURQUESA),
    ("[05/05] ESCRITA", "Consolidação visual e memorização.", ROXO)
]
top_y = 1.5
for tag, desc, col in met_logs:
    add_terminal_card(sl5, tag, desc, Inches(0.5), Inches(top_y), Inches(12.3), Inches(0.95), border_c=col, tit_c=col, desc_c=BRANCO, show_dots=False)
    top_y += 1.08

# SLIDES 6 A 8: EXPOSIÇÃO
for s_num in [6, 7, 8]:
    sl = base_slide(s_num)
    add_header(sl, f"GRAMMAR_EXPOSITION_0{s_num - 5}")
    add_terminal_card(sl, ">_ CORE_CONCEPT:", f"[Explicação e exemplo central da estrutura - Bloco {s_num - 5}]", Inches(0.5), Inches(1.5), Inches(7.5), Inches(5.2), border_c=ROXO, tit_size=20, desc_size=15)
    add_terminal_card(sl, "// CODE_EXAMPLES:", "[Exemplos adicionais e frases do dia a dia]", Inches(8.3), Inches(1.5), Inches(4.5), Inches(2.5), border_c=TURQUESA, tit_c=SLATE)
    add_terminal_card(sl, "// SYSTEM_EXCEPTIONS:", "[Exceções, nuances gramaticais e contexto]", Inches(8.3), Inches(4.2), Inches(4.5), Inches(2.5), border_c=ROXO, tit_c=SLATE)

# SLIDE 9: SÍNTESE + CULTURA
sl9 = base_slide(9)
add_header(sl9, "LESSON_SUMMARY_AND_MEDIA")
add_terminal_card(sl9, ">_ SUMMARY_LOG:", "[Resumo sintético dos pontos chave ensinados na aula]", Inches(0.5), Inches(1.5), Inches(6.0), Inches(5.2), border_c=ROXO)
add_terminal_card(sl9, "🎵 CULTURAL_RECOMMENDATION:", "[Nome da Música / Vídeo recomendado sobre o assunto].\n\nO professor pode rodar o áudio/vídeo ou deixar como recomendação extra para o aluno praticar.", Inches(6.8), Inches(1.5), Inches(6.0), Inches(5.3), border_c=TURQUESA)

# SLIDES 10 E 11: PRÁTICA GUIADA
sl10 = base_slide(10)
add_header(sl10, "PRACTICE_TASKS_01_05", "Correção de Vícios")
add_terminal_flashcard(sl10, "01", "[Lacunas contextualizadas com opções]", Inches(0.5), Inches(1.5), Inches(12.3), Inches(1.2))
add_terminal_flashcard(sl10, "02", "[Lacunas contextualizadas com opções]", Inches(0.5), Inches(2.9), Inches(6.0), Inches(1.8))
add_terminal_flashcard(sl10, "03", "[Preenchimento direto de verbos]", Inches(6.8), Inches(2.9), Inches(6.0), Inches(1.8))
add_terminal_flashcard(sl10, "04", "[Reestruturação / Transformação de frase]", Inches(0.5), Inches(4.9), Inches(6.0), Inches(1.9))
add_terminal_flashcard(sl10, "05", "[Reestruturação / Transformação de frase]", Inches(6.8), Inches(4.9), Inches(6.0), Inches(1.9))

sl11 = base_slide(11)
add_header(sl11, "PRACTICE_TASKS_06_10", "Desafios Práticos")
add_terminal_flashcard(sl11, "06", "[Preposição / Conector correto]", Inches(0.5), Inches(1.5), Inches(12.3), Inches(1.2))
add_terminal_flashcard(sl11, "07", "[Escolha da opção correta]", Inches(0.5), Inches(2.9), Inches(6.0), Inches(1.8))
add_terminal_flashcard(sl11, "08", "[Substituição por sinônimo formal]", Inches(6.8), Inches(2.9), Inches(6.0), Inches(1.8))
add_terminal_flashcard(sl11, "09", "[Identificação de erro e correção]", Inches(0.5), Inches(4.9), Inches(6.0), Inches(1.9))
add_terminal_flashcard(sl11, "10", "[Tradução de sentido sem uso literal]", Inches(6.8), Inches(4.9), Inches(6.0), Inches(1.9))

# ==========================================================
# SLIDES 12 A 19: CONVERSAÇÃO (INOVAÇÃO TERMINAL TREE INÉDITA!)
# ESQUERDA: ÁRVORE DE DIRETÓRIOS (/root/history/)
# DIREITA: EDITOR DE CÓDIGO DA PERGUNTA ATIVA + NOTA DE APOIO
# ==========================================================
perguntas_config = [
    "[Pergunta principal 1 sobre o tema]", "[Pergunta principal 2 sobre o tema]", "[Pergunta principal 3 sobre o tema]",
    "[Pergunta principal 4 sobre o tema]", "[Pergunta principal 5 sobre o tema]", "[Pergunta principal 6 sobre o tema]",
    "[Pergunta principal 7 sobre o tema]", "[Pergunta principal 8 sobre o tema]"
]

def add_terminal_directory_tree(sl, history_list):
    """Exibe o histórico de perguntas como arquivos de um diretório de terminal na coluna da esquerda"""
    count = len(history_list)
    if count == 0: return
    gap = 0.12
    avail_h = 5.2 - (gap * (count - 1))
    card_h = avail_h / count
    
    for i, h_txt in enumerate(history_list):
        y_pos = 1.5 + i * (card_h + gap)
        clean_t = h_txt[:28] + "..." if len(h_txt) > 28 else h_txt
        d_sz = 11 if count <= 3 else (10 if count <= 5 else 9)
        
        # Arquivo de histórico no estilo Terminal
        add_terminal_card(sl, f"📄 question_0{i+1}.sh", clean_t, Inches(0.5), Inches(y_pos), Inches(4.3), Inches(card_h), border_c=BORDA_DARK, tit_c=SLATE, desc_c=BRANCO, tit_size=11, desc_size=d_sz, show_dots=False, align_center_v=True)

for idx in range(8):
    s_num = 12 + idx
    sl = base_slide(s_num)
    add_header(sl, "SPEAKING_CONS_THREAD", f"Running Process [{idx + 1}/8]")
    
    p_desc = perguntas_config[idx]
    
    if idx == 0:
        # Pergunta 1 Ampla Centralizada no Editor
        add_terminal_card(sl, f">_ ACTIVE_QUESTION_0{idx+1}:", p_desc, Inches(0.5), Inches(1.5), Inches(12.3), Inches(2.7), border_c=ROXO, tit_c=TURQUESA, desc_c=BRANCO, tit_size=20, desc_size=18, align_center_v=True)
        add_terminal_card(sl, "// TEACHER_SUPPORT_NOTES:", "• Para expandir: 'Dê um exemplo recente...'\n• Para concordar/discordar: 'Qual é o outro lado dessa moeda?'", Inches(0.5), Inches(4.4), Inches(12.3), Inches(2.3), border_c=TURQUESA, tit_c=SLATE, desc_c=BRANCO, tit_size=16, desc_size=14, align_center_v=True)
    else:
        # Coluna da ESQUERDA: Árvore de Diretórios de Perguntas Anteriores (/root/history/)
        history_list = [perguntas_config[i] for i in range(idx)]
        add_terminal_directory_tree(sl, history_list)

        # Coluna da DIREITA: Editor de Código da Pergunta Ativa + Suporte ao Professor (TRAVADOS E FIXOS!)
        add_terminal_card(sl, f">_ ACTIVE_QUESTION_0{idx+1}:", p_desc, Inches(5.1), Inches(1.5), Inches(7.7), Inches(2.7), border_c=ROXO, tit_c=TURQUESA, desc_c=BRANCO, tit_size=18, desc_size=16, align_center_v=True)
        add_terminal_card(sl, "// TEACHER_SUPPORT_NOTES:", "• Para expandir: 'Dê um exemplo recente...'\n• Para discordar: 'Qual é o outro lado?'", Inches(5.1), Inches(4.4), Inches(7.7), Inches(2.3), border_c=TURQUESA, tit_c=SLATE, desc_c=BRANCO, tit_size=15, desc_size=14, align_center_v=True)

# SLIDE 20: ESCRITA
sl20 = base_slide(20)
add_header(sl20, "WRITING_TASK_EXECUTION")
add_terminal_card(sl20, ">_ RECOMMENDED_TASK:", "[Instrução detalhada da tarefa: escrever um texto, frases ou diálogo]", Inches(0.5), Inches(1.5), Inches(12.3), Inches(2.8), border_c=ROXO)
add_terminal_card(sl20, "// PEDAGOGICAL_RATIONALE:", "Identificamos que a escrita nesta etapa é fundamental para você fixar as estruturas aprendidas, expandir vocabulário ativo e diminuir a tradução mental automática do seu idioma nativo.", Inches(0.5), Inches(4.5), Inches(12.3), Inches(2.3), border_c=TURQUESA, tit_c=SLATE, align_center_v=True)

# SLIDE 21: TAREFA DE CASA
sl21 = base_slide(21)
add_header(sl21, "HOMEWORK_ASSIGNMENT")
add_terminal_card(sl21, ">_ CONTINUOUS_PRACTICE:", "• [Atividade prática 1 para realizar fora da aula]\n• [Atividade de fixação de vocabulário ou escuta]", Inches(0.5), Inches(1.5), Inches(12.3), Inches(5.2), border_c=ROXO)

# SLIDE 22: ENCERRAMENTO & FEEDBACK
sl22 = base_slide(22)
add_header(sl22, "SESSION_TERMINATED", "Considerações Finais")
add_terminal_card(sl22, ">_ THANK_YOU:", "Obrigado pela dedicação na aula de hoje!", Inches(0.5), Inches(1.5), Inches(12.3), Inches(2.0), border_c=TURQUESA, tit_size=20, desc_size=16)
add_terminal_card(sl22, "// QUICK_FEEDBACK_LOG:", "[Espaço destinado para o professor coletar dúvidas rápidas e considerações do aluno]", Inches(0.5), Inches(3.7), Inches(12.3), Inches(3.0), border_c=ROXO, tit_c=SLATE, tit_size=18, desc_size=15)

# SALVAR E COPIAR
out_local = "template-haas-modelo-17-terminal.pptx"
prs.save(out_local)

server_dir = "/var/www/haas-frontend-desk-mobile-oficial/public/temp-miniaturas/"
os.makedirs(server_dir, exist_ok=True)
out_server = os.path.join(server_dir, "template-haas-modelo-17-terminal.pptx")

import shutil
shutil.copyfile(out_local, out_server)
print(f"🚀 MODELO 17 TERMINAL GERADO EM: {out_server}")
