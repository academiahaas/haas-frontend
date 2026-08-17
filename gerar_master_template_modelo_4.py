import os
import urllib.request
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# PALETA OFICIAL HAAS LANGUAGE
BRANCO = RGBColor(0xFF, 0xFF, 0xFF)
ROXO = RGBColor(0x4F, 0x35, 0xE2)
AZUL = RGBColor(0x1E, 0x3A, 0x8A)
SLATE = RGBColor(0x33, 0x41, 0x55)
TURQUESA = RGBColor(0x14, 0xC8, 0xB0)
BORDA = RGBColor(0xCB, 0xD5, 0xE1)
FUNDO_CARD = RGBColor(0xF8, 0xFA, 0xFC)
ROXO_CLARO = RGBColor(0xED, 0xE9, 0xFE)
MUTED = RGBColor(0x94, 0xA3, 0xB8)

LOGO_URL = "https://jdppxfokfhqjudwfwckd.supabase.co/storage/v1/object/public/haas-academy/assignments/Design%20sem%20nome%20(4).png"
LOGO_PATH = "/tmp/haas-logo.png"

if not os.path.exists(LOGO_PATH):
    try:
        urllib.request.urlretrieve(LOGO_URL, LOGO_PATH)
    except Exception:
        pass

W, H = Inches(13.333), Inches(7.5)
TOTAL_SLIDES = 22

prs = Presentation()
prs.slide_width = W
prs.slide_height = H

def base_slide(num):
    sl = prs.slides.add_slide(prs.slide_layouts[6])
    sl.background.fill.solid()
    sl.background.fill.fore_color.rgb = BRANCO
    
    # MOLDURA MODELO 4: Barra Dupla Lateral Curva + Header Clean
    bar1 = sl.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.3), Inches(0.3), Inches(0.1), Inches(6.9))
    bar1.fill.solid(); bar1.fill.fore_color.rgb = ROXO; bar1.line.fill.background()
    
    # Logo Haas Language
    if os.path.exists(LOGO_PATH):
        sl.shapes.add_picture(LOGO_PATH, Inches(10.8), Inches(0.25), height=Inches(0.65))
        
    # Rodapé Flutuante Bento
    tb = sl.shapes.add_textbox(Inches(0.6), Inches(7.05), Inches(8), Inches(0.35))
    p = tb.text_frame.paragraphs[0]
    p.text = "Haas Language  •  Dynamic Studio"
    p.font.size = Pt(11); p.font.color.rgb = AZUL; p.font.bold = True
    
    tb2 = sl.shapes.add_textbox(Inches(10.8), Inches(7.05), Inches(2.0), Inches(0.35))
    p2 = tb2.text_frame.paragraphs[0]
    p2.text = f"{num} / {TOTAL_SLIDES}"
    p2.font.size = Pt(11); p2.font.color.rgb = ROXO; p2.font.bold = True; p2.alignment = PP_ALIGN.RIGHT
    return sl

def add_header(sl, tit, sub=""):
    tb = sl.shapes.add_textbox(Inches(0.6), Inches(0.35), Inches(10), Inches(0.65))
    tf = tb.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]; p.text = tit
    p.font.size = Pt(30); p.font.color.rgb = ROXO; p.font.bold = True
    
    if sub:
        tb_s = sl.shapes.add_textbox(Inches(0.6), Inches(1.0), Inches(10), Inches(0.4))
        p_s = tb_s.text_frame.paragraphs[0]
        p_s.text = sub
        p_s.font.size = Pt(17); p_s.font.color.rgb = AZUL; p_s.font.bold = True

def add_bento_block(sl, tit, desc, left, top, width, height, bg_color=FUNDO_CARD, border_color=BORDA, tit_color=ROXO, tit_size=16, desc_size=14):
    """Container Bento Grid de Canto Arredondado"""
    s = sl.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    s.fill.solid(); s.fill.fore_color.rgb = bg_color
    s.line.color.rgb = border_color; s.line.width = Pt(1)
    
    tf = s.text_frame; tf.word_wrap = True; tf.vertical_anchor = MSO_ANCHOR.TOP
    tf.margin_left = Inches(0.2); tf.margin_top = Inches(0.2); tf.margin_right = Inches(0.2); tf.margin_bottom = Inches(0.2)
    
    if tit:
        p = tf.paragraphs[0]; p.text = tit; p.font.size = Pt(tit_size); p.font.color.rgb = tit_color; p.font.bold = True
        if desc:
            p2 = tf.add_paragraph(); p2.text = desc; p2.font.size = Pt(desc_size); p2.font.color.rgb = SLATE; p2.space_before = Pt(6)
    else:
        if desc:
            p = tf.paragraphs[0]; p.text = desc; p.font.size = Pt(desc_size); p.font.color.rgb = SLATE

def add_quiz_card(sl, num_str, tit, desc, left, top, width, height, accent_color=ROXO):
    """Card de Exercício com Badge Numérico Estilo App"""
    # Card Base
    s = sl.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    s.fill.solid(); s.fill.fore_color.rgb = FUNDO_CARD
    s.line.color.rgb = BORDA; s.line.width = Pt(1)
    
    # Badge Numérico Circular
    circle = sl.shapes.add_shape(MSO_SHAPE.OVAL, left + Inches(0.15), top + Inches(0.15), Inches(0.45), Inches(0.45))
    circle.fill.solid(); circle.fill.fore_color.rgb = accent_color; circle.line.fill.background()
    tf_c = circle.text_frame; tf_c.vertical_anchor = MSO_ANCHOR.MIDDLE
    p_c = tf_c.paragraphs[0]; p_c.text = num_str; p_c.alignment = PP_ALIGN.CENTER
    p_c.font.size = Pt(12); p_c.font.color.rgb = BRANCO; p_c.font.bold = True
    
    # Texto
    tf = s.text_frame; tf.word_wrap = True; tf.vertical_anchor = MSO_ANCHOR.TOP
    tf.margin_left = Inches(0.75); tf.margin_top = Inches(0.15); tf.margin_right = Inches(0.15)
    p = tf.paragraphs[0]; p.text = tit; p.font.size = Pt(14); p.font.color.rgb = ROXO; p.font.bold = True
    if desc:
        p2 = tf.add_paragraph(); p2.text = desc; p2.font.size = Pt(13); p2.font.color.rgb = SLATE; p2.space_before = Pt(4)

def img_bento(sl, left, top, width, height, label="[Espaço de Imagem]"):
    s = sl.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    s.fill.solid(); s.fill.fore_color.rgb = ROXO_CLARO
    s.line.color.rgb = TURQUESA; s.line.width = Pt(1.5)
    tf = s.text_frame; tf.word_wrap = True; tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER
    p = tf.paragraphs[0]; p.text = label; p.font.size = Pt(15); p.font.color.rgb = ROXO; p.font.bold = True

# ==========================================================
# CONSTRUÇÃO DOS 22 SLIDES DO MODELO 4 (BENTO DYNAMIC)
# ==========================================================

# SLIDE 1: CAPA HERO BENTO
sl1 = base_slide(1)
add_bento_block(sl1, "[Título da Aula / Tema Principal]", "[Subtítulo / Descrição da Aula]", Inches(0.6), Inches(1.5), Inches(6.5), Inches(5.3), bg_color=ROXO_CLARO, border_color=ROXO, tit_size=36, desc_size=20)
img_bento(sl1, Inches(7.3), Inches(1.5), Inches(5.6), Inches(5.3), "[Imagem Hero da Capa]")

# SLIDE 2: PRESENÇA & INTEGRANTES (Mosaico Bento)
sl2 = base_slide(2)
add_header(sl2, "Presença & Integrantes")
add_bento_block(sl2, "Professor(a):", "[Nome do(a) Professor(a)]", Inches(0.6), Inches(1.6), Inches(4.0), Inches(5.2), bg_color=ROXO_CLARO, border_color=ROXO, tit_size=18, desc_size=16)
add_bento_block(sl2, "Aluno(s):", "[Nome do Aluno / Lista da Turma]", Inches(4.8), Inches(1.6), Inches(8.1), Inches(5.2), bg_color=FUNDO_CARD, border_color=BORDA, tit_size=18, desc_size=16)

# SLIDE 3: WARM-UP (Split Bento)
sl3 = base_slide(3)
add_header(sl3, "Aquecimento / Warm-up", "Foco 100% na Fala (3 Minutos)")
add_bento_block(sl3, "Pergunta Provocativa Rápida:", "[Insira aqui a pergunta provocativa de aquecimento]", Inches(0.6), Inches(1.6), Inches(6.5), Inches(3.2), bg_color=ROXO_CLARO, border_color=ROXO, tit_size=20, desc_size=16)
add_bento_block(sl3, "Fundamentação Neurocognitiva:", "Ativa o córtex pré-frontal e reduz a inibição linguística inicial.", Inches(0.6), Inches(5.0), Inches(6.5), Inches(1.8), bg_color=FUNDO_CARD, border_color=TURQUESA, tit_size=16, desc_size=14)
img_bento(sl3, Inches(7.3), Inches(1.6), Inches(5.6), Inches(5.2), "[Imagem de Aquecimento]")

# SLIDE 4: OBJETIVOS
sl4 = base_slide(4)
add_header(sl4, "Objetivos da Aula")
add_bento_block(sl4, "Metas da Aula:", "• [Objetivo 1]\n• [Objetivo 2]\n• [Objetivo 3]", Inches(0.6), Inches(1.6), Inches(12.3), Inches(3.2), tit_size=18, desc_size=16)
add_bento_block(sl4, "Nota de Escrita:", "Nesta aula daremos atenção à escrita para reforçar a memorização de estruturas/vocabulário e reduzir a interferência da língua materna.", Inches(0.6), Inches(5.0), Inches(12.3), Inches(1.8), bg_color=ROXO_CLARO, border_color=ROXO, tit_size=16, desc_size=14)

# SLIDE 5: METODOLOGIA (Bento Cards Escalonados)
sl5 = base_slide(5)
add_header(sl5, "A Nossa Metodologia", "Estrutura Fixa de Aprendizado Haas Language")
metodologias = [
    ("1. Fala Ativa", "Conversação inicial sem julgamentos."),
    ("2. Exposição", "Apresentação leve da gramática."),
    ("3. Prática Guiada", "Exercícios com suporte do professor."),
    ("4. Conversação", "Aplicação prática das estruturas."),
    ("5. Escrita", "Consolidação e redução de tradução.")
]
left_m = 0.6
for i, (tit_m, desc_m) in enumerate(metodologias):
    bg_col = ROXO_CLARO if i % 2 == 0 else FUNDO_CARD
    accent_b = ROXO if i % 2 == 0 else TURQUESA
    add_bento_block(sl5, tit_m, desc_m, Inches(left_m), Inches(1.6), Inches(2.25), Inches(5.2), bg_color=bg_col, border_color=accent_b, tit_size=16, desc_size=13)
    left_m += 2.45

# SLIDES 6 A 8: EXPOSIÇÃO (Layout Mosaico Bento)
for s_num in [6, 7, 8]:
    sl = base_slide(s_num)
    add_header(sl, f"Explicação do Tema — Parte {s_num - 5}")
    add_bento_block(sl, "Conceito Principal:", f"[Explicação e exemplo central da estrutura - Bloco {s_num - 5}]", Inches(0.6), Inches(1.6), Inches(12.3), Inches(2.5), bg_color=ROXO_CLARO, border_color=ROXO, tit_size=18, desc_size=15)
    add_bento_block(sl, "Exemplos Práticos:", "[Exemplos adicionais e frases do dia a dia]", Inches(0.6), Inches(4.3), Inches(6.0), Inches(2.5), tit_size=16, desc_size=14)
    add_bento_block(sl, "Notas & Detalhes:", "[Exceções, nuances gramaticais e contexto]", Inches(6.9), Inches(4.3), Inches(6.0), Inches(2.5), border_color=TURQUESA, tit_size=16, desc_size=14)

# SLIDE 9: SÍNTESE + CULTURA
sl9 = base_slide(9)
add_header(sl9, "Síntese & Recomendação Cultural")
add_bento_block(sl9, "Resumo do Tema:", "[Resumo sintético dos pontos chave ensinados na aula]", Inches(0.6), Inches(1.6), Inches(12.3), Inches(2.5), tit_size=18, desc_size=15)
add_bento_block(sl9, "Sugestão Cultural:", "🎵 [Nome da Música / Vídeo recomendado sobre o assunto].\nO professor pode rodar o áudio/vídeo ou deixar como recomendação extra.", Inches(0.6), Inches(4.3), Inches(12.3), Inches(2.5), bg_color=ROXO_CLARO, border_color=ROXO, tit_size=16, desc_size=14)

# SLIDES 10 E 11: PRÁTICA GUIADA (Quiz Cards Numéricos)
sl10 = base_slide(10)
add_header(sl10, "Prática Guiada — Exercícios 1 a 5", "Correção de Vícios")
exs_1_5 = [
    ("1", "Exercício 1", "[Lacunas contextualizadas com opções]"),
    ("2", "Exercício 2", "[Lacunas contextualizadas com opções]"),
    ("3", "Exercício 3", "[Preenchimento direto de verbos]"),
    ("4", "Exercício 4", "[Reestruturação / Transformação de frase]"),
    ("5", "Exercício 5", "[Reestruturação / Transformação de frase]")
]
for idx, (num_s, t_e, d_e) in enumerate(exs_1_5):
    col = idx % 2; row = idx // 2
    l = 0.6 if col == 0 else 6.9
    w = 6.0 if (idx < 4) else 12.3
    if idx == 4: l = 0.6
    t = 1.6 + row * 1.7
    add_quiz_card(sl10, num_s, t_e, d_e, Inches(l), Inches(t), Inches(w), Inches(1.5), accent_color=ROXO)

sl11 = base_slide(11)
add_header(sl11, "Prática Guiada — Exercícios 6 a 10", "Desafios Práticos")
exs_6_10 = [
    ("6", "Exercício 6", "[Preposição / Conector correto]"),
    ("7", "Exercício 7", "[Escolha da opção correta]"),
    ("8", "Exercício 8", "[Substituição por sinônimo formal]"),
    ("9", "Exercício 9", "[Identificação de erro e correção]"),
    ("10", "Exercício 10", "[Tradução de sentido sem uso literal]")
]
for idx, (num_s, t_e, d_e) in enumerate(exs_6_10):
    col = idx % 2; row = idx // 2
    l = 0.6 if col == 0 else 6.9
    w = 6.0 if (idx < 4) else 12.3
    if idx == 4: l = 0.6
    t = 1.6 + row * 1.7
    add_quiz_card(sl11, num_s, t_e, d_e, Inches(l), Inches(t), Inches(w), Inches(1.5), accent_color=TURQUESA)

# SLIDES 12 A 19: CONVERSAÇÃO (Spotlight Active Card)
perguntas_config = [
    "[Pergunta principal 1 sobre o tema]",
    "[Pergunta principal 2 sobre o tema]",
    "[Pergunta principal 3 sobre o tema]",
    "[Pergunta principal 4 sobre o tema]",
    "[Pergunta principal 5 sobre o tema]",
    "[Pergunta principal 6 sobre o tema]",
    "[Pergunta principal 7 sobre o tema]",
    "[Pergunta principal 8 sobre o tema]"
]

for idx in range(8):
    s_num = 12 + idx
    sl = base_slide(s_num)
    add_header(sl, "Conversação")
    
    if idx == 0:
        p_desc = perguntas_config[0]
        add_bento_block(sl, "", p_desc, Inches(0.6), Inches(1.6), Inches(12.3), Inches(2.7), bg_color=ROXO_CLARO, border_color=ROXO, desc_size=18)
        add_bento_block(sl, "", "• Para expandir: 'Dê um exemplo recente...'\n• Para concordar/discordar: 'Qual é o outro lado dessa moeda?'", Inches(0.6), Inches(4.5), Inches(12.3), Inches(2.3), desc_size=15)
    else:
        # Capsules Flutuantes de Histórico
        history_list = [perguntas_config[i] for i in range(idx)]
        count = len(history_list)
        gap = 0.15; avail_w = 12.3 - (gap * (count - 1))
        pill_w = avail_w / count
        
        for h_idx, h_txt in enumerate(history_list):
            l = 0.6 + h_idx * (pill_w + gap)
            clean_t = h_txt[:22] + "..." if len(h_txt) > 22 else h_txt
            add_bento_block(sl, "", clean_t, Inches(l), Inches(1.6), Inches(pill_w), Inches(0.6), bg_color=FUNDO_CARD, border_color=MUTED, desc_size=10)
        
        p_desc = perguntas_config[idx]
        add_bento_block(sl, "", p_desc, Inches(0.6), Inches(2.4), Inches(12.3), Inches(2.4), bg_color=ROXO_CLARO, border_color=ROXO, desc_size=17)
        add_bento_block(sl, "", "• Para expandir: 'Dê um exemplo recente...'\n• Para concordar/discordar: 'Qual é o outro lado dessa moeda?'", Inches(0.6), Inches(5.0), Inches(12.3), Inches(1.8), desc_size=14)

# SLIDE 20: ESCRITA
sl20 = base_slide(20)
add_header(sl20, "Atividade de Escrita", "Consolidação e Memorização")
add_bento_block(sl20, "Tarefa de Escrita Recomendada:", "[Instrução detalhada da tarefa: escrever um texto, frases ou diálogo]", Inches(0.6), Inches(1.6), Inches(12.3), Inches(2.8), tit_size=18, desc_size=15)
add_bento_block(sl20, "", "Identificamos que a escrita nesta etapa é fundamental para você fixar as estruturas aprendidas, expandir vocabulário ativo e diminuir a tradução mental automática do seu idioma nativo.", Inches(0.6), Inches(4.6), Inches(12.3), Inches(2.2), bg_color=ROXO_CLARO, border_color=ROXO, desc_size=14)

# SLIDE 21: TAREFA DE CASA
sl21 = base_slide(21)
add_header(sl21, "Encerramento — Tarefa de Casa", "Ação e Prática Contínua")
add_bento_block(sl21, "Tarefa de Casa:", "• [Atividade prática 1 para realizar fora da aula]\n• [Atividade de fixação de vocabulário ou escuta]", Inches(0.6), Inches(1.6), Inches(12.3), Inches(5.2), tit_size=18, desc_size=15)

# SLIDE 22: ENCERRAMENTO & FEEDBACK
sl22 = base_slide(22)
add_header(sl22, "Encerramento & Feedback Rápido", "Considerações Finais")
add_bento_block(sl22, "Agradecimento:", "Obrigado pela dedicação na aula de hoje!", Inches(0.6), Inches(1.6), Inches(12.3), Inches(2.0), tit_size=20, desc_size=16)
add_bento_block(sl22, "Espaço para Feedback Rápido:", "[Espaço destinado para o professor coletar dúvidas rápidas e considerações do aluno]", Inches(0.6), Inches(3.8), Inches(12.3), Inches(3.0), border_color=TURQUESA, tit_size=18, desc_size=15)

out_path = "/var/www/haas-frontend-desk-mobile-oficial/public/temp-miniaturas/template-haas-modelo-4.pptx"
prs.save(out_path)
print(f"🚀 MODELO 4 (BENTO DYNAMIC) COMPILADO COM SUCESSO EM: {out_path}")
