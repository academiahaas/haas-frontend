import os
import urllib.request
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# PALETA ESTRITAMENTE DUOTONE (+ BRANCO)
BRANCO = RGBColor(0xFF, 0xFF, 0xFF)
ROXO = RGBColor(0x4F, 0x35, 0xE2)
TURQUESA = RGBColor(0x14, 0xC8, 0xB0)

LOGO_URL = "https://jdppxfokfhqjudwfwckd.supabase.co/storage/v1/object/public/haas-academy/assignments/Design%20sem%20nome%20(4).png"
LOGO_PATH = "/tmp/haas-logo.png"

if not os.path.exists(LOGO_PATH):
    try:
        urllib.request.urlretrieve(LOGO_URL, LOGO_PATH)
    except Exception:
        pass

W, H = Inches(13.333), Inches(7.5)
TOTAL_SLIDES = 22

prs = Presentation()
prs.slide_width = W
prs.slide_height = H

def base_slide(num, bg_color=ROXO, hide_logo=False):
    sl = prs.slides.add_slide(prs.slide_layouts[6])
    sl.background.fill.solid()
    sl.background.fill.fore_color.rgb = bg_color
    
    tb = sl.shapes.add_textbox(Inches(0.5), Inches(7.1), Inches(8), Inches(0.35))
    p = tb.text_frame.paragraphs[0]
    p.text = "HAAS LANGUAGE  •  DUOTONE BRUTALIST"
    txt_color = BRANCO if bg_color == ROXO else ROXO
    p.font.size = Pt(11); p.font.color.rgb = txt_color; p.font.bold = True
    
    tb2 = sl.shapes.add_textbox(Inches(10.8), Inches(7.1), Inches(2.0), Inches(0.35))
    p2 = tb2.text_frame.paragraphs[0]
    p2.text = f"{num} // {TOTAL_SLIDES}"
    p2.font.size = Pt(11); p2.font.color.rgb = TURQUESA if bg_color == ROXO else ROXO; p2.font.bold = True; p2.alignment = PP_ALIGN.RIGHT
    
    if os.path.exists(LOGO_PATH) and not hide_logo:
        sl.shapes.add_picture(LOGO_PATH, Inches(10.8), Inches(0.2), height=Inches(0.65))
        
    return sl

def add_header(sl, tit, sub="", text_color=TURQUESA):
    tb = sl.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(10), Inches(0.8))
    tf = tb.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]; p.text = tit
    p.font.size = Pt(36); p.font.color.rgb = text_color; p.font.bold = True
    
    if sub:
        tb_s = sl.shapes.add_textbox(Inches(0.5), Inches(0.95), Inches(10), Inches(0.4))
        p_s = tb_s.text_frame.paragraphs[0]
        p_s.text = sub
        p_s.font.size = Pt(18); p_s.font.color.rgb = BRANCO; p_s.font.bold = True

def add_solid_block(sl, tit, desc, left, top, width, height, bg_c=BRANCO, tit_c=ROXO, desc_c=ROXO, tit_size=18, desc_size=15):
    s = sl.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    s.fill.solid(); s.fill.fore_color.rgb = bg_c
    s.line.fill.background()
        
    tf = s.text_frame; tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    tf.margin_left = Inches(0.3); tf.margin_top = Inches(0.15); tf.margin_right = Inches(0.3); tf.margin_bottom = Inches(0.15)
    
    if tit:
        p = tf.paragraphs[0]
        p.text = tit; p.font.size = Pt(tit_size); p.font.color.rgb = tit_c; p.font.bold = True
        if desc:
            p2 = tf.add_paragraph()
            p2.text = desc; p2.font.size = Pt(desc_size); p2.font.color.rgb = desc_c; p2.space_before = Pt(6)
    else:
        if desc:
            p = tf.paragraphs[0]
            p.text = desc; p.font.size = Pt(desc_size); p.font.color.rgb = desc_c; p.font.bold = True

def add_chat_bubble(sl, text, left, top, width, height, bg_c, txt_c, is_right=False, font_sz=13):
    s = sl.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    s.fill.solid(); s.fill.fore_color.rgb = bg_c
    s.line.fill.background()
    
    pt_w = Inches(0.25)
    pt_x = left + width - pt_w if is_right else left
    pt = sl.shapes.add_shape(MSO_SHAPE.RIGHT_TRIANGLE, pt_x, top + height - Inches(0.25), pt_w, Inches(0.25))
    pt.fill.solid(); pt.fill.fore_color.rgb = bg_c; pt.line.fill.background()
    if is_right: pt.rotation = 90
    else: pt.rotation = 180
    
    tf = s.text_frame; tf.word_wrap = True; tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    tf.margin_left = Inches(0.18); tf.margin_right = Inches(0.18)
    tf.margin_top = Inches(0.08); tf.margin_bottom = Inches(0.08)
    p = tf.paragraphs[0]; p.text = text; p.font.size = Pt(font_sz)
    p.font.color.rgb = txt_c; p.font.bold = True

# SLIDE 1: CAPA
sl1 = base_slide(1, bg_color=ROXO, hide_logo=True)
img_box = sl1.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), W, Inches(3.75))
img_box.fill.solid(); img_box.fill.fore_color.rgb = TURQUESA; img_box.line.fill.background()
tf_i = img_box.text_frame; tf_i.vertical_anchor = MSO_ANCHOR.MIDDLE
p_i = tf_i.paragraphs[0]; p_i.text = "[Imagem Horizontal Bleed]"; p_i.font.size = Pt(24); p_i.font.color.rgb = ROXO; p_i.font.bold = True; p_i.alignment = PP_ALIGN.CENTER
add_solid_block(sl1, "[Título da Aula / Tema Principal]", "[Subtítulo / Descrição da Aula]\n📅 Data da Aula: [DD/MM/AAAA]", Inches(0), Inches(3.75), W, Inches(3.3), bg_c=ROXO, tit_c=TURQUESA, desc_c=BRANCO, tit_size=44, desc_size=20)

# SLIDE 2: PRESENÇA
sl2 = base_slide(2, bg_color=BRANCO)
add_header(sl2, "Presença & Integrantes", text_color=ROXO)
add_solid_block(sl2, "Professor(a):", "[Nome do(a) Professor(a)]", Inches(0.5), Inches(1.8), Inches(6.0), Inches(4.5), bg_c=TURQUESA, tit_c=ROXO, desc_c=ROXO, tit_size=24, desc_size=18)
add_solid_block(sl2, "Aluno(s):", "[Nome do Aluno / Lista da Turma]", Inches(5.5), Inches(2.8), Inches(7.3), Inches(3.5), bg_c=ROXO, tit_c=TURQUESA, desc_c=BRANCO, tit_size=24, desc_size=18)

# SLIDE 3: WARM-UP
sl3 = base_slide(3, bg_color=TURQUESA)
add_header(sl3, "Aquecimento / Warm-up", "Foco 100% na Fala", text_color=ROXO)
img_w = sl3.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(1.5), W, Inches(2.5))
img_w.fill.solid(); img_w.fill.fore_color.rgb = BRANCO; img_w.line.fill.background()
tf_iw = img_w.text_frame; tf_iw.vertical_anchor = MSO_ANCHOR.MIDDLE
p_iw = tf_iw.paragraphs[0]; p_iw.text = "[Imagem Panorâmica]"; p_iw.font.size = Pt(20); p_iw.font.color.rgb = ROXO; p_iw.font.bold = True; p_iw.alignment = PP_ALIGN.CENTER
add_solid_block(sl3, "Pergunta Rápida:", "[Insira aqui a pergunta provocativa de aquecimento]", Inches(0), Inches(4.0), W, Inches(1.5), bg_c=ROXO, tit_c=TURQUESA, desc_c=BRANCO, tit_size=22, desc_size=18)
add_solid_block(sl3, "", "Neuro: Ativa o córtex pré-frontal e reduz a inibição.", Inches(0), Inches(5.5), W, Inches(1.0), bg_c=BRANCO, tit_c=ROXO, desc_c=ROXO, desc_size=16)

# SLIDE 4: OBJETIVOS
sl4 = base_slide(4, bg_color=ROXO)
add_header(sl4, "Objetivos da Aula", "O que vamos alcançar hoje", text_color=TURQUESA)
add_solid_block(sl4, "Metas:", "• [Objetivo 1]\n\n• [Objetivo 2]\n\n• [Objetivo 3]", Inches(0.5), Inches(1.5), Inches(7.5), Inches(5.0), bg_c=BRANCO, tit_c=ROXO, desc_c=ROXO, tit_size=24, desc_size=18)
add_solid_block(sl4, "Nota de Escrita:", "Atenção à escrita para reduzir a tradução da língua materna.", Inches(8.5), Inches(1.5), Inches(4.3), Inches(5.0), bg_c=TURQUESA, tit_c=ROXO, desc_c=ROXO, tit_size=20, desc_size=18)

# SLIDE 5: METODOLOGIA
sl5 = base_slide(5, bg_color=BRANCO)
add_header(sl5, "A Nossa Metodologia", text_color=ROXO)
metodologias = [
    ("1. Fala Ativa", "Conversação sem julgamentos.", ROXO, TURQUESA, BRANCO),
    ("2. Exposição", "Apresentação da gramática.", TURQUESA, ROXO, ROXO),
    ("3. Prática", "Exercícios guiados.", BRANCO, ROXO, ROXO),
    ("4. Conversa", "Aplicação prática.", ROXO, TURQUESA, BRANCO),
    ("5. Escrita", "Consolidação visual.", TURQUESA, ROXO, ROXO)
]
top_y = 1.5
for tit, desc, bg, t_c, d_c in metodologias:
    add_solid_block(sl5, tit, desc, Inches(0), Inches(top_y), W, Inches(1.1), bg_c=bg, tit_c=t_c, desc_c=d_c, tit_size=20, desc_size=16)
    top_y += 1.1

# SLIDES 6 A 8: EXPOSIÇÃO
for s_num in [6, 7, 8]:
    sl = base_slide(s_num, bg_color=BRANCO)
    add_header(sl, f"Explicação — Parte {s_num - 5}", text_color=ROXO)
    add_solid_block(sl, "Conteúdo:", f"[Explicação e exemplo central da estrutura - Bloco {s_num - 5}]", Inches(0), Inches(1.5), W, Inches(2.0), bg_c=ROXO, tit_c=TURQUESA, desc_c=BRANCO, tit_size=22, desc_size=18)
    add_solid_block(sl, "Prática:", "[Exemplos adicionais e frases do dia a dia]", Inches(0), Inches(3.5), W, Inches(1.7), bg_c=TURQUESA, tit_c=ROXO, desc_c=ROXO, tit_size=20, desc_size=16)
    add_solid_block(sl, "Atenção:", "[Exceções, nuances gramaticais e contexto de uso]", Inches(0), Inches(5.2), W, Inches(1.5), bg_c=BRANCO, tit_c=ROXO, desc_c=ROXO, tit_size=18, desc_size=16)

# SLIDE 9: SÍNTESE + CULTURA
sl9 = base_slide(9, bg_color=ROXO)
add_solid_block(sl9, "Resumo do Tema:", "[Resumo sintético dos pontos chave ensinados na aula]", Inches(0), Inches(0), Inches(6.666), H, bg_c=ROXO, tit_c=TURQUESA, desc_c=BRANCO, tit_size=26, desc_size=18)
add_solid_block(sl9, "Sugestão Cultural:", "🎵 [Nome da Música / Vídeo recomendado].\n\nO professor pode rodar o áudio ou deixar como extra.", Inches(6.666), Inches(0), Inches(6.666), H, bg_c=TURQUESA, tit_c=ROXO, desc_c=ROXO, tit_size=26, desc_size=18)

# SLIDES 10 E 11: PRÁTICA GUIADA ZIG-ZAG
def add_zigzag_card(sl, num, desc, left, top, width, bg_c, t_c):
    add_solid_block(sl, num, desc, Inches(left), Inches(top), Inches(width), Inches(1.0), bg_c=bg_c, tit_c=t_c, desc_c=t_c, tit_size=20, desc_size=16)

sl10 = base_slide(10, bg_color=ROXO)
add_header(sl10, "Prática Guiada 1-5", text_color=TURQUESA)
add_zigzag_card(sl10, "01", "[Lacunas contextualizadas com opções]", 0.5, 1.5, 8.0, BRANCO, ROXO)
add_zigzag_card(sl10, "02", "[Lacunas contextualizadas com opções]", 4.8, 2.6, 8.0, TURQUESA, ROXO)
add_zigzag_card(sl10, "03", "[Preenchimento direto de verbos]", 0.5, 3.7, 8.0, BRANCO, ROXO)
add_zigzag_card(sl10, "04", "[Reestruturação / Transformação de frase]", 4.8, 4.8, 8.0, TURQUESA, ROXO)
add_zigzag_card(sl10, "05", "[Reestruturação / Transformação de frase]", 0.5, 5.9, 12.3, BRANCO, ROXO)

sl11 = base_slide(11, bg_color=ROXO)
add_header(sl11, "Prática Guiada 6-10", text_color=TURQUESA)
add_zigzag_card(sl11, "06", "[Preposição / Conector correto]", 0.5, 1.5, 8.0, BRANCO, ROXO)
add_zigzag_card(sl11, "07", "[Escolha da opção correta]", 4.8, 2.6, 8.0, TURQUESA, ROXO)
add_zigzag_card(sl11, "08", "[Substituição por sinônimo formal]", 0.5, 3.7, 8.0, BRANCO, ROXO)
add_zigzag_card(sl11, "09", "[Identificação de erro e correção]", 4.8, 4.8, 8.0, TURQUESA, ROXO)
add_zigzag_card(sl11, "10", "[Tradução de sentido sem uso literal]", 0.5, 5.9, 12.3, BRANCO, ROXO)

# SLIDES 12 A 19: CONVERSAÇÃO (FIXO SEM ESTOURAR A TELA)
perguntas_config = [
    "[Pergunta principal 1 sobre o tema]", "[Pergunta principal 2 sobre o tema]", "[Pergunta principal 3 sobre o tema]",
    "[Pergunta principal 4 sobre o tema]", "[Pergunta principal 5 sobre o tema]", "[Pergunta principal 6 sobre o tema]",
    "[Pergunta principal 7 sobre o tema]", "[Pergunta principal 8 sobre o tema]"
]

for idx in range(8):
    sl = base_slide(12 + idx, bg_color=BRANCO)
    add_header(sl, "Conversação", text_color=ROXO)
    
    if idx == 0:
        # Pergunta 1 Ampla Centralizada
        add_chat_bubble(sl, perguntas_config[0], Inches(0.5), Inches(1.5), Inches(12.3), Inches(2.6), TURQUESA, ROXO, is_right=True, font_sz=20)
        add_solid_block(sl, "Apoio ao Professor:", "• Para expandir: \"Dê um exemplo recente...\"\n• Para concordar/discordar: \"Qual é o outro lado dessa moeda?\"", Inches(0.5), Inches(4.4), Inches(12.3), Inches(2.3), bg_c=ROXO, tit_c=TURQUESA, desc_c=BRANCO, tit_size=18, desc_size=15)
    else:
        # LADO ESQUERDO: Histórico que ESCALA DINAMICAMENTE sem estourar o slide
        tot_h = 5.2
        top_start = 1.5
        count = idx
        gap = 0.12
        avail_h = tot_h - (gap * (count - 1))
        bubble_h = avail_h / count
        f_size = 12 if count <= 3 else (11 if count <= 5 else 10)
        
        for i in range(idx):
            y_pos = top_start + i * (bubble_h + gap)
            h_txt = perguntas_config[i]
            clean_t = h_txt[:42] + "..." if len(h_txt) > 42 else h_txt
            add_chat_bubble(sl, clean_t, Inches(0.5), Inches(y_pos), Inches(4.5), Inches(bubble_h), ROXO, BRANCO, is_right=False, font_sz=f_size)
        
        # LADO DIREITO: Pergunta Ativa & Card Roxo de Apoio 100% TRAVADOS E FIXOS!
        p_desc = perguntas_config[idx]
        add_chat_bubble(sl, p_desc, Inches(5.3), Inches(1.5), Inches(7.5), Inches(2.6), TURQUESA, ROXO, is_right=True, font_sz=18)
        add_solid_block(sl, "Apoio ao Professor:", "• Para expandir: \"Dê um exemplo recente...\"\n• Para concordar/discordar: \"Qual é o outro lado dessa moeda?\"", Inches(5.3), Inches(4.4), Inches(7.5), Inches(2.3), bg_c=ROXO, tit_c=TURQUESA, desc_c=BRANCO, tit_size=18, desc_size=14)

# SLIDE 20: ESCRITA
sl20 = base_slide(20, bg_color=TURQUESA)
add_header(sl20, "Atividade de Escrita", text_color=ROXO)
add_solid_block(sl20, "Tarefa:", "[Instrução detalhada da tarefa: escrever um texto, frases ou diálogo]", Inches(0.5), Inches(1.5), Inches(12.3), Inches(3.5), bg_c=BRANCO, tit_c=ROXO, desc_c=ROXO, tit_size=24, desc_size=18)
add_solid_block(sl20, "", "A escrita é fundamental para fixar estruturas, expandir vocabulário e diminuir a tradução mental automática.", Inches(0), Inches(5.5), W, Inches(1.5), bg_c=ROXO, tit_c=BRANCO, desc_c=BRANCO, desc_size=16)

# SLIDE 21: TAREFA DE CASA
sl21 = base_slide(21, bg_color=ROXO)
add_header(sl21, "Tarefa de Casa", text_color=TURQUESA)
add_solid_block(sl21, "Ação Contínua:", "• [Atividade prática 1 fora da aula]\n\n• [Atividade de fixação de vocabulário ou escuta]", Inches(0), Inches(2.0), W, Inches(4.0), bg_c=ROXO, tit_c=BRANCO, desc_c=BRANCO, tit_size=36, desc_size=24)

# SLIDE 22: ENCERRAMENTO
sl22 = base_slide(22, bg_color=BRANCO)
add_header(sl22, "Encerramento & Feedback", text_color=ROXO)
add_solid_block(sl22, "Obrigado!", "Sua dedicação é o que faz a diferença.", Inches(0.5), Inches(1.5), Inches(6.0), Inches(5.0), bg_c=TURQUESA, tit_c=ROXO, desc_c=ROXO, tit_size=32, desc_size=20)
add_solid_block(sl22, "Feedback:", "[Professor: colete dúvidas rápidas e considerações do aluno aqui]", Inches(6.8), Inches(1.5), Inches(6.0), Inches(5.0), bg_c=ROXO, tit_c=TURQUESA, desc_c=BRANCO, tit_size=24, desc_size=18)

out_path = "/var/www/haas-frontend-desk-mobile-oficial/public/temp-miniaturas/template-haas-modelo-14-duotone.pptx"
prs.save(out_path)
print(f"✅ MODELO 14 CORRIGIDO E PERFEITO EM: {out_path}")
