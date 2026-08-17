import os
import urllib.request
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

# PALETA OFICIAL HAAS LANGUAGE
BRANCO = RGBColor(0xFF, 0xFF, 0xFF)
ROXO = RGBColor(0x4F, 0x35, 0xE2)
AZUL = RGBColor(0x1E, 0x3A, 0x8A)
SLATE = RGBColor(0x33, 0x41, 0x55)
TURQUESA = RGBColor(0x14, 0xC8, 0xB0)
BORDA = RGBColor(0xCB, 0xD5, 0xE1)
FUNDO_CARD = RGBColor(0xF8, 0xFA, 0xFC)
ROXO_CLARO = RGBColor(0xED, 0xE9, 0xFE)
MUTED = RGBColor(0x94, 0xA3, 0xB8)

# LOGO AUTOMÁTICA
LOGO_URL = "https://jdppxfokfhqjudwfwckd.supabase.co/storage/v1/object/public/haas-academy/assignments/Design%20sem%20nome%20(4).png"
LOGO_PATH = "/tmp/haas-logo.png"

if not os.path.exists(LOGO_PATH):
    try:
        urllib.request.urlretrieve(LOGO_URL, LOGO_PATH)
        print("✅ Logo baixada com sucesso.")
    except Exception as e:
        print(f"⚠️ Não foi possível baixar a logo: {e}")

W, H = Inches(13.333), Inches(7.5)

prs = Presentation()
prs.slide_width = W
prs.slide_height = H

def base_slide(num):
    sl = prs.slides.add_slide(prs.slide_layouts[6])
    sl.background.fill.solid()
    sl.background.fill.fore_color.rgb = BRANCO
    
    # Barra lateral roxa
    s = sl.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(0.12), H)
    s.fill.solid(); s.fill.fore_color.rgb = ROXO; s.line.fill.background()
    
    # Rodapé roxo
    b = sl.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(7.05), W, Inches(0.45))
    b.fill.solid(); b.fill.fore_color.rgb = ROXO; b.line.fill.background()
    
    # Inserção da Logo da Haas Language
    if os.path.exists(LOGO_PATH):
        sl.shapes.add_picture(LOGO_PATH, Inches(10.8), Inches(0.2), height=Inches(0.75))
    
    # Detalhe Turquesa
    det = sl.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.5), Inches(0.95), Inches(1.5), Pt(4))
    det.fill.solid(); det.fill.fore_color.rgb = TURQUESA; det.line.fill.background()
    
    # Rodapé Texto
    tb = sl.shapes.add_textbox(Inches(0.5), Inches(7.1), Inches(8), Inches(0.35))
    p = tb.text_frame.paragraphs[0]
    p.text = "Haas Language"
    p.font.size = Pt(12); p.font.color.rgb = BRANCO; p.font.bold = True
    
    tb2 = sl.shapes.add_textbox(Inches(11.2), Inches(7.1), Inches(1.6), Inches(0.35))
    p2 = tb2.text_frame.paragraphs[0]
    p2.text = f"Slide {num} / 18"
    p2.font.size = Pt(12); p2.font.color.rgb = BRANCO; p2.font.bold = True; p2.alignment = PP_ALIGN.RIGHT
    return sl

def add_header(sl, tit, sub=""):
    tb = sl.shapes.add_textbox(Inches(0.5), Inches(1.15), Inches(10), Inches(0.7))
    tf = tb.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]; p.text = tit
    p.font.size = Pt(32); p.font.color.rgb = ROXO; p.font.bold = True
    
    if sub:
        tb_s = sl.shapes.add_textbox(Inches(0.5), Inches(1.85), Inches(10), Inches(0.45))
        p_s = tb_s.text_frame.paragraphs[0]
        p_s.text = sub
        p_s.font.size = Pt(19); p_s.font.color.rgb = AZUL; p_s.font.bold = True

def add_card(sl, tit, desc, left, top, width, height, bg_color=FUNDO_CARD, border_color=BORDA, tit_size=16, desc_size=14):
    s = sl.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    s.fill.solid(); s.fill.fore_color.rgb = bg_color
    s.line.color.rgb = border_color; s.line.width = Pt(1)
    tf = s.text_frame; tf.word_wrap = True; tf.margin_left = Inches(0.2); tf.margin_top = Inches(0.2)
    if tit:
        p = tf.paragraphs[0]; p.text = tit; p.font.size = Pt(tit_size); p.font.color.rgb = ROXO; p.font.bold = True
        if desc:
            p2 = tf.add_paragraph(); p2.text = desc; p2.font.size = Pt(desc_size); p2.font.color.rgb = SLATE; p2.space_before = Pt(6)
    else:
        if desc:
            p = tf.paragraphs[0]; p.text = desc; p.font.size = Pt(desc_size); p.font.color.rgb = SLATE

def img_box(sl, left, top, width, height, label="[Espaço reservado para Imagem]"):
    s = sl.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    s.fill.solid(); s.fill.fore_color.rgb = ROXO_CLARO
    s.line.color.rgb = BORDA; s.line.width = Pt(1)
    tf = s.text_frame; tf.word_wrap = True
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER
    p = tf.paragraphs[0]; p.text = label
    p.font.size = Pt(15); p.font.color.rgb = ROXO; p.font.bold = True

# ==========================================================
# CONSTRUÇÃO DOS 18 SLIDES COM FONTES EXPANDIDAS
# ==========================================================

# SLIDE 1: CAPA (Apenas Título, Subtítulo e Imagem Principal)
sl1 = base_slide(1)
tb_c = sl1.shapes.add_textbox(Inches(0.5), Inches(2.2), Inches(6.5), Inches(1.5))
tf_c = tb_c.text_frame; tf_c.word_wrap = True
p = tf_c.paragraphs[0]; p.text = "[Título da Aula / Tema Principal]"
p.font.size = Pt(38); p.font.color.rgb = ROXO; p.font.bold = True

tb_sub = sl1.shapes.add_textbox(Inches(0.5), Inches(3.8), Inches(6.5), Inches(1.0))
tf_sub = tb_sub.text_frame; tf_sub.word_wrap = True
p_sub = tf_sub.paragraphs[0]; p_sub.text = "[Subtítulo / Descrição da Aula]"
p_sub.font.size = Pt(22); p_sub.font.color.rgb = AZUL; p_sub.font.bold = True

img_box(sl1, Inches(7.3), Inches(1.5), Inches(5.5), Inches(5.0), "[Imagem de Capa da Aula]")

# SLIDE 2: INTEGRANTES
sl2 = base_slide(2)
add_header(sl2, "Presença & Integrantes", "Lista de Alunos Previstos para esta Turma")
positions = [(0.5, 2.5), (4.7, 2.5), (8.9, 2.5), (0.5, 4.6), (4.7, 4.6), (8.9, 4.6)]
for i, (l, t) in enumerate(positions, 1):
    add_card(sl2, f"Aluno {i}:", "[Nome do Aluno]", Inches(l), Inches(t), Inches(3.9), Inches(1.8), tit_size=18, desc_size=15)

# SLIDE 3: WARM-UP (Com Imagem Reservada)
sl3 = base_slide(3)
add_header(sl3, "Aquecimento / Warm-up", "Foco 100% na Fala (3 Minutos)")
add_card(sl3, "Pergunta Provocativa Rápida:", "[Insira aqui a pergunta provocativa de aquecimento]", Inches(0.5), Inches(2.5), Inches(6.5), Inches(2.2), tit_size=18, desc_size=15)
add_card(sl3, "Fundamentação Neurocognitiva:", "Ativa o córtex pré-frontal e reduz a inibição linguística inicial.", Inches(0.5), Inches(4.9), Inches(6.5), Inches(1.8), bg_color=ROXO_CLARO, border_color=ROXO, tit_size=16, desc_size=14)
img_box(sl3, Inches(7.3), Inches(2.5), Inches(5.5), Inches(4.2), "[Imagem de Aquecimento / Warm-up]")

# SLIDE 4: OBJETIVOS
sl4 = base_slide(4)
add_header(sl4, "Objetivos da Aula", "O que vamos alcançar hoje")
add_card(sl4, "Metas da Aula:", "• [Objetivo 1]\n• [Objetivo 2]\n• [Objetivo 3]", Inches(0.5), Inches(2.5), Inches(12.3), Inches(2.6), tit_size=18, desc_size=15)
add_card(sl4, "Nota de Escrita:", "Nesta aula daremos atenção à escrita para reforçar a memorização de estruturas/vocabulário e reduzir a interferência da língua materna.", Inches(0.5), Inches(5.3), Inches(12.3), Inches(1.4), bg_color=FUNDO_CARD, border_color=TURQUESA, tit_size=16, desc_size=14)

# SLIDE 5: METODOLOGIA (Cards Verticais / Colunas)
sl5 = base_slide(5)
add_header(sl5, "A Nossa Metodologia", "Estrutura Fixa de Aprendizado Haas Language")
metodologias = [
    ("1. Fala Ativa", "Conversação inicial sem julgamentos."),
    ("2. Exposição", "Apresentação leve da gramática."),
    ("3. Prática Guiada", "Exercícios com suporte do professor."),
    ("4. Conversação", "Aplicação prática das estruturas."),
    ("5. Escrita", "Consolidação e redução de tradução.")
]
left_m = 0.5
for tit_m, desc_m in metodologias:
    add_card(sl5, tit_m, desc_m, Inches(left_m), Inches(2.5), Inches(2.25), Inches(4.2), tit_size=16, desc_size=13)
    left_m += 2.45

# SLIDES 6 A 8: EXPOSIÇÃO MULTINÍVEL
for s_num in [6, 7, 8]:
    sl = base_slide(s_num)
    add_header(sl, f"Explicação do Tema — Parte {s_num - 5}", "Estratégia Multinível")
    add_card(sl, "Conteúdo Principal (Nível Inicial / Básico):", f"[Explicação e exemplos simples para os níveis Inicial e Básico - Bloco {s_num - 5}]", Inches(0.5), Inches(2.5), Inches(7.8), Inches(4.2), tit_size=18, desc_size=15)
    add_card(sl, "Nuances & Exceções\n(Nível Independente / Avançado):", "[Nuances gramaticais detalhadas para alunos avançados]", Inches(8.5), Inches(2.5), Inches(4.3), Inches(4.2), bg_color=ROXO_CLARO, border_color=ROXO, tit_size=16, desc_size=14)

# SLIDE 9: SÍNTESE + CULTURA
sl9 = base_slide(9)
add_header(sl9, "Síntese do Conteúdo & Sugestão Cultural", "Resumo e Recomendação")
add_card(sl9, "Resumo do Tema:", "[Resumo sintético dos pontos chave ensinados na aula]", Inches(0.5), Inches(2.5), Inches(12.3), Inches(2.3), tit_size=18, desc_size=15)
add_card(sl9, "Sugestão Cultural (Música ou Vídeo):", "🎵 [Nome da Música / Vídeo recomendado sobre o assunto].\nO professor pode rodar o áudio/vídeo ou deixar como recomendação extra.", Inches(0.5), Inches(5.0), Inches(12.3), Inches(1.7), bg_color=FUNDO_CARD, border_color=AZUL, tit_size=16, desc_size=14)

# SLIDES 10 E 11: PRÁTICA GUIADA
sl10 = base_slide(10)
add_header(sl10, "Prática Guiada — Exercícios 1 a 5", "Máximo de 10 Exercícios Guiados (Correção de Vícios)")
exs_1_5 = [
    ("Exercício 1:", "[Lacunas contextualizadas com opções] — Inicial / Descoberta"),
    ("Exercício 2:", "[Lacunas contextualizadas com opções] — Básico / Conexão"),
    ("Exercício 3:", "[Preenchimento direto de verbos] — Básico / Conexão"),
    ("Exercício 4:", "[Reestruturação / Transformação] — Intermediário / Autônomo"),
    ("Exercício 5:", "[Reestruturação / Transformação] — Intermediário / Autônomo")
]
top_e = 2.4
for t_e, d_e in exs_1_5:
    add_card(sl10, t_e, d_e, Inches(0.5), Inches(top_e), Inches(12.3), Inches(0.85), tit_size=15, desc_size=14)
    top_e += 0.92

sl11 = base_slide(11)
add_header(sl11, "Prática Guiada — Exercícios 6 a 10", "Desafios para Níveis Superiores")
exs_6_10 = [
    ("Exercício 6:", "[Preposição / Conector correto] — Intermediário / Autônomo"),
    ("Exercício 7:", "[Escolha de opção correta] — Independente"),
    ("Exercício 8:", "[Substituição por sinônimo formal] — Independente"),
    ("Exercício 9:", "[Identificação de erro e correção] — Avançado / Especialista"),
    ("Exercício 10:", "[Tradução de sentido sem uso literal] — Avançado / Especialista")
]
top_e = 2.4
for t_e, d_e in exs_6_10:
    add_card(sl11, t_e, d_e, Inches(0.5), Inches(top_e), Inches(12.3), Inches(0.85), tit_size=15, desc_size=14)
    top_e += 0.92

# SLIDES 12 A 15: CONVERSAÇÃO CIRÚRGICA CUMULATIVA
sl12 = base_slide(12)
add_header(sl12, "Conversação Cirúrgica", "Efeito Cumulativo Progressivo")
add_card(sl12, "Pergunta 1 (Revelada):", "[Pergunta principal 1] — Inicial / Descoberta", Inches(0.5), Inches(2.5), Inches(12.3), Inches(2.1), bg_color=ROXO_CLARO, border_color=ROXO, tit_size=20, desc_size=16)
add_card(sl12, "Sub-prompts Visuais (Apoio):", "• Para expandir: 'Dê um exemplo recente...'\n• Para concordar/discordar: 'Qual é o outro lado dessa moeda?'", Inches(0.5), Inches(4.8), Inches(12.3), Inches(1.8), tit_size=16, desc_size=14)

sl13 = base_slide(13)
add_header(sl13, "Conversação Cirúrgica", "Efeito Cumulativo Progressivo")
add_card(sl13, "Histórico:", "1. [Pergunta principal 1]", Inches(0.5), Inches(2.4), Inches(12.3), Inches(0.9), bg_color=FUNDO_CARD, border_color=MUTED, tit_size=14, desc_size=13)
add_card(sl13, "Pergunta 2 (Nova):", "[Pergunta principal 2] — Básico / Conexão", Inches(0.5), Inches(3.5), Inches(12.3), Inches(1.9), bg_color=ROXO_CLARO, border_color=ROXO, tit_size=20, desc_size=16)
add_card(sl13, "Sub-prompts Visuais (Apoio):", "• Para expandir: 'Dê um exemplo recente...'\n• Para concordar/discordar: 'Qual é o outro lado dessa moeda?'", Inches(0.5), Inches(5.6), Inches(12.3), Inches(1.2), tit_size=15, desc_size=13)

sl14 = base_slide(14)
add_header(sl14, "Conversação Cirúrgica", "Efeito Cumulativo Progressivo")
add_card(sl14, "Histórico:", "1. [Pergunta 1]\n2. [Pergunta 2]", Inches(0.5), Inches(2.4), Inches(12.3), Inches(1.3), bg_color=FUNDO_CARD, border_color=MUTED, tit_size=14, desc_size=13)
add_card(sl14, "Pergunta 3 (Nova):", "[Pergunta principal 3] — Intermediário / Autônomo", Inches(0.5), Inches(3.9), Inches(12.3), Inches(1.8), bg_color=ROXO_CLARO, border_color=ROXO, tit_size=20, desc_size=16)
add_card(sl14, "Sub-prompts Visuais (Apoio):", "• Prompts de expansão para a Pergunta 3", Inches(0.5), Inches(5.9), Inches(12.3), Inches(0.9), tit_size=15, desc_size=13)

sl15 = base_slide(15)
add_header(sl15, "Conversação Cirúrgica", "Efeito Cumulativo Progressivo")
add_card(sl15, "Histórico Completo:", "1. [Pergunta 1]\n2. [Pergunta 2]\n3. [Pergunta 3]", Inches(0.5), Inches(2.4), Inches(12.3), Inches(1.6), bg_color=FUNDO_CARD, border_color=MUTED, tit_size=14, desc_size=13)
add_card(sl15, "Pergunta 4 (Nova):", "[Pergunta principal 4] — Avançado / Especialista", Inches(0.5), Inches(4.2), Inches(12.3), Inches(1.6), bg_color=ROXO_CLARO, border_color=ROXO, tit_size=20, desc_size=16)
add_card(sl15, "Sub-prompts Visuais (Apoio):", "• Prompts de análise crítica para a Pergunta 4", Inches(0.5), Inches(6.0), Inches(12.3), Inches(0.8), tit_size=14, desc_size=12)

# SLIDE 16: PARTE ESCRITA
sl16 = base_slide(16)
add_header(sl16, "Atividade de Escrita", "Consolidação e Memorização")
add_card(sl16, "Tarefa de Escrita Recomendada:", "[Instrução detalhada da tarefa: escrever um texto, frases ou diálogo dependendo do nível]", Inches(0.5), Inches(2.5), Inches(12.3), Inches(2.5), tit_size=18, desc_size=15)
add_card(sl16, "Razão Pedagógica (Rodapé Fixo):", "Identificamos que a escrita nesta etapa é fundamental para você fixar as estruturas aprendidas, expandir vocabulário ativo e diminuir a tradução mental automática do seu idioma nativo.", Inches(0.5), Inches(5.2), Inches(12.3), Inches(1.5), bg_color=ROXO_CLARO, border_color=ROXO, tit_size=16, desc_size=14)

# SLIDE 17: TAREFA DE CASA
sl17 = base_slide(17)
add_header(sl17, "Encerramento — Tarefa de Casa", "Ação e Prática Contínua")
add_card(sl17, "Tarefa de Casa:", "• [Atividade prática 1 para realizar fora da aula]\n• [Atividade de fixação de vocabulário ou escuta]", Inches(0.5), Inches(2.5), Inches(12.3), Inches(4.2), tit_size=18, desc_size=15)

# SLIDE 18: ENCERRAMENTO & FEEDBACK
sl18 = base_slide(18)
add_header(sl18, "Encerramento & Feedback Rápido", "Considerações Finais")
add_card(sl18, "Agradecimento:", "Obrigado pela dedicação na aula de hoje!", Inches(0.5), Inches(2.5), Inches(12.3), Inches(1.5), tit_size=20, desc_size=16)
add_card(sl18, "Espaço para Feedback Rápido:", "[Espaço destinado para o professor coletar dúvidas rápidas e considerações do aluno]", Inches(0.5), Inches(4.3), Inches(12.3), Inches(2.4), bg_color=FUNDO_CARD, border_color=TURQUESA, tit_size=18, desc_size=15)

# SALVAR
out_path = "/var/www/haas-frontend-desk-mobile-oficial/public/temp-miniaturas/template-haas-modelo.pptx"
prs.save(out_path)
print(f"✅ NOVO TEMPLATE HAAS LANGUAGE V3 CRIADO EM: {out_path}")
