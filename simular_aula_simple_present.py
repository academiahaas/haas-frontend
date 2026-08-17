import os
import urllib.request
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

# PALETA OFICIAL HAAS LANGUAGE
BRANCO = RGBColor(0xFF, 0xFF, 0xFF)
ROXO = RGBColor(0x4F, 0x35, 0xE2)
AZUL = RGBColor(0x1E, 0x3A, 0x8A)
SLATE = RGBColor(0x33, 0x41, 0x55)
TURQUESA = RGBColor(0x14, 0xC8, 0xB0)
BORDA = RGBColor(0xCB, 0xD5, 0xE1)
FUNDO_CARD = RGBColor(0xF8, 0xFA, 0xFC)
ROXO_CLARO = RGBColor(0xED, 0xE9, 0xFE)
MUTED = RGBColor(0x94, 0xA3, 0xB8)

LOGO_PATH = "/tmp/haas-logo.png"

W, H = Inches(13.333), Inches(7.5)
TOTAL_SLIDES = 22

prs = Presentation()
prs.slide_width = W
prs.slide_height = H

def base_slide(num):
    sl = prs.slides.add_slide(prs.slide_layouts[6])
    sl.background.fill.solid()
    sl.background.fill.fore_color.rgb = BRANCO
    
    # Barra lateral
    s = sl.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(0.12), H)
    s.fill.solid(); s.fill.fore_color.rgb = ROXO; s.line.fill.background()
    
    # Rodapé
    b = sl.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(7.05), W, Inches(0.45))
    b.fill.solid(); b.fill.fore_color.rgb = ROXO; b.line.fill.background()
    
    # Logo Haas Language
    if os.path.exists(LOGO_PATH):
        sl.shapes.add_picture(LOGO_PATH, Inches(10.8), Inches(0.2), height=Inches(0.75))
    
    # Detalhe Turquesa
    det = sl.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.5), Inches(0.95), Inches(1.5), Pt(4))
    det.fill.solid(); det.fill.fore_color.rgb = TURQUESA; det.line.fill.background()
    
    # Texto Rodapé (Español / Haas Language)
    tb = sl.shapes.add_textbox(Inches(0.5), Inches(7.1), Inches(8), Inches(0.35))
    p = tb.text_frame.paragraphs[0]
    p.text = "Haas Language"
    p.font.size = Pt(12); p.font.color.rgb = BRANCO; p.font.bold = True
    
    tb2 = sl.shapes.add_textbox(Inches(10.8), Inches(7.1), Inches(2.0), Inches(0.35))
    p2 = tb2.text_frame.paragraphs[0]
    p2.text = f"Slide {num} / {TOTAL_SLIDES}"
    p2.font.size = Pt(12); p2.font.color.rgb = BRANCO; p2.font.bold = True; p2.alignment = PP_ALIGN.RIGHT
    return sl

def add_header(sl, tit, sub=""):
    tb = sl.shapes.add_textbox(Inches(0.5), Inches(1.15), Inches(10), Inches(0.7))
    tf = tb.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]; p.text = tit
    p.font.size = Pt(32); p.font.color.rgb = ROXO; p.font.bold = True
    
    if sub:
        tb_s = sl.shapes.add_textbox(Inches(0.5), Inches(1.85), Inches(10), Inches(0.45))
        p_s = tb_s.text_frame.paragraphs[0]
        p_s.text = sub
        p_s.font.size = Pt(19); p_s.font.color.rgb = AZUL; p_s.font.bold = True

def add_card(sl, tit, desc, left, top, width, height, bg_color=FUNDO_CARD, border_color=BORDA, tit_size=16, desc_size=14):
    s = sl.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    s.fill.solid(); s.fill.fore_color.rgb = bg_color
    s.line.color.rgb = border_color; s.line.width = Pt(1)
    tf = s.text_frame; tf.word_wrap = True; tf.margin_left = Inches(0.15); tf.margin_top = Inches(0.15)
    if tit:
        p = tf.paragraphs[0]; p.text = tit; p.font.size = Pt(tit_size); p.font.color.rgb = ROXO; p.font.bold = True
        if desc:
            p2 = tf.add_paragraph(); p2.text = desc; p2.font.size = Pt(desc_size); p2.font.color.rgb = SLATE; p2.space_before = Pt(4)
    else:
        if desc:
            p = tf.paragraphs[0]; p.text = desc; p.font.size = Pt(desc_size); p.font.color.rgb = SLATE

def add_history_cards(sl, history_questions, top=2.25, height=1.05):
    count = len(history_questions)
    if count == 0:
        return
    total_w = 12.3
    margin_left = 0.5
    gap = 0.15 if count >= 5 else 0.25
    avail_w = total_w - (gap * (count - 1))
    card_w = avail_w / count
    
    for i, q_text in enumerate(history_questions):
        l = margin_left + i * (card_w + gap)
        q_num = i + 1
        desc_clean = q_text.split(" — ")[0]
        t_sz, d_sz = (13, 11) if count <= 2 else ((11, 10) if count <= 4 else (10, 9))
        add_card(sl, f"P{q_num}", desc_clean, Inches(l), Inches(top), Inches(card_w), Inches(height), 
                 bg_color=FUNDO_CARD, border_color=MUTED, tit_size=t_sz, desc_size=d_sz)

def img_box(sl, left, top, width, height, label="[Espacio reservado para Imagen]"):
    s = sl.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    s.fill.solid(); s.fill.fore_color.rgb = ROXO_CLARO
    s.line.color.rgb = BORDA; s.line.width = Pt(1)
    tf = s.text_frame; tf.word_wrap = True
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER
    p = tf.paragraphs[0]; p.text = label
    p.font.size = Pt(15); p.font.color.rgb = ROXO; p.font.bold = True

# ==========================================================
# CONSTRUCCIÓN DE LA CLASE: SIMPLE PRESENT (IDIOMA AUXILIAR: ESPAÑOL)
# ==========================================================

# SLIDE 1: CAPA
sl1 = base_slide(1)
tb_c = sl1.shapes.add_textbox(Inches(0.5), Inches(2.2), Inches(6.5), Inches(1.5))
tf_c = tb_c.text_frame; tf_c.word_wrap = True
p = tf_c.paragraphs[0]; p.text = "Simple Present Tense"
p.font.size = Pt(38); p.font.color.rgb = ROXO; p.font.bold = True

tb_sub = sl1.shapes.add_textbox(Inches(0.5), Inches(3.8), Inches(6.5), Inches(1.0))
tf_sub = tb_sub.text_frame; tf_sub.word_wrap = True
p_sub = tf_sub.paragraphs[0]; p_sub.text = "Habla sobre tu rutina y hechos reales en inglés"
p_sub.font.size = Pt(22); p_sub.font.color.rgb = AZUL; p_sub.font.bold = True

img_box(sl1, Inches(7.3), Inches(1.5), Inches(5.5), Inches(5.0), "[Imagen: Daily Routines & Habits]")

# SLIDE 2: PRESENCIA & INTEGRANTES
sl2 = base_slide(2)
add_header(sl2, "Presencia e Integrantes")
add_card(sl2, "Profesor(a):", "Sarah Miller", Inches(0.5), Inches(2.3), Inches(3.8), Inches(4.3), tit_size=18, desc_size=16)
add_card(sl2, "Alumno(s):", "Carlos Gómez", Inches(4.6), Inches(2.3), Inches(8.2), Inches(4.3), tit_size=18, desc_size=16)

# SLIDE 3: WARM-UP
sl3 = base_slide(3)
add_header(sl3, "Calentamiento / Warm-up", "Enfoque 100% en Habla (3 Minutos)")
add_card(sl3, "Pregunta Provocadora Rápida:", "What is the first thing you do when you wake up in the morning?", Inches(0.5), Inches(2.5), Inches(6.5), Inches(2.2), tit_size=18, desc_size=15)
add_card(sl3, "Fundamentación Neurocognitiva:", "Activa la corteza prefrontal y reduce la inhibición lingüística inicial.", Inches(0.5), Inches(4.9), Inches(6.5), Inches(1.8), bg_color=ROXO_CLARO, border_color=ROXO, tit_size=16, desc_size=14)
img_box(sl3, Inches(7.3), Inches(2.5), Inches(5.5), Inches(4.2), "[Imagen: Morning Routine]")

# SLIDE 4: OBJETIVOS
sl4 = base_slide(4)
add_header(sl4, "Objetivos de la Clase", "Lo que lograremos hoy")
add_card(sl4, "Metas de la Clase:", "• Expresar hábitos y rutinas diarias en inglés con fluidez.\n• Dominar las reglas de la 3ª persona del singular (he, she, it).\n• Formular preguntas y negaciones correctamente usando do/does.", Inches(0.5), Inches(2.5), Inches(12.3), Inches(2.6), tit_size=18, desc_size=15)
add_card(sl4, "Nota de Escritura:", "En esta clase prestaremos atención a la escritura para reforzar la memorización de estructuras/vocabulario y reducir la interferencia de la lengua materna.", Inches(0.5), Inches(5.3), Inches(12.3), Inches(1.4), bg_color=FUNDO_CARD, border_color=TURQUESA, tit_size=16, desc_size=14)

# SLIDE 5: METODOLOGÍA
sl5 = base_slide(5)
add_header(sl5, "Nuestra Metodología", "Estructura Fija de Aprendizaje Haas Language")
metodologias = [
    ("1. Habla Activa", "Conversación inicial sin juzgamientos."),
    ("2. Exposición", "Presentación ligera de la gramática."),
    ("3. Práctica Guiada", "Ejercicios con apoyo del profesor."),
    ("4. Conversación", "Aplicación práctica de estructuras."),
    ("5. Escritura", "Consolidación y reducción de traducción.")
]
left_m = 0.5
for tit_m, desc_m in metodologias:
    add_card(sl5, tit_m, desc_m, Inches(left_m), Inches(2.5), Inches(2.25), Inches(4.2), tit_size=16, desc_size=13)
    left_m += 2.45

# SLIDES 6 A 8: EXPOSICIÓN MULTINIVEL
# Slide 6
sl6 = base_slide(6)
add_header(sl6, "Simple Present — Estructura Básica", "Estrategia Multinivel")
add_card(sl6, "Contenido Principal (Inicial / Básico):", "Usamos el Simple Present para rutinas, hechos y verdades generales.\n\nEstructura:\n• Affirmative: I / You / We / They + Verb (base form)\n  Ex: I play tennis every Saturday.\n• Negative: Subject + don't + Verb\n  Ex: They don't eat meat.", Inches(0.5), Inches(2.5), Inches(7.8), Inches(4.2), tit_size=18, desc_size=14)
add_card(sl6, "Nuances & Excepciones\n(Independiente / Avanzado):", "Spelling rules para He / She / It:\n• Add -s: works, plays\n• Add -es (ch, sh, ss, x, o): watches, goes\n• Y preceded by consonant -> -ies: study -> studies.", Inches(8.5), Inches(2.5), Inches(4.3), Inches(4.2), bg_color=ROXO_CLARO, border_color=ROXO, tit_size=16, desc_size=13)

# Slide 7
sl7 = base_slide(7)
add_header(sl7, "Preguntas & Adverbios de Frecuencia", "Estrategia Multinivel")
add_card(sl7, "Contenido Principal (Inicial / Básico):", "Para hacer preguntas usamos DO o DOES:\n• Do you live in Spain?\n• Does she work here?\n\nAdverbios de Frecuencia:\nAlways (100%), Usually (80%), Sometimes (50%), Never (0%).\nEx: I always drink coffee in the morning.", Inches(0.5), Inches(2.5), Inches(7.8), Inches(4.2), tit_size=18, desc_size=14)
add_card(sl7, "Nuances & Excepciones\n(Independiente / Avanzado):", "Posición del Adverbio:\n• Antes del verbo principal: She seldom eats fast food.\n• Después del verbo TO BE: He is always late.\n• Word order con expresiones de tiempo: Every day va al final.", Inches(8.5), Inches(2.5), Inches(4.3), Inches(4.2), bg_color=ROXO_CLARO, border_color=ROXO, tit_size=16, desc_size=13)

# Slide 8
sl8 = base_slide(8)
add_header(sl8, "Stative Verbs vs Action Verbs", "Estrategia Multinivel")
add_card(sl8, "Contenido Principal (Inicial / Básico):", "Diferencia entre expresar acciones y estados de ánimo / posesión:\n\n• Action: I run, she works, we study.\n• State/Feeling: I love pizza, she needs help, they have a car.", Inches(0.5), Inches(2.5), Inches(7.8), Inches(4.2), tit_size=18, desc_size=14)
add_card(sl8, "Nuances & Excepciones\n(Independiente / Avanzado):", "Stative Verbs NO se usan normalmente en tiempos continuos (-ing):\n• Correct: I know the answer.\n• Incorrect: I am knowing the answer.\n• Future schedules: The train leaves at 8 PM.", Inches(8.5), Inches(2.5), Inches(4.3), Inches(4.2), bg_color=ROXO_CLARO, border_color=ROXO, tit_size=16, desc_size=13)

# SLIDE 9: SÍNTESIS + SUGESTIÓN CULTURAL
sl9 = base_slide(9)
add_header(sl9, "Síntense del Contenido & Sugerencia Cultural", "Resumen y Recomendación")
add_card(sl9, "Resumen de la Regla de Oro:", "• Rutina/Hecho -> Simple Present.\n• Recuerda agregar -s/-es a he/she/it.\n• Usa Do/Does para negar y preguntar.", Inches(0.5), Inches(2.5), Inches(12.3), Inches(2.3), tit_size=18, desc_size=15)
add_card(sl9, "Sugerencia Cultural:", "🎵 Canción recomendada: The Beatles - 'A Hard Day's Night'.\nPresta atención a cómo la letra usa el Simple Present para describir hábitos diarios.", Inches(0.5), Inches(5.0), Inches(12.3), Inches(1.7), bg_color=FUNDO_CARD, border_color=AZUL, tit_size=16, desc_size=14)

# SLIDES 10 E 11: PRÁTICA GUIADA
sl10 = base_slide(10)
add_header(sl10, "Práctica Guiada — Ejercicios 1 a 5", "Máximo 10 Ejercicios Guiados (Corrección de Vicios)")
exs_1_5 = [
    ("Ejercicio 1:", "She ______ (work / works) in a bank. — Inicial / Descubrimiento"),
    ("Ejercicio 2:", "They ______ (don't / doesn't) like tea. — Básico / Conexión"),
    ("Ejercicio 3:", "______ you speak French? (Do / Does) — Básico / Conexión"),
    ("Ejercicio 4:", "Rewrite: 'He study English on Mondays.' (Fix the verb) — Intermedio / Autónomo"),
    ("Ejercicio 5:", "Change to negative: 'My sister lives in London.' — Intermedio / Autónomo")
]
top_e = 2.4
for t_e, d_e in exs_1_5:
    add_card(sl10, t_e, d_e, Inches(0.5), Inches(top_e), Inches(12.3), Inches(0.85), tit_size=15, desc_size=14)
    top_e += 0.92

sl11 = base_slide(11)
add_header(sl11, "Prática Guiada — Ejercicios 6 a 10", "Desafíos para Niveles Superiores")
exs_6_10 = [
    ("Ejercicio 6:", "Place the adverb: 'She goes to the gym.' (always) — Intermedio / Autónomo"),
    ("Ejercicio 7:", "Which sentence is correct? A) He don't have time. B) He doesn't have time. — Independiente"),
    ("Ejercicio 8:", "Replace with a formal synonym: 'He starts the meeting at 9.' — Independente"),
    ("Ejercicio 9:", "Find and fix the error: 'I am understanding this rule now.' — Avanzado / Especialista"),
    ("Ejercicio 10:", "Explain the nuance: 'I think you are right' vs 'I am thinking about it'. — Avanzado / Especialista")
]
top_e = 2.4
for t_e, d_e in exs_6_10:
    add_card(sl11, t_e, d_e, Inches(0.5), Inches(top_e), Inches(12.3), Inches(0.85), tit_size=15, desc_size=14)
    top_e += 0.92

# SLIDES 12 A 19: CONVERSACIÓN (8 PREGUNTAS EN ESPAÑOL AUXILIAR)
perguntas_config = [
    ("Pregunta 1:", "What time do you usually wake up on weekends? — Inicial / Descubrimiento"),
    ("Pregunta 2:", "What do you usually eat for breakfast? — Básico / Conexión"),
    ("Pregunta 3:", "Where do you go when you want to relax? — Básico / Conexión"),
    ("Pregunta 4:", "How do you manage your daily routine when you have too much work? — Intermedio / Autónomo"),
    ("Pregunta 5:", "What habits do you think are essential for a healthy lifestyle? — Intermedio / Autónomo"),
    ("Pregunta 6:", "How does technology affect the daily routines of modern society? — Independiente"),
    ("Pregunta 7:", "In what ways do cultural differences shape people's daily habits around the world? — Independiente"),
    ("Pregunta 8:", "Debate: 'Habits define a person's destiny more than talent.' Do you agree? — Avanzado / Especialista"),
]

for idx in range(8):
    s_num = 12 + idx
    sl = base_slide(s_num)
    add_header(sl, "Conversación", f"Pregunta {idx + 1} de 8")
    
    if idx == 0:
        p_tit, p_desc = perguntas_config[0]
        add_card(sl, p_tit, p_desc, Inches(0.5), Inches(2.3), Inches(12.3), Inches(2.4), bg_color=ROXO_CLARO, border_color=ROXO, tit_size=22, desc_size=16)
        add_card(sl, "Sub-prompts Visuales (Apoyo al Alumno):", "• Para expandir: 'Give a recent example...'\n• Para estar de acuerdo/en desacuerdo: 'What is the other side of the coin?'", Inches(0.5), Inches(4.9), Inches(12.3), Inches(1.9), tit_size=16, desc_size=14)
    else:
        history_list = [perguntas_config[i][1] for i in range(idx)]
        add_history_cards(sl, history_list, top=2.25, height=1.05)
        
        p_tit, p_desc = perguntas_config[idx]
        add_card(sl, p_tit, p_desc, Inches(0.5), Inches(3.45), Inches(12.3), Inches(1.9), bg_color=ROXO_CLARO, border_color=ROXO, tit_size=18, desc_size=15)
        add_card(sl, "Sub-prompts Visuales (Apoyo al Alumno):", "• Para expandir: 'Give a recent example...'\n• Para estar de acuerdo/en desacuerdo: 'What is the other side of the coin?'", Inches(0.5), Inches(5.5), Inches(12.3), Inches(1.35), tit_size=15, desc_size=13)

# SLIDE 20: PARTE ESCRITA
sl20 = base_slide(20)
add_header(sl20, "Actividad de Escritura", "Consolidación y Memorización")
add_card(sl20, "Tarea de Escritura Recomendada:", "Escribe un párrafo corto (4 a 6 oraciones) describiendo tu rutina ideal de fin de semana.\nRegra: Utiliza al menos 4 verbos en Simple Present y 2 adverbios de frecuencia (always, usually, etc.).", Inches(0.5), Inches(2.4), Inches(12.3), Inches(2.6), tit_size=18, desc_size=15)
add_card(sl20, "", "Identificamos que la escritura en esta etapa es fundamental para fijar las estructuras aprendidas, expandir el vocabulario activo y disminuir la traducción mental automática de tu idioma nativo.", Inches(0.5), Inches(5.2), Inches(12.3), Inches(1.4), bg_color=ROXO_CLARO, border_color=ROXO, desc_size=14)

# SLIDE 21: TAREA PARA LA CASA
sl21 = base_slide(21)
add_header(sl21, "Cierre — Tarea para la Casa", "Acción y Práctica Continua")
add_card(sl21, "Tarea para la Casa:", "• Listening: Escucha una canción en inglés y escribe 3 oraciones que usen el Simple Present.\n• Speaking: Graba un audio de 1 minuto en WhatsApp describiendo la rutina de un familiar (usando He / She).", Inches(0.5), Inches(2.5), Inches(12.3), Inches(4.2), tit_size=18, desc_size=15)

# SLIDE 22: CIERRE & FEEDBACK
sl22 = base_slide(22)
add_header(sl22, "Cierre & Feedback Rápido", "Consideraciones Finales")
add_card(sl22, "Agradecimiento:", "¡Muchas gracias por tu dedicación en la clase de hoy!", Inches(0.5), Inches(2.5), Inches(12.3), Inches(1.5), tit_size=20, desc_size=16)
add_card(sl22, "Espacio para Feedback Rápido:", "[Espacio destinado para que el profesor recopile dudas rápidas y comentarios del alumno]", Inches(0.5), Inches(4.3), Inches(12.3), Inches(2.4), bg_color=FUNDO_CARD, border_color=TURQUESA, tit_size=18, desc_size=15)

# SALVAR SIMULAÇÃO
out_path = "/var/www/haas-frontend-desk-mobile-oficial/public/temp-miniaturas/aula-simple-present-espanol.pptx"
prs.save(out_path)
print(f"🚀 SIMULACIÓN COMPLETADA CON ÉXITO: {out_path}")
