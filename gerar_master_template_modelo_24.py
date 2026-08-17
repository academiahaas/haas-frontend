import os
import shutil
import urllib.request
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# PALETA OFICIAL HAAS LANGUAGE - MARQUEE RIBBON EDITION
BRANCO = RGBColor(0xFF, 0xFF, 0xFF)
ROXO = RGBColor(0x4F, 0x35, 0xE2)
AZUL = RGBColor(0x1E, 0x3A, 0x8A)
SLATE = RGBColor(0x33, 0x41, 0x55)
TURQUESA = RGBColor(0x14, 0xC8, 0xB0)
ROXO_CLARO = RGBColor(0xED, 0xE9, 0xFE)
TURQUESA_CLARO = RGBColor(0xE0, 0xF2, 0xF1)
FUNDO_BG = RGBColor(0xF8, 0xFA, 0xFC)

LOGO_URL = "https://jdppxfokfhqjudwfwckd.supabase.co/storage/v1/object/public/haas-academy/assignments/Design%20sem%20nome%20(4).png"
LOGO_PATH = "/tmp/haas-logo.png"

if not os.path.exists(LOGO_PATH):
    try:
        urllib.request.urlretrieve(LOGO_URL, LOGO_PATH)
    except Exception:
        pass

W, H = Inches(13.333), Inches(7.5)
TOTAL_SLIDES = 22

prs = Presentation()
prs.slide_width = W
prs.slide_height = H

def base_slide(num, hide_logo=False):
    sl = prs.slides.add_slide(prs.slide_layouts[6])
    sl.background.fill.solid()
    sl.background.fill.fore_color.rgb = FUNDO_BG

    if os.path.exists(LOGO_PATH) and not hide_logo:
        sl.shapes.add_picture(LOGO_PATH, Inches(10.8), Inches(0.2), height=Inches(0.65))
        
    tb = sl.shapes.add_textbox(Inches(0.5), Inches(7.12), Inches(8), Inches(0.35))
    p = tb.text_frame.paragraphs[0]
    p.text = "HAAS LANGUAGE  •  MARQUEE RIBBON EDITION 🎟️"
    p.font.size = Pt(10); p.font.color.rgb = ROXO; p.font.bold = True
    
    tb2 = sl.shapes.add_textbox(Inches(10.8), Inches(7.12), Inches(2.0), Inches(0.35))
    p2 = tb2.text_frame.paragraphs[0]
    p2.text = f"{num:02d} // {TOTAL_SLIDES:02d}"
    p2.font.size = Pt(10); p2.font.color.rgb = AZUL; p2.font.bold = True; p2.alignment = PP_ALIGN.RIGHT
    return sl

def add_header(sl, tit, sub=""):
    tb = sl.shapes.add_textbox(Inches(0.5), Inches(0.25), Inches(10), Inches(0.65))
    tf = tb.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]; p.text = tit
    p.font.size = Pt(30); p.font.color.rgb = ROXO; p.font.bold = True
    
    if sub:
        tb_s = sl.shapes.add_textbox(Inches(0.5), Inches(0.95), Inches(10), Inches(0.4))
        p_s = tb_s.text_frame.paragraphs[0]
        p_s.text = sub
        p_s.font.size = Pt(16); p_s.font.color.rgb = AZUL; p_s.font.bold = True

def add_marquee_strip(sl, text, top, height=0.42, bg_c=ROXO, txt_c=TURQUESA, font_sz=11):
    """Cria uma faixa de transmissão continua estilo Marquee Ticker Tape de ponta a ponta"""
    s = sl.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(top), W, Inches(height))
    s.fill.solid(); s.fill.fore_color.rgb = bg_c
    s.line.fill.background()
    
    tf = s.text_frame; tf.word_wrap = False
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    tf.margin_left = Inches(0.1); tf.margin_right = Inches(0.1)
    
    repeated = (text + "   ★   ") * 5
    p = tf.paragraphs[0]
    p.text = repeated
    p.font.size = Pt(font_sz); p.font.color.rgb = txt_c; p.font.bold = True; p.alignment = PP_ALIGN.LEFT

def add_marquee_card(sl, tit, desc, left, top, width, height, bg_c=BRANCO, border_c=ROXO, tit_c=ROXO, desc_c=SLATE, tit_sz=16, desc_sz=14, ribbon_label="", align_center_v=False):
    s = sl.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(left), Inches(top), Inches(width), Inches(height))
    s.fill.solid(); s.fill.fore_color.rgb = bg_c
    s.line.color.rgb = border_c; s.line.width = Pt(1.5)
    
    if ribbon_label:
        rib = sl.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(left), Inches(top), Inches(width), Inches(0.32))
        rib.fill.solid(); rib.fill.fore_color.rgb = border_c; rib.line.fill.background()
        tf_r = rib.text_frame; tf_r.vertical_anchor = MSO_ANCHOR.MIDDLE
        p_r = tf_r.paragraphs[0]; p_r.text = f">>> {ribbon_label.upper()} >>>"
        p_r.font.size = Pt(10); p_r.font.color.rgb = BRANCO; p_r.font.bold = True
        
    tf = s.text_frame; tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE if align_center_v else MSO_ANCHOR.TOP
    margin_t = Inches(0.4) if ribbon_label else Inches(0.2)
    tf.margin_left = Inches(0.25); tf.margin_top = margin_t; tf.margin_right = Inches(0.25); tf.margin_bottom = Inches(0.2)
    
    if tit:
        p = tf.paragraphs[0]
        p.text = tit; p.font.size = Pt(tit_sz); p.font.color.rgb = tit_c; p.font.bold = True
        if desc:
            p2 = tf.add_paragraph()
            p2.text = desc; p2.font.size = Pt(desc_sz); p2.font.color.rgb = desc_c; p2.space_before = Pt(6)
    else:
        if desc:
            p = tf.paragraphs[0]
            p.text = desc; p.font.size = Pt(desc_sz); p.font.color.rgb = desc_c

def img_marquee_box(sl, left, top, width, height, label="[MEDIA TICKER]"):
    s = sl.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(left), Inches(top), Inches(width), Inches(height))
    s.fill.solid(); s.fill.fore_color.rgb = ROXO_CLARO
    s.line.color.rgb = TURQUESA; s.line.width = Pt(2.0)
    
    # Header de esteira na imagem
    h_rib = sl.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(left), Inches(top), Inches(width), Inches(0.3))
    h_rib.fill.solid(); h_rib.fill.fore_color.rgb = TURQUESA; h_rib.line.fill.background()
    tf_h = h_rib.text_frame; tf_h.vertical_anchor = MSO_ANCHOR.MIDDLE
    p_h = tf_h.paragraphs[0]; p_h.text = "/// MARQUEE MEDIA STREAM ///"; p_h.alignment = PP_ALIGN.CENTER
    p_h.font.size = Pt(9); p_h.font.color.rgb = AZUL; p_h.font.bold = True
    
    tf = s.text_frame; tf.word_wrap = True; tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER
    p = tf.paragraphs[0]; p.text = label; p.font.size = Pt(16); p.font.color.rgb = ROXO; p.font.bold = True

def add_marquee_flashcard(sl, num_str, desc, left, top, width, height, border_c=ROXO):
    s = sl.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(left), Inches(top), Inches(width), Inches(height))
    s.fill.solid(); s.fill.fore_color.rgb = BRANCO
    s.line.color.rgb = border_c; s.line.width = Pt(1.5)

    chip = sl.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(left), Inches(top), Inches(width), Inches(0.28))
    chip.fill.solid(); chip.fill.fore_color.rgb = border_c; chip.line.fill.background()
    tf_c = chip.text_frame; tf_c.vertical_anchor = MSO_ANCHOR.MIDDLE
    p_c = tf_c.paragraphs[0]; p_c.text = f">>> TASK #{num_str} >>> FLUENCY CHECK >>>"
    p_c.font.size = Pt(9); p_c.font.color.rgb = BRANCO; p_c.font.bold = True

    tf = s.text_frame; tf.word_wrap = True; tf.vertical_anchor = MSO_ANCHOR.TOP
    tf.margin_left = Inches(0.25); tf.margin_top = Inches(0.35); tf.margin_right = Inches(0.15); tf.margin_bottom = Inches(0.1)
    p = tf.paragraphs[0]; p.text = desc; p.font.size = Pt(14); p.font.color.rgb = SLATE; p.font.bold = False

# ==========================================================
# CONSTRUÇÃO DOS 22 SLIDES - MODELO 24 (MARQUEE RIBBON)
# ==========================================================

# SLIDE 1: CAPA
sl1 = base_slide(1, hide_logo=True)
add_marquee_strip(sl1, "HAAS LANGUAGE ★ HIGH-LEVEL EDUCATION ★ SPEAKING MODULE ★ ACTIVE FLUENCY", 0.9, height=0.45, bg_c=TURQUESA, txt_c=AZUL)
add_marquee_card(sl1, "[Título da Aula / Tema Principal]", "[Subtítulo / Descrição da Aula]\n\n📅 Data da Aula: [DD/MM/AAAA]", 0.5, 1.6, 7.2, 5.0, bg_c=BRANCO, border_c=ROXO, tit_sz=36, desc_sz=18, ribbon_label="MAIN BROADCAST")
img_marquee_box(sl1, 8.0, 1.6, 4.8, 5.0, "[Poster Marquee Hero]")

# SLIDE 2: PRESENÇA
sl2 = base_slide(2)
add_header(sl2, "Presença & Integrantes", "Membros em Transmissão")
add_marquee_strip(sl2, "SESSION PARTICIPANTS /// ACTIVE TEACHERS AND STUDENTS /// HAAS ACADEMY", 1.4, height=0.35, bg_c=AZUL, txt_c=TURQUESA)
add_marquee_card(sl2, "Professor(a):", "[Nome do(a) Professor(a)]", 0.5, 1.9, 12.333, 2.3, bg_c=BRANCO, border_c=ROXO, tit_sz=20, desc_sz=16, ribbon_label="INSTRUCTOR LOG")
add_marquee_card(sl2, "Aluno(s):", "[Nome do Aluno / Lista da Turma]", 0.5, 4.4, 12.333, 2.5, bg_c=TURQUESA_CLARO, border_c=TURQUESA, tit_sz=20, desc_sz=16, ribbon_label="STUDENTS STREAM")

# SLIDE 3: WARM-UP
sl3 = base_slide(3)
add_header(sl3, "Aquecimento / Warm-up", "Foco 100% na Fala (3 Minutos)")
add_marquee_card(sl3, "Pergunta Provocativa Rápida:", "[Insira aqui a pergunta provocativa de aquecimento]", 0.5, 1.5, 7.0, 2.6, bg_c=BRANCO, border_c=ROXO, tit_sz=20, desc_sz=16, ribbon_label="SPEECH PROMPT")
add_marquee_card(sl3, "Fundamentação Neurocognitiva:", "Ativa o córtex pré-frontal e reduz a inibição linguística inicial.", 0.5, 4.3, 7.0, 2.4, bg_c=TURQUESA_CLARO, border_c=TURQUESA, tit_sz=18, desc_sz=15, ribbon_label="NEUROSCIENCE LOG", align_center_v=True)
img_marquee_box(sl3, 7.8, 1.5, 5.0, 5.2, "[Imagem Warm-up]")

# SLIDE 4: OBJETIVOS
sl4 = base_slide(4)
add_header(sl4, "Objetivos da Aula", "Metas da Transmissão")
add_marquee_card(sl4, "Goal #01", "[Objetivo 1]", 0.5, 1.5, 7.0, 1.5, bg_c=BRANCO, border_c=ROXO, tit_sz=14, desc_sz=15, align_center_v=True)
add_marquee_card(sl4, "Goal #02", "[Objetivo 2]", 0.5, 3.2, 7.0, 1.5, bg_c=BRANCO, border_c=TURQUESA, tit_sz=14, desc_sz=15, align_center_v=True)
add_marquee_card(sl4, "Goal #03", "[Objetivo 3]", 0.5, 4.9, 7.0, 1.8, bg_c=BRANCO, border_c=AZUL, tit_sz=14, desc_sz=15, align_center_v=True)
add_marquee_card(sl4, "Nota de Escrita:", "Nesta aula daremos atenção à escrita para reforçar a memorização de estruturas/vocabulário e reduzir a interferência da língua materna.", 7.8, 1.5, 5.0, 5.2, bg_c=ROXO_CLARO, border_c=ROXO, tit_sz=18, desc_sz=15, ribbon_label="WRITING FOCUS")

# SLIDE 5: METODOLOGIA (5 MARQUEE CONVEYOR TRACKS)
sl5 = base_slide(5)
add_header(sl5, "A Nossa Metodologia", "Esteiras de Aprendizado Haas Language")
met_tracks = [
    ("01. FALA ATIVA", "Conversação inicial sem julgamentos.", ROXO, TURQUESA),
    ("02. EXPOSIÇÃO", "Apresentação leve da gramática.", AZUL, TURQUESA),
    ("03. PRÁTICA GUIADA", "Exercícios com suporte do professor.", TURQUESA, AZUL),
    ("04. CONVERSAÇÃO", "Aplicação prática das estruturas em contexto real.", ROXO, BRANCO),
    ("05. ESCRITA", "Consolidação visual e memorização de vocabulário.", AZUL, TURQUESA)
]
top_y = 1.5
for tit, desc, bg, txt_c in met_tracks:
    add_marquee_strip(sl5, f">>> {tit} /// {desc} /// HAAS PIPELINE", top_y, height=0.9, bg_c=bg, txt_c=txt_c, font_sz=13)
    top_y += 1.08

# SLIDES 6 A 8: EXPOSIÇÃO
for s_num in [6, 7, 8]:
    sl = base_slide(s_num)
    add_header(sl, f"Explicação do Tema — Parte {s_num - 5}")
    add_marquee_card(sl, "Conteúdo Principal:", f"[Explicação e exemplo central da estrutura - Bloco {s_num - 5}]", 0.5, 1.5, 7.0, 5.2, bg_c=BRANCO, border_c=ROXO, tit_sz=20, desc_sz=15, ribbon_label="CORE STRUCTURE")
    add_marquee_card(sl, "Exemplos Práticos:", "[Exemplos adicionais e frases do dia a dia]", 7.8, 1.5, 5.0, 2.5, bg_c=TURQUESA_CLARO, border_c=TURQUESA, ribbon_label="EXAMPLES STREAM")
    add_marquee_card(sl, "Notas & Detalhes:", "[Exceções, nuances gramaticais e contexto de uso]", 7.8, 4.2, 5.0, 2.5, bg_c=ROXO_CLARO, border_c=AZUL, ribbon_label="RULES LOG")

# SLIDE 9: SÍNTESE + CULTURA
sl9 = base_slide(9)
add_header(sl9, "Síntese do Conteúdo & Sugestão Cultural")
add_marquee_card(sl9, "Resumo do Tema:", "[Resumo sintético dos pontos chave ensinados na aula]", 0.5, 1.5, 5.8, 5.2, bg_c=BRANCO, border_c=ROXO, tit_sz=18, desc_sz=15, ribbon_label="SUMMARY TICKER")
add_marquee_card(sl9, "Sugestão Cultural:", "🎵 [Nome da Música / Vídeo recomendado sobre o assunto].\n\nO professor pode rodar o áudio/vídeo ou deixar como recomendação extra para o aluno praticar.", 6.6, 1.5, 6.2, 5.2, bg_c=TURQUESA_CLARO, border_c=TURQUESA, tit_sz=18, desc_sz=15, ribbon_label="CULTURAL MEDIA STREAM")

# SLIDES 10 E 11: PRÁTICA GUIADA
sl10 = base_slide(10)
add_header(sl10, "Prática Guiada — Blocos 1 a 5", "Correção de Vícios")
add_marquee_flashcard(sl10, "01", "[Lacunas contextualizadas com opções]", 0.5, 1.5, 12.333, 1.2, border_c=ROXO)
add_marquee_flashcard(sl10, "02", "[Lacunas contextualizadas com opções]", 0.5, 2.9, 6.0, 1.8, border_c=TURQUESA)
add_marquee_flashcard(sl10, "03", "[Preenchimento direto de verbos]", 6.8, 2.9, 6.0, 1.8, border_c=TURQUESA)
add_marquee_flashcard(sl10, "04", "[Reestruturação / Transformação de frase]", 0.5, 4.9, 6.0, 1.8, border_c=ROXO)
add_marquee_flashcard(sl10, "05", "[Reestruturação / Transformação de frase]", 6.8, 4.9, 6.0, 1.8, border_c=ROXO)

sl11 = base_slide(11)
add_header(sl11, "Prática Guiada — Blocos 6 a 10", "Desafios Práticos")
add_marquee_flashcard(sl11, "06", "[Preposição / Conector correto]", 0.5, 1.5, 12.333, 1.2, border_c=TURQUESA)
add_marquee_flashcard(sl11, "07", "[Escolha da opção correta]", 0.5, 2.9, 6.0, 1.8, border_c=ROXO)
add_marquee_flashcard(sl11, "08", "[Substituição por sinônimo formal]", 6.8, 2.9, 6.0, 1.8, border_c=ROXO)
add_marquee_flashcard(sl11, "09", "[Identificação de erro e correção]", 0.5, 4.9, 6.0, 1.8, border_c=TURQUESA)
add_marquee_flashcard(sl11, "10", "[Tradução de sentido sem uso literal]", 6.8, 4.9, 6.0, 1.8, border_c=TURQUESA)

# ==========================================================
# SLIDES 12 A 19: CONVERSAÇÃO (INOVAÇÃO MARQUEE STREAM STACK INÉDITA!)
# TOPO: FAIXA MARQUEE DE TRANSMISSÃO DA PERGUNTA ATIVA
# ESQUERDA: HERO ACTIVE QUESTION (X = 0.5", W = 6.8") + SUPPORT PROMPTS
# DIREITA: PILHA DE FITA MARQUEE DE HISTÓRICO (X = 7.6", W = 5.2")
# ==========================================================
perguntas_config = [
    "[Pergunta principal 1 sobre o tema]", "[Pergunta principal 2 sobre o tema]", "[Pergunta principal 3 sobre o tema]",
    "[Pergunta principal 4 sobre o tema]", "[Pergunta principal 5 sobre o tema]", "[Pergunta principal 6 sobre o tema]",
    "[Pergunta principal 7 sobre o tema]", "[Pergunta principal 8 sobre o tema]"
]

def add_marquee_history_ribbon_stack(sl, history_list):
    """Exibe o histórico de perguntas passadas como faixas de esteira Marquee empilhadas na direita"""
    count = len(history_list)
    if count == 0: return
    
    gap = 0.12
    avail_h = 5.2 - (gap * (count - 1))
    card_h = avail_h / count
    
    bg_colors = [ROXO, TURQUESA, AZUL, ROXO, TURQUESA, AZUL, ROXO]
    txt_colors = [TURQUESA, AZUL, BRANCO, TURQUESA, AZUL, BRANCO, TURQUESA]
    
    for i, h_txt in enumerate(history_list):
        y_pos = 1.5 + i * (card_h + gap)
        clean_t = h_txt[:32] + "..." if len(h_txt) > 32 else h_txt
        bg_c = bg_colors[i % len(bg_colors)]
        tx_c = txt_colors[i % len(txt_colors)]
        d_sz = 11 if count <= 3 else (10 if count <= 5 else 9)
        
        # Fita Marquee de Histórico
        add_marquee_card(sl, f"Q{i+1:02d} STREAM", clean_t, 7.6, y_pos, 5.2, card_h, bg_c=BRANCO, border_c=bg_c, tit_sz=11, desc_sz=d_sz, align_center_v=True, ribbon_label=f"PASSED Q0{i+1}")

for idx in range(8):
    s_num = 12 + idx
    sl = base_slide(s_num)
    add_header(sl, "Conversação", f"Thread {idx + 1} de 8")
    p_desc = perguntas_config[idx]
    
    if idx == 0:
        # Pergunta 1 Ampla Hero
        add_marquee_card(sl, f"PERGUNTA ATIVA Q0{idx+1}:", p_desc, 0.5, 1.5, 12.333, 2.6, bg_c=BRANCO, border_c=ROXO, tit_sz=18, desc_sz=18, ribbon_label="LIVE BROADCAST", align_center_v=True)
        add_marquee_card(sl, "PROMPTS DE APOIO AO PROFESSOR:", "• Para expandir: 'Dê um exemplo recente...'\n• Para concordar/discordar: 'Qual é o outro lado dessa moeda?'", 0.5, 4.3, 12.333, 2.4, bg_c=TURQUESA_CLARO, border_c=TURQUESA, tit_sz=16, desc_sz=14, ribbon_label="TEACHER STREAM", align_center_v=True)
    else:
        # LADO ESQUERDO: Pergunta Ativa + Prompts
        add_marquee_card(sl, f"PERGUNTA ATIVA Q0{idx+1}:", p_desc, 0.5, 1.5, 6.8, 2.6, bg_c=BRANCO, border_c=ROXO, tit_sz=18, desc_sz=16, ribbon_label=f"LIVE BROADCAST Q0{idx+1}", align_center_v=True)
        add_marquee_card(sl, "PROMPTS DE APOIO:", "• Para expandir: 'Dê um exemplo recente...'\n• Para discordar: 'Qual é o outro lado?'", 0.5, 4.3, 6.8, 2.4, bg_c=TURQUESA_CLARO, border_c=TURQUESA, tit_sz=15, desc_sz=14, ribbon_label="TEACHER STREAM", align_center_v=True)
        
        # LADO DIREITO: Esteira Marquee de Histórico Empilhada!
        history_list = [perguntas_config[i] for i in range(idx)]
        add_marquee_history_ribbon_stack(sl, history_list)

# SLIDE 20: ESCRITA
sl20 = base_slide(20)
add_header(sl20, "Atividade de Escrita", "Consolidação e Memorização")
add_marquee_card(sl20, "Tarefa de Escrita Recomendada:", "[Instrução detalhada da tarefa: escrever um texto, frases ou diálogo]", 0.5, 1.5, 12.333, 2.6, bg_c=BRANCO, border_c=ROXO, tit_sz=18, desc_sz=15, ribbon_label="WRITING TASK STREAM")
add_marquee_card(sl20, "Fundamentação Pedagógica:", "Identificamos que a escrita nesta etapa é fundamental para você fixar as estruturas aprendidas, expandir vocabulário ativo e diminuir a tradução mental automática do seu idioma nativo.", 0.5, 4.3, 12.333, 2.4, bg_c=TURQUESA_CLARO, border_c=TURQUESA, desc_sz=14, ribbon_label="PEDAGOGICAL STREAM", align_center_v=True)

# SLIDE 21: TAREFA DE CASA
sl21 = base_slide(21)
add_header(sl21, "Encerramento — Tarefa de Casa", "Ação e Prática Contínua")
add_marquee_card(sl21, "Tarefa de Casa:", "• [Atividade prática 1 para realizar fora da aula]\n• [Atividade de fixação de vocabulário ou escuta]", 0.5, 1.5, 12.333, 5.2, bg_c=BRANCO, border_c=ROXO, tit_sz=18, desc_sz=15, ribbon_label="HOMEWORK STREAM")

# SLIDE 22: ENCERRAMENTO & FEEDBACK
sl22 = base_slide(22)
add_header(sl22, "Encerramento & Feedback Rápido", "Considerações Finais")
add_marquee_card(sl22, "Agradecimento:", "Obrigado pela dedicação na aula de hoje!", 0.5, 1.5, 12.333, 2.2, bg_c=ROXO_CLARO, border_c=ROXO, tit_sz=20, desc_sz=16, ribbon_label="CLOSING STREAM")
add_marquee_card(sl22, "Espaço para Feedback Rápido:", "[Espaço destinado para o professor coletar dúvidas rápidas e considerações do aluno]", 0.5, 3.9, 12.333, 2.8, bg_c=TURQUESA_CLARO, border_c=TURQUESA, tit_sz=18, desc_sz=15, ribbon_label="FEEDBACK STREAM")

# SALVAR E COPIAR
out_path_local = "template-haas-modelo-24-marquee.pptx"
prs.save(out_path_local)

server_dir = "/var/www/haas-frontend-desk-mobile-oficial/public/temp-miniaturas/"
os.makedirs(server_dir, exist_ok=True)
out_path_server = os.path.join(server_dir, "template-haas-modelo-24-marquee.pptx")

shutil.copyfile(out_path_local, out_path_server)
print(f"🎉 MODELO 24 MARQUEE RIBBON GERADO EM: {out_path_server}")
