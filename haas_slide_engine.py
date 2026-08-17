from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
import os
import json

# PALETA OFICIAL HAAS ACADEMY
BRANCO = RGBColor(0xFF, 0xFF, 0xFF)
ROXO = RGBColor(0x4F, 0x35, 0xE2)
AZUL = RGBColor(0x1E, 0x3A, 0x8A)
SLATE = RGBColor(0x33, 0x41, 0x55)
TURQUESA = RGBColor(0x14, 0xC8, 0xB0)
BORDA = RGBColor(0xCB, 0xD5, 0xE1)
FUNDO_CARD = RGBColor(0xF8, 0xFA, 0xFC)
ROXO_CLARO = RGBColor(0xED, 0xE9, 0xFE)
LOGO = "/tmp/haas-logo.png"

W, H = Inches(13.333), Inches(7.5)

class HaasPresentationEngine:
    def __init__(self):
        self.prs = Presentation()
        self.prs.slide_width = W
        self.prs.slide_height = H

    def _slide_base(self, numero_slide):
        sl = self.prs.slides.add_slide(self.prs.slide_layouts[6])
        
        # Fundo Branco
        sl.background.fill.solid()
        sl.background.fill.fore_color.rgb = BRANCO
        
        # Lateral Roxo Dominante
        s = sl.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(0.12), H)
        s.fill.solid(); s.fill.fore_color.rgb = ROXO; s.line.fill.background()
        
        # Barra Rodapé
        b = sl.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(7.1), W, Inches(0.4))
        b.fill.solid(); b.fill.fore_color.rgb = ROXO; b.line.fill.background()
        
        # Logo Haas
        if os.path.exists(LOGO):
            sl.shapes.add_picture(LOGO, Inches(0.4), Inches(0.25), width=Inches(1.5))
            
        # Detalhe Turquesa (5% Acento)
        det = sl.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.5), Inches(1.1), Inches(1.2), Pt(4))
        det.fill.solid(); det.fill.fore_color.rgb = TURQUESA; det.line.fill.background()
        
        # Círculo decorativo
        c = sl.shapes.add_shape(MSO_SHAPE.OVAL, Inches(12.5), Inches(0.5), Inches(0.4), Inches(0.4))
        c.fill.solid(); c.fill.fore_color.rgb = TURQUESA; c.line.fill.background()
        
        # Rodapé Texto + Contador de Slides
        tb = sl.shapes.add_textbox(Inches(0.5), Inches(7.15), Inches(10), Inches(0.3))
        tb.text_frame.paragraphs[0].text = "Haas Academia de Idiomas"
        tb.text_frame.paragraphs[0].font.size = Pt(10)
        tb.text_frame.paragraphs[0].font.color.rgb = BRANCO
        
        tb_num = sl.shapes.add_textbox(Inches(11.5), Inches(7.15), Inches(1.3), Inches(0.3))
        p = tb_num.text_frame.paragraphs[0]
        p.text = f"{numero_slide} / 18"
        p.font.size = Pt(10); p.font.color.rgb = BRANCO; p.alignment = PP_ALIGN.RIGHT
        
        return sl

    def add_header(self, slide, titulo_txt, subtitulo_txt=""):
        # Título Roxo
        tb = slide.shapes.add_textbox(Inches(0.5), Inches(1.3), Inches(12), Inches(0.7))
        tf = tb.text_frame; tf.word_wrap = True
        p = tf.paragraphs[0]; p.text = titulo_txt
        p.font.size = Pt(32); p.font.color.rgb = ROXO; p.font.bold = True
        
        # Subtítulo Azul Profundo
        if subtitulo_txt:
            tb_s = slide.shapes.add_textbox(Inches(0.5), Inches(2.0), Inches(12), Inches(0.5))
            tf_s = tb_s.text_frame; tf_s.word_wrap = True
            p_s = tf_s.paragraphs[0]; p_s.text = subtitulo_txt
            p_s.font.size = Pt(18); p_s.font.color.rgb = AZUL; p_s.font.bold = True

    def add_card(self, slide, tit, corpo, left, top, width, height):
        s = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
        s.fill.solid(); s.fill.fore_color.rgb = FUNDO_CARD
        s.line.color.rgb = BORDA; s.line.width = Pt(1)
        tf = s.text_frame; tf.word_wrap = True; tf.margin_left = Inches(0.2); tf.margin_top = Inches(0.15)
        p = tf.paragraphs[0]; p.text = tit; p.font.size = Pt(14); p.font.color.rgb = ROXO; p.font.bold = True
        p2 = tf.add_paragraph(); p2.text = corpo; p2.font.size = Pt(12); p2.font.color.rgb = SLATE; p2.space_before = Pt(4)

    def add_bullet(self, slide, texto, top, level_tag=""):
        tb = slide.shapes.add_textbox(Inches(0.5), top, Inches(12), Inches(0.4))
        tf = tb.text_frame; tf.word_wrap = True; p = tf.paragraphs[0]
        
        r1 = p.add_run(); r1.text = "● "
        r1.font.size = Pt(14); r1.font.color.rgb = TURQUESA
        
        r2 = p.add_run(); r2.text = texto
        r2.font.size = Pt(14); r2.font.color.rgb = SLATE
        
        if level_tag:
            r3 = p.add_run(); r3.text = f" — [{level_tag}]"
            r3.font.size = Pt(12); r3.font.color.rgb = TURQUESA; r3.font.bold = True

    def build_presentation(self, payload_aula, output_path):
        """Preenche os 18 slides dinamicamente a partir de um dicionário/JSON"""
        
        for num in range(1, 19):
            sl = self._slide_base(num)
            slide_key = f"slide_{num}"
            data = payload_aula.get(slide_key, {})
            
            tit = data.get("titulo", f"Slide {num}")
            sub = data.get("subtitulo", "")
            self.add_header(sl, tit, sub)
            
            # Renderização dinâmica por tipo de bloco
            conteudo = data.get("conteudo", [])
            
            if isinstance(conteudo, list):
                top_offset = Inches(2.7)
                for item in conteudo:
                    if isinstance(item, dict):
                        self.add_bullet(sl, item.get("texto", ""), top_offset, item.get("nivel", ""))
                    else:
                        self.add_bullet(sl, str(item), top_offset)
                    top_offset += Inches(0.5)
            elif isinstance(conteudo, str):
                tb = sl.shapes.add_textbox(Inches(0.5), Inches(2.7), Inches(12), Inches(3.5))
                tf = tb.text_frame; tf.word_wrap = True
                p = tf.paragraphs[0]; p.text = conteudo
                p.font.size = Pt(16); p.font.color.rgb = SLATE
                
            # Card inferior pedagógico (se houver)
            if "nota_pedagogica" in data:
                self.add_card(sl, "Nota Pedagógica:", data["nota_pedagogica"], Inches(0.5), Inches(5.5), Inches(12.3), Inches(1.2))

        self.prs.save(output_path)
        print(f"✅ BASE HAAS COMPILADA COM SUCESSO: {output_path}")

