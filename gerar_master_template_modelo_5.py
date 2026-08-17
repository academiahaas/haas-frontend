import os
import urllib.request
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# PALETA OFICIAL HAAS LANGUAGE
BRANCO = RGBColor(0xFF, 0xFF, 0xFF)
ROXO = RGBColor(0x4F, 0x35, 0xE2)
AZUL = RGBColor(0x1E, 0x3A, 0x8A)
SLATE = RGBColor(0x33, 0x41, 0x55)
TURQUESA = RGBColor(0x14, 0xC8, 0xB0)
BORDA = RGBColor(0xCB, 0xD5, 0xE1)
FUNDO_CARD = RGBColor(0xF8, 0xFA, 0xFC)
ROXO_CLARO = RGBColor(0xED, 0xE9, 0xFE)
MUTED = RGBColor(0x94, 0xA3, 0xB8)

LOGO_URL = "https://jdppxfokfhqjudwfwckd.supabase.co/storage/v1/object/public/haas-academy/assignments/Design%20sem%20nome%20(4).png"
LOGO_PATH = "/tmp/haas-logo.png"

if not os.path.exists(LOGO_PATH):
    try:
        urllib.request.urlretrieve(LOGO_URL, LOGO_PATH)
    except Exception:
        pass

W, H = Inches(13.333), Inches(7.5)
TOTAL_SLIDES = 22

prs = Presentation()
prs.slide_width = W
prs.slide_height = H

def base_slide(num):
    sl = prs.slides.add_slide(prs.slide_layouts[6])
    sl.background.fill.solid()
    sl.background.fill.fore_color.rgb = BRANCO
    
    # MOLDURA MODELO 5: Elemento de Topo Editorial Dual
    top_line = sl.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.5), Inches(0.2), Inches(2.0), Pt(4))
    top_line.fill.solid(); top_line.fill.fore_color.rgb = ROXO; top_line.line.fill.background()
    
    top_line2 = sl.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(2.6), Inches(0.2), Inches(0.6), Pt(4))
    top_line2.fill.solid(); top_line2.fill.fore_color.rgb = TURQUESA; top_line2.line.fill.background()
    
    # Logo Haas Language
    if os.path.exists(LOGO_PATH):
        sl.shapes.add_picture(LOGO_PATH, Inches(10.8), Inches(0.2), height=Inches(0.65))
        
    # Rodapé Editorial
    tb = sl.shapes.add_textbox(Inches(0.5), Inches(7.05), Inches(8), Inches(0.35))
    p = tb.text_frame.paragraphs[0]
    p.text = "HAAS LANGUAGE  •  STUDIO BOLD EDITION"
    p.font.size = Pt(10); p.font.color.rgb = AZUL; p.font.bold = True
    
    tb2 = sl.shapes.add_textbox(Inches(10.8), Inches(7.05), Inches(2.0), Inches(0.35))
    p2 = tb2.text_frame.paragraphs[0]
    p2.text = f"SLIDE {num} / {TOTAL_SLIDES}"
    p2.font.size = Pt(10); p2.font.color.rgb = ROXO; p2.font.bold = True; p2.alignment = PP_ALIGN.RIGHT
    return sl

def add_header(sl, tit, sub=""):
    tb = sl.shapes.add_textbox(Inches(0.5), Inches(0.35), Inches(10), Inches(0.65))
    tf = tb.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]; p.text = tit
    p.font.size = Pt(30); p.font.color.rgb = ROXO; p.font.bold = True
    
    if sub:
        tb_s = sl.shapes.add_textbox(Inches(0.5), Inches(0.98), Inches(10), Inches(0.4))
        p_s = tb_s.text_frame.paragraphs[0]
        p_s.text = sub
        p_s.font.size = Pt(17); p_s.font.color.rgb = AZUL; p_s.font.bold = True

def add_bold_card(sl, tit, desc, left, top, width, height, bg_color=FUNDO_CARD, border_color=BORDA, tit_size=16, desc_size=14, badge_txt=""):
    s = sl.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    s.fill.solid(); s.fill.fore_color.rgb = bg_color
    s.line.color.rgb = border_color; s.line.width = Pt(1.5)
    
    tf = s.text_frame; tf.word_wrap = True; tf.vertical_anchor = MSO_ANCHOR.TOP
    
    # Preenchimento interno adaptativo para evitar transbordamento
    margin_val = Inches(0.1) if height < Inches(1.0) else Inches(0.18)
    tf.margin_left = margin_val; tf.margin_top = margin_val; tf.margin_right = margin_val; tf.margin_bottom = margin_val
    
    if badge_txt:
        p_b = tf.paragraphs[0]
        p_b.text = badge_txt.upper()
        p_b.font.size = Pt(10); p_b.font.color.rgb = TURQUESA; p_b.font.bold = True
        
        if tit:
            p = tf.add_paragraph()
            p.text = tit; p.font.size = Pt(tit_size); p.font.color.rgb = ROXO; p.font.bold = True; p.space_before = Pt(4)
    else:
        if tit:
            p = tf.paragraphs[0]
            p.text = tit; p.font.size = Pt(tit_size); p.font.color.rgb = ROXO; p.font.bold = True
            
    if desc:
        p2 = tf.add_paragraph() if (badge_txt or tit) else tf.paragraphs[0]
        p2.text = desc; p2.font.size = Pt(desc_size); p2.font.color.rgb = SLATE
        if badge_txt or tit:
            p2.space_before = Pt(4)

def add_history_pill(sl, text, left, top, width, height):
    """Gera cartões de histórico perfeitamente ajustados, impedindo vazamento de texto"""
    s = sl.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    s.fill.solid(); s.fill.fore_color.rgb = FUNDO_CARD
    s.line.color.rgb = MUTED; s.line.width = Pt(1)
    
    tf = s.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    tf.margin_left = Inches(0.08)
    tf.margin_right = Inches(0.08)
    tf.margin_top = Inches(0.04)
    tf.margin_bottom = Inches(0.04)
    
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(9)
    p.font.color.rgb = SLATE
    p.alignment = PP_ALIGN.CENTER

def add_flashcard(sl, num_str, tit, desc, left, top, width, height):
    s = sl.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    s.fill.solid(); s.fill.fore_color.rgb = BRANCO
    s.line.color.rgb = BORDA; s.line.width = Pt(1.5)
    
    top_bar = sl.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, Inches(0.12))
    top_bar.fill.solid(); top_bar.fill.fore_color.rgb = ROXO; top_bar.line.fill.background()
    
    tag = sl.shapes.add_shape(MSO_SHAPE.RECTANGLE, left + Inches(0.15), top + Inches(0.25), Inches(0.5), Inches(0.3))
    tag.fill.solid(); tag.fill.fore_color.rgb = ROXO_CLARO; tag.line.fill.background()
    tf_t = tag.text_frame; tf_t.vertical_anchor = MSO_ANCHOR.MIDDLE
    p_t = tf_t.paragraphs[0]; p_t.text = num_str; p_t.alignment = PP_ALIGN.CENTER
    p_t.font.size = Pt(11); p_t.font.color.rgb = ROXO; p_t.font.bold = True
    
    tf = s.text_frame; tf.word_wrap = True; tf.vertical_anchor = MSO_ANCHOR.TOP
    tf.margin_left = Inches(0.75); tf.margin_top = Inches(0.2); tf.margin_right = Inches(0.2); tf.margin_bottom = Inches(0.1)
    p = tf.paragraphs[0]; p.text = tit; p.font.size = Pt(14); p.font.color.rgb = ROXO; p.font.bold = True
    if desc:
        p2 = tf.add_paragraph(); p2.text = desc; p2.font.size = Pt(13); p2.font.color.rgb = SLATE; p2.space_before = Pt(4)

def img_bold(sl, left, top, width, height, label="[Espaço de Imagem Editorial]"):
    s = sl.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    s.fill.solid(); s.fill.fore_color.rgb = ROXO_CLARO
    s.line.color.rgb = TURQUESA; s.line.width = Pt(2)
    tf = s.text_frame; tf.word_wrap = True; tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER
    p = tf.paragraphs[0]; p.text = label; p.font.size = Pt(16); p.font.color.rgb = ROXO; p.font.bold = True

# ==========================================================
# CONSTRUÇÃO DOS 22 SLIDES REVISADOS
# ==========================================================

# SLIDE 1: CAPA EDITORIAL + CAMPO DE DATA
sl1 = base_slide(1)

tb_c = sl1.shapes.add_textbox(Inches(0.5), Inches(1.8), Inches(6.5), Inches(1.8))
tf_c = tb_c.text_frame; tf_c.word_wrap = True
p = tf_c.paragraphs[0]; p.text = "[Título da Aula / Tema Principal]"
p.font.size = Pt(38); p.font.color.rgb = ROXO; p.font.bold = True

tb_sub = sl1.shapes.add_textbox(Inches(0.5), Inches(3.7), Inches(6.5), Inches(1.2))
tf_sub = tb_sub.text_frame; tf_sub.word_wrap = True
p_sub = tf_sub.paragraphs[0]; p_sub.text = "[Subtítulo / Descrição da Aula]"
p_sub.font.size = Pt(22); p_sub.font.color.rgb = AZUL; p_sub.font.bold = True

# Campo de Data Personalizada
tb_date = sl1.shapes.add_textbox(Inches(0.5), Inches(5.2), Inches(6.5), Inches(0.6))
tf_date = tb_date.text_frame; tf_date.word_wrap = True
p_date = tf_date.paragraphs[0]
p_date.text = "📅 Data da Aula: [DD/MM/AAAA]"
p_date.font.size = Pt(15); p_date.font.color.rgb = TURQUESA; p_date.font.bold = True

img_bold(sl1, Inches(7.3), Inches(1.4), Inches(5.5), Inches(5.3), "[Imagem Editorial da Capa]")

# SLIDE 2: PRESENÇA & INTEGRANTES
sl2 = base_slide(2)
add_header(sl2, "Presença & Integrantes")
add_bold_card(sl2, "Professor(a):", "[Nome do(a) Professor(a)]", Inches(0.5), Inches(1.5), Inches(12.3), Inches(1.5), bg_color=ROXO_CLARO, border_color=ROXO, tit_size=18, desc_size=16)
add_bold_card(sl2, "Aluno(s):", "[Nome do Aluno / Lista da Turma]", Inches(0.5), Inches(3.2), Inches(12.3), Inches(3.6), bg_color=FUNDO_CARD, border_color=BORDA, tit_size=18, desc_size=16)

# SLIDE 3: WARM-UP
sl3 = base_slide(3)
add_header(sl3, "Aquecimento / Warm-up", "Foco 100% na Fala (3 Minutos)")
add_bold_card(sl3, "Pergunta Provocativa Rápida:", "[Insira aqui a pergunta provocativa de aquecimento]", Inches(0.5), Inches(1.5), Inches(6.5), Inches(3.2), bg_color=ROXO_CLARO, border_color=ROXO, tit_size=20, desc_size=16)
add_bold_card(sl3, "", "Ativa o córtex pré-frontal e reduz a inibição linguística inicial.", Inches(0.5), Inches(4.9), Inches(6.5), Inches(1.9), bg_color=FUNDO_CARD, border_color=TURQUESA, desc_size=14)
img_bold(sl3, Inches(7.3), Inches(1.5), Inches(5.5), Inches(5.3), "[Imagem de Aquecimento]")

# SLIDE 4: OBJETIVOS
sl4 = base_slide(4)
add_header(sl4, "Objetivos da Aula")
add_bold_card(sl4, "Metas da Aula:", "• [Objetivo 1]\n• [Objetivo 2]\n• [Objetivo 3]", Inches(0.5), Inches(1.5), Inches(12.3), Inches(3.3), tit_size=18, desc_size=16)
add_bold_card(sl4, "", "Nesta aula daremos atenção à escrita para reforçar a memorização de estruturas/vocabulário e reduzir a interferência da língua materna.", Inches(0.5), Inches(5.0), Inches(12.3), Inches(1.8), bg_color=ROXO_CLARO, border_color=ROXO, desc_size=14)

# SLIDE 5: METODOLOGIA EM ESCADA BOLD
sl5 = base_slide(5)
add_header(sl5, "A Nossa Metodologia", "Estrutura Fixa de Aprendizado Haas Language")
metodologias = [
    ("1. Fala Ativa", "Conversação inicial sem julgamentos."),
    ("2. Exposição", "Apresentação leve da gramática."),
    ("3. Prática Guiada", "Exercícios com suporte do professor."),
    ("4. Conversação", "Aplicação prática das estruturas."),
    ("5. Escrita", "Consolidação e redução de tradução.")
]
left_m = 0.5
for i, (tit_m, desc_m) in enumerate(metodologias, 1):
    bg_col = ROXO_CLARO if i % 2 != 0 else FUNDO_CARD
    b_col = ROXO if i % 2 != 0 else TURQUESA
    add_bold_card(sl5, tit_m, desc_m, Inches(left_m), Inches(1.5), Inches(2.25), Inches(5.3), bg_color=bg_col, border_color=b_col, tit_size=16, desc_size=13)
    left_m += 2.45

# SLIDES 6 A 8: EXPOSIÇÃO
for s_num in [6, 7, 8]:
    sl = base_slide(s_num)
    add_header(sl, f"Explicação do Tema — Parte {s_num - 5}")
    add_bold_card(sl, "Conteúdo Principal:", f"[Explicação e exemplo central da estrutura - Bloco {s_num - 5}]", Inches(0.5), Inches(1.5), Inches(7.5), Inches(5.3), bg_color=ROXO_CLARO, border_color=ROXO, tit_size=20, desc_size=15)
    add_bold_card(sl, "Exemplos Práticos:", "[Exemplos adicionais e frases do dia a dia]", Inches(8.3), Inches(1.5), Inches(4.5), Inches(2.5), tit_size=16, desc_size=13)
    add_bold_card(sl, "Notas & Detalhes:", "[Exceções, nuances gramaticais e contexto]", Inches(8.3), Inches(4.2), Inches(4.5), Inches(2.6), border_color=TURQUESA, tit_size=16, desc_size=13)

# SLIDE 9: SÍNTESE + CULTURA
sl9 = base_slide(9)
add_header(sl9, "Síntese do Conteúdo & Sugestão Cultural")
add_bold_card(sl9, "Resumo do Tema:", "[Resumo sintético dos pontos chave ensinados na aula]", Inches(0.5), Inches(1.5), Inches(12.3), Inches(2.6), tit_size=18, desc_size=15)
add_bold_card(sl9, "Sugestão Cultural:", "🎵 [Nome da Música / Vídeo recomendado sobre o assunto].\nO professor pode rodar o áudio/vídeo ou deixar como recomendação extra.", Inches(0.5), Inches(4.3), Inches(12.3), Inches(2.5), bg_color=ROXO_CLARO, border_color=ROXO, tit_size=16, desc_size=14)

# SLIDES 10 E 11: PRÁTICA GUIADA
sl10 = base_slide(10)
add_header(sl10, "Prática Guiada — Exercícios 1 a 5", "Correção de Vícios")
exs_1_5 = [
    ("01", "Exercício 1", "[Lacunas contextualizadas com opções]"),
    ("02", "Exercício 2", "[Lacunas contextualizadas com opções]"),
    ("03", "Exercício 3", "[Preenchimento direto de verbos]"),
    ("04", "Exercício 4", "[Reestruturação / Transformação de frase]"),
    ("05", "Exercício 5", "[Reestruturação / Transformação de frase]")
]
for idx, (num_s, t_e, d_e) in enumerate(exs_1_5):
    col = idx % 2; row = idx // 2
    l = 0.5 if col == 0 else 6.8
    w = 6.0 if (idx < 4) else 12.3
    if idx == 4: l = 0.5
    t = 1.5 + row * 1.7
    add_flashcard(sl10, num_s, t_e, d_e, Inches(l), Inches(t), Inches(w), Inches(1.5))

sl11 = base_slide(11)
add_header(sl11, "Prática Guiada — Exercícios 6 a 10", "Desafios Práticos")
exs_6_10 = [
    ("06", "Exercício 6", "[Preposição / Conector correto]"),
    ("07", "Exercício 7", "[Escolha da opção correta]"),
    ("08", "Exercício 8", "[Substituição por sinônimo formal]"),
    ("09", "Exercício 9", "[Identificação de erro e correção]"),
    ("10", "Exercício 10", "[Tradução de sentido sem uso literal]")
]
for idx, (num_s, t_e, d_e) in enumerate(exs_6_10):
    col = idx % 2; row = idx // 2
    l = 0.5 if col == 0 else 6.8
    w = 6.0 if (idx < 4) else 12.3
    if idx == 4: l = 0.5
    t = 1.5 + row * 1.7
    add_flashcard(sl11, num_s, t_e, d_e, Inches(l), Inches(t), Inches(w), Inches(1.5))

# SLIDES 12 A 19: CONVERSAÇÃO (Sem rótulos confusos, sem texto vazando)
perguntas_config = [
    "[Pergunta principal 1 sobre o tema]",
    "[Pergunta principal 2 sobre o tema]",
    "[Pergunta principal 3 sobre o tema]",
    "[Pergunta principal 4 sobre o tema]",
    "[Pergunta principal 5 sobre o tema]",
    "[Pergunta principal 6 sobre o tema]",
    "[Pergunta principal 7 sobre o tema]",
    "[Pergunta principal 8 sobre o tema]"
]

for idx in range(8):
    s_num = 12 + idx
    sl = base_slide(s_num)
    add_header(sl, "Conversação")
    
    if idx == 0:
        p_desc = perguntas_config[0]
        add_bold_card(sl, "", p_desc, Inches(0.5), Inches(1.5), Inches(12.3), Inches(2.7), bg_color=ROXO_CLARO, border_color=ROXO, desc_size=18)
        add_bold_card(sl, "", "• Para expandir: 'Dê um exemplo recente...'\n• Para concordar/discordar: 'Qual é o outro lado dessa moeda?'", Inches(0.5), Inches(4.4), Inches(12.3), Inches(2.4), desc_size=15)
    else:
        history_list = [perguntas_config[i] for i in range(idx)]
        count = len(history_list)
        gap = 0.15; avail_w = 12.3 - (gap * (count - 1))
        pill_w = avail_w / count
        
        # Cápsulas de histórico perfeitamente ajustadas
        for h_idx, h_txt in enumerate(history_list):
            l = 0.5 + h_idx * (pill_w + gap)
            clean_t = h_txt[:20] + "..." if len(h_txt) > 20 else h_txt
            add_history_pill(sl, clean_t, Inches(l), Inches(1.5), Inches(pill_w), Inches(0.75))
        
        p_desc = perguntas_config[idx]
        add_bold_card(sl, "", p_desc, Inches(0.5), Inches(2.45), Inches(12.3), Inches(2.5), bg_color=ROXO_CLARO, border_color=ROXO, desc_size=17)
        add_bold_card(sl, "", "• Para expandir: 'Dê um exemplo recente...'\n• Para concordar/discordar: 'Qual é o outro lado dessa moeda?'", Inches(0.5), Inches(5.1), Inches(12.3), Inches(1.7), desc_size=14)

# SLIDE 20: ESCRITA
sl20 = base_slide(20)
add_header(sl20, "Atividade de Escrita", "Consolidação e Memorização")
add_bold_card(sl20, "Tarefa de Escrita Recomendada:", "[Instrução detalhada da tarefa: escrever um texto, frases ou diálogo]", Inches(0.5), Inches(1.5), Inches(12.3), Inches(2.8), tit_size=18, desc_size=15)
add_bold_card(sl20, "", "Identificamos que a escrita nesta etapa é fundamental para você fixar as estruturas aprendidas, expandir vocabulário ativo e diminuir a tradução mental automática do seu idioma nativo.", Inches(0.5), Inches(4.5), Inches(12.3), Inches(2.3), bg_color=ROXO_CLARO, border_color=ROXO, desc_size=14)

# SLIDE 21: TAREFA DE CASA
sl21 = base_slide(21)
add_header(sl21, "Encerramento — Tarefa de Casa", "Ação e Prática Contínua")
add_bold_card(sl21, "Tarefa de Casa:", "• [Atividade prática 1 para realizar fora da aula]\n• [Atividade de fixação de vocabulário ou escuta]", Inches(0.5), Inches(1.5), Inches(12.3), Inches(5.2), tit_size=18, desc_size=15)

# SLIDE 22: ENCERRAMENTO & FEEDBACK
sl22 = base_slide(22)
add_header(sl22, "Encerramento & Feedback Rápido", "Considerações Finais")
add_bold_card(sl22, "Agradecimento:", "Obrigado pela dedicação na aula de hoje!", Inches(0.5), Inches(1.5), Inches(12.3), Inches(2.0), tit_size=20, desc_size=16)
add_bold_card(sl22, "Espaço para Feedback Rápido:", "[Espaço destinado para o professor coletar dúvidas rápidas e considerações do aluno]", Inches(0.5), Inches(3.7), Inches(12.3), Inches(3.0), border_color=TURQUESA, tit_size=18, desc_size=15)

# SALVAR NOVO ARQUIVO REVISADO
out_path = "/var/www/haas-frontend-desk-mobile-oficial/public/temp-miniaturas/template-haas-modelo-5.pptx"
prs.save(out_path)
print(f"✅ MODELO 5 ATUALIZADO COM SUCESSO EM: {out_path}")
