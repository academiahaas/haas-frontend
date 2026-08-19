import json
import re
import random
import urllib.request
import urllib.parse
from datetime import datetime, timedelta, timezone
from pptx import Presentation

# ===== CONFIGURAÇÃO =====
SUPABASE_URL = "https://jdppxfokfhqjudwfwckd.supabase.co"
SERVICE_KEY = "sb_secret_PngpeWMeQNFg1OcKpl0xOw_ZVZ8PMVz"
DEEPSEEK_KEY = "sk-e426fa20f2c64907bb550d7eccf1261f"
BUCKET = "haas-academy"
PASTA_MODELOS = "Untitled folder/slides_modelos"

HEADERS = {"apikey": SERVICE_KEY, "Authorization": f"Bearer {SERVICE_KEY}", "Content-Type": "application/json"}


def supabase_get(path):
    req = urllib.request.Request(f"{SUPABASE_URL}/rest/v1/{path}", headers=HEADERS)
    with urllib.request.urlopen(req, timeout=30) as r:
        return json.loads(r.read().decode())


def supabase_patch(table, filtro, dados):
    req = urllib.request.Request(
        f"{SUPABASE_URL}/rest/v1/{table}?{filtro}",
        data=json.dumps(dados).encode(),
        headers=HEADERS, method="PATCH",
    )
    with urllib.request.urlopen(req, timeout=30) as r:
        return r.read().decode()


def storage_list(prefix):
    body = json.dumps({"prefix": prefix, "limit": 100}).encode()
    req = urllib.request.Request(
        f"{SUPABASE_URL}/storage/v1/object/list/{BUCKET}",
        data=body, headers=HEADERS, method="POST",
    )
    with urllib.request.urlopen(req, timeout=30) as r:
        return json.loads(r.read().decode())


def storage_download(caminho, destino_local):
    url = f"{SUPABASE_URL}/storage/v1/object/public/{BUCKET}/{urllib.parse.quote(caminho)}"
    urllib.request.urlretrieve(url, destino_local)


def storage_upload(caminho, arquivo_local):
    with open(arquivo_local, "rb") as f:
        conteudo = f.read()
    if arquivo_local.endswith(".pdf"):
        content_type = "application/pdf"
    else:
        content_type = "application/vnd.openxmlformats-officedocument.presentationml.presentation"
    req = urllib.request.Request(
        f"{SUPABASE_URL}/storage/v1/object/{BUCKET}/{urllib.parse.quote(caminho)}",
        data=conteudo,
        headers={"apikey": SERVICE_KEY, "Authorization": f"Bearer {SERVICE_KEY}",
                 "Content-Type": content_type,
                 "x-upsert": "true"},
        method="POST",
    )
    urllib.request.urlopen(req, timeout=60)
    return f"{SUPABASE_URL}/storage/v1/object/public/{BUCKET}/{urllib.parse.quote(caminho)}"


def cortar_seguro(texto, limite):
    if limite is None or len(texto) <= limite:
        return texto
    cortado = texto[:limite].rsplit(" ", 1)[0]
    return cortado.rstrip(".,;:—-") + "…"


print("=" * 60)
print("MOTOR DE GERAÇÃO DE SLIDES - iniciando busca")
print("=" * 60)

agora = datetime.now(timezone.utc)
janela_inicio = agora.isoformat()
janela_fim = (agora + timedelta(hours=10)).isoformat()

aulas = supabase_get(
    f"aulas_disponiveis?data_hora_inicio=gte.{urllib.parse.quote(janela_inicio)}&data_hora_inicio=lte.{urllib.parse.quote(janela_fim)}"
    f"&slides_status=is.null&select=id,tipo_aula,idioma,data_hora_inicio"
)

print(f"Aulas encontradas na janela de 10h: {len(aulas)}")

for aula in aulas:
    aula_id = aula["id"]
    print(f"\n--- Processando aula {aula_id} ({aula['tipo_aula']}) ---")

    try:
        # Marca como "gerando" pra não processar duas vezes
        supabase_patch("aulas_disponiveis", f"id=eq.{aula_id}", {"slides_status": "gerando"})

        # 1. Busca alunos matriculados
        matriculas = supabase_get(f"aula_matriculas?aula_id=eq.{aula_id}&select=user_id")
        user_ids = [m["user_id"] for m in matriculas]
        print(f"Alunos matriculados: {len(user_ids)}")

        if not user_ids:
            raise Exception("Nenhum aluno matriculado nessa aula")

        # 2. Busca dados de cada aluno
        alunos_info = []
        for uid in user_ids:
            u = supabase_get(f"users?id=eq.{uid}&select=id,name,current_unit_id,topico_deficitario,course_language,native_language")
            if u:
                alunos_info.append(u[0])

        # 3. Verifica se algum é corporativo
        eh_corporativo = False
        for uid in user_ids:
            subs = supabase_get(f"user_subscriptions?user_id=eq.{uid}&select=plan_category&order=created_at.desc&limit=1")
            if subs:
                plano = supabase_get(f"master_plans?plan_category=eq.{urllib.parse.quote(subs[0]['plan_category'])}&select=is_corporate")
                if plano and plano[0].get("is_corporate"):
                    eh_corporativo = True
                    break

        # 4. Monta o tema: pega a unidade do primeiro aluno como referência + junta dificuldades
        unit_id_referencia = alunos_info[0]["current_unit_id"]
        unidade_info = supabase_get(f"units?id=eq.{unit_id_referencia}&select=unit_title,hidden_grammatical_structure")
        tema = unidade_info[0]["unit_title"] if unidade_info else "Revisão Geral"

        dificuldades = list({a["topico_deficitario"] for a in alunos_info if a.get("topico_deficitario")})
        idioma_aprendido = alunos_info[0].get("course_language") or "Português"
        idioma_auxiliar = alunos_info[0].get("native_language") or "Espanhol"

        print(f"Tema escolhido: {tema}")
        print(f"Dificuldades da turma: {dificuldades}")
        print(f"Pacote corporativo: {eh_corporativo}")

        # 5. Sorteia um template
        arquivos = storage_list(PASTA_MODELOS)
        candidatos = [a["name"] for a in arquivos if a["name"].endswith(".pptx")]
        if not candidatos:
            raise Exception("Nenhum template encontrado na pasta slides_modelos")
        template_escolhido = random.choice(candidatos)
        print(f"Template sorteado: {template_escolhido}")

        storage_download(f"{PASTA_MODELOS}/{template_escolhido}", "/tmp/motor_template.pptx")

        # 6. Gera o conteúdo (reaproveitando a mesma lógica testada)
        prs = Presentation("/tmp/motor_template.pptx")
        padrao = re.compile(r"\[([A-Z_0-9]+):(\d+)\]")
        estrutura_real = {}
        ocorrencias = []

        for slide_idx, slide in enumerate(prs.slides, 1):
            for shape in slide.shapes:
                if not shape.has_text_frame:
                    continue
                for paragraph in shape.text_frame.paragraphs:
                    for run in paragraph.runs:
                        for match in padrao.finditer(run.text):
                            nome, limite = match.group(1), int(match.group(2))
                            ocorrencias.append({"slide": slide_idx, "run": run, "marcador": match.group(0), "nome": nome, "limite": limite})
                            entry = estrutura_real.setdefault(slide_idx, {"titulo_limite": None, "itens": []})
                            if nome == "TITULO_SLIDE":
                                entry["titulo_limite"] = limite
                            elif nome.startswith("ITEM_"):
                                entry["itens"].append(limite)

        # Numera os itens pela ORDEM DE APARICÇÃO no slide, não pelo número escrito
        # (protege contra o Canva às vezes repetir/pular números de item)
        contador_por_slide = {}
        for oc in ocorrencias:
            if oc["nome"].startswith("ITEM_"):
                s = oc["slide"]
                contador_por_slide.setdefault(s, 0)
                oc["indice_real"] = contador_por_slide[s]
                contador_por_slide[s] += 1

        SLIDE_ALUNOS = 2
        ultimo_slide = max(estrutura_real.keys()) if estrutura_real else None

        contexto_corporativo = "IMPORTANTE: este aluno está num pacote CORPORATIVO/PROFISSIONAL. Use vocabulário e situações de trabalho (reuniões, e-mails, apresentações profissionais) em vez de temas do dia a dia genéricos." if eh_corporativo else ""

        descricao_estrutura = "\n".join(
            f'"{s}": título com até {d["titulo_limite"] or 0} caracteres, e {len(d["itens"])} itens com limites {d["itens"]}'
            for s, d in sorted(estrutura_real.items()) if s != SLIDE_ALUNOS
        )

        system_prompt = f"""Você é um gerador de conteúdo pedagógico para aulas de idiomas.
{contexto_corporativo}
Estrutura: capa, lista alunos, aquecimento, objetivos, metodologia, conteúdo do tema, resumo+cultura, exercícios, conversação progressiva, escrita, tarefa, agradecimento.
Dificuldades específicas dos alunos dessa turma pra reforçar durante a aula: {", ".join(dificuldades) if dificuldades else "nenhuma específica"}.
REGRAS: NUNCA use códigos CEFR (A1,A2,B1,B2,C1). Use: Inicial, Básico, Intermediário, Independente, Avançado. Texto sempre pro ALUNO. Responda APENAS JSON válido."""

        user_prompt = f"""Aula de {idioma_aprendido} sobre "{tema}", com {idioma_auxiliar} como apoio.
Estrutura exata: {descricao_estrutura}
Retorne JSON: {{"1": {{"titulo": "...", "itens": [...]}}, ...}}
Use PERTO do máximo de caracteres de cada campo, sem ultrapassar. Preencha TODOS os slides listados."""

        payload = json.dumps({
            "model": "deepseek-v4-flash",
            "messages": [{"role": "system", "content": system_prompt}, {"role": "user", "content": user_prompt}],
            "response_format": {"type": "json_object"},
            "max_tokens": 16000,
            "thinking": {"type": "disabled"},
        }).encode("utf-8")

        req = urllib.request.Request(
            "https://api.deepseek.com/chat/completions", data=payload,
            headers={"Authorization": f"Bearer {DEEPSEEK_KEY}", "Content-Type": "application/json"}, method="POST",
        )
        with urllib.request.urlopen(req, timeout=180) as response:
            resp_data = json.loads(response.read().decode("utf-8"))

        content_json = json.loads(resp_data["choices"][0]["message"]["content"])

        for oc in ocorrencias:
            slide_num = oc["slide"]
            dados = content_json.get(str(slide_num), {})
            if oc["nome"] == "TITULO_SLIDE":
                valor = "Presença" if slide_num == SLIDE_ALUNOS else dados.get("titulo", "")
            elif oc["nome"].startswith("ITEM_"):
                idx = oc["indice_real"]
                if slide_num == SLIDE_ALUNOS:
                    nomes = ", ".join(a["name"] for a in alunos_info)
                    valor = nomes if idx == 0 else ""
                elif slide_num == ultimo_slide and idx == 1:
                    valor = "[Feedback baseado no desempenho do aluno — conectar dados reais]"
                else:
                    itens = dados.get("itens", [])
                    valor = itens[idx] if idx < len(itens) else ""
            else:
                valor = ""
            oc["run"].text = oc["run"].text.replace(oc["marcador"], cortar_seguro(str(valor), oc["limite"]))

        arquivo_final = f"/tmp/motor_saida_{aula_id}.pptx"
        prs.save(arquivo_final)

        # 7. Converte também para PDF usando o LibreOffice
        import subprocess
        subprocess.run(
            ["soffice", "--headless", "--convert-to", "pdf", "--outdir", "/tmp", arquivo_final],
            timeout=120, check=True, capture_output=True
        )
        arquivo_pdf = f"/tmp/motor_saida_{aula_id}.pdf"

        # 8. Sobe os dois pro Supabase e atualiza o status
        url_pptx = storage_upload(f"aulas_geradas/{aula_id}.pptx", arquivo_final)
        url_pdf = storage_upload(f"aulas_geradas/{aula_id}.pdf", arquivo_pdf)

        supabase_patch("aulas_disponiveis", f"id=eq.{aula_id}", {
            "slides_status": "pronto",
            "slides_pptx_path": url_pptx,
            "slides_pdf_path": url_pdf,
            "slides_generated_at": datetime.now(timezone.utc).isoformat(),
        })

        print(f"✅ SUCESSO! PPTX: {url_pptx}")
        print(f"✅ SUCESSO! PDF: {url_pdf}")

    except Exception as e:
        print(f"❌ ERRO ao processar aula {aula_id}: {e}")
        supabase_patch("aulas_disponiveis", f"id=eq.{aula_id}", {"slides_status": "falhou"})

print("\n" + "=" * 60)
print("MOTOR FINALIZADO")
print("=" * 60)
