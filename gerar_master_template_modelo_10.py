import os
import urllib.request
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml import parse_xml

# PALETA OFICIAL HAAS LANGUAGE - MODELO 10 (TECH BLUE GRADIENT)
BRANCO = RGBColor(0xFF, 0xFF, 0xFF)
ROXO = RGBColor(0x4F, 0x35, 0xE2)
AZUL = RGBColor(0x1E, 0x3A, 0x8A)
AZUL_TECH = RGBColor(0x25, 0x63, 0xEB) # Azul Vibrante Tech
SLATE = RGBColor(0x33, 0x41, 0x55)
TURQUESA = RGBColor(0x14, 0xC8, 0xB0)
BORDA = RGBColor(0xCB, 0xD5, 0xE1)

LOGO_URL = "https://jdppxfokfhqjudwfwckd.supabase.co/storage/v1/object/public/haas-academy/assignments/Design%20sem%20nome%20(4).png"
LOGO_PATH = "/tmp/haas-logo.png"

if not os.path.exists(LOGO_PATH):
    try:
        urllib.request.urlretrieve(LOGO_URL, LOGO_PATH)
    except Exception:
        pass

W, H = Inches(13.333), Inches(7.5)
TOTAL_SLIDES = 22

prs = Presentation()
prs.slide_width = W
prs.slide_height = H

def apply_bg_gradient(sl):
    """Degradê Tecnológico e Profissional (Azul Noturno para Azul Tech Vibrante)"""
    bg = sl.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), W, H)
    spPr = bg.element.spPr
    grad_xml = """
    <a:gradFill xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main">
      <a:gsLst>
        <a:gs pos="0"><a:srgbClr val="0F172A"/></a:gs>
        <a:gs pos="100000"><a:srgbClr val="1D4ED8"/></a:gs>
      </a:gsLst>
      <a:lin ang="8100000" scaled="1"/>
    </a:gradFill>
    """
    gradFill = parse_xml(grad_xml)
    solidFill = spPr.find('{http://schemas.openxmlformats.org/drawingml/2006/main}solidFill')
    if solidFill is not None:
        spPr.remove(solidFill)
    spPr.append(gradFill)
    bg.line.fill.background()

def base_slide(num, hide_logo=False):
    sl = prs.slides.add_slide(prs.slide_layouts[6])
    
    apply_bg_gradient(sl)

    if os.path.exists(LOGO_PATH) and not hide_logo:
        sl.shapes.add_picture(LOGO_PATH, Inches(10.8), Inches(0.2), height=Inches(0.65))
        
    tb = sl.shapes.add_textbox(Inches(0.5), Inches(7.12), Inches(8), Inches(0.35))
    p = tb.text_frame.paragraphs[0]
    p.text = "HAAS LANGUAGE  •  TECH BLUE EDITION"
    p.font.size = Pt(10); p.font.color.rgb = BRANCO; p.font.bold = True
    
    tb2 = sl.shapes.add_textbox(Inches(10.8), Inches(7.12), Inches(2.0), Inches(0.35))
    p2 = tb2.text_frame.paragraphs[0]
    p2.text = f"SLIDE {num} / {TOTAL_SLIDES}"
    p2.font.size = Pt(10); p2.font.color.rgb = BRANCO; p2.font.bold = True; p2.alignment = PP_ALIGN.RIGHT
    return sl

def add_header(sl, tit, sub=""):
    tb = sl.shapes.add_textbox(Inches(0.5), Inches(0.35), Inches(10), Inches(0.65))
    tf = tb.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]; p.text = tit
    p.font.size = Pt(30); p.font.color.rgb = BRANCO; p.font.bold = True
    
    if sub:
        tb_s = sl.shapes.add_textbox(Inches(0.5), Inches(0.98), Inches(10), Inches(0.4))
        p_s = tb_s.text_frame.paragraphs[0]
        p_s.text = sub
        p_s.font.size = Pt(17); p_s.font.color.rgb = TURQUESA; p_s.font.bold = True

def add_solid_card(sl, tit, desc, left, top, width, height, bg_color=BRANCO, border_color=TURQUESA, tit_color=AZUL, desc_color=SLATE, tit_size=16, desc_size=14, align_center_v=False):
    """Cartões Brancos de Alto Contraste sobre Fundo Escuro"""
    s = sl.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    s.fill.solid(); s.fill.fore_color.rgb = bg_color
    if border_color:
        s.line.color.rgb = border_color; s.line.width = Pt(1.5)
    else:
        s.line.fill.background()
        
    tf = s.text_frame; tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE if align_center_v else MSO_ANCHOR.TOP
    
    margin_val = Inches(0.12) if height < Inches(1.0) else Inches(0.2)
    tf.margin_left = margin_val; tf.margin_top = margin_val; tf.margin_right = margin_val; tf.margin_bottom = margin_val
    
    if tit:
        p = tf.paragraphs[0]
        p.text = tit; p.font.size = Pt(tit_size); p.font.color.rgb = tit_color; p.font.bold = True
        if desc:
            p2 = tf.add_paragraph()
            p2.text = desc; p2.font.size = Pt(desc_size); p2.font.color.rgb = desc_color; p2.space_before = Pt(6)
    else:
        if desc:
            p = tf.paragraphs[0]
            p.text = desc; p.font.size = Pt(desc_size); p.font.color.rgb = desc_color

def img_box_bleed(sl, left, top, width, height, label="[Espaço de Imagem]", bg_color=AZUL, text_color=BRANCO):
    s = sl.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    s.fill.solid(); s.fill.fore_color.rgb = bg_color; s.line.fill.background()
    tf = s.text_frame; tf.word_wrap = True; tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER
    p = tf.paragraphs[0]; p.text = label; p.font.size = Pt(18); p.font.color.rgb = text_color; p.font.bold = True

def add_vertical_history_right(sl, history_list, left=7.8, top=1.5, width=5.0, total_height=5.3):
    count = len(history_list)
    if count == 0:
        return
    gap = 0.12
    avail_h = total_height - (gap * (count - 1))
    card_h = avail_h / count
    
    for i, h_txt in enumerate(history_list):
        y = top + i * (card_h + gap)
        clean_t = h_txt[:40] + "..." if len(h_txt) > 40 else h_txt
        d_sz = 12 if count <= 3 else (11 if count <= 5 else 10)
        # Cartões brancos e com borda clean contrastam incrivelmente com o fundo tech blue
        add_solid_card(sl, "", clean_t, Inches(left), Inches(y), Inches(width), Inches(card_h), 
                       bg_color=BRANCO, border_color=BORDA, desc_color=SLATE, desc_size=d_sz, align_center_v=True)

def add_clean_flashcard(sl, num_str, desc, left, top, width, height, b_c=AZUL_TECH):
    s = sl.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    s.fill.solid(); s.fill.fore_color.rgb = BRANCO; s.line.color.rgb = b_c; s.line.width = Pt(1.5)
    
    chip = sl.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left + Inches(0.12), top + Inches(0.12), Inches(0.55), Inches(0.35))
    chip.fill.solid(); chip.fill.fore_color.rgb = b_c; chip.line.fill.background()
    tf_chip = chip.text_frame; tf_chip.vertical_anchor = MSO_ANCHOR.MIDDLE
    p_chip = tf_chip.paragraphs[0]; p_chip.text = num_str; p_chip.alignment = PP_ALIGN.CENTER
    p_chip.font.size = Pt(11); p_chip.font.color.rgb = BRANCO; p_chip.font.bold = True
    
    tf = s.text_frame; tf.word_wrap = True; tf.vertical_anchor = MSO_ANCHOR.TOP
    tf.margin_left = Inches(0.8); tf.margin_top = Inches(0.15); tf.margin_right = Inches(0.2); tf.margin_bottom = Inches(0.12)
    p = tf.paragraphs[0]; p.text = desc; p.font.size = Pt(14); p.font.color.rgb = SLATE; p.font.bold = False

# ==========================================================
# CONSTRUÇÃO DOS 22 SLIDES DO MODELO 10 (TECH BLUE)
# ==========================================================

# SLIDE 1: CAPA
sl1 = base_slide(1, hide_logo=True)
img_box_bleed(sl1, Inches(0), Inches(0), W, Inches(3.5), "[Imagem de Capa / Herói no Topo]", bg_color=AZUL)

if os.path.exists(LOGO_PATH):
    sl1.shapes.add_picture(LOGO_PATH, Inches(0.5), Inches(3.8), height=Inches(0.65))

tb_c = sl1.shapes.add_textbox(Inches(0.4), Inches(4.5), Inches(12.3), Inches(1.0))
tf_c = tb_c.text_frame; tf_c.word_wrap = True
p = tf_c.paragraphs[0]; p.text = "[Título da Aula / Tema Principal]"
p.font.size = Pt(38); p.font.color.rgb = BRANCO; p.font.bold = True

tb_sub = sl1.shapes.add_textbox(Inches(0.4), Inches(5.5), Inches(12.3), Inches(0.6))
tf_sub = tb_sub.text_frame; tf_sub.word_wrap = True
p_sub = tf_sub.paragraphs[0]; p_sub.text = "[Subtítulo / Descrição da Aula]"
p_sub.font.size = Pt(22); p_sub.font.color.rgb = TURQUESA; p_sub.font.bold = True

tb_date = sl1.shapes.add_textbox(Inches(0.4), Inches(6.1), Inches(12.3), Inches(0.6))
tf_date = tb_date.text_frame; tf_date.word_wrap = True
p_date = tf_date.paragraphs[0]
p_date.text = "📅 Data da Aula: [DD/MM/AAAA]"
p_date.font.size = Pt(16); p_date.font.color.rgb = TURQUESA; p_date.font.bold = True

# SLIDE 2: PRESENÇA & INTEGRANTES
sl2 = base_slide(2)
add_header(sl2, "Presença & Integrantes")
add_solid_card(sl2, "Professor(a):", "[Nome do(a) Professor(a)]", Inches(0.5), Inches(1.5), Inches(12.3), Inches(1.5), border_color=AZUL_TECH)
add_solid_card(sl2, "Aluno(s):", "[Nome do Aluno / Lista da Turma]", Inches(0.5), Inches(3.2), Inches(12.3), Inches(3.6), border_color=TURQUESA)

# SLIDE 3: WARM-UP
sl3 = base_slide(3)
add_header(sl3, "Aquecimento / Warm-up", "Foco 100% na Fala (3 Minutos)")
add_solid_card(sl3, "Pergunta Provocativa Rápida:", "[Insira aqui a pergunta provocativa de aquecimento]", Inches(0.5), Inches(1.5), Inches(7.8), Inches(2.2), border_color=AZUL_TECH, tit_size=20, desc_size=16)
add_solid_card(sl3, "", "Ativa o córtex pré-frontal e reduz a inibição linguística inicial.", Inches(8.5), Inches(1.5), Inches(4.3), Inches(2.2), border_color=TURQUESA, desc_color=SLATE, desc_size=14, align_center_v=True)
img_box_bleed(sl3, Inches(0), Inches(4.2), W, Inches(2.8), "[Imagem Warm-up (Rodapé)]", bg_color=AZUL)

# SLIDE 4: OBJETIVOS
sl4 = base_slide(4)
add_header(sl4, "Objetivos da Aula")
add_solid_card(sl4, "Metas da Aula:", "• [Objetivo 1]\n• [Objetivo 2]\n• [Objetivo 3]", Inches(0.5), Inches(1.5), Inches(7.8), Inches(5.3), border_color=AZUL_TECH)
add_solid_card(sl4, "Nota de Escrita:", "Nesta aula daremos atenção à escrita para reforçar a memorização de estruturas/vocabulário e reduzir a interferência da língua materna.", Inches(8.5), Inches(1.5), Inches(4.3), Inches(5.3), border_color=TURQUESA)

# SLIDE 5: METODOLOGIA
sl5 = base_slide(5)
add_header(sl5, "A Nossa Metodologia", "Estrutura Fixa de Aprendizado Haas Language")
add_solid_card(sl5, "1. Fala Ativa", "Conversação inicial sem julgamentos.", Inches(0.5), Inches(1.5), Inches(3.8), Inches(5.3), border_color=AZUL_TECH)

met_grid = [
    ("2. Exposição", "Apresentação leve da gramática."),
    ("3. Prática Guiada", "Exercícios com suporte."),
    ("4. Conversação", "Aplicação prática das estruturas."),
    ("5. Escrita", "Consolidação visual.")
]
for idx, (t, d) in enumerate(met_grid):
    col = idx % 2; row = idx // 2
    l = 4.5 + (col * 4.0); t_pos = 1.5 + (row * 2.7)
    add_solid_card(sl5, t, d, Inches(l), Inches(t_pos), Inches(3.8), Inches(2.5), border_color=TURQUESA)

# SLIDES 6 A 8: EXPOSIÇÃO
for s_num in [6, 7, 8]:
    sl = base_slide(s_num)
    add_header(sl, f"Explicação do Tema — Parte {s_num - 5}")
    add_solid_card(sl, "Conteúdo Principal:", f"[Explicação e exemplo central da estrutura - Bloco {s_num - 5}]", Inches(0.5), Inches(1.5), Inches(12.3), Inches(2.5), border_color=AZUL_TECH, tit_size=20, desc_size=15)
    add_solid_card(sl, "Exemplos Práticos:", "[Exemplos adicionais e frases do dia a dia]", Inches(0.5), Inches(4.2), Inches(6.0), Inches(2.6), border_color=TURQUESA)
    add_solid_card(sl, "Notas & Detalhes:", "[Exceções, nuances gramaticais e contexto de uso]", Inches(6.8), Inches(4.2), Inches(6.0), Inches(2.6), border_color=ROXO)

# SLIDE 9: SÍNTESE + CULTURA
sl9 = base_slide(9)
add_header(sl9, "Síntese do Conteúdo & Sugestão Cultural")
add_solid_card(sl9, "Resumo do Tema:", "[Resumo sintético dos pontos chave ensinados na aula]", Inches(0.5), Inches(1.5), Inches(6.0), Inches(5.3), border_color=AZUL_TECH)
add_solid_card(sl9, "Sugestão Cultural:", "🎵 [Nome da Música / Vídeo recomendado sobre o assunto].\n\nO professor pode rodar o áudio/vídeo ou deixar como recomendação extra para o aluno praticar.", Inches(6.8), Inches(1.5), Inches(6.0), Inches(5.3), border_color=TURQUESA)

# SLIDES 10 E 11: PRÁTICA GUIADA
sl10 = base_slide(10)
add_header(sl10, "Prática Guiada — Exercícios 1 a 5", "Correção de Vícios")
add_clean_flashcard(sl10, "01", "[Lacunas contextualizadas com opções]", Inches(0.5), Inches(1.5), Inches(12.3), Inches(1.2), b_c=AZUL_TECH)
add_clean_flashcard(sl10, "02", "[Lacunas contextualizadas com opções]", Inches(0.5), Inches(2.9), Inches(6.0), Inches(1.8), b_c=TURQUESA)
add_clean_flashcard(sl10, "03", "[Preenchimento direto de verbos]", Inches(6.8), Inches(2.9), Inches(6.0), Inches(1.8), b_c=TURQUESA)
add_clean_flashcard(sl10, "04", "[Reestruturação / Transformação de frase]", Inches(0.5), Inches(4.9), Inches(6.0), Inches(1.9), b_c=AZUL_TECH)
add_clean_flashcard(sl10, "05", "[Reestruturação / Transformação de frase]", Inches(6.8), Inches(4.9), Inches(6.0), Inches(1.9), b_c=AZUL_TECH)

sl11 = base_slide(11)
add_header(sl11, "Prática Guiada — Exercícios 6 a 10", "Desafios Práticos")
add_clean_flashcard(sl11, "06", "[Preposição / Conector correto]", Inches(0.5), Inches(1.5), Inches(12.3), Inches(1.2), b_c=TURQUESA)
add_clean_flashcard(sl11, "07", "[Escolha da opção correta]", Inches(0.5), Inches(2.9), Inches(6.0), Inches(1.8), b_c=AZUL_TECH)
add_clean_flashcard(sl11, "08", "[Substituição por sinônimo formal]", Inches(6.8), Inches(2.9), Inches(6.0), Inches(1.8), b_c=AZUL_TECH)
add_clean_flashcard(sl11, "09", "[Identificação de erro e correção]", Inches(0.5), Inches(4.9), Inches(6.0), Inches(1.9), b_c=TURQUESA)
add_clean_flashcard(sl11, "10", "[Tradução de sentido sem uso literal]", Inches(6.8), Inches(4.9), Inches(6.0), Inches(1.9), b_c=TURQUESA)

# SLIDES 12 A 19: CONVERSAÇÃO (CRESCENDO NA DIREITA)
perguntas_config = [
    "[Pergunta principal 1 sobre o tema]",
    "[Pergunta principal 2 sobre o tema]",
    "[Pergunta principal 3 sobre o tema]",
    "[Pergunta principal 4 sobre o tema]",
    "[Pergunta principal 5 sobre o tema]",
    "[Pergunta principal 6 sobre o tema]",
    "[Pergunta principal 7 sobre o tema]",
    "[Pergunta principal 8 sobre o tema]"
]

for idx in range(8):
    s_num = 12 + idx
    sl = base_slide(s_num)
    add_header(sl, "Conversação")
    
    if idx == 0:
        p_desc = perguntas_config[0]
        # Layout Amplo Central
        add_solid_card(sl, "", p_desc, Inches(0.5), Inches(1.5), Inches(12.3), Inches(2.7), border_color=AZUL_TECH, desc_color=AZUL, desc_size=18, align_center_v=True)
        add_solid_card(sl, "", "• Para expandir: 'Dê um exemplo recente...'\n• Para concordar/discordar: 'Qual é o outro lado dessa moeda?'", Inches(0.5), Inches(4.4), Inches(12.3), Inches(2.4), border_color=TURQUESA, desc_color=SLATE, desc_size=15, align_center_v=True)
    else:
        # Coluna ESQUERDA: Pergunta Ativa e Prompts
        p_desc = perguntas_config[idx]
        add_solid_card(sl, "", p_desc, Inches(0.5), Inches(1.5), Inches(7.0), Inches(2.7), border_color=AZUL_TECH, desc_color=AZUL, desc_size=17, align_center_v=True)
        add_solid_card(sl, "", "• Para expandir: 'Dê um exemplo recente...'\n• Para concordar/discordar: 'Qual é o outro lado dessa moeda?'", Inches(0.5), Inches(4.4), Inches(7.0), Inches(2.4), border_color=TURQUESA, desc_color=SLATE, desc_size=14, align_center_v=True)
        
        # Coluna DIREITA: Histórico Empilhado que cresce dinamicamente
        history_list = [perguntas_config[i] for i in range(idx)]
        add_vertical_history_right(sl, history_list, left=7.8, top=1.5, width=5.0, total_height=5.3)

# SLIDE 20: ESCRITA
sl20 = base_slide(20)
add_header(sl20, "Atividade de Escrita", "Consolidação e Memorização")
add_solid_card(sl20, "Tarefa de Escrita Recomendada:", "[Instrução detalhada da tarefa: escrever um texto, frases ou diálogo]", Inches(0.5), Inches(1.5), Inches(12.3), Inches(2.8), border_color=AZUL_TECH)
add_solid_card(sl20, "", "Identificamos que a escrita nesta etapa é fundamental para você fixar as estruturas aprendidas, expandir vocabulário ativo e diminuir a tradução mental automática do seu idioma nativo.", Inches(0.5), Inches(4.5), Inches(12.3), Inches(2.3), border_color=TURQUESA, align_center_v=True)

# SLIDE 21: TAREFA DE CASA
sl21 = base_slide(21)
add_header(sl21, "Encerramento — Tarefa de Casa", "Ação e Prática Contínua")
add_solid_card(sl21, "Tarefa de Casa:", "• [Atividade prática 1 para realizar fora da aula]\n• [Atividade de fixação de vocabulário ou escuta]", Inches(0.5), Inches(1.5), Inches(12.3), Inches(5.2), border_color=AZUL_TECH)

# SLIDE 22: ENCERRAMENTO & FEEDBACK
sl22 = base_slide(22)
add_header(sl22, "Encerramento & Feedback Rápido", "Considerações Finais")
add_solid_card(sl22, "Agradecimento:", "Obrigado pela dedicação na aula de hoje!", Inches(0.5), Inches(1.5), Inches(12.3), Inches(2.0), border_color=TURQUESA, tit_size=20, desc_size=16)
add_solid_card(sl22, "Espaço para Feedback Rápido:", "[Espaço destinado para o professor coletar dúvidas rápidas e considerações do aluno]", Inches(0.5), Inches(3.7), Inches(12.3), Inches(3.0), border_color=AZUL_TECH, tit_size=18, desc_size=15)

out_path = "/var/www/haas-frontend-desk-mobile-oficial/public/temp-miniaturas/template-haas-modelo-10.pptx"
prs.save(out_path)
print(f"✅ MODELO 10 FINALIZADO: {out_path}")
