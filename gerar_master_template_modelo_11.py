import os
import urllib.request
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# PALETA OFICIAL HAAS LANGUAGE - PICASSO CUBISM EDITION
BRANCO = RGBColor(0xFF, 0xFF, 0xFF)
ROXO = RGBColor(0x4F, 0x35, 0xE2)
AZUL = RGBColor(0x1E, 0x3A, 0x8A)
SLATE = RGBColor(0x33, 0x41, 0x55)
TURQUESA = RGBColor(0x14, 0xC8, 0xB0)
BORDA = RGBColor(0xCB, 0xD5, 0xE1)
FUNDO_CARD = RGBColor(0xF8, 0xFA, 0xFC)
MUTED = RGBColor(0x94, 0xA3, 0xB8)
AZUL_CLARO = RGBColor(0xE0, 0xF2, 0xFE)
ROXO_CLARO = RGBColor(0xED, 0xE9, 0xFE)

LOGO_URL = "https://jdppxfokfhqjudwfwckd.supabase.co/storage/v1/object/public/haas-academy/assignments/Design%20sem%20nome%20(4).png"
LOGO_PATH = "/tmp/haas-logo.png"

if not os.path.exists(LOGO_PATH):
    try:
        urllib.request.urlretrieve(LOGO_URL, LOGO_PATH)
    except Exception:
        pass

W, H = Inches(13.333), Inches(7.5)
TOTAL_SLIDES = 22

prs = Presentation()
prs.slide_width = W
prs.slide_height = H

def base_slide(num, hide_logo=False):
    sl = prs.slides.add_slide(prs.slide_layouts[6])
    sl.background.fill.solid()
    sl.background.fill.fore_color.rgb = BRANCO
    
    # ARTE CUBISTA DE FUNDO / MOLDURA
    # Triângulo Roxo Topo-Esquerdo
    t1 = sl.shapes.add_shape(MSO_SHAPE.RIGHT_TRIANGLE, Inches(0), Inches(0), Inches(1.5), Inches(1.5))
    t1.fill.solid(); t1.fill.fore_color.rgb = ROXO; t1.line.fill.background()
    t1.rotation = 270
    
    # Bloco Retangular Assimétrico Base-Esquerda
    r1 = sl.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(6.8), Inches(4.5), Inches(0.7))
    r1.fill.solid(); r1.fill.fore_color.rgb = AZUL; r1.line.fill.background()
    
    # Triângulo Turquesa Base-Direita (recortando o rodapé)
    t2 = sl.shapes.add_shape(MSO_SHAPE.RIGHT_TRIANGLE, Inches(12.333), Inches(6.5), Inches(1.0), Inches(1.0))
    t2.fill.solid(); t2.fill.fore_color.rgb = TURQUESA; t2.line.fill.background()
    t2.rotation = 90
    
    # Linha Diagonal Sutil Topo-Direita
    l1 = sl.shapes.add_shape(MSO_SHAPE.ISOSCELES_TRIANGLE, Inches(11.0), Inches(0), Inches(2.333), Inches(0.8))
    l1.fill.solid(); l1.fill.fore_color.rgb = AZUL_CLARO; l1.line.fill.background()
    l1.rotation = 180

    if os.path.exists(LOGO_PATH) and not hide_logo:
        sl.shapes.add_picture(LOGO_PATH, Inches(10.8), Inches(0.2), height=Inches(0.65))

    # Rodapé Textos
    tb = sl.shapes.add_textbox(Inches(0.5), Inches(7.1), Inches(8), Inches(0.35))
    p = tb.text_frame.paragraphs[0]
    p.text = "HAAS LANGUAGE  •  PICASSO CUBISM EDITION"
    p.font.size = Pt(10); p.font.color.rgb = BRANCO; p.font.bold = True
    
    tb2 = sl.shapes.add_textbox(Inches(10.8), Inches(7.05), Inches(2.0), Inches(0.35))
    p2 = tb2.text_frame.paragraphs[0]
    p2.text = f"SLIDE {num} / {TOTAL_SLIDES}"
    p2.font.size = Pt(10); p2.font.color.rgb = SLATE; p2.font.bold = True; p2.alignment = PP_ALIGN.RIGHT
    return sl

def add_header(sl, tit, sub=""):
    tb = sl.shapes.add_textbox(Inches(1.8), Inches(0.2), Inches(9), Inches(0.65))
    tf = tb.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]; p.text = tit
    p.font.size = Pt(32); p.font.color.rgb = ROXO; p.font.bold = True
    
    if sub:
        tb_s = sl.shapes.add_textbox(Inches(1.8), Inches(0.85), Inches(9), Inches(0.4))
        p_s = tb_s.text_frame.paragraphs[0]
        p_s.text = sub
        p_s.font.size = Pt(18); p_s.font.color.rgb = AZUL; p_s.font.bold = True

def add_cubist_block(sl, tit, desc, left, top, width, height, bg_color=FUNDO_CARD, border_color=BORDA, tit_color=ROXO, desc_color=SLATE, tit_size=16, desc_size=14, align_center_v=False, shadow_color=None):
    """Blocos Cubistas Rígidos (Cantos Retos) com Sombra Deslocada Geométrica"""
    
    if shadow_color:
        sh_offset = Inches(0.1)
        sh = sl.shapes.add_shape(MSO_SHAPE.RECTANGLE, left + sh_offset, top + sh_offset, width, height)
        sh.fill.solid(); sh.fill.fore_color.rgb = shadow_color; sh.line.fill.background()

    s = sl.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    s.fill.solid(); s.fill.fore_color.rgb = bg_color
    s.line.color.rgb = border_color; s.line.width = Pt(1.5)
        
    tf = s.text_frame; tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE if align_center_v else MSO_ANCHOR.TOP
    margin_val = Inches(0.12) if height < Inches(1.0) else Inches(0.2)
    tf.margin_left = margin_val; tf.margin_top = margin_val; tf.margin_right = margin_val; tf.margin_bottom = margin_val
    
    if tit:
        p = tf.paragraphs[0]
        p.text = tit; p.font.size = Pt(tit_size); p.font.color.rgb = tit_color; p.font.bold = True
        if desc:
            p2 = tf.add_paragraph()
            p2.text = desc; p2.font.size = Pt(desc_size); p2.font.color.rgb = desc_color; p2.space_before = Pt(6)
    else:
        if desc:
            p = tf.paragraphs[0]
            p.text = desc; p.font.size = Pt(desc_size); p.font.color.rgb = desc_color

def img_box_cubist(sl, left, top, width, height, label="[Espaço de Imagem]"):
    sh = sl.shapes.add_shape(MSO_SHAPE.RECTANGLE, left + Inches(0.15), top + Inches(0.15), width, height)
    sh.fill.solid(); sh.fill.fore_color.rgb = TURQUESA; sh.line.fill.background()
    
    s = sl.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    s.fill.solid(); s.fill.fore_color.rgb = ROXO; s.line.color.rgb = BRANCO; s.line.width = Pt(2)
    tf = s.text_frame; tf.word_wrap = True; tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER
    p = tf.paragraphs[0]; p.text = label; p.font.size = Pt(18); p.font.color.rgb = BRANCO; p.font.bold = True

def add_vertical_history_cubist(sl, history_list, left=7.8, top=1.5, width=5.0, total_height=5.3):
    count = len(history_list)
    if count == 0: return
    gap = 0.12
    avail_h = total_height - (gap * (count - 1))
    card_h = avail_h / count
    for i, h_txt in enumerate(history_list):
        y = top + i * (card_h + gap)
        clean_t = h_txt[:40] + "..." if len(h_txt) > 40 else h_txt
        d_sz = 12 if count <= 3 else (11 if count <= 5 else 10)
        bg_col = ROXO_CLARO if i % 2 == 0 else AZUL_CLARO
        b_col = ROXO if i % 2 == 0 else AZUL
        add_cubist_block(sl, "", clean_t, Inches(left), Inches(y), Inches(width), Inches(card_h), 
                         bg_color=bg_col, border_color=b_col, desc_color=SLATE, desc_size=d_sz, align_center_v=True)

def add_cubist_flashcard(sl, num_str, desc, left, top, width, height, b_c=AZUL):
    s = sl.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    s.fill.solid(); s.fill.fore_color.rgb = BRANCO; s.line.color.rgb = b_c; s.line.width = Pt(1.5)
    
    chip = sl.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, Inches(0.55), Inches(0.45))
    chip.fill.solid(); chip.fill.fore_color.rgb = b_c; chip.line.fill.background()
    tf_chip = chip.text_frame; tf_chip.vertical_anchor = MSO_ANCHOR.MIDDLE
    p_chip = tf_chip.paragraphs[0]; p_chip.text = num_str; p_chip.alignment = PP_ALIGN.CENTER
    p_chip.font.size = Pt(12); p_chip.font.color.rgb = BRANCO; p_chip.font.bold = True
    
    tf = s.text_frame; tf.word_wrap = True; tf.vertical_anchor = MSO_ANCHOR.TOP
    tf.margin_left = Inches(0.8); tf.margin_top = Inches(0.15); tf.margin_right = Inches(0.2); tf.margin_bottom = Inches(0.12)
    p = tf.paragraphs[0]; p.text = desc; p.font.size = Pt(14); p.font.color.rgb = SLATE; p.font.bold = False

# SLIDE 1: CAPA CUBISTA (Blocos Assimétricos)
sl1 = base_slide(1, hide_logo=True)
img_box_cubist(sl1, Inches(0.5), Inches(1.5), Inches(5.5), Inches(5.2), "[Imagem Cubista / Arte]")
add_cubist_block(sl1, "[Título da Aula / Tema Principal]", "", Inches(6.5), Inches(1.5), Inches(6.3), Inches(2.0), bg_color=BRANCO, border_color=ROXO, tit_size=36, shadow_color=AZUL_CLARO, align_center_v=True)
add_cubist_block(sl1, "[Subtítulo / Descrição da Aula]", "", Inches(6.5), Inches(3.8), Inches(6.3), Inches(1.2), bg_color=AZUL, border_color=AZUL, tit_color=BRANCO, tit_size=20, shadow_color=TURQUESA, align_center_v=True)
add_cubist_block(sl1, "📅 Data da Aula: [DD/MM/AAAA]", "", Inches(6.5), Inches(5.3), Inches(6.3), Inches(0.8), bg_color=BRANCO, border_color=TURQUESA, tit_color=TURQUESA, tit_size=16, shadow_color=ROXO_CLARO, align_center_v=True)

# SLIDE 2: PRESENÇA & INTEGRANTES
sl2 = base_slide(2)
add_header(sl2, "Presença & Integrantes")
add_cubist_block(sl2, "Professor(a):", "[Nome do(a) Professor(a)]", Inches(0.5), Inches(1.5), Inches(12.3), Inches(1.5), shadow_color=TURQUESA)
add_cubist_block(sl2, "Aluno(s):", "[Nome do Aluno / Lista da Turma]", Inches(0.5), Inches(3.4), Inches(12.3), Inches(3.4), shadow_color=ROXO_CLARO)

# SLIDE 3: WARM-UP (Blocos Empilhados + Sombra)
sl3 = base_slide(3)
add_header(sl3, "Aquecimento / Warm-up", "Foco 100% na Fala (3 Minutos)")
add_cubist_block(sl3, "Pergunta Provocativa Rápida:", "[Insira aqui a pergunta provocativa de aquecimento]", Inches(0.5), Inches(1.5), Inches(7.0), Inches(2.4), bg_color=ROXO, tit_color=TURQUESA, desc_color=BRANCO, tit_size=20, desc_size=16, shadow_color=TURQUESA)
add_cubist_block(sl3, "", "Ativa o córtex pré-frontal e reduz a inibição linguística inicial.", Inches(0.5), Inches(4.3), Inches(7.0), Inches(1.8), bg_color=BRANCO, border_color=AZUL, desc_color=SLATE, desc_size=14, shadow_color=AZUL_CLARO)
img_box_cubist(sl3, Inches(8.0), Inches(1.5), Inches(4.8), Inches(4.6), "[Imagem Warm-up]")

# SLIDE 4: OBJETIVOS
sl4 = base_slide(4)
add_header(sl4, "Objetivos da Aula")
add_cubist_block(sl4, "Metas da Aula:", "• [Objetivo 1]\n• [Objetivo 2]\n• [Objetivo 3]", Inches(0.5), Inches(1.5), Inches(7.5), Inches(5.3), shadow_color=ROXO_CLARO)
add_cubist_block(sl4, "Nota de Escrita:", "Nesta aula daremos atenção à escrita para reforçar a memorização de estruturas/vocabulário e reduzir a interferência da língua materna.", Inches(8.4), Inches(1.5), Inches(4.4), Inches(5.3), shadow_color=TURQUESA)

# SLIDE 5: METODOLOGIA (Painel Piet Mondrian)
sl5 = base_slide(5)
add_header(sl5, "A Nossa Metodologia", "Estrutura Fixa de Aprendizado Haas Language")
add_cubist_block(sl5, "1. Fala Ativa", "Conversação inicial sem julgamentos.", Inches(0.5), Inches(1.5), Inches(5.8), Inches(2.5), bg_color=ROXO, tit_color=TURQUESA, desc_color=BRANCO)
add_cubist_block(sl5, "2. Exposição", "Apresentação leve da gramática.", Inches(6.5), Inches(1.5), Inches(6.3), Inches(2.5), bg_color=TURQUESA, tit_color=ROXO, desc_color=SLATE)
add_cubist_block(sl5, "3. Prática", "Exercícios com suporte.", Inches(0.5), Inches(4.2), Inches(3.8), Inches(2.6), bg_color=AZUL, tit_color=TURQUESA, desc_color=BRANCO)
add_cubist_block(sl5, "4. Conversação", "Aplicação prática das estruturas.", Inches(4.5), Inches(4.2), Inches(4.0), Inches(2.6), bg_color=ROXO_CLARO, tit_color=ROXO, desc_color=SLATE)
add_cubist_block(sl5, "5. Escrita", "Consolidação visual.", Inches(8.7), Inches(4.2), Inches(4.1), Inches(2.6), bg_color=BRANCO, border_color=AZUL, tit_color=AZUL, desc_color=SLATE)

# SLIDES 6 A 8: EXPOSIÇÃO
for s_num in [6, 7, 8]:
    sl = base_slide(s_num)
    add_header(sl, f"Explicação do Tema — Parte {s_num - 5}")
    add_cubist_block(sl, "Conteúdo Principal:", f"[Explicação e exemplo central da estrutura - Bloco {s_num - 5}]", Inches(0.5), Inches(1.5), Inches(12.3), Inches(2.5), shadow_color=ROXO_CLARO, tit_size=20, desc_size=15)
    add_cubist_block(sl, "Exemplos Práticos:", "[Exemplos adicionais e frases do dia a dia]", Inches(0.5), Inches(4.4), Inches(6.0), Inches(2.4), shadow_color=TURQUESA)
    add_cubist_block(sl, "Notas & Detalhes:", "[Exceções, nuances gramaticais e contexto de uso]", Inches(6.8), Inches(4.4), Inches(6.0), Inches(2.4), shadow_color=AZUL_CLARO)

# SLIDE 9: SÍNTESE + CULTURA
sl9 = base_slide(9)
add_header(sl9, "Síntese do Conteúdo & Sugestão Cultural")
add_cubist_block(sl9, "Resumo do Tema:", "[Resumo sintético dos pontos chave ensinados na aula]", Inches(0.5), Inches(1.5), Inches(6.0), Inches(5.3), shadow_color=ROXO_CLARO)
add_cubist_block(sl9, "Sugestão Cultural:", "🎵 [Nome da Música / Vídeo recomendado sobre o assunto].\n\nO professor pode rodar o áudio/vídeo ou deixar como recomendação extra para o aluno praticar.", Inches(6.8), Inches(1.5), Inches(6.0), Inches(5.3), shadow_color=TURQUESA)

# SLIDES 10 E 11: PRÁTICA GUIADA
sl10 = base_slide(10)
add_header(sl10, "Prática Guiada — Exercícios 1 a 5", "Correção de Vícios")
add_cubist_flashcard(sl10, "01", "[Lacunas contextualizadas com opções]", Inches(0.5), Inches(1.5), Inches(12.3), Inches(1.2), b_c=ROXO)
add_cubist_flashcard(sl10, "02", "[Lacunas contextualizadas com opções]", Inches(0.5), Inches(2.9), Inches(6.0), Inches(1.8), b_c=AZUL)
add_cubist_flashcard(sl10, "03", "[Preenchimento direto de verbos]", Inches(6.8), Inches(2.9), Inches(6.0), Inches(1.8), b_c=AZUL)
add_cubist_flashcard(sl10, "04", "[Reestruturação / Transformação de frase]", Inches(0.5), Inches(4.9), Inches(6.0), Inches(1.9), b_c=TURQUESA)
add_cubist_flashcard(sl10, "05", "[Reestruturação / Transformação de frase]", Inches(6.8), Inches(4.9), Inches(6.0), Inches(1.9), b_c=TURQUESA)

sl11 = base_slide(11)
add_header(sl11, "Prática Guiada — Exercícios 6 a 10", "Desafios Práticos")
add_cubist_flashcard(sl11, "06", "[Preposição / Conector correto]", Inches(0.5), Inches(1.5), Inches(12.3), Inches(1.2), b_c=TURQUESA)
add_cubist_flashcard(sl11, "07", "[Escolha da opção correta]", Inches(0.5), Inches(2.9), Inches(6.0), Inches(1.8), b_c=ROXO)
add_cubist_flashcard(sl11, "08", "[Substituição por sinônimo formal]", Inches(6.8), Inches(2.9), Inches(6.0), Inches(1.8), b_c=ROXO)
add_cubist_flashcard(sl11, "09", "[Identificação de erro e correção]", Inches(0.5), Inches(4.9), Inches(6.0), Inches(1.9), b_c=AZUL)
add_cubist_flashcard(sl11, "10", "[Tradução de sentido sem uso literal]", Inches(6.8), Inches(4.9), Inches(6.0), Inches(1.9), b_c=AZUL)

# SLIDES 12 A 19: CONVERSAÇÃO (Histórico na direita com blocos geométricos)
perguntas_config = [
    "[Pergunta principal 1 sobre o tema]", "[Pergunta principal 2 sobre o tema]", "[Pergunta principal 3 sobre o tema]",
    "[Pergunta principal 4 sobre o tema]", "[Pergunta principal 5 sobre o tema]", "[Pergunta principal 6 sobre o tema]",
    "[Pergunta principal 7 sobre o tema]", "[Pergunta principal 8 sobre o tema]"
]
for idx in range(8):
    s_num = 12 + idx
    sl = base_slide(s_num)
    add_header(sl, "Conversação")
    
    if idx == 0:
        add_cubist_block(sl, "", perguntas_config[0], Inches(0.5), Inches(1.5), Inches(12.3), Inches(2.7), shadow_color=ROXO_CLARO, desc_size=18, align_center_v=True)
        add_cubist_block(sl, "", "• Para expandir: 'Dê um exemplo recente...'\n• Para concordar/discordar: 'Qual é o outro lado dessa moeda?'", Inches(0.5), Inches(4.5), Inches(12.3), Inches(2.3), bg_color=ROXO, desc_color=BRANCO, shadow_color=TURQUESA, desc_size=15, align_center_v=True)
    else:
        add_cubist_block(sl, "", perguntas_config[idx], Inches(0.5), Inches(1.5), Inches(7.0), Inches(2.7), shadow_color=ROXO_CLARO, desc_size=17, align_center_v=True)
        add_cubist_block(sl, "", "• Para expandir: 'Dê um exemplo recente...'\n• Para concordar/discordar: 'Qual é o outro lado dessa moeda?'", Inches(0.5), Inches(4.5), Inches(7.0), Inches(2.3), bg_color=ROXO, desc_color=BRANCO, shadow_color=TURQUESA, desc_size=14, align_center_v=True)
        
        history_list = [perguntas_config[i] for i in range(idx)]
        add_vertical_history_cubist(sl, history_list, left=7.8, top=1.5, width=5.0, total_height=5.3)

# SLIDE 20: ESCRITA
sl20 = base_slide(20)
add_header(sl20, "Atividade de Escrita", "Consolidação e Memorização")
add_cubist_block(sl20, "Tarefa de Escrita Recomendada:", "[Instrução detalhada da tarefa: escrever um texto, frases ou diálogo]", Inches(0.5), Inches(1.5), Inches(12.3), Inches(2.8), shadow_color=ROXO_CLARO)
add_cubist_block(sl20, "", "Identificamos que a escrita nesta etapa é fundamental para você fixar as estruturas aprendidas, expandir vocabulário ativo e diminuir a tradução mental automática do seu idioma nativo.", Inches(0.5), Inches(4.6), Inches(12.3), Inches(2.2), bg_color=TURQUESA, desc_color=SLATE, align_center_v=True)

# SLIDE 21: TAREFA DE CASA
sl21 = base_slide(21)
add_header(sl21, "Encerramento — Tarefa de Casa", "Ação e Prática Contínua")
add_cubist_block(sl21, "Tarefa de Casa:", "• [Atividade prática 1 para realizar fora da aula]\n• [Atividade de fixação de vocabulário ou escuta]", Inches(0.5), Inches(1.5), Inches(12.3), Inches(5.2), shadow_color=AZUL_CLARO)

# SLIDE 22: ENCERRAMENTO & FEEDBACK
sl22 = base_slide(22)
add_header(sl22, "Encerramento & Feedback Rápido", "Considerações Finais")
add_cubist_block(sl22, "Agradecimento:", "Obrigado pela dedicação na aula de hoje!", Inches(0.5), Inches(1.5), Inches(12.3), Inches(2.0), shadow_color=TURQUESA, tit_size=20, desc_size=16)
add_cubist_block(sl22, "Espaço para Feedback Rápido:", "[Espaço destinado para o professor coletar dúvidas rápidas e considerações do aluno]", Inches(0.5), Inches(3.8), Inches(12.3), Inches(2.9), shadow_color=ROXO_CLARO, tit_size=18, desc_size=15)

out_path = "/var/www/haas-frontend-desk-mobile-oficial/public/temp-miniaturas/template-haas-modelo-11-cubist-picasso.pptx"
prs.save(out_path)
print(f"✅ MODELO 11 CUBISMO PURO (PICASSO) GERADO EM: {out_path}")
