import os
import urllib.request
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml import parse_xml

# PALETA OFICIAL HAAS LANGUAGE - OUTLINE TYPOGRAPHY EDITION
BRANCO = RGBColor(0xFF, 0xFF, 0xFF)
ROXO = RGBColor(0x4F, 0x35, 0xE2)
AZUL = RGBColor(0x1E, 0x3A, 0x8A)
SLATE = RGBColor(0x33, 0x41, 0x55)
TURQUESA = RGBColor(0x14, 0xC8, 0xB0)
BORDA = RGBColor(0xCB, 0xD5, 0xE1)
FUNDO_OUTLINE = RGBColor(0xF8, 0xFA, 0xFC)
ROXO_CLARO = RGBColor(0xED, 0xE9, 0xFE)
TURQUESA_CLARO = RGBColor(0xE0, 0xF2, 0xF1)

LOGO_URL = "https://jdppxfokfhqjudwfwckd.supabase.co/storage/v1/object/public/haas-academy/assignments/Design%20sem%20nome%20(4).png"
LOGO_PATH = "/tmp/haas-logo.png"

if not os.path.exists(LOGO_PATH):
    try:
        urllib.request.urlretrieve(LOGO_URL, LOGO_PATH)
    except Exception:
        pass

W, H = Inches(13.333), Inches(7.5)
TOTAL_SLIDES = 22

prs = Presentation()
prs.slide_width = W
prs.slide_height = H

def apply_outline_to_run(run, outline_hex="4F35E2", width_pt=1.5):
    """Aplica efeito de fonte vazada (Outline/Stroke) no texto via OXML"""
    rPr = run._r.get_or_add_rPr()
    for fill_tag in ['solidFill', 'gradFill', 'pattFill']:
        f = rPr.find(f'{{http://schemas.openxmlformats.org/drawingml/2006/main}}{fill_tag}')
        if f is not None:
            rPr.remove(f)
            
    noFill = parse_xml('<a:noFill xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main"/>')
    rPr.append(noFill)
    
    w_emu = int(width_pt * 12700)
    ln_xml = f"""
    <a:ln xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main" w="{w_emu}">
      <a:solidFill>
        <a:srgbClr val="{outline_hex}"/>
      </a:solidFill>
    </a:ln>
    """
    ln = parse_xml(ln_xml)
    rPr.append(ln)

def base_slide(num, hide_logo=False):
    sl = prs.slides.add_slide(prs.slide_layouts[6])
    sl.background.fill.solid()
    sl.background.fill.fore_color.rgb = FUNDO_OUTLINE

    frame = sl.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.5), Inches(0.2), Pt(2), Inches(6.8))
    frame.fill.solid(); frame.fill.fore_color.rgb = ROXO; frame.line.fill.background()

    if os.path.exists(LOGO_PATH) and not hide_logo:
        sl.shapes.add_picture(LOGO_PATH, Inches(10.8), Inches(0.2), height=Inches(0.65))
        
    tb = sl.shapes.add_textbox(Inches(0.7), Inches(7.1), Inches(8), Inches(0.35))
    p = tb.text_frame.paragraphs[0]
    p.text = "HAAS LANGUAGE  •  OUTLINE TYPOGRAPHY EDITION"
    p.font.size = Pt(10); p.font.color.rgb = AZUL; p.font.bold = True
    
    tb2 = sl.shapes.add_textbox(Inches(10.8), Inches(7.1), Inches(2.0), Inches(0.35))
    p2 = tb2.text_frame.paragraphs[0]
    p2.text = f"{num:02d} // {TOTAL_SLIDES:02d}"
    p2.font.size = Pt(10); p2.font.color.rgb = ROXO; p2.font.bold = True; p2.alignment = PP_ALIGN.RIGHT
    return sl

def add_header(sl, tit, sub=""):
    tb = sl.shapes.add_textbox(Inches(0.7), Inches(0.25), Inches(9.8), Inches(0.75))
    tf = tb.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]; p.text = tit
    p.font.size = Pt(30); p.font.color.rgb = ROXO; p.font.bold = True
    
    if sub:
        tb_s = sl.shapes.add_textbox(Inches(0.7), Inches(0.98), Inches(9.8), Inches(0.4))
        p_s = tb_s.text_frame.paragraphs[0]
        p_s.text = sub
        p_s.font.size = Pt(16); p_s.font.color.rgb = AZUL; p_s.font.bold = True

def add_outline_card(sl, tit, desc, left, top, width, height, bg_color=BRANCO, border_color=ROXO, tit_color=ROXO, desc_color=SLATE, tit_size=16, desc_size=14, align_center_v=False, outline_word=""):
    s = sl.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    s.fill.solid(); s.fill.fore_color.rgb = bg_color
    s.line.color.rgb = border_color; s.line.width = Pt(1.5)
        
    tf = s.text_frame; tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE if align_center_v else MSO_ANCHOR.TOP
    tf.margin_left = Inches(0.25); tf.margin_top = Inches(0.2); tf.margin_right = Inches(0.25); tf.margin_bottom = Inches(0.2)
    
    if outline_word:
        p_o = tf.paragraphs[0]
        r_o = p_o.add_run()
        r_o.text = outline_word.upper()
        r_o.font.size = Pt(28)
        r_o.font.bold = True
        apply_outline_to_run(r_o, "4F35E2", width_pt=1.5)
        
        if tit:
            p = tf.add_paragraph()
            p.text = tit; p.font.size = Pt(tit_size); p.font.color.rgb = tit_color; p.font.bold = True; p.space_before = Pt(4)
    else:
        if tit:
            p = tf.paragraphs[0]
            p.text = tit; p.font.size = Pt(tit_size); p.font.color.rgb = tit_color; p.font.bold = True
            
    if desc:
        p2 = tf.add_paragraph() if (outline_word or tit) else tf.paragraphs[0]
        p2.text = desc; p2.font.size = Pt(desc_size); p2.font.color.rgb = desc_color
        if outline_word or tit:
            p2.space_before = Pt(4)

def img_outline_box(sl, left, top, width, height, label="[IMAGEM EDITORIAL]"):
    s = sl.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    s.fill.solid(); s.fill.fore_color.rgb = BRANCO
    s.line.color.rgb = TURQUESA; s.line.width = Pt(2)
    tf = s.text_frame; tf.word_wrap = True; tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER
    p = tf.paragraphs[0]; p.text = label; p.font.size = Pt(16); p.font.color.rgb = ROXO; p.font.bold = True

def add_outline_flashcard(sl, num_str, desc, left, top, width, height, b_c=ROXO):
    s = sl.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    s.fill.solid(); s.fill.fore_color.rgb = BRANCO; s.line.color.rgb = b_c; s.line.width = Pt(1.5)
    
    tf_tag = s.text_frame; tf_tag.word_wrap = True; tf_tag.vertical_anchor = MSO_ANCHOR.MIDDLE
    tf_tag.margin_left = Inches(0.2); tf_tag.margin_top = Inches(0.12); tf_tag.margin_right = Inches(0.2)
    
    p = tf_tag.paragraphs[0]
    r_n = p.add_run()
    r_n.text = f"#{num_str}  "
    r_n.font.size = Pt(22)
    r_n.font.bold = True
    apply_outline_to_run(r_n, "14C8B0" if b_c == TURQUESA else "4F35E2", width_pt=1.5)
    
    r_d = p.add_run()
    r_d.text = desc
    r_d.font.size = Pt(14)
    r_d.font.color.rgb = SLATE

# SLIDE 1
sl1 = base_slide(1, hide_logo=True)
tb_c = sl1.shapes.add_textbox(Inches(0.7), Inches(1.2), Inches(7.0), Inches(2.2))
tf_c = tb_c.text_frame; tf_c.word_wrap = True
p = tf_c.paragraphs[0]
r_out = p.add_run()
r_out.text = "LANGUAGE\n"
r_out.font.size = Pt(54)
r_out.font.bold = True
apply_outline_to_run(r_out, "4F35E2", width_pt=2.0)

r_sol = p.add_run()
r_sol.text = "[Título da Aula / Tema Principal]"
r_sol.font.size = Pt(22)
r_sol.font.color.rgb = AZUL
r_sol.font.bold = True

tb_sub = sl1.shapes.add_textbox(Inches(0.7), Inches(3.8), Inches(7.0), Inches(1.8))
tf_sub = tb_sub.text_frame; tf_sub.word_wrap = True
p_sub = tf_sub.paragraphs[0]; p_sub.text = "[Subtítulo / Descrição da Aula]"
p_sub.font.size = Pt(18); p_sub.font.color.rgb = SLATE; p_sub.font.bold = True

p_date = tf_sub.add_paragraph(); p_date.text = "📅 Data da Aula: [DD/MM/AAAA]"
p_date.font.size = Pt(15); p_date.font.color.rgb = TURQUESA; p_date.font.bold = True; p_date.space_before = Pt(12)

img_outline_box(sl1, Inches(8.0), Inches(1.2), Inches(4.8), Inches(5.5), "[IMAGEM STREETWEAR]")

# SLIDE 2
sl2 = base_slide(2)
add_header(sl2, "Presença & Integrantes", "Membros Registrados")
add_outline_card(sl2, "Professor(a):", "[Nome do(a) Professor(a)]", Inches(0.7), Inches(1.5), Inches(5.8), Inches(5.3), outline_word="TEACHER", border_color=ROXO)
add_outline_card(sl2, "Aluno(s):", "[Nome do Aluno / Lista da Turma]", Inches(6.8), Inches(1.5), Inches(6.0), Inches(5.3), outline_word="STUDENTS", border_color=TURQUESA)

# SLIDE 3
sl3 = base_slide(3)
add_header(sl3, "Aquecimento / Warm-up", "Foco 100% na Fala (3 Minutos)")
add_outline_card(sl3, "Pergunta Provocativa Rápida:", "[Insira aqui a pergunta provocativa de aquecimento]", Inches(0.7), Inches(1.5), Inches(6.8), Inches(2.6), outline_word="WARMUP", border_color=ROXO)
add_outline_card(sl3, "", "Ativa o córtex pré-frontal e reduz a inibição linguística inicial.", Inches(0.7), Inches(4.3), Inches(6.8), Inches(2.5), outline_word="NEURO", border_color=TURQUESA)
img_outline_box(sl3, Inches(7.8), Inches(1.5), Inches(5.0), Inches(5.3), "[IMAGEM WARMUP]")

# SLIDE 4
sl4 = base_slide(4)
add_header(sl4, "Objetivos da Aula", "PONTOS DE FOCO")
add_outline_card(sl4, "Metas da Aula:", "• [Objetivo 1]\n• [Objetivo 2]\n• [Objetivo 3]", Inches(0.7), Inches(1.5), Inches(7.2), Inches(5.3), outline_word="GOALS", border_color=ROXO)
add_outline_card(sl4, "Nota de Escrita:", "Nesta aula daremos atenção à escrita para reforçar a memorização de estruturas/vocabulário e reduzir a interferência da língua materna.", Inches(8.2), Inches(1.5), Inches(4.6), Inches(5.3), outline_word="WRITE", border_color=TURQUESA)

# SLIDE 5
sl5 = base_slide(5)
add_header(sl5, "A Nossa Metodologia", "Estrutura Haas Language")
metodologias = [
    ("01", "Fala Ativa", "Conversação inicial sem julgamentos.", ROXO),
    ("02", "Exposição", "Apresentação leve da gramática.", TURQUESA),
    ("03", "Prática", "Exercícios com suporte.", AZUL),
    ("04", "Conversa", "Aplicação prática real.", ROXO),
    ("05", "Escrita", "Consolidação e memorização.", TURQUESA)
]
top_m = 1.5
for num_o, tit_m, desc_m, col_m in metodologias:
    add_outline_card(sl5, tit_m, desc_m, Inches(0.7), Inches(top_m), Inches(12.1), Inches(0.95), outline_word=f"STEP {num_o}", border_color=col_m, tit_size=15, desc_size=13)
    top_m += 1.08

# SLIDES 6 A 8
for s_num in [6, 7, 8]:
    sl = base_slide(s_num)
    add_header(sl, f"Explicação do Tema — Parte {s_num - 5}")
    add_outline_card(sl, "Conceito Principal:", f"[Explicação e exemplo central da estrutura - Bloco {s_num - 5}]", Inches(0.7), Inches(1.5), Inches(12.1), Inches(2.6), outline_word="STRUCTURE", border_color=ROXO, tit_size=18, desc_size=15)
    add_outline_card(sl, "Aplicações Práticas:", "[Exemplos adicionais e frases do dia a dia]", Inches(0.7), Inches(4.3), Inches(5.9), Inches(2.5), outline_word="EXAMP", border_color=TURQUESA, tit_size=15, desc_size=13)
    add_outline_card(sl, "Notas & Detalhes:", "[Exceções, nuances gramaticais e contexto]", Inches(6.9), Inches(4.3), Inches(5.9), Inches(2.5), outline_word="RULES", border_color=AZUL, tit_size=15, desc_size=13)

# SLIDE 9
sl9 = base_slide(9)
add_header(sl9, "Síntese do Conteúdo & Sugestão Cultural")
add_outline_card(sl9, "Resumo do Tema:", "[Resumo sintético dos pontos chave ensinados na aula]", Inches(0.7), Inches(1.5), Inches(5.9), Inches(5.3), outline_word="SUMMARY", border_color=ROXO)
add_outline_card(sl9, "Mídia Recomendada:", "🎵 [Nome da Música / Vídeo recomendado sobre o assunto].\n\nO professor pode rodar o áudio/vídeo ou deixar como recomendação extra para o aluno praticar.", Inches(6.9), Inches(1.5), Inches(5.9), Inches(5.3), outline_word="MEDIA", border_color=TURQUESA)

# SLIDES 10 E 11
sl10 = base_slide(10)
add_header(sl10, "Prática Guiada — Blocos 1 a 5", "Correção de Vícios")
exs_1_5 = [
    ("01", "[Lacunas contextualizadas com opções]"),
    ("02", "[Lacunas contextualizadas com opções]"),
    ("03", "[Preenchimento direto de verbos]"),
    ("04", "[Reestruturação / Transformação de frase]"),
    ("05", "[Reestruturação / Transformação de frase]")
]
for idx, (num_s, d_e) in enumerate(exs_1_5):
    col = idx % 2; row = idx // 2
    l = 0.7 if col == 0 else 6.9
    w = 5.9 if (idx < 4 or len(exs_1_5) % 2 == 0) else 12.1
    if idx == 4: l = 0.7
    t = 1.5 + row * 1.7
    add_outline_flashcard(sl10, num_s, d_e, Inches(l), Inches(t), Inches(w), Inches(1.5), b_c=ROXO if idx%2==0 else TURQUESA)

sl11 = base_slide(11)
add_header(sl11, "Prática Guiada — Blocos 6 a 10", "Desafios Práticos")
exs_6_10 = [
    ("06", "[Preposição / Conector correto]"),
    ("07", "[Escolha da opção correta]"),
    ("08", "[Substituição por sinônimo formal]"),
    ("09", "[Identificação de erro e correção]"),
    ("10", "[Tradução de sentido sem uso literal]")
]
for idx, (num_s, d_e) in enumerate(exs_6_10):
    col = idx % 2; row = idx // 2
    l = 0.7 if col == 0 else 6.9
    w = 5.9 if (idx < 4 or len(exs_6_10) % 2 == 0) else 12.1
    if idx == 4: l = 0.7
    t = 1.5 + row * 1.7
    add_outline_flashcard(sl11, num_s, d_e, Inches(l), Inches(t), Inches(w), Inches(1.5), b_c=TURQUESA if idx%2==0 else ROXO)

# SLIDES 12 A 19: CONVERSAÇÃO (OUTLINE STAIRCASE INÉDITA)
perguntas_config = [
    "[Pergunta principal 1 sobre o tema]", "[Pergunta principal 2 sobre o tema]", "[Pergunta principal 3 sobre o tema]",
    "[Pergunta principal 4 sobre o tema]", "[Pergunta principal 5 sobre o tema]", "[Pergunta principal 6 sobre o tema]",
    "[Pergunta principal 7 sobre o tema]", "[Pergunta principal 8 sobre o tema]"
]

def add_outline_history_staircase(sl, history_list):
    count = len(history_list)
    if count == 0: return
    gap = 0.12
    avail_h = 5.3 - (gap * (count - 1))
    card_h = avail_h / count
    
    for i, h_txt in enumerate(history_list):
        y_pos = 1.5 + i * (card_h + gap)
        clean_t = h_txt[:32] + "..." if len(h_txt) > 32 else h_txt
        d_sz = 11 if count <= 3 else (10 if count <= 5 else 9)
        bg_col = ROXO_CLARO if i % 2 == 0 else TURQUESA_CLARO
        border_c = ROXO if i % 2 == 0 else TURQUESA
        add_outline_card(sl, "", clean_t, Inches(8.2), Inches(y_pos), Inches(4.6), Inches(card_h), bg_color=bg_col, border_color=border_c, outline_word=f"Q0{i+1}", desc_size=d_sz)

for idx in range(8):
    s_num = 12 + idx
    sl = base_slide(s_num)
    add_header(sl, "Conversação", f"Sessão Ativa [{idx + 1}/8]")
    p_desc = perguntas_config[idx]
    
    if idx == 0:
        add_outline_card(sl, "Pergunta Principal:", p_desc, Inches(0.7), Inches(1.5), Inches(12.1), Inches(2.7), outline_word=f"ACTIVE Q0{idx+1}", border_color=ROXO, tit_size=18, desc_size=18, align_center_v=True)
        add_outline_card(sl, "Apoio ao Professor:", "• Para expandir: 'Dê um exemplo recente...'\n• Para concordar/discordar: 'Qual é o outro lado dessa moeda?'", Inches(0.7), Inches(4.4), Inches(12.1), Inches(2.4), outline_word="PROMPTS", border_color=TURQUESA, tit_size=16, desc_size=14, align_center_v=True)
    else:
        add_outline_card(sl, "Pergunta Principal:", p_desc, Inches(0.7), Inches(1.5), Inches(7.2), Inches(2.7), outline_word=f"ACTIVE Q0{idx+1}", border_color=ROXO, tit_size=18, desc_size=16, align_center_v=True)
        add_outline_card(sl, "Apoio ao Professor:", "• Para expandir: 'Dê um exemplo recente...'\n• Para discordar: 'Qual é o outro lado?'", Inches(0.7), Inches(4.4), Inches(7.2), Inches(2.4), outline_word="PROMPTS", border_color=TURQUESA, tit_size=15, desc_size=14, align_center_v=True)
        history_list = [perguntas_config[i] for i in range(idx)]
        add_outline_history_staircase(sl, history_list)

# SLIDE 20
sl20 = base_slide(20)
add_header(sl20, "Atividade de Escrita", "Consolidação e Memorização")
add_outline_card(sl20, "Tarefa de Escrita Recomendada:", "[Instrução detalhada da tarefa: escrever um texto, frases ou diálogo]", Inches(0.7), Inches(1.5), Inches(12.1), Inches(2.8), outline_word="TASK", border_color=ROXO, tit_size=18, desc_size=15)
add_outline_card(sl20, "Fundamentação Pedagógica:", "Identificamos que a escrita nesta etapa é fundamental para você fixar as estruturas aprendidas, expandir vocabulário ativo e diminuir a tradução mental automática do seu idioma nativo.", Inches(0.7), Inches(4.5), Inches(12.1), Inches(2.3), outline_word="WHY", border_color=TURQUESA, desc_size=14, align_center_v=True)

# SLIDE 21
sl21 = base_slide(21)
add_header(sl21, "Encerramento — Tarefa de Casa", "Ação e Prática Contínua")
add_outline_card(sl21, "Tarefa de Casa:", "• [Atividade prática 1 para realizar fora da aula]\n• [Atividade de fixação de vocabulário ou escuta]", Inches(0.7), Inches(1.5), Inches(12.1), Inches(5.3), outline_word="HOMEWORK", border_color=ROXO, tit_size=18, desc_size=15)

# SLIDE 22
sl22 = base_slide(22)
add_header(sl22, "Encerramento & Feedback Rápido", "Considerações Finais")
add_outline_card(sl22, "AGRADECIMENTO:", "Obrigado pela dedicação na aula de hoje!", Inches(0.7), Inches(1.5), Inches(12.1), Inches(2.0), outline_word="THANKS", border_color=TURQUESA, tit_size=20, desc_size=16)
add_outline_card(sl22, "ESPAÇO PARA FEEDBACK RÁPIDO:", "[Espaço destinado para o professor coletar dúvidas rápidas e considerações do aluno]", Inches(0.7), Inches(3.7), Inches(12.1), Inches(3.1), outline_word="FEEDBACK", border_color=ROXO, tit_size=18, desc_size=15)

out_path_local = "template-haas-modelo-19-outline.pptx"
prs.save(out_path_local)

server_dir = "/var/www/haas-frontend-desk-mobile-oficial/public/temp-miniaturas/"
os.makedirs(server_dir, exist_ok=True)
out_path_server = os.path.join(server_dir, "template-haas-modelo-19-outline.pptx")

import shutil
shutil.copyfile(out_path_local, out_path_server)
print(f"SUCCESSFULLY GENERATED OUTLINE MODEL 19: {out_path_server}")
