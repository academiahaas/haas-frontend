import os
import urllib.request
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# PALETA OFICIAL HAAS LANGUAGE - MODELO 15 (BUBBLE CHAT EDITION)
BRANCO = RGBColor(0xFF, 0xFF, 0xFF)
ROXO = RGBColor(0x4F, 0x35, 0xE2)
AZUL = RGBColor(0x1E, 0x3A, 0x8A)
SLATE = RGBColor(0x33, 0x41, 0x55)
TURQUESA = RGBColor(0x14, 0xC8, 0xB0)
BORDA = RGBColor(0xCB, 0xD5, 0xE1)
ROXO_CLARO = RGBColor(0xED, 0xE9, 0xFE)
TURQUESA_CLARO = RGBColor(0xE0, 0xF2, 0xF1)
FUNDO_CHAT = RGBColor(0xF1, 0xF5, 0xF9)

LOGO_URL = "https://jdppxfokfhqjudwfwckd.supabase.co/storage/v1/object/public/haas-academy/assignments/Design%20sem%20nome%20(4).png"
LOGO_PATH = "/tmp/haas-logo.png"

if not os.path.exists(LOGO_PATH):
    try:
        urllib.request.urlretrieve(LOGO_URL, LOGO_PATH)
    except Exception:
        pass

W, H = Inches(13.333), Inches(7.5)
TOTAL_SLIDES = 22

prs = Presentation()
prs.slide_width = W
prs.slide_height = H

def base_slide(num, hide_logo=False):
    sl = prs.slides.add_slide(prs.slide_layouts[6])
    sl.background.fill.solid()
    sl.background.fill.fore_color.rgb = FUNDO_CHAT

    top_bar = sl.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), W, Inches(0.12))
    top_bar.fill.solid(); top_bar.fill.fore_color.rgb = ROXO; top_bar.line.fill.background()

    if os.path.exists(LOGO_PATH) and not hide_logo:
        sl.shapes.add_picture(LOGO_PATH, Inches(10.8), Inches(0.2), height=Inches(0.65))
        
    tb = sl.shapes.add_textbox(Inches(0.5), Inches(7.12), Inches(8), Inches(0.35))
    p = tb.text_frame.paragraphs[0]
    p.text = "HAAS LANGUAGE  •  BUBBLE CHAT EDITION 💬"
    p.font.size = Pt(10); p.font.color.rgb = ROXO; p.font.bold = True
    
    tb2 = sl.shapes.add_textbox(Inches(10.8), Inches(7.12), Inches(2.0), Inches(0.35))
    p2 = tb2.text_frame.paragraphs[0]
    p2.text = f"SLIDE {num} / {TOTAL_SLIDES}"
    p2.font.size = Pt(10); p2.font.color.rgb = AZUL; p2.font.bold = True; p2.alignment = PP_ALIGN.RIGHT
    return sl

def add_header(sl, tit, sub=""):
    tb = sl.shapes.add_textbox(Inches(0.5), Inches(0.35), Inches(10), Inches(0.65))
    tf = tb.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]; p.text = tit
    p.font.size = Pt(30); p.font.color.rgb = ROXO; p.font.bold = True
    
    if sub:
        tb_s = sl.shapes.add_textbox(Inches(0.5), Inches(0.98), Inches(10), Inches(0.4))
        p_s = tb_s.text_frame.paragraphs[0]
        p_s.text = sub
        p_s.font.size = Pt(17); p_s.font.color.rgb = AZUL; p_s.font.bold = True

def add_chat_bubble(sl, tit, desc, left, top, width, height, bg_color=BRANCO, border_color=BORDA, tit_color=ROXO, desc_color=SLATE, tail_dir='left_top', tit_size=16, desc_size=14, align_center_v=False, avatar_label=""):
    s = sl.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    s.fill.solid(); s.fill.fore_color.rgb = bg_color
    if border_color:
        s.line.color.rgb = border_color; s.line.width = Pt(1.5)
    else:
        s.line.fill.background()

    tail_w, tail_h = Inches(0.2), Inches(0.2)
    if tail_dir == 'left_top':
        t = sl.shapes.add_shape(MSO_SHAPE.RIGHT_TRIANGLE, left - Inches(0.12), top + Inches(0.15), tail_w, tail_h)
        t.fill.solid(); t.fill.fore_color.rgb = bg_color
        if border_color: t.line.color.rgb = border_color; t.line.width = Pt(1.5)
        else: t.line.fill.background()
        t.rotation = 180
    elif tail_dir == 'right_bottom':
        t = sl.shapes.add_shape(MSO_SHAPE.RIGHT_TRIANGLE, left + width - Inches(0.08), top + height - Inches(0.35), tail_w, tail_h)
        t.fill.solid(); t.fill.fore_color.rgb = bg_color
        if border_color: t.line.color.rgb = border_color; t.line.width = Pt(1.5)
        else: t.line.fill.background()
        t.rotation = 0

    if avatar_label:
        av = sl.shapes.add_shape(MSO_SHAPE.OVAL, left - Inches(0.55), top + Inches(0.1), Inches(0.45), Inches(0.45))
        av.fill.solid(); av.fill.fore_color.rgb = ROXO if bg_color != ROXO else TURQUESA
        av.line.fill.background()
        tf_a = av.text_frame; tf_a.vertical_anchor = MSO_ANCHOR.MIDDLE
        p_a = tf_a.paragraphs[0]; p_a.text = avatar_label; p_a.alignment = PP_ALIGN.CENTER
        p_a.font.size = Pt(11); p_a.font.color.rgb = BRANCO; p_a.font.bold = True
        
    tf = s.text_frame; tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE if align_center_v else MSO_ANCHOR.TOP
    margin_val = Inches(0.12) if height < Inches(1.0) else Inches(0.2)
    tf.margin_left = margin_val; tf.margin_top = margin_val; tf.margin_right = margin_val; tf.margin_bottom = margin_val
    
    if tit:
        p = tf.paragraphs[0]
        p.text = tit; p.font.size = Pt(tit_size); p.font.color.rgb = tit_color; p.font.bold = True
        if desc:
            p2 = tf.add_paragraph()
            p2.text = desc; p2.font.size = Pt(desc_size); p2.font.color.rgb = desc_color; p2.space_before = Pt(6)
    else:
        if desc:
            p = tf.paragraphs[0]
            p.text = desc; p.font.size = Pt(desc_size); p.font.color.rgb = desc_color

def img_chat_media(sl, left, top, width, height, label="[Anexo de Mídia / Imagem]"):
    s = sl.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    s.fill.solid(); s.fill.fore_color.rgb = ROXO_CLARO
    s.line.color.rgb = TURQUESA; s.line.width = Pt(2)
    tf = s.text_frame; tf.word_wrap = True; tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER
    p = tf.paragraphs[0]; p.text = f"📷 {label}"; p.font.size = Pt(16); p.font.color.rgb = ROXO; p.font.bold = True

def add_chat_flashcard(sl, num_str, desc, left, top, width, height, bg_c=BRANCO, b_c=TURQUESA):
    s = sl.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    s.fill.solid(); s.fill.fore_color.rgb = bg_c; s.line.color.rgb = b_c; s.line.width = Pt(1.5)
    
    chip = sl.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left + Inches(0.12), top + Inches(0.12), Inches(0.55), Inches(0.35))
    chip.fill.solid(); chip.fill.fore_color.rgb = b_c; chip.line.fill.background()
    tf_chip = chip.text_frame; tf_chip.vertical_anchor = MSO_ANCHOR.MIDDLE
    p_chip = tf_chip.paragraphs[0]; p_chip.text = f"#{num_str}"; p_chip.alignment = PP_ALIGN.CENTER
    p_chip.font.size = Pt(11); p_chip.font.color.rgb = BRANCO; p_chip.font.bold = True
    
    tf = s.text_frame; tf.word_wrap = True; tf.vertical_anchor = MSO_ANCHOR.TOP
    tf.margin_left = Inches(0.8); tf.margin_top = Inches(0.12); tf.margin_right = Inches(0.15); tf.margin_bottom = Inches(0.1)
    p = tf.paragraphs[0]; p.text = desc; p.font.size = Pt(14); p.font.color.rgb = SLATE; p.font.bold = False

# SLIDE 1: CAPA
sl1 = base_slide(1, hide_logo=True)
add_chat_bubble(sl1, "[Título da Aula / Tema Principal]", "[Subtítulo / Descrição da Aula]\n\n📅 Data da Aula: [DD/MM/AAAA]", Inches(0.8), Inches(1.5), Inches(6.5), Inches(5.2), bg_color=ROXO, border_color=None, tit_color=BRANCO, desc_color=TURQUESA_CLARO, tail_dir='left_top', tit_size=36, desc_size=18, avatar_label="HL")
img_chat_media(sl1, Inches(7.8), Inches(1.5), Inches(5.0), Inches(5.2), "[Imagem da Capa]")

# SLIDE 2: PRESENÇA
sl2 = base_slide(2)
add_header(sl2, "Presença & Integrantes", "Membros da Sessão de Hoje")
add_chat_bubble(sl2, "Professor(a):", "[Nome do(a) Professor(a)]", Inches(0.8), Inches(1.6), Inches(11.5), Inches(1.6), bg_color=ROXO_CLARO, border_color=ROXO, tail_dir='left_top', avatar_label="P")
add_chat_bubble(sl2, "Aluno(s):", "[Nome do Aluno / Lista da Turma]", Inches(1.2), Inches(3.5), Inches(11.5), Inches(3.2), bg_color=TURQUESA_CLARO, border_color=TURQUESA, tail_dir='right_bottom', avatar_label="A")

# SLIDE 3: WARM-UP
sl3 = base_slide(3)
add_header(sl3, "Aquecimento / Warm-up", "Foco 100% na Fala (3 Minutos)")
add_chat_bubble(sl3, "Pergunta Provocativa Rápida:", "[Insira aqui a pergunta provocativa de aquecimento]", Inches(0.8), Inches(1.5), Inches(6.8), Inches(2.6), bg_color=BRANCO, border_color=ROXO, tit_size=18, desc_size=15, tail_dir='left_top')
add_chat_bubble(sl3, "🎙️ Áudio Neurocognitivo:", "Ativa o córtex pré-frontal e reduz a inibição linguística inicial.", Inches(0.8), Inches(4.3), Inches(6.8), Inches(2.3), bg_color=TURQUESA_CLARO, border_color=TURQUESA, tit_color=AZUL, desc_color=SLATE, tit_size=15, desc_size=14, tail_dir='left_top')
img_chat_media(sl3, Inches(8.0), Inches(1.5), Inches(4.8), Inches(5.1), "[Imagem de Aquecimento]")

# SLIDE 4: OBJETIVOS
sl4 = base_slide(4)
add_header(sl4, "Objetivos da Aula", "📌 Mensagens Fixadas da Aula")
add_chat_bubble(sl4, "Metas da Aula:", "• [Objetivo 1]\n• [Objetivo 2]\n• [Objetivo 3]", Inches(0.8), Inches(1.5), Inches(7.2), Inches(5.2), bg_color=BRANCO, border_color=ROXO, tail_dir='left_top')
add_chat_bubble(sl4, "Nota de Escrita:", "Nesta aula daremos atenção à escrita para reforçar a memorização de estruturas/vocabulário e reduzir a interferência da língua materna.", Inches(8.3), Inches(1.5), Inches(4.5), Inches(5.2), bg_color=ROXO_CLARO, border_color=ROXO, tail_dir='right_bottom')

# SLIDE 5: METODOLOGIA
sl5 = base_slide(5)
add_header(sl5, "A Nossa Metodologia", "Thread de Aprendizado Haas Language")
met_bubbles = [
    ("1. Fala Ativa", "Conversação inicial sem julgamentos.", 0.8, 1.5, 5.8, ROXO_CLARO, ROXO, 'left_top'),
    ("2. Exposição", "Apresentação leve da gramática.", 6.8, 1.5, 6.0, TURQUESA_CLARO, TURQUESA, 'right_bottom'),
    ("3. Prática Guiada", "Exercícios com suporte.", 0.8, 4.2, 3.8, BRANCO, ROXO, 'left_top'),
    ("4. Conversação", "Aplicação prática das estruturas.", 4.8, 4.2, 4.0, TURQUESA_CLARO, TURQUESA, 'left_top'),
    ("5. Escrita", "Consolidação e redução de tradução.", 9.0, 4.2, 3.8, ROXO_CLARO, ROXO, 'right_bottom')
]
for tit, desc, l, t, w, bg_c, b_c, t_dir in met_bubbles:
    add_chat_bubble(sl5, tit, desc, Inches(l), Inches(t), Inches(w), Inches(2.4), bg_color=bg_c, border_color=b_c, tail_dir=t_dir)

# SLIDES 6 A 8: EXPOSIÇÃO
for s_num in [6, 7, 8]:
    sl = base_slide(s_num)
    add_header(sl, f"Explicação do Tema — Parte {s_num - 5}")
    add_chat_bubble(sl, "Conceito Principal:", f"[Explicação e exemplo central da estrutura - Bloco {s_num - 5}]", Inches(0.8), Inches(1.5), Inches(7.5), Inches(5.2), bg_color=BRANCO, border_color=ROXO, tit_size=20, desc_size=15, tail_dir='left_top')
    add_chat_bubble(sl, "Exemplos Práticos:", "[Exemplos adicionais e frases do dia a dia]", Inches(8.5), Inches(1.5), Inches(4.3), Inches(2.5), bg_color=TURQUESA_CLARO, border_color=TURQUESA, tail_dir='right_bottom')
    add_chat_bubble(sl, "Notas & Detalhes:", "[Exceções, nuances gramaticais e contexto]", Inches(8.5), Inches(4.2), Inches(4.3), Inches(2.5), bg_color=ROXO_CLARO, border_color=ROXO, tail_dir='right_bottom')

# SLIDE 9: SÍNTESE + CULTURA
sl9 = base_slide(9)
add_header(sl9, "Síntese do Conteúdo & Sugestão Cultural")
add_chat_bubble(sl9, "Resumo do Tema:", "[Resumo sintético dos pontos chave ensinados na aula]", Inches(0.8), Inches(1.5), Inches(6.0), Inches(5.2), bg_color=BRANCO, border_color=ROXO, tail_dir='left_top')
add_chat_bubble(sl9, "🎵 Mídia Cultural Encodada:", "[Nome da Música / Vídeo recomendado sobre o assunto].\n\nO professor pode rodar o áudio/vídeo ou deixar como recomendação extra para o aluno praticar.", Inches(7.1), Inches(1.5), Inches(5.7), Inches(5.2), bg_color=TURQUESA_CLARO, border_color=TURQUESA, tit_color=AZUL, tail_dir='right_bottom')

# SLIDES 10 E 11: PRÁTICA GUIADA
sl10 = base_slide(10)
add_header(sl10, "Prática Guiada — Blocos 1 a 5", "Correção de Vícios")
add_chat_flashcard(sl10, "01", "[Lacunas contextualizadas com opções]", Inches(0.5), Inches(1.5), Inches(12.3), Inches(1.2), bg_c=BRANCO, b_c=ROXO)
add_chat_flashcard(sl10, "02", "[Lacunas contextualizadas com opções]", Inches(0.5), Inches(2.9), Inches(6.0), Inches(1.8), bg_c=TURQUESA_CLARO, b_c=TURQUESA)
add_chat_flashcard(sl10, "03", "[Preenchimento direto de verbos]", Inches(6.8), Inches(2.9), Inches(6.0), Inches(1.8), bg_c=TURQUESA_CLARO, b_c=TURQUESA)
add_chat_flashcard(sl10, "04", "[Reestruturação / Transformação de frase]", Inches(0.5), Inches(4.9), Inches(6.0), Inches(1.9), bg_c=ROXO_CLARO, b_c=ROXO)
add_chat_flashcard(sl10, "05", "[Reestruturação / Transformação de frase]", Inches(6.8), Inches(4.9), Inches(6.0), Inches(1.9), bg_c=ROXO_CLARO, b_c=ROXO)

sl11 = base_slide(11)
add_header(sl11, "Prática Guiada — Blocos 6 a 10", "Desafios Práticos")
add_chat_flashcard(sl11, "06", "[Preposição / Conector correto]", Inches(0.5), Inches(1.5), Inches(12.3), Inches(1.2), bg_c=BRANCO, b_c=TURQUESA)
add_chat_flashcard(sl11, "07", "[Escolha da opção correta]", Inches(0.5), Inches(2.9), Inches(6.0), Inches(1.8), bg_c=ROXO_CLARO, b_c=ROXO)
add_chat_flashcard(sl11, "08", "[Substituição por sinônimo formal]", Inches(6.8), Inches(2.9), Inches(6.0), Inches(1.8), bg_c=ROXO_CLARO, b_c=ROXO)
add_chat_flashcard(sl11, "09", "[Identificação de erro e correção]", Inches(0.5), Inches(4.9), Inches(6.0), Inches(1.9), bg_c=TURQUESA_CLARO, b_c=TURQUESA)
add_chat_flashcard(sl11, "10", "[Tradução de sentido sem uso literal]", Inches(6.8), Inches(4.9), Inches(6.0), Inches(1.9), bg_c=TURQUESA_CLARO, b_c=TURQUESA)

# SLIDES 12 A 19: CONVERSAÇÃO (INÉDITO: PERGUNTA ATIVA NO TOPO + HISTÓRICO HORIZONTAL NO RODAPÉ)
perguntas_config = [
    "[Pergunta principal 1 sobre o tema]", "[Pergunta principal 2 sobre o tema]", "[Pergunta principal 3 sobre o tema]",
    "[Pergunta principal 4 sobre o tema]", "[Pergunta principal 5 sobre o tema]", "[Pergunta principal 6 sobre o tema]",
    "[Pergunta principal 7 sobre o tema]", "[Pergunta principal 8 sobre o tema]"
]

def add_chat_history_bottom_bar(sl, history_list):
    count = len(history_list)
    if count == 0: return
    gap = 0.12
    avail_w = 12.3 - (gap * (count - 1))
    bubble_w = avail_w / count
    
    for i, h_txt in enumerate(history_list):
        left_x = 0.5 + i * (bubble_w + gap)
        clean_t = h_txt[:28] + "..." if len(h_txt) > 28 else h_txt
        d_sz = 11 if count <= 3 else (10 if count <= 5 else 9)
        bg_col = ROXO_CLARO if i % 2 == 0 else TURQUESA_CLARO
        b_col = ROXO if i % 2 == 0 else TURQUESA
        
        add_chat_bubble(sl, f"💬 Q{i+1}", clean_t, Inches(left_x), Inches(5.4), Inches(bubble_w), Inches(1.4), bg_color=bg_col, border_color=b_col, tit_size=11, desc_size=d_sz, align_center_v=False, tail_dir='left_top')

for idx in range(8):
    s_num = 12 + idx
    sl = base_slide(s_num)
    add_header(sl, "Conversação", f"Thread de Pergunta {idx + 1} de 8")
    
    p_desc = perguntas_config[idx]
    
    if idx == 0:
        add_chat_bubble(sl, f"Pergunta {idx+1}:", p_desc, Inches(0.8), Inches(1.5), Inches(12.0), Inches(2.7), bg_color=BRANCO, border_color=ROXO, tit_size=20, desc_size=18, tail_dir='left_top', avatar_label="Q1")
        add_chat_bubble(sl, "💡 Dicas de Apoio para a Resposta:", "• Para expandir: 'Dê um exemplo recente...'\n• Para concordar/discordar: 'Qual é o outro lado dessa moeda?'", Inches(0.8), Inches(4.4), Inches(12.0), Inches(2.3), bg_color=ROXO_CLARO, border_color=ROXO, tit_color=AZUL, tit_size=16, desc_size=14, tail_dir='right_bottom')
    else:
        add_chat_bubble(sl, f"Pergunta Ativa {idx+1}:", p_desc, Inches(0.8), Inches(1.5), Inches(12.0), Inches(2.2), bg_color=BRANCO, border_color=ROXO, tit_size=18, desc_size=16, tail_dir='left_top', avatar_label=f"Q{idx+1}")
        add_chat_bubble(sl, "💡 Dicas de Apoio:", "• Para expandir: 'Dê um exemplo recente...'\n• Para discordar: 'Qual é o outro lado?'", Inches(0.8), Inches(3.85), Inches(12.0), Inches(1.35), bg_color=ROXO_CLARO, border_color=ROXO, tit_color=AZUL, tit_size=14, desc_size=13, tail_dir='right_bottom')
        
        history_list = [perguntas_config[i] for i in range(idx)]
        add_chat_history_bottom_bar(sl, history_list)

# SLIDE 20: ESCRITA
sl20 = base_slide(20)
add_header(sl20, "Atividade de Escrita", "Consolidação e Memorização")
add_chat_bubble(sl20, "Tarefa de Escrita Recomendada:", "[Instrução detalhada da tarefa: escrever um texto, frases ou diálogo]", Inches(0.8), Inches(1.5), Inches(12.0), Inches(2.8), bg_color=BRANCO, border_color=ROXO, tail_dir='left_top')
add_chat_bubble(sl20, "Fundamentação:", "Identificamos que a escrita nesta etapa é fundamental para você fixar as estruturas aprendidas, expandir vocabulário ativo e diminuir a tradução mental automática do seu idioma nativo.", Inches(0.8), Inches(4.6), Inches(12.0), Inches(2.2), bg_color=TURQUESA_CLARO, border_color=TURQUESA, tit_color=AZUL, desc_size=14, tail_dir='right_bottom')

# SLIDE 21: TAREFA DE CASA
sl21 = base_slide(21)
add_header(sl21, "Encerramento — Tarefa de Casa", "Ação e Prática Contínua")
add_chat_bubble(sl21, "Tarefa de Casa:", "• [Atividade prática 1 para realizar fora da aula]\n• [Atividade de fixação de vocabulário ou escuta]", Inches(0.8), Inches(1.5), Inches(12.0), Inches(5.2), bg_color=BRANCO, border_color=ROXO, tail_dir='left_top')

# SLIDE 22: ENCERRAMENTO & FEEDBACK
sl22 = base_slide(22)
add_header(sl22, "Encerramento & Feedback Rápido", "Considerações Finais")
add_chat_bubble(sl22, "Agradecimento:", "Obrigado pela dedicação na aula de hoje!", Inches(0.8), Inches(1.5), Inches(12.0), Inches(2.0), bg_color=ROXO_CLARO, border_color=ROXO, tit_size=20, desc_size=16, tail_dir='left_top')
add_chat_bubble(sl22, "Espaço para Feedback Rápido:", "[Espaço destinado para o professor coletar dúvidas rápidas e considerações do aluno]", Inches(0.8), Inches(3.8), Inches(12.0), Inches(2.9), bg_color=TURQUESA_CLARO, border_color=TURQUESA, tit_size=18, desc_size=15, tail_dir='right_bottom')

# SALVAR NO SERVIDOR E DISPONIBILIZAR LINK
out_path_server = "/var/www/haas-frontend-desk-mobile-oficial/public/temp-miniaturas/template-haas-modelo-15-bubble-chat.pptx"
prs.save(out_path_server)
print(f"🎉 MODELO 15 BUBBLE CHAT GERADO EM: {out_path_server}")
